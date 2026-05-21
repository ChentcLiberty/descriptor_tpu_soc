from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

import render_breath_soc_rtl_arch_ppt as arch


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
BASE_DECK = DOCS / "陈韦东tinytpusoc_呼吸识别SoC更新版.pptx"
OUT_STANDALONE = DOCS / "20_CPU_TPU_呼吸识别SoC_RTL架构图_6p_含讲解.pptx"
OUT_APPENDED = DOCS / "陈韦东tinytpusoc_呼吸识别SoC更新版_v3_RTL架构图_含讲解.pptx"
OUT_NOTES = DOCS / "CPU_TPU_呼吸识别SoC_RTL架构图_讲解稿_20260419.md"


def add_talk_track_slide(prs: Presentation, page: str = "P23"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    arch.title(
        slide,
        page,
        "RTL 架构图讲解顺序",
        "先讲系统分工，再讲控制面和数据面，最后落到 Q8.8 MLP RTL 精度修正",
    )
    arch.pill(slide, 10.26, 0.34, 2.36, 0.34, "建议 2-3 分钟讲完", arch.GREEN_FILL, arch.GREEN_LINE, 10.8, True)

    arch.round_box(slide, 0.44, 1.22, 12.44, 5.98, arch.WHITE, arch.PANEL, 1.1)
    arch.module(slide, 0.78, 1.68, 2.10, 1.02, "1. 先定边界", "CPU 跑 CNN/FiLM 前端\nTPU RTL 跑 MLP/Classifier\nshared SRAM 做数据交换区", arch.BLUE_FILL, arch.BLUE_LINE, 10.8)
    arch.module(slide, 3.54, 1.68, 2.10, 1.02, "2. 讲控制面", "CPU 通过 DBUS 访问 TPU_CTRL\nsplitter 区分 UART/APB 和 TPU\nlaunch/status/perf 是软件可见接口", arch.GREEN_FILL, arch.GREEN_LINE, 10.4)
    arch.module(slide, 6.30, 1.68, 2.10, 1.02, "3. 讲数据面", "descriptor 放在 shared SRAM\nDMA 拉 input/param\nMLP 输出再写回 output blob", arch.ORANGE_FILL, arch.ORANGE_LINE, 10.4)
    arch.module(slide, 9.06, 1.68, 2.10, 1.02, "4. 讲验证点", "Q8.8 MAC 先完整累加\n最后 round shift\nRTL 输出与 Python golden 对齐", arch.RED_FILL, arch.RED_LINE, 10.4)

    arch.poly_arrow(slide, [(2.88, 2.20), (3.54, 2.20)], arch.BLUE_LINE, 1.8)
    arch.poly_arrow(slide, [(5.64, 2.20), (6.30, 2.20)], arch.GREEN_LINE, 1.8)
    arch.poly_arrow(slide, [(8.40, 2.20), (9.06, 2.20)], arch.ORANGE_LINE, 1.8)

    arch.module(slide, 0.78, 3.42, 3.20, 1.46, "30 秒版本", "这个 SoC 现在已经把 CPU 启动、TPU 控制寄存器、descriptor DMA、shared SRAM 和 Q8.8 MLP RTL 串起来了。CPU 负责前端算法和调度，TPU 负责最后的 MLP/分类器，数据通过 shared SRAM 交换。", arch.INK_FILL, arch.INK_LINE, 10.7)
    arch.module(slide, 4.38, 3.42, 3.20, 1.46, "2 分钟版本", "按 P19 到 P22 讲：总图说明模块分工；控制面说明 memory-mapped TPU_CTRL；数据面说明 descriptor DMA 读写 shared SRAM；最后说明 Q8.8 MAC 精度修正和 RTL 验证结果。", arch.INK_FILL, arch.INK_LINE, 10.7)
    arch.module(slide, 7.98, 3.42, 3.20, 1.46, "被追问时", "明确说当前不是完整 CNN 都映射到 TPU。已经验证的是 TPU MLP/Classifier RTL 主路径；CNN/FiLM 前端先跑 CPU，后续若性能不足再考虑下沉到 TPU 或新增专用 CNN/VPU。", arch.INK_FILL, arch.INK_LINE, 10.5)

    arch.pill(slide, 0.90, 5.62, 1.38, 0.28, "不要过度承诺", arch.RED_FILL, arch.RED_LINE, 9.0, True)
    arch.textbox(slide, 2.42, 5.56, 8.90, 0.46, "推荐口径：已经完成 TPU MLP/Classifier 的 RTL 验证链路，以及 CPU+TPU 端到端数据闭环；完整算法端到端结果用 Python/fixed golden 和 CPU 软件路径验证。", size=10.0, color=arch.TEXT)
    arch.textbox(slide, 0.62, 6.90, 11.70, 0.18, "这页可以作为答辩讲稿索引，不需要展开每个信号名；追问时再回到对应图页。", size=9.4, color=arch.SUB)


def add_generated_code_slide(prs: Presentation, page: str = "P24"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    arch.title(
        slide,
        page,
        "新生成代码与文件说明",
        "本轮只新增文档/绘图脚本和 PPT/PNG 产物，没有改动 RTL 功能逻辑",
    )
    arch.pill(slide, 10.12, 0.34, 2.54, 0.34, "代码边界说明", arch.ORANGE_FILL, arch.ORANGE_LINE, 10.8, True)

    arch.round_box(slide, 0.44, 1.22, 12.44, 5.98, arch.WHITE, arch.PANEL, 1.1)
    arch.module(slide, 0.78, 1.62, 3.36, 1.16, "render_breath_soc_rtl_arch_ppt.py", "作用：用 python-pptx 画 4 页 RTL 架构图\n输入：人工抽象后的现有 RTL 结构\n输出：4p standalone PPT + 追加到 18p 总 PPT 的 v2", arch.BLUE_FILL, arch.BLUE_LINE, 10.3)
    arch.module(slide, 4.92, 1.62, 3.36, 1.16, "render_breath_soc_rtl_arch_explained_ppt.py", "作用：复用前一个绘图脚本，再追加 2 页讲解\n输出：6p 含讲解 PPT + v3 总 PPT\n同时生成 Markdown 讲稿", arch.GREEN_FILL, arch.GREEN_LINE, 10.1)
    arch.module(slide, 9.06, 1.62, 2.44, 1.16, "讲稿 Markdown", "逐页讲解口径\n新文件清单\nRTL 映射关系\n当前边界和追问回答", arch.ORANGE_FILL, arch.ORANGE_LINE, 10.2)

    arch.module(slide, 0.78, 3.34, 2.58, 1.34, "图页对应的 RTL", "P19: panda_soc_stage2_base_top.v\nP20: cpu_tpu_axil_splitter.v / tpu_ctrl_axil_regs.v\nP21: tpu_desc_fetch_dma_stub.v / shared_mem_subsys.v\nP22: tpu_mlp_compute_stub.v", arch.INK_FILL, arch.INK_LINE, 9.4)
    arch.module(slide, 3.78, 3.34, 2.58, 1.34, "生成物", "20_CPU_TPU_..._6p_含讲解.pptx\n陈韦东..._v3_RTL架构图_含讲解.pptx\nPNG/PDF 预览图\n讲解稿 md", arch.INK_FILL, arch.INK_LINE, 9.4)
    arch.module(slide, 6.78, 3.34, 2.58, 1.34, "未做的事", "没有把 CNN/FiLM 下沉到 TPU\n没有新增真实 CNN/VPU RTL\n没有改变 CPU 软件算法\n没有修改已验证 MLP RTL 行为", arch.RED_FILL, arch.RED_LINE, 9.6)
    arch.module(slide, 9.78, 3.34, 2.00, 1.34, "后续怎么改", "改图：编辑两个 render 脚本\n改讲稿：编辑 md\n改硬件：另行修改 RTL，再同步更新图", arch.GREEN_FILL, arch.GREEN_LINE, 9.4)

    arch.pill(slide, 1.02, 5.36, 2.00, 0.30, "这不是自动 RTL 解析器", arch.RED_FILL, arch.RED_LINE, 9.2, True)
    arch.textbox(slide, 3.22, 5.32, 8.30, 0.58, "绘图脚本是面试表达用的可编辑 PPT 生成器。它根据已经读过的 RTL 结构手工建模成图，不参与综合、不参与仿真，也不会影响验证结果。", size=10.2, color=arch.TEXT)
    arch.textbox(slide, 0.62, 6.90, 11.70, 0.18, "面试里可以说：图是我根据 RTL 层次重新整理的架构表达，真实验证依据仍然是 VCS log、Python golden 和 RTL testbench。", size=9.4, color=arch.SUB)


def write_notes():
    text = """# CPU+TPU 呼吸识别 SoC RTL 架构图讲解稿

生成时间：2026-04-19

## 一句话总览

当前硬件验证重点是 TPU MLP/Classifier RTL 链路。CPU 负责呼吸算法前端和任务调度，TPU 控制寄存器负责 launch/status，descriptor DMA 负责从 shared SRAM 拉取 input/param 并写回 output，Q8.8 MLP compute block 负责最后的线性层/分类器计算。

需要明确边界：完整 CNN/FiLM 前端当前没有全部映射到 TPU，先由 CPU 软件路径承担；TPU 已验证的是 MLP/Classifier 加速链路和 CPU+TPU 数据闭环。

## 新生成代码有哪些

1. `docs/render_breath_soc_rtl_arch_ppt.py`

   这是第一版 RTL 架构图生成脚本。它使用 `python-pptx` 画 4 页可编辑图：SoC 总图、TPU 控制面、descriptor DMA/shared SRAM 数据面、Q8.8 MLP compute 细节。脚本的 helper 函数包括 `module`、`pill`、`poly_arrow`、`title` 等，只负责画 PPT 元素。

   这个脚本不会参与 RTL 综合或仿真，也不会修改硬件逻辑。它的输入是我根据当前 RTL 文件整理出来的结构关系，输出是 PPT/PNG/PDF 文档产物。

2. `docs/render_breath_soc_rtl_arch_explained_ppt.py`

   这是本轮新增的改进脚本。它复用 `render_breath_soc_rtl_arch_ppt.py` 里的绘图函数和 4 页架构图，再追加 2 页：`RTL 架构图讲解顺序` 和 `新生成代码与文件说明`。同时它会写出这一份 Markdown 讲稿。

   输出文件是 6 页 standalone PPT 和追加到原 18 页总 PPT 后面的 v3 总 PPT。

3. `docs/CPU_TPU_呼吸识别SoC_RTL架构图_讲解稿_20260419.md`

   这是面试讲稿和文件说明。它把每一页图的讲法、对应 RTL、当前边界、常见追问回答整理成文字。

## 新生成文件不是 RTL 功能代码

本轮新增的是文档生成脚本和 PPT/PNG/PDF 产物，不是新的硬件功能 RTL。已经画进去的硬件模块来自现有工程：

- `panda_soc_stage2_base_top.v`
- `cpu_tpu_axil_splitter.v`
- `tpu_ctrl_axil_regs.v`
- `tpu_desc_fetch_dma_stub.v`
- `panda_soc_shared_mem_subsys.v`
- `tpu_mlp_compute_stub.v`

## 每页怎么讲

### P19 CPU+TPU 呼吸识别 SoC RTL 架构总图

先讲系统分工：Panda RISC-V CPU 是调度者，负责启动程序、写 TPU 控制寄存器、准备 descriptor/input/param，并通过 dcache AXI 访问 shared SRAM。TPU 侧由 `TPU_CTRL Regs` 接收 launch，再由 `TPU Desc DMA` 从 shared SRAM 读 descriptor、input、param，最后把 MLP 输出写回 output 区。

这一页要强调两条路径：控制面是 CPU DBUS -> AXI-Lite splitter -> TPU_CTRL；数据面是 CPU dcache AXI 和 TPU AXI master 共同访问 shared SRAM。

推荐表达：现在已经闭合的是 CPU 写 TPU_CTRL、TPU DMA 读写 shared SRAM、Q8.8 MLP 输出写回这条硬件主路径。

### P20 TPU 控制面 RTL 细节

讲 `cpu_tpu_axil_splitter.v` 怎么按地址窗口分流：`0x4000_4000` 起的 4KB 走 TPU_CTRL，其他 legacy 外设地址走 AXI-APB bridge，例如 UART。写通道用 `write_busy_reg` 保持 AW 和 W 属于同一路，读通道用 `read_busy_reg` 保持返回数据来自同一路。

讲 `tpu_ctrl_axil_regs.v` 的寄存器：`CTRL` 产生 `launch_pulse` 和 `soft_reset_pulse`，`STATUS` 回读 busy/done/error，`MODE/NET_ID/DESC_LO/DESC_HI` 描述任务，`PERF_CYCLE` 在 busy 时计数。

### P21 Descriptor DMA 与 Shared SRAM 数据面

讲 descriptor 驱动的数据搬运：CPU 先把 8-word descriptor 放到 shared SRAM，descriptor 里包含 net_id、input/output/param/scratch 地址、input/output word 数和 flags。`tpu_desc_fetch_dma_stub.v` 收到 launch 后依次进入 `DESC / INPUT / PARAM / OUTPUT` phase，先读 descriptor，再读 input，再读 param，最后写 output。

这里要强调它已经不是单纯 done 延迟 stub，而是实际发 AXI transaction 读写 shared SRAM，并把 input/param 以 valid 流送进 compute block。

### P22 Q8.8 MLP Compute RTL 细节

讲 `tpu_mlp_compute_stub.v` 的核心：input word 和 param word 流式写入内部 memory。`flags[16]` 置位时进入 tile2x2 Q8.8 线性层模式，每个 output word 包两个 int16 Q8.8 输出。

精度修正点是最重要的：RTL 不再对每个乘积单独右移，而是先用 Q16.16 精度累加 `x_q8.8 * w_q8.8`，再加 `bias_q8.8 << 8`，最后统一 round shift 回 Q8.8，并做 int16 饱和和可选 ReLU。这一点已经和 Python golden 对齐。

### P23 RTL 架构图讲解顺序

这一页是讲稿索引。快速讲时按四步走：先讲系统边界，再讲控制面，再讲数据面，最后讲 Q8.8 MLP RTL 和验证证据。

被追问“是不是整个算法都进 TPU”时，要明确回答：不是。当前 TPU RTL 验证的是 MLP/Classifier，CNN/FiLM 前端先在 CPU 软件路径执行；如果性能不够，下一步再考虑把 CNN/FiLM 下沉到 TPU 或新增 CNN/VPU 单元。

### P24 新生成代码与文件说明

这一页说明本轮新增的都是文档和绘图生成代码，不是硬件功能 RTL。真实硬件依据仍然是 RTL 文件、VCS testbench、Python golden 和仿真 log。

## 常见追问回答

问：现在跑通的是整个呼吸识别算法的 RTL 验证吗？

答：不是全算法 RTL。当前跑通的是 TPU MLP/Classifier 的 RTL 验证链路，以及 CPU+TPU 数据闭环；完整算法端到端结果用 CPU 前端软件路径和 Python/fixed golden 对齐验证。

问：TPU 有没有加速剩下的 CNN/FiLM？

答：当前没有。CNN/FiLM 先放在 CPU 上跑，TPU 加速最后的 MLP/Classifier。后续如果性能不够，可以把 CNN/FiLM 下沉到 TPU 或新增专用 CNN/VPU 单元。

问：为什么要画 shared SRAM？

答：这是 CPU 和 TPU 之间的核心交换区。CPU 写 descriptor/input/param，TPU DMA 读这些数据并写回 output，CPU 或 testbench 再读结果验证。

问：这次新增代码会影响硬件验证结果吗？

答：不会。这次新增的是 PPT 绘图脚本和讲稿，不参与综合和仿真。硬件验证结果来自之前的 RTL/testbench/golden。
"""
    OUT_NOTES.write_text(text, encoding="utf-8")


def build_six_page_deck(prs: Presentation):
    arch.build_deck(prs)
    add_talk_track_slide(prs, "P23")
    add_generated_code_slide(prs, "P24")


def main():
    DOCS.mkdir(parents=True, exist_ok=True)

    standalone = Presentation()
    standalone.slide_width = arch.SLIDE_W
    standalone.slide_height = arch.SLIDE_H
    build_six_page_deck(standalone)
    standalone.save(OUT_STANDALONE)

    if BASE_DECK.exists():
        appended = Presentation(str(BASE_DECK))
        build_six_page_deck(appended)
        appended.save(OUT_APPENDED)

    write_notes()
    print(OUT_STANDALONE)
    if OUT_APPENDED.exists():
        print(OUT_APPENDED)
    print(OUT_NOTES)


if __name__ == "__main__":
    main()
