from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
BASE_DECK = DOCS / "陈韦东tinytpusoc_呼吸识别SoC更新版.pptx"
OUT_STANDALONE = DOCS / "19_CPU_TPU_呼吸识别SoC_RTL架构图_4p_可编辑.pptx"
OUT_APPENDED = DOCS / "陈韦东tinytpusoc_呼吸识别SoC更新版_v2_RTL架构图.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = "F7F8F6"
TEXT = "243241"
SUB = "68778A"
PANEL = "D7DFE7"
WHITE = "FEFEFD"
BLUE_FILL = "EAF3FA"
BLUE_LINE = "2C7DA8"
GREEN_FILL = "E7F2EB"
GREEN_LINE = "5E946F"
ORANGE_FILL = "F8EBDD"
ORANGE_LINE = "C9733D"
RED_FILL = "F8E4E4"
RED_LINE = "B65D5D"
INK_FILL = "ECF0F3"
INK_LINE = "627185"
FONT = "WenQuanYi Zen Hei"


def rgb(v: str) -> RGBColor:
    v = v.replace("#", "")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def I(v: float):
    return Inches(v)


def set_fill(shape, color: str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.05):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width_pt)


def write_text(
    shape,
    text: str,
    *,
    size: float = 12.0,
    bold: bool = False,
    color: str = TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    ml: float = 0.08,
    mr: float = 0.08,
    mt: float = 0.04,
    mb: float = 0.03,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def textbox(slide, x, y, w, h, text, **kwargs):
    sp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    write_text(sp, text, **kwargs)
    return sp


def round_box(slide, x, y, w, h, fill=WHITE, line=PANEL, width_pt=1.05):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    set_fill(sp, fill)
    set_line(sp, line, width_pt)
    return sp


def pill(slide, x, y, w, h, text, fill, color, size=9.4, bold=False):
    sp = round_box(slide, x, y, w, h, fill, fill, 0.6)
    write_text(
        sp,
        text,
        size=size,
        bold=bold,
        color=color,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        ml=0.03,
        mr=0.03,
        mt=0.01,
        mb=0.01,
    )
    return sp


def module(slide, x, y, w, h, title, subtitle="", fill=WHITE, line=PANEL, title_size=12.2):
    round_box(slide, x, y, w, h, fill, line, 1.05)
    if subtitle:
        textbox(slide, x + 0.10, y + 0.10, w - 0.20, 0.22, title, size=title_size, bold=True)
        textbox(slide, x + 0.10, y + 0.38, w - 0.20, h - 0.44, subtitle, size=8.6, color=SUB)
    else:
        textbox(
            slide,
            x + 0.10,
            y + 0.08,
            w - 0.20,
            h - 0.14,
            title,
            size=title_size,
            bold=True,
            align=PP_ALIGN.CENTER,
            valign=MSO_ANCHOR.MIDDLE,
        )


def connector(slide, x1, y1, x2, y2, color, width_pt=1.8):
    sp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    sp.line.color.rgb = rgb(color)
    sp.line.width = Pt(width_pt)
    return sp


def arrow_head(slide, x, y, direction, color, size=0.065):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, I(x - size), I(y - size), I(size * 2), I(size * 2))
    set_fill(sp, color)
    sp.line.fill.background()
    sp.rotation = {"up": 0, "right": 90, "down": 180, "left": 270}[direction]


def poly_arrow(slide, pts, color, width_pt=1.8):
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        connector(slide, x1, y1, x2, y2, color, width_pt)
    (x1, y1), (x2, y2) = pts[-2], pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = "right" if x2 > x1 else "left"
    else:
        direction = "down" if y2 > y1 else "up"
    arrow_head(slide, x2, y2, direction, color)


def title(slide, page, title_text, subtitle):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)
    textbox(slide, 0.52, 0.18, 0.52, 0.20, page, size=11.0, bold=True, color=BLUE_LINE)
    textbox(slide, 0.52, 0.45, 5.90, 0.38, title_text, size=21.5, bold=True)
    textbox(slide, 0.52, 0.84, 8.60, 0.22, subtitle, size=11.2, color=SUB)


def add_top_arch(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "P19", "CPU+TPU 呼吸识别 SoC RTL 架构总图", "panda_soc_stage2_base_top.v：CPU 控制面、TPU 控制/数据面、shared SRAM 骨架")
    pill(slide, 10.05, 0.34, 2.70, 0.34, "当前已验证主路径", GREEN_FILL, GREEN_LINE, 11.2, True)

    round_box(slide, 0.40, 1.26, 12.52, 5.88, WHITE, PANEL, 1.1)
    pill(slide, 0.58, 1.38, 2.34, 0.24, "top: panda_soc_stage2_base_top.v", BLUE_FILL, BLUE_LINE, 9.6, True)

    module(slide, 0.74, 1.88, 2.20, 1.26, "Panda RISC-V CPU", "panda_risc_v_min_proc_sys\nIMEM boot @0x0000_0800\nDBUS + DCACHE AXI masters", BLUE_FILL, BLUE_LINE, 12.0)
    module(slide, 3.48, 1.64, 2.08, 0.80, "AXI-Lite Splitter", "TPU window: 0x4000_4000\n其他地址走 legacy APB", INK_FILL, INK_LINE, 11.2)
    module(slide, 6.12, 1.52, 2.18, 1.05, "TPU_CTRL Regs", "CTRL / STATUS / MODE\nNET_ID / DESC_LO / HI\nPERF_CYCLE", GREEN_FILL, GREEN_LINE, 11.4)
    module(slide, 9.02, 1.52, 2.60, 1.16, "TPU Desc DMA", "读 descriptor / input / param\n写 output\n内含 Q8.8 MLP compute", ORANGE_FILL, ORANGE_LINE, 11.4)

    module(slide, 3.48, 3.22, 2.08, 0.80, "AXI-APB Bridge", "legacy 外设桥\n当前保留 UART APB #0", INK_FILL, INK_LINE, 11.2)
    module(slide, 6.08, 3.35, 1.78, 0.76, "APB UART", "0x4000_3000\nboot log / printf", BLUE_FILL, BLUE_LINE, 11.0)
    module(slide, 4.50, 4.70, 3.00, 1.18, "Shared SRAM Subsys", "CPU dcache AXI + TPU AXI\n2-master -> 1-slave axi_interconnect\nshared SRAM base 0x6000_0000", GREEN_FILL, GREEN_LINE, 11.6)
    module(slide, 8.20, 4.70, 2.05, 1.18, "axi_ram", "32-bit AXI RAM\naddr width 23\n预加载 raw/param/desc", WHITE, GREEN_LINE, 11.6)
    module(slide, 10.92, 4.18, 1.52, 1.92, "Testbench\n观测点", "done / busy / error\noutput words\nUART print", RED_FILL, RED_LINE, 11.0)

    poly_arrow(slide, [(2.94, 2.20), (3.48, 2.20)], BLUE_LINE, 2.0)
    pill(slide, 3.02, 1.94, 0.36, 0.20, "dbus", BLUE_FILL, BLUE_LINE, 8.0)
    poly_arrow(slide, [(5.56, 1.98), (6.12, 1.98)], GREEN_LINE, 2.0)
    pill(slide, 5.62, 1.72, 0.42, 0.20, "m1", GREEN_FILL, GREEN_LINE, 8.0)
    poly_arrow(slide, [(5.56, 2.25), (6.72, 2.25), (6.72, 2.57)], GREEN_LINE, 1.4)
    poly_arrow(slide, [(5.56, 2.20), (5.84, 2.20), (5.84, 3.54), (5.56, 3.54)], INK_LINE, 1.8)
    pill(slide, 5.70, 2.88, 0.42, 0.20, "m0", INK_FILL, INK_LINE, 8.0)
    poly_arrow(slide, [(5.56, 3.62), (6.08, 3.62)], BLUE_LINE, 1.8)

    poly_arrow(slide, [(7.20, 2.57), (7.20, 3.02), (9.55, 3.02), (9.55, 2.68)], ORANGE_LINE, 2.0)
    pill(slide, 7.74, 2.78, 1.16, 0.22, "launch / desc", ORANGE_FILL, ORANGE_LINE, 8.6)
    poly_arrow(slide, [(9.02, 2.12), (8.30, 2.12)], GREEN_LINE, 1.6)
    pill(slide, 8.42, 1.86, 0.48, 0.20, "status", GREEN_FILL, GREEN_LINE, 8.0)

    poly_arrow(slide, [(1.84, 3.14), (1.84, 5.28), (4.50, 5.28)], BLUE_LINE, 2.0)
    pill(slide, 2.50, 5.00, 0.86, 0.22, "dcache AXI", BLUE_FILL, BLUE_LINE, 8.4)
    poly_arrow(slide, [(10.08, 2.68), (10.08, 4.70)], ORANGE_LINE, 2.0)
    pill(slide, 9.72, 3.48, 0.76, 0.22, "TPU AXI", ORANGE_FILL, ORANGE_LINE, 8.4)
    poly_arrow(slide, [(7.50, 5.28), (8.20, 5.28)], GREEN_LINE, 2.0)
    poly_arrow(slide, [(10.25, 5.28), (10.92, 5.28)], RED_LINE, 1.8)
    poly_arrow(slide, [(7.00, 5.88), (7.00, 6.38), (11.68, 6.38), (11.68, 6.10)], RED_LINE, 1.3)
    pill(slide, 9.76, 6.14, 0.72, 0.20, "scoreboard", RED_FILL, RED_LINE, 7.8)

    textbox(slide, 0.58, 6.92, 11.80, 0.18, "边界说明：当前真实 RTL 闭合的是 CPU 写 TPU_CTRL、DMA 读写 shared SRAM、Q8.8 MLP 输出写回；完整 CNN/FiLM 仍在 CPU 软件路径。", size=9.4, color=SUB)


def add_control_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "P20", "TPU 控制面 RTL 细节", "cpu_tpu_axil_splitter.v + tpu_ctrl_axil_regs.v：地址分流、寄存器、状态回读")
    pill(slide, 10.24, 0.34, 2.42, 0.34, "CPU 可编程控制入口", BLUE_FILL, BLUE_LINE, 11.0, True)

    round_box(slide, 0.44, 1.25, 12.44, 5.94, WHITE, PANEL, 1.1)
    module(slide, 0.78, 1.76, 1.95, 1.12, "CPU DBUS AXI-Lite", "AW/W/B + AR/R\n来自 m_axi_dbus_*", BLUE_FILL, BLUE_LINE, 11.0)
    module(slide, 3.32, 1.52, 2.32, 1.58, "cpu_tpu_axil_splitter", "aw_sel_tpu / ar_sel_tpu\nwrite_busy / read_busy\n保持读写响应返回同一路", INK_FILL, INK_LINE, 11.2)
    module(slide, 6.40, 1.44, 2.20, 1.72, "m1: TPU_CTRL", "0x4000_4000 + 4KB\nCTRL launch / reset / irq\nSTATUS busy / done / error", GREEN_FILL, GREEN_LINE, 11.2)
    module(slide, 9.20, 1.44, 2.40, 1.72, "任务控制输出", "launch_pulse\nsoft_reset_pulse\nmode / net_id\ndesc_lo / desc_hi", ORANGE_FILL, ORANGE_LINE, 11.2)

    module(slide, 6.40, 3.72, 2.20, 1.40, "状态/性能回读", "status_busy / done / error\nperf_cycle_reg\nbusy 时自动累加", GREEN_FILL, GREEN_LINE, 11.2)
    module(slide, 9.20, 3.72, 2.40, 1.40, "DMA/Stub 状态源选择", "优先 desc_dma_stub\n否则 task_stub\n否则外部真实 TPU wrapper", ORANGE_FILL, ORANGE_LINE, 11.0)
    module(slide, 3.32, 4.18, 2.32, 1.16, "m0: legacy APB", "0x4000_3000 UART\n其他 legacy 外设保留", INK_FILL, INK_LINE, 11.2)
    module(slide, 0.78, 4.18, 1.95, 1.16, "仿真旁路", "SIM_TPU_CTRL_AXIL_BYPASS\nTB 可直接访问 regs", RED_FILL, RED_LINE, 10.6)

    poly_arrow(slide, [(2.73, 2.24), (3.32, 2.24)], BLUE_LINE, 2.0)
    poly_arrow(slide, [(5.64, 2.02), (6.40, 2.02)], GREEN_LINE, 2.0)
    poly_arrow(slide, [(8.60, 2.10), (9.20, 2.10)], ORANGE_LINE, 2.0)
    poly_arrow(slide, [(10.40, 3.72), (10.40, 3.16)], ORANGE_LINE, 1.8)
    poly_arrow(slide, [(9.20, 4.42), (8.60, 4.42)], GREEN_LINE, 2.0)
    poly_arrow(slide, [(7.44, 3.72), (7.44, 3.16)], GREEN_LINE, 1.8)
    poly_arrow(slide, [(5.64, 2.74), (5.90, 2.74), (5.90, 4.76), (5.64, 4.76)], INK_LINE, 1.8)
    poly_arrow(slide, [(2.73, 4.76), (3.32, 4.76)], RED_LINE, 1.6)
    poly_arrow(slide, [(1.76, 4.18), (1.76, 3.36), (6.98, 3.36), (6.98, 3.16)], RED_LINE, 1.4)

    pill(slide, 3.60, 3.35, 1.76, 0.28, "地址窗口判断", INK_FILL, INK_LINE, 10.0, True)
    pill(slide, 6.62, 5.45, 0.76, 0.24, "CTRL 0x00", GREEN_FILL, GREEN_LINE, 8.4)
    pill(slide, 7.46, 5.45, 0.84, 0.24, "STATUS 0x04", GREEN_FILL, GREEN_LINE, 8.4)
    pill(slide, 8.38, 5.45, 0.78, 0.24, "MODE 0x08", GREEN_FILL, GREEN_LINE, 8.4)
    pill(slide, 9.24, 5.45, 0.86, 0.24, "NET_ID 0x0C", GREEN_FILL, GREEN_LINE, 8.4)
    pill(slide, 10.18, 5.45, 0.92, 0.24, "DESC 0x10/14", GREEN_FILL, GREEN_LINE, 8.4)
    pill(slide, 11.18, 5.45, 0.86, 0.24, "PERF 0x18", GREEN_FILL, GREEN_LINE, 8.4)

    textbox(slide, 0.62, 6.90, 11.70, 0.18, "面试讲法：CPU 只看到 memory-mapped TPU_CTRL，RTL 里用 splitter 把控制面和 UART/APB 分开，launch 后由 DMA 子系统接管数据搬运。", size=9.4, color=SUB)


def add_dma_shared_mem_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "P21", "Descriptor DMA 与 Shared SRAM 数据面", "tpu_desc_fetch_dma_stub.v + panda_soc_shared_mem_subsys.v：descriptor 驱动的数据搬运闭环")
    pill(slide, 9.92, 0.34, 2.74, 0.34, "已闭合 CPU -> TPU -> SRAM", GREEN_FILL, GREEN_LINE, 11.0, True)

    round_box(slide, 0.44, 1.22, 12.44, 5.98, WHITE, PANEL, 1.1)
    module(slide, 0.78, 1.70, 2.12, 1.18, "Descriptor in SRAM", "8 words:\nnet/input/output/param/scratch\ninput_words/output_words/flags", BLUE_FILL, BLUE_LINE, 10.6)
    module(slide, 3.34, 1.42, 2.70, 1.86, "tpu_desc_fetch_dma_stub", "FSM: IDLE -> AR -> R -> AW -> W -> B -> DONE\nphase: DESC / INPUT / PARAM / OUTPUT", ORANGE_FILL, ORANGE_LINE, 11.0)
    module(slide, 6.68, 1.70, 2.14, 1.18, "MLP Compute Block", "input stream\nparam stream\noutput_word(index)", GREEN_FILL, GREEN_LINE, 11.0)
    module(slide, 9.58, 1.70, 2.08, 1.18, "Output Blob", "DMA 写回 shared SRAM\nTB/CPU 读取结果", BLUE_FILL, BLUE_LINE, 10.8)

    module(slide, 0.78, 4.26, 2.20, 1.28, "CPU dcache AXI", "CPU 软件写 descriptor\n写 input / param\n读 output", BLUE_FILL, BLUE_LINE, 11.0)
    module(slide, 3.58, 4.08, 2.24, 1.58, "2-to-1 AXI Interconnect", "S_COUNT=2, M_COUNT=1\nM_BASE_ADDR=0x6000_0000\nread/write connect = 2'b11", INK_FILL, INK_LINE, 11.0)
    module(slide, 6.56, 4.16, 2.08, 1.42, "AXI RAM", "DATA_WIDTH=32\nADDR_WIDTH=23\n预加载样本/参数池", GREEN_FILL, GREEN_LINE, 11.0)
    module(slide, 9.28, 4.16, 2.40, 1.42, "TPU AXI Master Path", "当前默认 desc_dma_stub\n未来可换真实 TPU DMA\n顶层 mux 保留外部 tpu_axi_*", ORANGE_FILL, ORANGE_LINE, 10.8)

    poly_arrow(slide, [(2.90, 2.26), (3.34, 2.26)], BLUE_LINE, 2.0)
    poly_arrow(slide, [(6.04, 2.10), (6.68, 2.10)], GREEN_LINE, 2.0)
    pill(slide, 6.14, 1.84, 0.42, 0.20, "valid", GREEN_FILL, GREEN_LINE, 8.0)
    poly_arrow(slide, [(8.82, 2.26), (9.58, 2.26)], BLUE_LINE, 2.0)
    poly_arrow(slide, [(4.70, 3.28), (4.70, 4.08)], ORANGE_LINE, 2.0)
    poly_arrow(slide, [(2.98, 4.90), (3.58, 4.90)], BLUE_LINE, 2.0)
    poly_arrow(slide, [(5.82, 4.90), (6.56, 4.90)], GREEN_LINE, 2.0)
    poly_arrow(slide, [(10.48, 4.16), (10.48, 3.30), (5.64, 3.30)], ORANGE_LINE, 1.8)
    poly_arrow(slide, [(8.64, 4.90), (9.28, 4.90)], ORANGE_LINE, 1.8)
    poly_arrow(slide, [(7.60, 4.16), (7.60, 3.30), (9.95, 3.30), (9.95, 2.88)], BLUE_LINE, 1.6)

    pill(slide, 3.46, 3.50, 0.92, 0.24, "read desc", ORANGE_FILL, ORANGE_LINE, 8.6)
    pill(slide, 4.48, 3.50, 0.92, 0.24, "read input", ORANGE_FILL, ORANGE_LINE, 8.6)
    pill(slide, 5.50, 3.50, 0.92, 0.24, "read param", ORANGE_FILL, ORANGE_LINE, 8.6)
    pill(slide, 6.52, 3.50, 0.92, 0.24, "write out", ORANGE_FILL, ORANGE_LINE, 8.6)

    textbox(slide, 0.62, 6.90, 11.70, 0.18, "关键点：DMA stub 已不是纯 done 延迟占位，它真正按 descriptor 从 shared SRAM 拉 input/param，并把 MLP 输出写回 output 区。", size=9.4, color=SUB)


def add_mlp_detail(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title(slide, "P22", "Q8.8 MLP Compute RTL 细节", "tpu_mlp_compute_stub.v：流式装载、2-output tile、Q16.16 累加后 Q8.8 舍入")
    pill(slide, 9.92, 0.34, 2.74, 0.34, "TPU RTL 验证核心", ORANGE_FILL, ORANGE_LINE, 11.0, True)

    round_box(slide, 0.44, 1.22, 12.44, 5.98, WHITE, PANEL, 1.1)
    module(slide, 0.78, 1.74, 2.02, 1.16, "Input Stream", "input_word_valid\n32-bit word packs x1:x0\n写入 input_mem[]", BLUE_FILL, BLUE_LINE, 11.0)
    module(slide, 0.78, 4.18, 2.02, 1.16, "Param Stream", "param_word_valid\nw00/w01/w10/w11\nbias word", BLUE_FILL, BLUE_LINE, 11.0)
    module(slide, 3.38, 1.62, 2.18, 1.34, "Input/Param SRAM", "input_mem[256]\nparam_mem[8192]\n同时保留 checksum/last_word", GREEN_FILL, GREEN_LINE, 11.0)
    module(slide, 6.12, 1.58, 2.32, 1.42, "Tile2x2 Addressing", "linear_stride_words = input_words*2 + 1\nbase = output_idx * stride", INK_FILL, INK_LINE, 10.6)
    module(slide, 6.12, 3.64, 2.32, 1.54, "Q8.8 MAC", "acc_q16.16 += x_q8.8 * w_q8.8\nbias_q8.8 << 8\n最后 round >>> 8", ORANGE_FILL, ORANGE_LINE, 10.6)
    module(slide, 9.16, 2.46, 2.38, 1.62, "Output Word", "saturate to int16\n可选 ReLU flags[0]\n{y1_q8.8, y0_q8.8}", GREEN_FILL, GREEN_LINE, 11.0)
    module(slide, 9.16, 4.72, 2.38, 0.88, "Fallback Path", "未置 flags[16] 时输出 checksum 兼容旧 smoke test", INK_FILL, INK_LINE, 10.2)

    poly_arrow(slide, [(2.80, 2.32), (3.38, 2.32)], BLUE_LINE, 2.0)
    poly_arrow(slide, [(2.80, 4.76), (3.38, 4.76), (3.38, 2.66)], BLUE_LINE, 2.0)
    poly_arrow(slide, [(5.56, 2.34), (6.12, 2.34)], GREEN_LINE, 2.0)
    poly_arrow(slide, [(7.28, 3.00), (7.28, 3.64)], ORANGE_LINE, 2.0)
    poly_arrow(slide, [(8.44, 4.40), (9.16, 3.26)], GREEN_LINE, 2.0)
    poly_arrow(slide, [(10.34, 4.08), (10.34, 4.72)], INK_LINE, 1.6)

    pill(slide, 3.70, 3.26, 1.48, 0.28, "flags[16] = tile2x2_q8_8", ORANGE_FILL, ORANGE_LINE, 8.8, True)
    pill(slide, 5.96, 5.70, 1.18, 0.26, "Q16.16 acc", ORANGE_FILL, ORANGE_LINE, 9.0, True)
    pill(slide, 7.24, 5.70, 1.18, 0.26, "round + sat", ORANGE_FILL, ORANGE_LINE, 9.0, True)
    pill(slide, 8.52, 5.70, 0.86, 0.26, "Q8.8 out", GREEN_FILL, GREEN_LINE, 9.0, True)

    textbox(slide, 0.72, 6.20, 11.70, 0.36, "精度修正点：不再对每个乘积单独右移；RTL 先完整累加 x*w 和 bias<<8，再一次性 round shift 到 Q8.8，与 Python golden 对齐。", size=10.0, color=TEXT)
    textbox(slide, 0.62, 6.92, 11.70, 0.18, "面试讲法：这一页是 TPU 已验证的 MLP/Classifier RTL 核心，CNN/FiLM 前端当前由 CPU 软件承担。", size=9.4, color=SUB)


def build_deck(prs):
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    add_top_arch(prs)
    add_control_detail(prs)
    add_dma_shared_mem_detail(prs)
    add_mlp_detail(prs)


def main():
    DOCS.mkdir(parents=True, exist_ok=True)

    standalone = Presentation()
    build_deck(standalone)
    standalone.save(OUT_STANDALONE)

    if BASE_DECK.exists():
        appended = Presentation(str(BASE_DECK))
        build_deck(appended)
        appended.save(OUT_APPENDED)

    print(OUT_STANDALONE)
    if OUT_APPENDED.exists():
        print(OUT_APPENDED)


if __name__ == "__main__":
    main()
