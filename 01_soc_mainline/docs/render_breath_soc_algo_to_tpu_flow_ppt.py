from __future__ import annotations

from pathlib import Path

from pptx import Presentation

import render_breath_soc_rtl_arch_ppt as arch
import render_breath_soc_rtl_arch_explained_ppt as explained


ROOT = Path(__file__).resolve().parents[1]
DOCS = ROOT / "docs"
BASE_V3 = DOCS / "陈韦东tinytpusoc_呼吸识别SoC更新版_v3_RTL架构图_含讲解.pptx"
OUT_3P = DOCS / "21_CPU_TPU_呼吸识别_算法拆分_CPU发送TPU_3p.pptx"
OUT_9P = DOCS / "22_CPU_TPU_呼吸识别_算法到RTL_9p_讲解顺序版.pptx"
OUT_V4 = DOCS / "陈韦东tinytpusoc_呼吸识别SoC更新版_v4_算法到CPU发送TPU.pptx"
OUT_NOTES = DOCS / "CPU_TPU_呼吸识别_算法拆分_CPU发送TPU_讲解稿_20260419.md"


def add_algo_split_slide(prs: Presentation, page: str = "A1"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    arch.title(slide, page, "先从算法拆分讲起", "呼吸识别模型不是一次性全部丢进 TPU，而是拆成 CPU 前端 + TPU MLP/Classifier 任务序列")
    arch.pill(slide, 9.92, 0.34, 2.76, 0.34, "讲解主线第一步", arch.BLUE_FILL, arch.BLUE_LINE, 10.8, True)

    arch.round_box(slide, 0.44, 1.22, 12.44, 5.98, arch.WHITE, arch.PANEL, 1.1)
    arch.module(slide, 0.78, 1.78, 1.70, 0.92, "Raw Window", "原始呼吸窗口\n样本/传感器输入", arch.BLUE_FILL, arch.BLUE_LINE, 10.6)
    arch.module(slide, 3.00, 1.56, 2.24, 1.36, "CPU Front-end", "定点预处理\nCNN/FiLM/特征前端\n组织 fusion 输入", arch.GREEN_FILL, arch.GREEN_LINE, 10.8)
    arch.module(slide, 5.86, 1.56, 2.34, 1.36, "TPU MLP Branches", "MLP_KEY 多层 linear\nMLP_OTHER 多层 linear\nQ8.8 tile2x2", arch.ORANGE_FILL, arch.ORANGE_LINE, 10.6)
    arch.module(slide, 8.84, 1.56, 2.24, 1.36, "TPU Classifier", "融合输入 161 words\nclassifier linear layers\n输出 logits", arch.ORANGE_FILL, arch.ORANGE_LINE, 10.6)
    arch.module(slide, 11.62, 1.78, 0.84, 0.92, "Class", "预测类别\nlogits", arch.GREEN_FILL, arch.GREEN_LINE, 10.4)

    arch.poly_arrow(slide, [(2.48, 2.24), (3.00, 2.24)], arch.BLUE_LINE, 1.9)
    arch.poly_arrow(slide, [(5.24, 2.24), (5.86, 2.24)], arch.GREEN_LINE, 1.9)
    arch.poly_arrow(slide, [(8.20, 2.24), (8.84, 2.24)], arch.ORANGE_LINE, 1.9)
    arch.poly_arrow(slide, [(11.08, 2.24), (11.62, 2.24)], arch.ORANGE_LINE, 1.9)

    arch.module(slide, 0.96, 3.64, 2.72, 1.38, "为什么这样拆", "TPU 当前已经验证的是 Q8.8 MLP/Classifier RTL。CNN/FiLM 前端先放 CPU，避免把尚未验证的复杂卷积控制和数据复用一次性塞进硬件。", arch.INK_FILL, arch.INK_LINE, 9.8)
    arch.module(slide, 4.10, 3.64, 2.72, 1.38, "CPU 做什么", "准备输入窗口和前端特征；把每个 TPU stage 的 input/param/output 地址写成 descriptor；通过 MMIO 发 launch。", arch.INK_FILL, arch.INK_LINE, 9.8)
    arch.module(slide, 7.24, 3.64, 2.72, 1.38, "TPU 做什么", "按 descriptor 读 shared SRAM，把 input/param 流送入 MLP compute block，完成 linear/ReLU/round/saturate，再写回 output。", arch.INK_FILL, arch.INK_LINE, 9.8)
    arch.module(slide, 10.38, 3.64, 1.78, 1.38, "任务粒度", "不是一次 launch 全模型。当前 demo 是多个 linear stage 依次 launch。", arch.RED_FILL, arch.RED_LINE, 9.4)

    arch.pill(slide, 0.96, 5.66, 1.36, 0.28, "一句话", arch.RED_FILL, arch.RED_LINE, 9.0, True)
    arch.textbox(slide, 2.52, 5.61, 8.80, 0.42, "算法层面先把 CNN/FiLM 前端留在 CPU，把最后 MLP/Classifier 拆成一串 Q8.8 linear 任务交给 TPU；CPU 负责调度和发 descriptor。", size=10.4, color=arch.TEXT)
    arch.textbox(slide, 0.62, 6.90, 11.70, 0.18, "这样讲可以先回答“算法怎么拆”，再自然过渡到“CPU 怎么发给 TPU”。", size=9.4, color=arch.SUB)


def add_cpu_send_slide(prs: Presentation, page: str = "A2"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    arch.title(slide, page, "CPU 怎么把一个 TPU 任务发出去", "软件不是直接驱动 TPU 内部信号，而是写 shared SRAM 里的 descriptor/data，再写 TPU_CTRL MMIO 寄存器")
    arch.pill(slide, 9.86, 0.34, 2.84, 0.34, "descriptor + MMIO", arch.GREEN_FILL, arch.GREEN_LINE, 10.8, True)

    arch.round_box(slide, 0.44, 1.22, 12.44, 5.98, arch.WHITE, arch.PANEL, 1.1)
    arch.module(slide, 0.74, 1.58, 2.08, 1.20, "1. 准备数据", "CPU 把 input vector 写到 shared SRAM\n参数池 param pool 已预加载或由 CPU 写入\n输出区先清零", arch.BLUE_FILL, arch.BLUE_LINE, 9.8)
    arch.module(slide, 3.12, 1.58, 2.36, 1.20, "2. 写 descriptor", "8 words:\nnet_id / input_addr / output_addr / param_addr\nscratch_addr / input_words / output_words / flags", arch.GREEN_FILL, arch.GREEN_LINE, 9.2)
    arch.module(slide, 5.82, 1.58, 2.02, 1.20, "3. 同步内存", "fence rw,rw\ncache evict\nsettle loops\n保证 TPU DMA 能看到最新数据", arch.INK_FILL, arch.INK_LINE, 9.4)
    arch.module(slide, 8.18, 1.58, 2.02, 1.20, "4. MMIO launch", "写 MODE / NET_ID\n写 DESC_LO = desc addr\n写 CTRL.START", arch.ORANGE_FILL, arch.ORANGE_LINE, 9.8)
    arch.module(slide, 10.52, 1.58, 1.74, 1.20, "5. 等结果", "轮询 STATUS\ndone/error\n读 output words", arch.RED_FILL, arch.RED_LINE, 9.8)

    for x1, x2, color in [(2.82, 3.12, arch.BLUE_LINE), (5.48, 5.82, arch.GREEN_LINE), (7.84, 8.18, arch.INK_LINE), (10.20, 10.52, arch.ORANGE_LINE)]:
        arch.poly_arrow(slide, [(x1, 2.18), (x2, 2.18)], color, 1.8)

    arch.module(slide, 0.86, 3.40, 3.00, 1.42, "descriptor 真实结构", "TPUDesc {\n  net_id, input_addr, output_addr, param_addr,\n  scratch_addr, input_words, output_words, flags\n}\nflags[16] = TILE2X2_Q8_8, flags[0] = ReLU", arch.INK_FILL, arch.INK_LINE, 8.8)
    arch.module(slide, 4.22, 3.40, 2.82, 1.42, "MMIO 真实寄存器", "TPU_CTRL_BASE = 0x4000_4000\nDESC_LO @ 0x10\nDESC_HI @ 0x14\nCTRL.START bit0\nSTATUS busy/done/error", arch.INK_FILL, arch.INK_LINE, 9.0)
    arch.module(slide, 7.40, 3.40, 3.86, 1.42, "对应 C runtime 调用", "tpu_select_desc_buffer()\ntpu_build_desc()\ntpu_submit_desc(): 写 desc_slot[0..7] + 写 TPU_REG_DESC_LO + 写 TPU_REG_CTRL START\ntpu_wait_done(): 轮询 STATUS", arch.INK_FILL, arch.INK_LINE, 8.8)

    arch.pill(slide, 0.96, 5.62, 1.74, 0.28, "重点口径", arch.RED_FILL, arch.RED_LINE, 9.0, True)
    arch.textbox(slide, 2.90, 5.55, 8.50, 0.50, "CPU 发给 TPU 的核心不是函数调用，而是“共享内存地址 + descriptor + MMIO START”。descriptor 告诉 TPU 去哪里读 input/param、去哪里写 output、这一层有多少 word。", size=10.2, color=arch.TEXT)
    arch.textbox(slide, 0.62, 6.90, 11.70, 0.18, "这页可以直接回答“CPU 怎么发送到 TPU”：先写 SRAM，再写寄存器启动。", size=9.4, color=arch.SUB)


def add_tpu_receive_slide(prs: Presentation, page: str = "A3"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    arch.title(slide, page, "TPU 收到 launch 后怎么执行", "TPU_CTRL 只负责启动和状态；真正搬数据的是 descriptor DMA，真正算的是 Q8.8 MLP compute block")
    arch.pill(slide, 9.84, 0.34, 2.86, 0.34, "从 MMIO 到 AXI DMA", arch.ORANGE_FILL, arch.ORANGE_LINE, 10.8, True)

    arch.round_box(slide, 0.44, 1.22, 12.44, 5.98, arch.WHITE, arch.PANEL, 1.1)
    arch.module(slide, 0.76, 1.70, 1.80, 0.96, "CPU MMIO", "CTRL.START\nDESC_LO", arch.BLUE_FILL, arch.BLUE_LINE, 10.4)
    arch.module(slide, 3.08, 1.58, 2.10, 1.20, "TPU_CTRL Regs", "锁存 descriptor 地址\n产生 launch_pulse\n回读 busy/done/error", arch.GREEN_FILL, arch.GREEN_LINE, 10.2)
    arch.module(slide, 5.72, 1.58, 2.30, 1.20, "Desc DMA FSM", "读 8-word descriptor\n读 input blob\n读 param blob\n写 output blob", arch.ORANGE_FILL, arch.ORANGE_LINE, 10.0)
    arch.module(slide, 8.56, 1.58, 2.18, 1.20, "Q8.8 MLP Compute", "input_word_valid\nparam_word_valid\noutput_word(index)", arch.GREEN_FILL, arch.GREEN_LINE, 10.0)
    arch.module(slide, 11.22, 1.70, 1.08, 0.96, "Done", "STATUS\noutput", arch.RED_FILL, arch.RED_LINE, 10.4)

    arch.poly_arrow(slide, [(2.56, 2.18), (3.08, 2.18)], arch.BLUE_LINE, 1.9)
    arch.poly_arrow(slide, [(5.18, 2.18), (5.72, 2.18)], arch.GREEN_LINE, 1.9)
    arch.poly_arrow(slide, [(8.02, 2.18), (8.56, 2.18)], arch.ORANGE_LINE, 1.9)
    arch.poly_arrow(slide, [(10.74, 2.18), (11.22, 2.18)], arch.GREEN_LINE, 1.9)

    arch.module(slide, 1.02, 3.54, 2.28, 1.30, "控制面路径", "CPU DBUS -> splitter -> TPU_CTRL\n地址窗口 0x4000_4000\n只传控制寄存器和状态", arch.INK_FILL, arch.INK_LINE, 9.4)
    arch.module(slide, 3.82, 3.54, 2.28, 1.30, "数据面路径", "TPU DMA 作为 AXI master\n通过 shared_mem_subsys\n访问 0x6000_0000 shared SRAM", arch.INK_FILL, arch.INK_LINE, 9.4)
    arch.module(slide, 6.62, 3.54, 2.28, 1.30, "计算路径", "DMA 把 input/param 变成 valid 流\ncompute block 内部累加 Q16.16\n输出 packed Q8.8 word", arch.INK_FILL, arch.INK_LINE, 9.4)
    arch.module(slide, 9.42, 3.54, 2.08, 1.30, "返回路径", "DMA 写 output_addr\nSTATUS.done 置位\nCPU poll 后读 output", arch.INK_FILL, arch.INK_LINE, 9.4)

    arch.pill(slide, 1.04, 5.60, 1.58, 0.28, "讲解结论", arch.RED_FILL, arch.RED_LINE, 9.0, True)
    arch.textbox(slide, 2.82, 5.54, 8.60, 0.50, "所以这套接口本质上是一个轻量任务队列：CPU 写 descriptor，TPU DMA 按 descriptor 自己取数和写回，CPU 只负责启动、等待和读结果。", size=10.2, color=arch.TEXT)
    arch.textbox(slide, 0.62, 6.90, 11.70, 0.18, "这页之后再接 RTL 总图，就能把模块名和前面的算法/软件流程对应起来。", size=9.4, color=arch.SUB)


def build_algo_flow(prs: Presentation, page_prefix: str = "A"):
    prs.slide_width = arch.SLIDE_W
    prs.slide_height = arch.SLIDE_H
    add_algo_split_slide(prs, f"{page_prefix}1")
    add_cpu_send_slide(prs, f"{page_prefix}2")
    add_tpu_receive_slide(prs, f"{page_prefix}3")


def write_notes():
    text = """# 从算法拆分到 CPU 发送 TPU 的讲解稿

生成时间：2026-04-19

## 推荐讲法

先不要从 RTL 模块名开始讲。更顺的主线是：算法怎么拆，CPU 做什么，TPU 做什么，CPU 怎么把任务发给 TPU，TPU 收到后怎么自己搬数和计算。

可以这样说：

> 这个呼吸识别算法我没有一开始就把全部计算都塞进 TPU。当前拆法是 CPU 负责前端，TPU 负责后面的 MLP/Classifier。CPU 先完成定点预处理、CNN/FiLM 或特征前端，整理出要送给 MLP/Classifier 的输入向量；然后 CPU 把输入、参数地址和输出地址写成一个 descriptor，放到 shared SRAM；最后通过 memory-mapped TPU_CTRL 寄存器写 `DESC_LO` 和 `CTRL.START`，把这个任务 launch 给 TPU。TPU 收到 launch 后，descriptor DMA 会自己从 shared SRAM 读 descriptor、读 input、读 param，把数据送进 Q8.8 MLP compute block，算完再把 output 写回 shared SRAM。CPU 只需要轮询 STATUS.done，然后读 output。

## 算法怎么拆

当前口径：

- CPU：原始窗口预处理、定点化、CNN/FiLM/特征前端、TPU 任务调度。
- TPU：MLP_KEY 分支、MLP_OTHER 分支、Classifier head，对应一串 Q8.8 linear tile 任务。
- Shared SRAM：CPU 和 TPU 之间的数据交换区。
- TPU_CTRL：CPU 发起任务、TPU 返回状态的 MMIO 控制口。

要明确：当前不是完整 CNN/FiLM 都进 TPU。当前已经做 RTL 验证的是 MLP/Classifier 任务链和 CPU+TPU 数据闭环。

## CPU 怎么发送给 TPU

CPU 发送一个 TPU stage 的动作可以拆成 5 步：

1. CPU 选择双缓冲区，例如 `TPU_DESC0_BASE / TPU_IN_BUF0_BASE / TPU_OUT_BUF0_BASE / TPU_SCRATCH0_BASE`。
2. CPU 把 input vector 写入 shared SRAM，把参数池 param pool 放在固定 shared SRAM 地址。
3. CPU 写 8-word descriptor：

```c
TPUDesc {
    net_id,
    input_addr,
    output_addr,
    param_addr,
    scratch_addr,
    input_words,
    output_words,
    flags
};
```

4. CPU 做 `fence`、dcache eviction 和 settle，保证 TPU DMA 看到最新 shared SRAM 内容。
5. CPU 通过 MMIO 写 TPU 控制寄存器：

```c
tpu_reg_write(runtime, TPU_REG_CTRL, TPU_CTRL_SOFT_RESET_MASK);
tpu_reg_write(runtime, TPU_REG_MODE, TPU_MODE_INFER);
tpu_reg_write(runtime, TPU_REG_NET_ID, desc->net_id);
tpu_reg_write(runtime, TPU_REG_DESC_LO, runtime->active_desc_addr);
tpu_reg_write(runtime, TPU_REG_DESC_HI, 0u);
tpu_reg_write(runtime, TPU_REG_CTRL, TPU_CTRL_START_MASK);
```

寄存器基地址是 `TPU_CTRL_BASEADDR = 0x4000_4000`。`DESC_LO` 偏移是 `0x10`，`CTRL.START` 是 bit0，`STATUS.done` 是 bit1。

## TPU 收到任务后怎么跑

TPU 内部不是 CPU 一拍一拍喂数据，而是 DMA 自己跑：

1. `tpu_ctrl_axil_regs.v` 收到 `CTRL.START`，产生 `launch_pulse`。
2. `tpu_desc_fetch_dma_stub.v` 进入 FSM，先从 `DESC_LO` 指向的 shared SRAM 地址读 8-word descriptor。
3. DMA 根据 descriptor 里的 `input_addr/input_words` 读 input blob。
4. DMA 根据 `param_addr` 和 `flags/input_words/output_words` 读 param blob。
5. DMA 把 input/param 以 `valid` 流送给 `tpu_mlp_compute_stub.v`。
6. compute block 生成 packed Q8.8 output word。
7. DMA 把 output word 写回 descriptor 里的 `output_addr`。
8. `STATUS.done` 置位，CPU 轮询到 done 后读 output。

## 和代码对应关系

- `work/600_competition_5stage/software/include/tpu_desc.h`：descriptor 字段、shared SRAM buffer 地址、net_id、flags。
- `work/600_competition_5stage/software/include/tpu_regs.h`：TPU_CTRL MMIO 地址和寄存器 bit。
- `work/600_competition_5stage/software/lib/tpu_runtime.c`：CPU 侧 `tpu_build_desc()`、`tpu_submit_desc()`、`tpu_wait_done()`。
- `work/600_competition_5stage/software/test/breath_tpu_soc_demo/main.c`：按 stage 依次 launch MLP_KEY、MLP_OTHER、classifier schedule。
- `work/600_competition_5stage/fpga/panda_soc_eva/rtl/tpu_ctrl_axil_regs.v`：MMIO 控制寄存器硬件。
- `work/600_competition_5stage/fpga/panda_soc_eva/rtl/tpu_desc_fetch_dma_stub.v`：descriptor DMA FSM。
- `work/600_competition_5stage/fpga/panda_soc_eva/rtl/tpu_mlp_compute_stub.v`：Q8.8 MLP compute block。

## 被追问时的边界回答

问：这是不是整个算法都映射到 TPU？

答：不是。当前拆法是 CPU 跑前端，TPU 跑 MLP/Classifier。TPU RTL 已经验证的是 descriptor DMA + Q8.8 MLP/Classifier 链路。

问：CPU 发给 TPU 的到底是什么？

答：CPU 发的不是算法源码，也不是直接拉 RTL 内部 wire，而是 shared SRAM 中的 descriptor 和数据地址，再通过 MMIO 写 `CTRL.START`。TPU 根据 descriptor 自己发 AXI 读写 shared SRAM。

问：为什么用 descriptor？

答：descriptor 把每个 TPU stage 的 input/output/param 地址、word 数和 flags 固定下来。这样 CPU 只需要按 stage 改 descriptor，TPU 侧 DMA/compute 逻辑可以复用。
"""
    OUT_NOTES.write_text(text, encoding="utf-8")


def main():
    DOCS.mkdir(parents=True, exist_ok=True)

    deck3 = Presentation()
    build_algo_flow(deck3, "A")
    deck3.save(OUT_3P)

    deck9 = Presentation()
    build_algo_flow(deck9, "A")
    explained.build_six_page_deck(deck9)
    deck9.save(OUT_9P)

    if BASE_V3.exists():
        v4 = Presentation(str(BASE_V3))
    else:
        v4 = Presentation()
    add_algo_split_slide(v4, "P25")
    add_cpu_send_slide(v4, "P26")
    add_tpu_receive_slide(v4, "P27")
    v4.save(OUT_V4)

    write_notes()
    print(OUT_3P)
    print(OUT_9P)
    print(OUT_V4)
    print(OUT_NOTES)


if __name__ == "__main__":
    main()
