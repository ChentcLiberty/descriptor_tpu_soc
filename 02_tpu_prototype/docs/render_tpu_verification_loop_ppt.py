from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/home/jjt/tpu-soc")
PACK = ROOT / "career" / "hisilicon_25min_interview_20260410"
OUT_PPT = PACK / "09_TinyTPU_验证闭环_可编辑.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = "F8F6F1"
TEXT = "243241"
SUB = "718094"
PANEL = "D7DFE7"
WHITE = "FEFEFD"
BLUE_FILL = "EAF3FA"
BLUE_LINE = "2C7DA8"
ORANGE_FILL = "F7EBDD"
ORANGE_LINE = "C9733D"
GREEN_FILL = "E6F2E9"
GREEN_LINE = "5E946F"
PURPLE_FILL = "ECE7F8"
PURPLE_LINE = "8770D0"
PILL_BLUE = "E7F0F6"
FONT = "WenQuanYi Zen Hei"


def rgb(v: str) -> RGBColor:
    v = v.replace("#", "")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def I(v: float):
    return Inches(v)


def set_fill(shape, color: str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.15):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width_pt)


def write_text(shape, text: str, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
               valign=MSO_ANCHOR.TOP, ml=0.08, mr=0.08, mt=0.05, mb=0.03):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def round_box(slide, x, y, w, h, fill, line, width_pt=1.15):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    set_fill(sp, fill)
    set_line(sp, line, width_pt)
    return sp


def textbox(slide, x, y, w, h, text, **kwargs):
    sp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    write_text(sp, text, **kwargs)
    return sp


def pill(slide, x, y, w, h, text, fill, color, size=10.0, bold=False):
    sp = round_box(slide, x, y, w, h, fill, fill, 0.6)
    write_text(sp, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER,
               valign=MSO_ANCHOR.MIDDLE, ml=0.03, mr=0.03, mt=0.01, mb=0.01)
    return sp


def card(slide, x, y, w, h, title, subtitle, fill, accent, *, body_size=10.2):
    round_box(slide, x, y, w, h, fill, PANEL, 1.0)
    pill(slide, x + 0.10, y + 0.10, min(w - 0.20, 1.10), 0.22, title, fill, accent, 8.8, True)
    textbox(slide, x + 0.14, y + 0.40, w - 0.28, h - 0.48, subtitle, size=body_size, valign=MSO_ANCHOR.TOP)


def line(slide, x1, y1, x2, y2, color, width_pt=2.0):
    sp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    sp.line.color.rgb = rgb(color)
    sp.line.width = Pt(width_pt)
    return sp


def arrow_head(slide, x, y, direction, color, size=0.07):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, I(x - size), I(y - size), I(size * 2), I(size * 2))
    set_fill(sp, color)
    sp.line.fill.background()
    sp.rotation = {"up": 0, "right": 90, "down": 180, "left": 270}[direction]


def poly_arrow(slide, pts, color, width_pt=2.0):
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line(slide, x1, y1, x2, y2, color, width_pt)
    (x1, y1), (x2, y2) = pts[-2], pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = "right" if x2 > x1 else "left"
    else:
        direction = "down" if y2 > y1 else "up"
    arrow_head(slide, x2, y2, direction, color)


def tag(slide, x, y, w, text, fill, color):
    pill(slide, x, y, w, 0.21, text, fill, color, 8.6)


def main():
    PACK.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    textbox(slide, 0.56, 0.18, 0.42, 0.18, "P3", size=12.0, color=BLUE_LINE)
    textbox(slide, 0.56, 0.45, 4.10, 0.38, "AXI-Lite 多 epoch 验证闭环", size=22.0, bold=True)
    textbox(slide, 0.56, 0.84, 7.80, 0.22, "host preload / imem load / start loop / ref model / assertion", size=12.0, color=SUB)
    pill(slide, 10.16, 0.34, 2.50, 0.34, "可编辑单页验证图", PILL_BLUE, BLUE_LINE, 11.2)

    round_box(slide, 0.40, 1.28, 12.52, 5.74, WHITE, PANEL, 1.1)

    round_box(slide, 0.64, 1.60, 6.00, 2.32, WHITE, PANEL, 1.0)
    textbox(slide, 0.90, 1.80, 2.30, 0.24, "Stimulus / Load", size=16.4, bold=True)
    textbox(slide, 0.90, 2.06, 3.40, 0.18, "参数、UB、IMEM 一次性装载，然后重复 start", size=9.4, color=SUB)

    round_box(slide, 6.90, 1.60, 5.74, 2.32, WHITE, PANEL, 1.0)
    textbox(slide, 7.18, 1.80, 2.30, 0.24, "Epoch Run / Observe", size=16.4, bold=True)
    textbox(slide, 7.18, 2.06, 3.10, 0.18, "轮询状态、读回参数、在线观察关键链路", size=9.4, color=SUB)

    round_box(slide, 0.64, 4.18, 12.00, 1.42, WHITE, PANEL, 1.0)
    textbox(slide, 0.90, 4.40, 2.20, 0.24, "Check / Result", size=16.0, bold=True)
    textbox(slide, 0.90, 4.66, 3.60, 0.18, "参考模型对齐、loss 收敛、pred 正确、assert 通过", size=9.4, color=SUB)

    card(slide, 0.92, 2.44, 1.30, 1.02, "cfg", "0x050 LEAK\n0x054 INV_N2\n0x058 LR", BLUE_FILL, BLUE_LINE, body_size=9.6)
    card(slide, 2.44, 2.44, 1.70, 1.02, "UB preload", "load_all_data\n0x020 / 0x028 /\n0x024", GREEN_FILL, GREEN_LINE, body_size=9.4)
    card(slide, 4.40, 2.44, 1.68, 1.02, "IMEM load", "imem_load()\n0x030 / 0x034 / 0x040 / 0x044", ORANGE_FILL, ORANGE_LINE, body_size=9.2)

    poly_arrow(slide, [(2.22, 2.95), (2.44, 2.95)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(4.14, 2.95), (4.40, 2.95)], GREEN_LINE, 1.8)

    tag(slide, 0.98, 3.58, 1.14, "X@0  8w", BLUE_FILL, BLUE_LINE)
    tag(slide, 2.18, 3.58, 1.10, "Y@8  4w", GREEN_FILL, GREEN_LINE)
    tag(slide, 3.34, 3.58, 1.16, "W1@12  4w", ORANGE_FILL, ORANGE_LINE)
    tag(slide, 4.58, 3.58, 1.10, "B1@16  2w", PURPLE_FILL, PURPLE_LINE)
    tag(slide, 5.74, 3.58, 1.00, "W2@18  2w", BLUE_FILL, BLUE_LINE)

    card(slide, 7.22, 2.42, 1.38, 1.06, "epoch loop", "run_one_epoch\nAXI write\n0x000=0x2", BLUE_FILL, BLUE_LINE, body_size=9.5)
    card(slide, 8.88, 2.42, 1.42, 1.06, "status", "read 0x004\nbusy clear => done", ORANGE_FILL, ORANGE_LINE, body_size=9.6)
    card(slide, 10.58, 2.42, 1.54, 1.06, "readback", "read_hw_params\nW1/B1/W2/B2\nfrom UB", GREEN_FILL, GREEN_LINE, body_size=9.2)

    poly_arrow(slide, [(8.60, 2.95), (8.88, 2.95)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(10.30, 2.95), (10.58, 2.95)], ORANGE_LINE, 1.8)

    tag(slide, 7.30, 3.60, 1.36, "TRAIN_EPOCHS=12", BLUE_FILL, BLUE_LINE)
    tag(slide, 8.86, 3.60, 1.54, "500 cycles poll", ORANGE_FILL, ORANGE_LINE)
    tag(slide, 10.62, 3.60, 1.30, "wr_ptr check", GREEN_FILL, GREEN_LINE)

    card(slide, 1.00, 4.96, 2.18, 0.88, "ref model", "forward_model() / backward_model() / update_model()\nNumPy Q8.8", PURPLE_FILL, PURPLE_LINE, body_size=9.6)
    card(slide, 3.44, 4.96, 1.96, 0.88, "online mon", "monitor_lrd()\nmonitor_gd()", BLUE_FILL, BLUE_LINE, body_size=10.0)
    card(slide, 5.64, 4.96, 2.04, 0.88, "UB probe", "epoch1 dump dZ2@29 / dZ1@33\n定位 writeback 与 update", GREEN_FILL, GREEN_LINE, body_size=9.4)
    card(slide, 7.94, 4.96, 1.42, 0.88, "assert", "pred == (0,1,1,0)", ORANGE_FILL, ORANGE_LINE, body_size=10.0)
    card(slide, 9.60, 4.96, 1.32, 0.88, "assert", "loss < 0.21", ORANGE_FILL, ORANGE_LINE, body_size=10.0)
    card(slide, 11.14, 4.96, 1.18, 0.88, "assert", "loss\nimproves", ORANGE_FILL, ORANGE_LINE, body_size=10.0)

    poly_arrow(slide, [(3.18, 5.40), (3.44, 5.40)], PURPLE_LINE, 1.7)
    poly_arrow(slide, [(5.40, 5.40), (5.64, 5.40)], BLUE_LINE, 1.7)
    poly_arrow(slide, [(7.68, 5.40), (7.94, 5.40)], GREEN_LINE, 1.7)
    poly_arrow(slide, [(9.36, 5.40), (9.60, 5.40)], ORANGE_LINE, 1.7)
    poly_arrow(slide, [(10.92, 5.40), (11.14, 5.40)], ORANGE_LINE, 1.7)

    pill(slide, 0.98, 6.10, 1.18, 0.24, "init_loss 0.2529", BLUE_FILL, BLUE_LINE, 9.0)
    pill(slide, 2.34, 6.10, 1.20, 0.24, "12 epoch", GREEN_FILL, GREEN_LINE, 9.0)
    pill(slide, 3.74, 6.10, 1.30, 0.24, "final_loss 0.1807", ORANGE_FILL, ORANGE_LINE, 9.0)
    pill(slide, 5.24, 6.10, 1.42, 0.24, "pred (0,1,1,0)", PURPLE_FILL, PURPLE_LINE, 9.0)
    pill(slide, 6.86, 6.10, 1.52, 0.24, "single load, multi start", BLUE_FILL, BLUE_LINE, 9.0)
    pill(slide, 8.60, 6.10, 1.34, 0.24, "Q8.8 ref model", GREEN_FILL, GREEN_LINE, 9.0)
    pill(slide, 10.14, 6.10, 1.66, 0.24, "HW / SW loop aligned", ORANGE_FILL, ORANGE_LINE, 9.0)

    textbox(slide, 0.56, 6.74, 6.00, 0.18, "主叙事：一次 preload + 一次 imem load，之后只反复 start，观察收敛闭环。", size=10.0, color=SUB)
    textbox(slide, 8.28, 6.74, 4.26, 0.18, "代码入口：train_convergence.py", size=10.0, color=SUB, align=PP_ALIGN.RIGHT)

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    main()
