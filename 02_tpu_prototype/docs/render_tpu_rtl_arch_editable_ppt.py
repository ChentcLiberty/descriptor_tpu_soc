from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/home/jjt/tpu-soc")
PACK = ROOT / "career" / "hisilicon_25min_interview_20260410"
OUT_PPT = PACK / "07_TinyTPU_RTL架构图_整合版_可编辑.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = "F8F6F1"
TEXT = "243241"
SUB = "718094"
PANEL = "D7DFE7"
WHITE = "FEFEFD"
BLUE_FILL = "EAF3FA"
BLUE_LINE = "2C7DA8"
ORANGE_FILL = "F7EBDD"
ORANGE_LINE = "C9733D"
GREEN_FILL = "E6F2E9"
GREEN_LINE = "5E946F"
PURPLE_FILL = "ECE7F8"
PURPLE_LINE = "8770D0"
PILL_BLUE = "E7F0F6"
FONT = "WenQuanYi Zen Hei"


def rgb(v: str) -> RGBColor:
    v = v.replace('#', '')
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def I(v: float):
    return Inches(v)


def set_fill(shape, color: str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.15):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width_pt)


def write_text(shape, text: str, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP, ml=0.08, mr=0.08, mt=0.05, mb=0.03):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def round_box(slide, x, y, w, h, fill, line, width_pt=1.15):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    set_fill(sp, fill)
    set_line(sp, line, width_pt)
    return sp


def textbox(slide, x, y, w, h, text, **kwargs):
    sp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    write_text(sp, text, **kwargs)
    return sp


def pill(slide, x, y, w, h, text, fill, color, size=10.0, bold=False):
    sp = round_box(slide, x, y, w, h, fill, fill, 0.6)
    write_text(sp, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, ml=0.03, mr=0.03, mt=0.01, mb=0.01)
    return sp


def module(slide, x, y, w, h, title, fill=WHITE, line=PANEL, title_size=12.0):
    round_box(slide, x, y, w, h, fill, line, 1.05)
    textbox(slide, x + 0.12, y + 0.11, w - 0.24, h - 0.18, title, size=title_size, bold=True, valign=MSO_ANCHOR.MIDDLE)


def line(slide, x1, y1, x2, y2, color, width_pt=2.0):
    sp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    sp.line.color.rgb = rgb(color)
    sp.line.width = Pt(width_pt)
    return sp


def arrow_head(slide, x, y, direction, color, size=0.07):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, I(x - size), I(y - size), I(size * 2), I(size * 2))
    set_fill(sp, color)
    sp.line.fill.background()
    sp.rotation = {"up": 0, "right": 90, "down": 180, "left": 270}[direction]


def poly_arrow(slide, pts, color, width_pt=2.0):
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line(slide, x1, y1, x2, y2, color, width_pt)
    (x1, y1), (x2, y2) = pts[-2], pts[-1]
    direction = 'right' if abs(x2 - x1) >= abs(y2 - y1) and x2 > x1 else 'left' if abs(x2 - x1) >= abs(y2 - y1) else 'down' if y2 > y1 else 'up'
    arrow_head(slide, x2, y2, direction, color)


def flow(slide, x, y, w, text, fill, color):
    pill(slide, x, y, w, 0.21, text, fill, color, size=9.0)


def main():
    PACK.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    textbox(slide, 0.56, 0.18, 0.40, 0.18, 'P2', size=12.0, color=BLUE_LINE)
    textbox(slide, 0.56, 0.45, 2.2, 0.38, 'RTL 架构图', size=22.0, bold=True)
    textbox(slide, 0.56, 0.84, 6.6, 0.22, '单页整合：top wrapper + frontend + TPU core + observability', size=12.0, color=SUB)
    pill(slide, 10.30, 0.34, 2.30, 0.34, '可编辑单页细化版', PILL_BLUE, BLUE_LINE, 11.5)

    round_box(slide, 0.40, 1.30, 12.50, 5.62, WHITE, PANEL, 1.1)
    pill(slide, 0.56, 1.40, 1.10, 0.24, 'top wrapper', PILL_BLUE, BLUE_LINE, 10.6, True)
    textbox(slide, 1.82, 1.39, 1.60, 0.18, 'src_axi/tpu_soc.sv', size=10.2, color=SUB)

    round_box(slide, 0.66, 1.78, 4.50, 4.46, WHITE, PANEL, 1.05)
    textbox(slide, 0.94, 2.00, 2.40, 0.26, 'Frontend 控制域', size=16.6, bold=True)
    textbox(slide, 0.94, 2.28, 3.80, 0.20, 'AXI-Lite / CSR / IMEM / sequencer / decode', size=9.4, color=SUB)
    pill(slide, 0.98, 2.58, 1.06, 0.22, 'frontend_axil.sv', BLUE_FILL, BLUE_LINE, 8.9)
    pill(slide, 2.12, 2.58, 0.66, 0.22, 'cu.sv', ORANGE_FILL, ORANGE_LINE, 8.9)

    module(slide, 0.98, 3.06, 1.08, 0.58, 'AXI IF', BLUE_FILL, PANEL, 11.2)
    module(slide, 2.16, 3.06, 1.08, 0.58, 'Reg Map', BLUE_FILL, PANEL, 11.2)
    module(slide, 3.36, 3.06, 1.10, 0.58, 'Cfg Regs', GREEN_FILL, PANEL, 11.0)
    module(slide, 0.98, 3.96, 0.92, 0.62, 'IMEM', ORANGE_FILL, PANEL, 12.0)
    module(slide, 2.02, 3.96, 1.40, 0.74, 'Sequencer', ORANGE_FILL, PANEL, 11.6)
    module(slide, 3.58, 3.96, 0.88, 0.62, 'CU', ORANGE_FILL, PANEL, 12.0)
    module(slide, 0.98, 5.16, 1.76, 0.68, 'Host Write', GREEN_FILL, PANEL, 11.4)
    module(slide, 2.92, 5.16, 1.54, 0.68, 'Step / Status', BLUE_FILL, PANEL, 10.8)

    round_box(slide, 5.48, 1.78, 5.26, 4.46, WHITE, PANEL, 1.05)
    textbox(slide, 5.76, 2.00, 2.60, 0.26, 'TPU Core 数据域', size=16.6, bold=True)
    textbox(slide, 5.76, 2.28, 3.40, 0.20, 'UB / SA / VPU / writeback', size=9.4, color=SUB)
    pill(slide, 5.80, 2.58, 0.52, 0.22, 'tpu.sv', BLUE_FILL, BLUE_LINE, 8.9)

    module(slide, 5.82, 3.02, 2.08, 2.16, 'Unified Buffer', WHITE, PANEL, 12.8)
    pill(slide, 6.10, 4.62, 1.10, 0.22, 'input / weight', BLUE_FILL, BLUE_LINE, 9.6)
    pill(slide, 6.10, 4.96, 1.06, 0.22, 'bias / Y / H', GREEN_FILL, GREEN_LINE, 9.6)
    module(slide, 8.28, 2.90, 2.18, 1.14, 'Systolic Array', WHITE, PANEL, 13.0)
    module(slide, 8.28, 4.28, 2.18, 1.62, 'VPU', WHITE, PANEL, 13.2)
    pill(slide, 8.54, 4.60, 0.64, 0.22, 'vpu_inst', PILL_BLUE, BLUE_LINE, 8.8)
    pill(slide, 8.50, 5.00, 0.70, 0.36, 'Bias', BLUE_FILL, BLUE_LINE, 11.8, True)
    pill(slide, 9.54, 5.00, 0.68, 0.36, 'Loss', ORANGE_FILL, ORANGE_LINE, 11.8, True)
    pill(slide, 8.50, 5.42, 0.70, 0.36, 'LReLU', GREEN_FILL, GREEN_LINE, 11.2, True)
    pill(slide, 9.54, 5.42, 0.68, 0.36, 'dLReLU', PURPLE_FILL, PURPLE_LINE, 10.4, True)

    round_box(slide, 11.08, 2.24, 1.54, 3.42, WHITE, PANEL, 1.05)
    textbox(slide, 11.30, 2.50, 1.10, 0.24, '对外观测', size=15.2, bold=True, align=PP_ALIGN.CENTER)
    textbox(slide, 11.22, 2.86, 1.18, 0.20, 'host / scoreboard', size=8.8, color=SUB, align=PP_ALIGN.CENTER)
    pill(slide, 11.36, 3.94, 0.92, 0.24, 'sys_out', BLUE_FILL, BLUE_LINE, 10.6)
    pill(slide, 11.34, 4.42, 0.96, 0.24, 'vpu_out', ORANGE_FILL, ORANGE_LINE, 10.6)
    pill(slide, 11.38, 4.90, 0.88, 0.24, 'valid', GREEN_FILL, GREEN_LINE, 10.6)

    poly_arrow(slide, [(2.06, 3.35), (2.16, 3.35)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(3.24, 3.35), (3.36, 3.35)], GREEN_LINE, 1.8)
    poly_arrow(slide, [(2.54, 3.64), (2.54, 3.82), (1.44, 3.82), (1.44, 3.96)], ORANGE_LINE, 1.7)
    poly_arrow(slide, [(2.48, 3.64), (1.96, 3.64), (1.96, 4.86), (1.86, 4.86), (1.86, 5.16)], GREEN_LINE, 1.7)
    poly_arrow(slide, [(2.92, 3.64), (3.52, 3.64), (3.52, 4.94), (3.69, 4.94), (3.69, 5.16)], BLUE_LINE, 1.7)
    poly_arrow(slide, [(1.90, 4.27), (2.02, 4.27)], ORANGE_LINE, 1.8)
    poly_arrow(slide, [(3.38, 5.16), (3.38, 4.78), (2.72, 4.78)], BLUE_LINE, 1.7)
    poly_arrow(slide, [(3.42, 4.27), (3.58, 4.27)], ORANGE_LINE, 1.8)

    poly_arrow(slide, [(4.46, 3.18), (5.48, 3.18)], GREEN_LINE, 1.9)
    flow(slide, 4.74, 2.95, 0.44, 'cfg', GREEN_FILL, GREEN_LINE)
    poly_arrow(slide, [(4.46, 4.12), (5.02, 4.12), (5.02, 2.72), (8.08, 2.72), (8.08, 3.18), (8.28, 3.18)], ORANGE_LINE, 1.9)
    flow(slide, 7.16, 2.49, 0.98, 'switch -> SA', ORANGE_FILL, ORANGE_LINE)
    poly_arrow(slide, [(4.46, 4.26), (5.82, 4.26)], ORANGE_LINE, 2.1)
    flow(slide, 4.64, 4.02, 0.72, 'ub_rd_*', ORANGE_FILL, ORANGE_LINE)
    poly_arrow(slide, [(4.46, 4.42), (5.06, 4.42), (5.06, 5.30), (8.10, 5.30), (8.10, 4.72), (8.28, 4.72)], ORANGE_LINE, 1.8)
    flow(slide, 5.18, 5.04, 0.74, 'pathway', ORANGE_FILL, ORANGE_LINE)
    poly_arrow(slide, [(4.46, 5.50), (5.82, 5.50)], GREEN_LINE, 2.1)
    flow(slide, 4.62, 5.28, 0.78, 'host load', GREEN_FILL, GREEN_LINE)
    poly_arrow(slide, [(7.90, 3.46), (8.28, 3.46)], BLUE_LINE, 2.1)
    poly_arrow(slide, [(7.90, 5.08), (8.28, 5.08)], GREEN_LINE, 2.1)
    poly_arrow(slide, [(9.36, 4.04), (9.36, 4.28)], BLUE_LINE, 2.1)
    poly_arrow(slide, [(8.28, 5.74), (7.56, 5.74), (7.56, 5.50), (7.90, 5.50)], PURPLE_LINE, 2.1)
    flow(slide, 7.30, 5.98, 0.42, 'wb', PURPLE_FILL, PURPLE_LINE)
    poly_arrow(slide, [(10.46, 3.46), (11.08, 3.46)], BLUE_LINE, 2.0)
    poly_arrow(slide, [(10.46, 5.06), (11.08, 5.06)], ORANGE_LINE, 2.0)
    poly_arrow(slide, [(11.08, 5.02), (10.84, 5.02), (10.84, 6.14), (2.88, 6.14), (2.88, 4.78), (2.72, 4.78)], GREEN_LINE, 1.6)
    flow(slide, 5.92, 5.88, 0.76, 'valid fb', GREEN_FILL, GREEN_LINE)

    pill(slide, 5.54, 6.48, 1.86, 0.24, 'ctrl: cfg / ub_rd / switch', ORANGE_FILL, ORANGE_LINE, 8.6)
    pill(slide, 7.46, 6.48, 1.26, 0.24, 'data: UB -> SA -> VPU', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 8.82, 6.48, 1.18, 0.24, 'wb: VPU -> UB', GREEN_FILL, GREEN_LINE, 8.8)
    textbox(slide, 0.56, 7.00, 6.20, 0.18, '只保留面试可讲的真实 RTL 层次和主链路。', size=10.0, color=SUB)

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == '__main__':
    main()
