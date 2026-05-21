from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/home/jjt/tpu-soc")
PACK = ROOT / "career" / "hisilicon_25min_interview_20260410"
OUT_PPT = PACK / "11_TinyTPU_UB128_to_64_DC_可编辑.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = "F8F6F1"
TEXT = "243241"
SUB = "718094"
PANEL = "D7DFE7"
WHITE = "FEFEFD"
BLUE_FILL = "EAF3FA"
BLUE_LINE = "2C7DA8"
ORANGE_FILL = "F7EBDD"
ORANGE_LINE = "C9733D"
GREEN_FILL = "E6F2E9"
GREEN_LINE = "5E946F"
PURPLE_FILL = "ECE7F8"
PURPLE_LINE = "8770D0"
PILL_BLUE = "E7F0F6"
GRAY_FILL = "EFF1F4"
FONT = "Microsoft YaHei"


def rgb(v: str) -> RGBColor:
    v = v.replace("#", "")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def I(v: float):
    return Inches(v)


def set_fill(shape, color: str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.15):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width_pt)


def write_text(
    shape,
    text: str,
    *,
    size=14,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    ml=0.08,
    mr=0.08,
    mt=0.05,
    mb=0.03,
):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    tf.vertical_anchor = valign
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = 1.05
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def round_box(slide, x, y, w, h, fill, line, width_pt=1.15):
    sp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h)
    )
    set_fill(sp, fill)
    set_line(sp, line, width_pt)
    return sp


def textbox(slide, x, y, w, h, text, **kwargs):
    sp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    write_text(sp, text, **kwargs)
    return sp


def pill(slide, x, y, w, h, text, fill, color, size=10.0, bold=False):
    sp = round_box(slide, x, y, w, h, fill, fill, 0.6)
    write_text(
        sp,
        text,
        size=size,
        bold=bold,
        color=color,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        ml=0.03,
        mr=0.03,
        mt=0.01,
        mb=0.01,
    )
    return sp


def card(
    slide,
    x,
    y,
    w,
    h,
    title,
    subtitle,
    fill,
    accent,
    *,
    body_size=10.2,
    label_w=1.08,
    align=PP_ALIGN.LEFT,
):
    round_box(slide, x, y, w, h, fill, PANEL, 1.0)
    pill(
        slide,
        x + 0.10,
        y + 0.10,
        min(w - 0.20, label_w),
        0.22,
        title,
        fill,
        accent,
        8.8,
        True,
    )
    textbox(
        slide,
        x + 0.14,
        y + 0.38,
        w - 0.28,
        h - 0.44,
        subtitle,
        size=body_size,
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
        ml=0.05,
        mr=0.05,
    )


def line(slide, x1, y1, x2, y2, color, width_pt=2.0):
    sp = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2)
    )
    sp.line.color.rgb = rgb(color)
    sp.line.width = Pt(width_pt)
    return sp


def arrow_head(slide, x, y, direction, color, size=0.07):
    sp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        I(x - size),
        I(y - size),
        I(size * 2),
        I(size * 2),
    )
    set_fill(sp, color)
    sp.line.fill.background()
    sp.rotation = {"up": 0, "right": 90, "down": 180, "left": 270}[direction]


def poly_arrow(slide, pts, color, width_pt=2.0):
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line(slide, x1, y1, x2, y2, color, width_pt)
    (x1, y1), (x2, y2) = pts[-2], pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = "right" if x2 > x1 else "left"
    else:
        direction = "down" if y2 > y1 else "up"
    arrow_head(slide, x2, y2, direction, color)


def metric(slide, x, y, w, h, title, value, fill, accent):
    round_box(slide, x, y, w, h, fill, PANEL, 0.9)
    textbox(
        slide,
        x + 0.08,
        y + 0.08,
        w - 0.16,
        0.14,
        title,
        size=8.4,
        color=SUB,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        x + 0.08,
        y + 0.23,
        w - 0.16,
        h - 0.26,
        value,
        size=11.8,
        bold=True,
        color=accent,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def table_cell(
    slide,
    x,
    y,
    w,
    h,
    text,
    fill,
    line_color,
    *,
    size=9.2,
    bold=False,
    color=TEXT,
    align=PP_ALIGN.CENTER,
):
    sp = round_box(slide, x, y, w, h, fill, line_color, 0.8)
    write_text(
        sp,
        text,
        size=size,
        bold=bold,
        color=color,
        align=align,
        valign=MSO_ANCHOR.MIDDLE,
        ml=0.04,
        mr=0.04,
        mt=0.02,
        mb=0.02,
    )
    return sp


def main():
    PACK.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    textbox(slide, 0.56, 0.18, 0.46, 0.18, "P4", size=12.0, color=BLUE_LINE)
    textbox(slide, 0.56, 0.45, 5.40, 0.38, "UB depth128 -> 64 DC 综合对比", size=21.0, bold=True)
    textbox(
        slide,
        0.56,
        0.84,
        9.20,
        0.22,
        "Synopsys Design Compiler / SMIC180 TT 1.8V 25C / 10ns = 100MHz / pre-layout PPA",
        size=11.0,
        color=SUB,
    )
    pill(slide, 10.02, 0.34, 2.64, 0.34, "面试一页直讲版", PILL_BLUE, BLUE_LINE, 10.2)

    round_box(slide, 0.40, 1.28, 12.52, 5.74, WHITE, PANEL, 1.1)

    round_box(slide, 0.64, 1.58, 2.76, 2.12, WHITE, PANEL, 1.0)
    textbox(slide, 0.90, 1.78, 1.70, 0.24, "工具 / 约束", size=15.0, bold=True)
    textbox(slide, 0.90, 2.04, 2.20, 0.18, "同一 flow 对比 depth128 与 depth64", size=8.8, color=SUB)
    card(
        slide,
        0.90,
        2.36,
        0.74,
        0.92,
        "tool",
        "dc_shell\nT-2022.03-SP2",
        BLUE_FILL,
        BLUE_LINE,
        body_size=8.2,
        label_w=0.52,
        align=PP_ALIGN.CENTER,
    )
    card(
        slide,
        1.74,
        2.36,
        0.88,
        0.92,
        "PDK",
        "SMIC180\nTT 1.8V\n25C",
        GREEN_FILL,
        GREEN_LINE,
        body_size=8.1,
        label_w=0.56,
        align=PP_ALIGN.CENTER,
    )
    card(
        slide,
        2.72,
        2.36,
        0.46,
        0.92,
        "SDC",
        "10ns\n100MHz",
        ORANGE_FILL,
        ORANGE_LINE,
        body_size=8.1,
        label_w=0.40,
        align=PP_ALIGN.CENTER,
    )
    textbox(
        slide,
        0.92,
        3.36,
        2.18,
        0.24,
        "证据：command.log 显示 Running dc_shell Version T-2022.03-SP2",
        size=7.5,
        color=SUB,
    )

    round_box(slide, 3.64, 1.58, 3.12, 2.12, WHITE, PANEL, 1.0)
    textbox(slide, 3.90, 1.78, 1.26, 0.24, "DC 流程", size=15.0, bold=True)
    textbox(slide, 3.90, 2.04, 2.30, 0.18, "同一脚本，不同 UB depth 参数快照", size=8.8, color=SUB)
    card(
        slide,
        3.90,
        2.42,
        0.58,
        0.82,
        "1",
        "analyze\n+ elaborate",
        BLUE_FILL,
        BLUE_LINE,
        body_size=7.7,
        label_w=0.30,
        align=PP_ALIGN.CENTER,
    )
    card(
        slide,
        4.70,
        2.42,
        0.66,
        0.82,
        "2",
        "source\nconstraints.sdc",
        GREEN_FILL,
        GREEN_LINE,
        body_size=7.5,
        label_w=0.30,
        align=PP_ALIGN.CENTER,
    )
    card(
        slide,
        5.60,
        2.42,
        0.62,
        0.82,
        "3",
        "compile_ultra\n-gate_clock",
        ORANGE_FILL,
        ORANGE_LINE,
        body_size=7.5,
        label_w=0.30,
        align=PP_ALIGN.CENTER,
    )
    card(
        slide,
        6.42,
        2.42,
        0.20,
        0.82,
        "4",
        "qor\narea\npower\ncg",
        PURPLE_FILL,
        PURPLE_LINE,
        body_size=7.0,
        label_w=0.20,
        align=PP_ALIGN.CENTER,
    )
    poly_arrow(slide, [(4.48, 2.83), (4.70, 2.83)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(5.36, 2.83), (5.60, 2.83)], GREEN_LINE, 1.8)
    poly_arrow(slide, [(6.22, 2.83), (6.42, 2.83)], ORANGE_LINE, 1.8)
    pill(slide, 4.00, 3.32, 1.16, 0.22, "dc_shell -f dc_script.tcl", PILL_BLUE, BLUE_LINE, 7.8, False)
    pill(slide, 5.28, 3.32, 1.08, 0.22, "clock gating on", GREEN_FILL, GREEN_LINE, 7.8, False)

    round_box(slide, 7.00, 1.58, 5.64, 2.12, WHITE, PANEL, 1.0)
    textbox(slide, 7.26, 1.78, 1.44, 0.24, "结果对比", size=15.0, bold=True)
    textbox(slide, 7.26, 2.04, 2.86, 0.18, "两版都满足 10ns；重点比较面积 / 功耗 / ICG", size=8.8, color=SUB)
    pill(slide, 10.48, 1.76, 1.84, 0.22, "WNS/TNS = 0/0 for both", GREEN_FILL, GREEN_LINE, 7.8, False)

    x0 = 7.22
    y0 = 2.36
    row_h = 0.26
    c0 = 1.28
    c1 = 1.18
    c2 = 1.18
    c3 = 1.12
    table_cell(slide, x0, y0, c0, row_h, "metric", GRAY_FILL, PANEL, size=8.8, bold=True)
    table_cell(slide, x0 + c0, y0, c1, row_h, "depth128", BLUE_FILL, PANEL, size=8.8, bold=True, color=BLUE_LINE)
    table_cell(slide, x0 + c0 + c1, y0, c2, row_h, "depth64", GREEN_FILL, PANEL, size=8.8, bold=True, color=GREEN_LINE)
    table_cell(slide, x0 + c0 + c1 + c2, y0, c3, row_h, "delta", ORANGE_FILL, PANEL, size=8.8, bold=True, color=ORANGE_LINE)

    rows = [
        ("area", "770.9K", "518.5K", "-32.7%"),
        ("power", "15.020mW", "13.842mW", "-7.8%"),
        ("ICG cells", "163", "99", "-39.3%"),
        ("gated regs", "2773", "1741", "-37.2%"),
        ("leaf cells", "49012", "31314", "-36.1%"),
    ]
    for idx, (name, v128, v64, delta) in enumerate(rows, start=1):
        yy = y0 + idx * row_h
        table_cell(slide, x0, yy, c0, row_h, name, WHITE, PANEL, size=8.6, bold=True, align=PP_ALIGN.LEFT)
        table_cell(slide, x0 + c0, yy, c1, row_h, v128, WHITE, PANEL, size=8.6)
        table_cell(slide, x0 + c0 + c1, yy, c2, row_h, v64, WHITE, PANEL, size=8.6)
        table_cell(slide, x0 + c0 + c1 + c2, yy, c3, row_h, delta, WHITE, PANEL, size=8.6, bold=True, color=ORANGE_LINE)

    round_box(slide, 0.64, 3.98, 3.44, 2.76, WHITE, PANEL, 1.0)
    textbox(slide, 0.90, 4.18, 1.66, 0.24, "结论怎么讲", size=15.0, bold=True)
    textbox(slide, 0.90, 4.44, 2.30, 0.18, "一页里只保留最稳口径", size=8.8, color=SUB)
    card(
        slide,
        0.90,
        4.78,
        2.92,
        1.40,
        "answer line",
        "在相同 10ns / 100MHz 约束下，UB depth 从 128 降到 64 后，面积从 770.9K 降到 518.5K，功耗从 15.020mW 降到 13.842mW。两版均满足 WNS/TNS=0/0，因此后续选 depth64 作为更均衡的 PPA 配置。",
        BLUE_FILL,
        BLUE_LINE,
        body_size=8.4,
        label_w=0.86,
    )
    pill(slide, 0.92, 6.26, 1.02, 0.22, "same flow", GREEN_FILL, GREEN_LINE, 7.8, False)
    pill(slide, 2.04, 6.26, 0.94, 0.22, "same 100MHz", ORANGE_FILL, ORANGE_LINE, 7.8, False)
    pill(slide, 3.08, 6.26, 0.50, 0.22, "A/B", PURPLE_FILL, PURPLE_LINE, 7.8, False)

    round_box(slide, 4.24, 3.98, 8.40, 2.76, WHITE, PANEL, 1.0)
    textbox(slide, 4.50, 4.18, 2.34, 0.24, "报告位置 / 截图清单", size=15.0, bold=True)
    textbox(slide, 4.50, 4.44, 3.40, 0.18, "先截约束和脚本，再截 qor / power / clock_gating summary", size=8.8, color=SUB)

    shot_w = 2.46
    shot_h = 0.88
    gap_x = 0.20
    gap_y = 0.16
    sx = 4.50
    sy = 4.78
    cards = [
        (
            "1 constraints.sdc",
            "tpu_depth64.../constraints.sdc\n截：CLK_PERIOD=10.0\ncreate_clock -period 10.0",
            BLUE_FILL,
            BLUE_LINE,
        ),
        (
            "2 dc_script.tcl",
            "tpu_depth64.../dc_script.tcl\n截：source constraints.sdc\ncompile_ultra -gate_clock",
            GREEN_FILL,
            GREEN_LINE,
        ),
        (
            "3 depth128 qor.rpt",
            "tpu_depth128.../reports/qor.rpt\n截：clk 10.00 / slack 0.00\narea 770899.356977",
            ORANGE_FILL,
            ORANGE_LINE,
        ),
        (
            "4 depth64 qor.rpt",
            "tpu_depth64.../reports/qor.rpt\n截：clk 10.00 / slack 0.00\narea 518543.557142",
            PURPLE_FILL,
            PURPLE_LINE,
        ),
        (
            "5 power.rpt",
            "分别看两版第一页顶层总功耗\n128 = 15.020mW\n64 = 13.842mW",
            BLUE_FILL,
            BLUE_LINE,
        ),
        (
            "6 clock_gating.rpt",
            "summary 区域\n128: 163 ICG / 2773 gated regs\n64: 99 ICG / 1741 gated regs",
            GREEN_FILL,
            GREEN_LINE,
        ),
    ]
    for idx, (title, body, fill, accent) in enumerate(cards):
        cx = sx + (idx % 3) * (shot_w + gap_x)
        cy = sy + (idx // 3) * (shot_h + gap_y)
        card(
            slide,
            cx,
            cy,
            shot_w,
            shot_h,
            title,
            body,
            fill,
            accent,
            body_size=6.8,
            label_w=1.18,
        )

    textbox(
        slide,
        0.56,
        6.86,
        6.10,
        0.18,
        "讲法：同一 DC flow、同一 100MHz 约束下，直接比较 depth128 和 depth64 的面积功耗差异。",
        size=8.8,
        color=SUB,
    )
    textbox(
        slide,
        8.12,
        6.86,
        4.52,
        0.18,
        "输出：career/hisilicon_25min_interview_20260410/11_TinyTPU_UB128_to_64_DC_可编辑.pptx",
        size=8.4,
        color=SUB,
        align=PP_ALIGN.RIGHT,
    )

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    main()
