from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path('/home/jjt/tpu-soc')
PACK = ROOT / 'career' / 'hisilicon_25min_interview_20260410'
OUT_PPT = PACK / '11_TinyTPU_实现边界与交接_可编辑.pptx'

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = 'F8F6F1'
TEXT = '243241'
SUB = '718094'
PANEL = 'D7DFE7'
WHITE = 'FEFEFD'
BLUE_FILL = 'EAF3FA'
BLUE_LINE = '2C7DA8'
ORANGE_FILL = 'F7EBDD'
ORANGE_LINE = 'C9733D'
GREEN_FILL = 'E6F2E9'
GREEN_LINE = '5E946F'
PURPLE_FILL = 'ECE7F8'
PURPLE_LINE = '8770D0'
PILL_BLUE = 'E7F0F6'
GRAY_FILL = 'EFF1F4'
GRAY_LINE = '6A7482'
FONT = 'WenQuanYi Zen Hei'


def rgb(v: str) -> RGBColor:
    v = v.replace('#', '')
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def I(v: float):
    return Inches(v)


def set_fill(shape, color: str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.15):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width_pt)


def write_text(shape, text: str, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
               valign=MSO_ANCHOR.TOP, ml=0.08, mr=0.08, mt=0.05, mb=0.03):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def round_box(slide, x, y, w, h, fill, line, width_pt=1.15):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    set_fill(sp, fill)
    set_line(sp, line, width_pt)
    return sp


def textbox(slide, x, y, w, h, text, **kwargs):
    sp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    write_text(sp, text, **kwargs)
    return sp


def pill(slide, x, y, w, h, text, fill, color, size=10.0, bold=False):
    sp = round_box(slide, x, y, w, h, fill, fill, 0.6)
    write_text(sp, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER,
               valign=MSO_ANCHOR.MIDDLE, ml=0.03, mr=0.03, mt=0.01, mb=0.01)
    return sp


def card(slide, x, y, w, h, title, subtitle, fill, accent, *, body_size=10.0, label_w=1.08,
         align=PP_ALIGN.LEFT):
    round_box(slide, x, y, w, h, fill, PANEL, 1.0)
    pill(slide, x + 0.10, y + 0.10, min(w - 0.20, label_w), 0.22, title, fill, accent, 8.8, True)
    textbox(slide, x + 0.14, y + 0.38, w - 0.28, h - 0.44, subtitle, size=body_size,
            align=align, valign=MSO_ANCHOR.MIDDLE, ml=0.05, mr=0.05)


def line(slide, x1, y1, x2, y2, color, width_pt=2.0):
    sp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    sp.line.color.rgb = rgb(color)
    sp.line.width = Pt(width_pt)
    return sp


def arrow_head(slide, x, y, direction, color, size=0.07):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, I(x - size), I(y - size), I(size * 2), I(size * 2))
    set_fill(sp, color)
    sp.line.fill.background()
    sp.rotation = {'up': 0, 'right': 90, 'down': 180, 'left': 270}[direction]


def poly_arrow(slide, pts, color, width_pt=2.0):
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line(slide, x1, y1, x2, y2, color, width_pt)
    (x1, y1), (x2, y2) = pts[-2], pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = 'right' if x2 > x1 else 'left'
    else:
        direction = 'down' if y2 > y1 else 'up'
    arrow_head(slide, x2, y2, direction, color)


def metric(slide, x, y, w, h, title, value, fill, accent):
    round_box(slide, x, y, w, h, fill, PANEL, 0.9)
    textbox(slide, x + 0.08, y + 0.08, w - 0.16, 0.14, title, size=8.4, color=SUB, align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.08, y + 0.22, w - 0.16, h - 0.24, value, size=11.4, bold=True, color=accent,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def main():
    PACK.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    textbox(slide, 0.56, 0.18, 0.42, 0.18, 'P5', size=12.0, color=BLUE_LINE)
    textbox(slide, 0.56, 0.45, 5.10, 0.38, '实现边界 / Reset / SDC 交接', size=22.0, bold=True)
    textbox(slide, 0.56, 0.84, 8.70, 0.22, 'AXI wrapper boundary / core synthesis target / timing assumptions / handoff to P&R', size=12.0, color=SUB)
    pill(slide, 10.06, 0.34, 2.60, 0.34, '可编辑单页实现图', PILL_BLUE, BLUE_LINE, 11.2)

    round_box(slide, 0.40, 1.28, 12.52, 5.74, WHITE, PANEL, 1.1)

    round_box(slide, 0.64, 1.60, 5.78, 2.34, WHITE, PANEL, 1.0)
    textbox(slide, 0.90, 1.80, 2.30, 0.24, '层次分区 / floorplan 入口', size=16.0, bold=True)
    textbox(slide, 0.90, 2.06, 3.20, 0.18, '综合保留层次，后续物理实现可以按模块切区', size=9.2, color=SUB)

    round_box(slide, 0.96, 2.38, 5.10, 1.26, WHITE, PANEL, 0.9)
    pill(slide, 1.08, 2.50, 0.92, 0.22, 'tpu_soc', BLUE_FILL, BLUE_LINE, 8.8, True)
    round_box(slide, 1.14, 2.78, 0.94, 0.60, BLUE_FILL, PANEL, 0.8)
    textbox(slide, 1.18, 2.92, 0.86, 0.22, 'frontend\nAXI-Lite', size=9.2, bold=True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
    round_box(slide, 2.30, 2.68, 3.46, 0.82, WHITE, PANEL, 0.9)
    pill(slide, 2.40, 2.50, 0.92, 0.22, 'tpu core', GREEN_FILL, GREEN_LINE, 8.8, True)
    round_box(slide, 2.48, 2.92, 0.88, 0.38, ORANGE_FILL, PANEL, 0.8)
    textbox(slide, 2.52, 3.02, 0.80, 0.12, 'systolic', size=8.8, bold=True, align=PP_ALIGN.CENTER)
    round_box(slide, 3.54, 2.92, 0.84, 0.38, PURPLE_FILL, PANEL, 0.8)
    textbox(slide, 3.58, 3.02, 0.76, 0.12, 'VPU', size=8.8, bold=True, align=PP_ALIGN.CENTER)
    round_box(slide, 4.56, 2.92, 0.96, 0.38, BLUE_FILL, PANEL, 0.8)
    textbox(slide, 4.60, 3.02, 0.88, 0.12, 'UB', size=8.8, bold=True, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(2.08, 3.08), (2.30, 3.08)], BLUE_LINE, 1.7)
    pill(slide, 1.02, 3.54, 1.34, 0.22, '-no_autoungroup', GREEN_FILL, GREEN_LINE, 8.8)
    pill(slide, 2.52, 3.54, 1.46, 0.22, 'hierarchy kept for APR', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 4.18, 3.54, 1.36, 0.22, 'good for floorplan', ORANGE_FILL, ORANGE_LINE, 8.8)

    round_box(slide, 6.66, 1.60, 5.98, 2.34, WHITE, PANEL, 1.0)
    textbox(slide, 6.92, 1.80, 2.40, 0.24, 'clock / reset 边界', size=16.0, bold=True)
    textbox(slide, 6.92, 2.06, 3.20, 0.18, 'SoC wrapper 把 AXI 域时钟复位映射到 core clk/rst', size=9.2, color=SUB)

    card(slide, 6.98, 2.42, 1.48, 1.00, 'AXI ports', 's_axil_aclk\ns_axil_aresetn\nactive-low', BLUE_FILL, BLUE_LINE, body_size=9.6, label_w=0.78, align=PP_ALIGN.CENTER)
    card(slide, 8.78, 2.42, 1.64, 1.00, 'frontend', 'clk_out = aclk\nrst_out = ~aresetn', GREEN_FILL, GREEN_LINE, body_size=9.4, label_w=0.78, align=PP_ALIGN.CENTER)
    card(slide, 10.74, 2.42, 1.54, 1.00, 'core pins', 'clk\nrst\nactive-high', ORANGE_FILL, ORANGE_LINE, body_size=9.6, label_w=0.72, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(8.46, 2.92), (8.78, 2.92)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(10.42, 2.92), (10.74, 2.92)], GREEN_LINE, 1.8)
    pill(slide, 6.98, 3.56, 1.52, 0.22, 'current RTL = single clk domain', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 8.72, 3.56, 1.24, 0.22, 'set_ideal_network rst', PURPLE_FILL, PURPLE_LINE, 8.8)
    pill(slide, 10.16, 3.56, 1.28, 0.22, 'false_path from rst', ORANGE_FILL, ORANGE_LINE, 8.8)

    round_box(slide, 0.64, 4.18, 5.60, 2.18, WHITE, PANEL, 1.0)
    textbox(slide, 0.90, 4.38, 2.20, 0.24, '真实 SDC 假设', size=16.0, bold=True)
    textbox(slide, 0.90, 4.64, 3.10, 0.18, '来自 constraints.sdc / write_sdc', size=9.2, color=SUB)

    metric(slide, 0.94, 4.94, 1.02, 0.54, 'clock', '10ns\n100MHz', BLUE_FILL, BLUE_LINE)
    metric(slide, 2.08, 4.94, 1.02, 0.54, 'setup unc', '0.3ns', GREEN_FILL, GREEN_LINE)
    metric(slide, 3.22, 4.94, 1.02, 0.54, 'hold unc', '0.1ns', ORANGE_FILL, ORANGE_LINE)
    metric(slide, 4.36, 4.94, 1.02, 0.54, 'clk tran', '0.15ns', PURPLE_FILL, PURPLE_LINE)
    metric(slide, 0.94, 5.56, 1.02, 0.54, 'I/O delay', '3ns\n30%', ORANGE_FILL, ORANGE_LINE)
    metric(slide, 2.08, 5.56, 1.02, 0.54, 'max tran', '1.0', BLUE_FILL, BLUE_LINE)
    metric(slide, 3.22, 5.56, 1.02, 0.54, 'fanout', '16', GREEN_FILL, GREEN_LINE)
    metric(slide, 4.36, 5.56, 1.02, 0.54, 'DRC', '0', PURPLE_FILL, PURPLE_LINE)
    pill(slide, 0.96, 6.18, 1.24, 0.22, 'write_sdc -> P&R', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 2.38, 6.18, 1.34, 0.22, 'constraints expanded', GREEN_FILL, GREEN_LINE, 8.8)
    pill(slide, 3.92, 6.18, 1.10, 0.22, 'core target = clk/rst', ORANGE_FILL, ORANGE_LINE, 8.4)

    round_box(slide, 6.48, 4.18, 6.16, 2.18, WHITE, PANEL, 1.0)
    textbox(slide, 6.74, 4.38, 2.60, 0.24, '交接物 / 当前边界', size=16.0, bold=True)
    textbox(slide, 6.74, 4.64, 3.40, 0.18, '已经能交给 APR 的内容，以及还没做的实现步骤', size=9.2, color=SUB)

    card(slide, 6.76, 4.96, 1.52, 0.92, 'handoff', 'tpu_syn.v\ntpu_syn.ddc\ntpu_syn.sdc\ntpu_syn.sdf', BLUE_FILL, BLUE_LINE, body_size=9.2, label_w=0.82, align=PP_ALIGN.CENTER)
    card(slide, 8.52, 4.96, 1.56, 0.92, 'gating', 'clock_gating.rpt\n99 banks\nreported', GREEN_FILL, GREEN_LINE, body_size=9.2, label_w=0.68, align=PP_ALIGN.CENTER)
    card(slide, 10.32, 4.96, 1.44, 0.92, 'current', 'pre-layout\nideal clock\nno CTS', ORANGE_FILL, ORANGE_LINE, body_size=9.2, label_w=0.70, align=PP_ALIGN.CENTER)
    card(slide, 11.76, 4.96, 0.68, 0.92, 'next', 'APR', PURPLE_FILL, PURPLE_LINE, body_size=10.2, label_w=0.38, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(11.58, 5.42), (11.76, 5.42)], ORANGE_LINE, 1.8)

    pill(slide, 6.80, 6.02, 1.24, 0.22, 'floorplan / CTS / route', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 8.20, 6.02, 1.40, 0.22, 'sign-off STA', GREEN_FILL, GREEN_LINE, 8.8)
    pill(slide, 9.78, 6.02, 1.28, 0.22, 'VCD + PT PX', PURPLE_FILL, PURPLE_LINE, 8.8)
    pill(slide, 11.22, 6.02, 1.10, 0.22, 'SS corner', ORANGE_FILL, ORANGE_LINE, 8.8)

    textbox(slide, 0.56, 6.82, 6.10, 0.18, '讲法：SoC wrapper 定边界，core netlist 做综合交接，约束和 reset 语义一起交给后端。', size=9.4, color=SUB)
    textbox(slide, 8.04, 6.82, 4.60, 0.18, '来源：constraints.sdc / clock_gating.rpt / tpu_syn.sdc / tpu_soc.sv', size=9.4, color=SUB, align=PP_ALIGN.RIGHT)

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == '__main__':
    main()
