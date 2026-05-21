from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path("/home/jjt/tpu-soc")
PACK = ROOT / "career" / "hisilicon_25min_interview_20260410"
OUT_PPT = PACK / "08_TinyTPU_编译到IMEM控制链_可编辑.pptx"

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = "F8F6F1"
TEXT = "243241"
SUB = "718094"
PANEL = "D7DFE7"
WHITE = "FEFEFD"
BLUE_FILL = "EAF3FA"
BLUE_LINE = "2C7DA8"
ORANGE_FILL = "F7EBDD"
ORANGE_LINE = "C9733D"
GREEN_FILL = "E6F2E9"
GREEN_LINE = "5E946F"
PURPLE_FILL = "ECE7F8"
PURPLE_LINE = "8770D0"
PILL_BLUE = "E7F0F6"
FONT = "WenQuanYi Zen Hei"


def rgb(v: str) -> RGBColor:
    v = v.replace("#", "")
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def I(v: float):
    return Inches(v)


def set_fill(shape, color: str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.15):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width_pt)


def write_text(shape, text: str, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
               valign=MSO_ANCHOR.TOP, ml=0.08, mr=0.08, mt=0.05, mb=0.03):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def round_box(slide, x, y, w, h, fill, line, width_pt=1.15):
    sp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h)
    )
    set_fill(sp, fill)
    set_line(sp, line, width_pt)
    return sp


def textbox(slide, x, y, w, h, text, **kwargs):
    sp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    write_text(sp, text, **kwargs)
    return sp


def pill(slide, x, y, w, h, text, fill, color, size=10.0, bold=False):
    sp = round_box(slide, x, y, w, h, fill, fill, 0.6)
    write_text(
        sp,
        text,
        size=size,
        bold=bold,
        color=color,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
        ml=0.03,
        mr=0.03,
        mt=0.01,
        mb=0.01,
    )
    return sp


def card(slide, x, y, w, h, title, subtitle, fill, accent, *, body_size=10.4):
    round_box(slide, x, y, w, h, fill, PANEL, 1.0)
    pill(slide, x + 0.10, y + 0.10, min(w - 0.20, 1.08), 0.22, title, fill, accent, 8.8, True)
    textbox(slide, x + 0.14, y + 0.40, w - 0.28, h - 0.48, subtitle, size=body_size, valign=MSO_ANCHOR.TOP)


def line(slide, x1, y1, x2, y2, color, width_pt=2.0):
    sp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    sp.line.color.rgb = rgb(color)
    sp.line.width = Pt(width_pt)
    return sp


def arrow_head(slide, x, y, direction, color, size=0.07):
    sp = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE,
        I(x - size),
        I(y - size),
        I(size * 2),
        I(size * 2),
    )
    set_fill(sp, color)
    sp.line.fill.background()
    sp.rotation = {"up": 0, "right": 90, "down": 180, "left": 270}[direction]


def poly_arrow(slide, pts, color, width_pt=2.0):
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line(slide, x1, y1, x2, y2, color, width_pt)
    (x1, y1), (x2, y2) = pts[-2], pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = "right" if x2 > x1 else "left"
    else:
        direction = "down" if y2 > y1 else "up"
    arrow_head(slide, x2, y2, direction, color)


def tag(slide, x, y, w, text, fill, color):
    pill(slide, x, y, w, 0.21, text, fill, color, 8.6)


def main():
    PACK.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    textbox(slide, 0.56, 0.18, 0.42, 0.18, "P1", size=12.0, color=BLUE_LINE)
    textbox(slide, 0.56, 0.45, 4.20, 0.38, "编译到 IMEM 控制链", size=22.0, bold=True)
    textbox(
        slide,
        0.56,
        0.84,
        7.60,
        0.22,
        "model spec / scheduler / encoder / AXI imem_load / sequencer / control_unit",
        size=12.0,
        color=SUB,
    )
    pill(slide, 10.22, 0.34, 2.42, 0.34, "可编辑单页链路图", PILL_BLUE, BLUE_LINE, 11.2)

    round_box(slide, 0.40, 1.28, 12.52, 5.72, WHITE, PANEL, 1.1)

    round_box(slide, 0.64, 1.62, 5.56, 2.12, WHITE, PANEL, 1.0)
    textbox(slide, 0.88, 1.82, 2.00, 0.24, "Compiler 产物域", size=16.4, bold=True)
    textbox(slide, 0.88, 2.08, 2.60, 0.18, "输入规格 -> UB 分配 -> stage schedule -> opcode", size=9.4, color=SUB)

    round_box(slide, 6.48, 1.62, 5.96, 3.60, WHITE, PANEL, 1.0)
    textbox(slide, 6.74, 1.82, 2.20, 0.24, "Runtime 控制域", size=16.4, bold=True)
    textbox(slide, 6.74, 2.08, 3.50, 0.18, "AXI-Lite load -> IMEM -> sequencer -> CU -> core", size=9.4, color=SUB)

    round_box(slide, 0.64, 4.02, 12.00, 1.14, WHITE, PANEL, 1.0)
    textbox(slide, 0.88, 4.22, 2.10, 0.24, "阶段与关键约束", size=16.0, bold=True)
    textbox(slide, 0.88, 4.48, 3.40, 0.18, "面试时一口气讲清的关键信息", size=9.4, color=SUB)

    card(slide, 0.92, 2.42, 1.18, 0.86, "输入", "MLP spec\n2-2-1", BLUE_FILL, BLUE_LINE, body_size=10.4)
    card(slide, 2.28, 2.42, 1.18, 0.86, "分配", "UB map\naddr/shape", GREEN_FILL, GREEN_LINE, body_size=10.2)
    card(slide, 3.64, 2.42, 1.22, 0.86, "调度", "stage cmds\nhost load", ORANGE_FILL, ORANGE_LINE, body_size=10.2)
    card(slide, 5.04, 2.42, 0.90, 0.86, "编码", "imem hex", PURPLE_FILL, PURPLE_LINE, body_size=10.2)

    poly_arrow(slide, [(2.10, 2.84), (2.28, 2.84)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(3.46, 2.84), (3.64, 2.84)], GREEN_LINE, 1.8)
    poly_arrow(slide, [(4.86, 2.84), (5.04, 2.84)], ORANGE_LINE, 1.8)

    tag(slide, 0.98, 3.42, 1.40, "2 layer MLP", BLUE_FILL, BLUE_LINE)
    tag(slide, 2.54, 3.42, 1.60, "2x2 / 2-lane", GREEN_FILL, GREEN_LINE)
    tag(slide, 4.36, 3.42, 1.36, "59 instr", ORANGE_FILL, ORANGE_LINE)

    card(slide, 6.84, 2.40, 1.36, 1.16, "产物", "schedule.json\nimem.hex\nimem.txt", PURPLE_FILL, PURPLE_LINE, body_size=10.0)
    card(slide, 8.48, 2.40, 1.42, 1.16, "装载", "cocotb\nimem_load()", BLUE_FILL, BLUE_LINE, body_size=10.2)
    card(slide, 10.18, 2.40, 1.84, 1.16, "AXI regs", "0x030 IMEM_ADDR\n0x034 IMEM_W0\n0x040 IMEM_WE\n0x044 IMEM_LEN", GREEN_FILL, GREEN_LINE, body_size=9.3)
    card(slide, 7.16, 3.86, 1.16, 0.92, "IMEM", "imem[0:63]\npc=0..len-1", ORANGE_FILL, ORANGE_LINE, body_size=9.8)
    card(slide, 8.68, 3.86, 1.30, 0.92, "SEQ", "IDLE\nDISPATCH\nWAIT\nADVANCE", ORANGE_FILL, ORANGE_LINE, body_size=9.3)
    card(slide, 10.30, 3.86, 1.34, 0.92, "CU", "decode opcode\nfield split", BLUE_FILL, BLUE_LINE, body_size=10.0)
    card(slide, 11.86, 3.86, 0.78, 0.92, "Core", "TPU\ncore", WHITE, PANEL, body_size=10.8)

    poly_arrow(slide, [(5.94, 2.84), (6.84, 2.84)], PURPLE_LINE, 1.9)
    poly_arrow(slide, [(8.20, 2.98), (8.48, 2.98)], BLUE_LINE, 1.9)
    poly_arrow(slide, [(9.90, 2.98), (10.18, 2.98)], GREEN_LINE, 1.9)
    poly_arrow(slide, [(11.10, 3.56), (11.10, 3.86), (7.74, 3.86)], ORANGE_LINE, 1.8)
    poly_arrow(slide, [(8.32, 4.32), (8.68, 4.32)], ORANGE_LINE, 1.8)
    poly_arrow(slide, [(9.98, 4.32), (10.30, 4.32)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(11.64, 4.32), (11.94, 4.32)], BLUE_LINE, 1.8)

    tag(slide, 7.08, 5.00, 1.18, "opcode[2:0]", BLUE_FILL, BLUE_LINE)
    tag(slide, 8.44, 5.00, 1.42, "wait_after[23]", ORANGE_FILL, ORANGE_LINE)
    tag(slide, 10.08, 5.00, 1.42, "vpu_drain", GREEN_FILL, GREEN_LINE)
    tag(slide, 11.66, 5.00, 0.82, "decode", PURPLE_FILL, PURPLE_LINE)

    poly_arrow(slide, [(10.96, 4.78), (10.96, 4.96)], BLUE_LINE, 1.5)
    poly_arrow(slide, [(9.36, 4.78), (9.36, 4.96)], ORANGE_LINE, 1.5)

    tag(slide, 0.96, 4.70, 1.56, "forward_layer1", BLUE_FILL, BLUE_LINE)
    tag(slide, 2.68, 4.70, 1.72, "transition_layer2", ORANGE_FILL, ORANGE_LINE)
    tag(slide, 4.56, 4.70, 1.70, "backward_layer1", GREEN_FILL, GREEN_LINE)
    tag(slide, 6.42, 4.70, 1.54, "update_w1_tile_*", PURPLE_FILL, PURPLE_LINE)
    tag(slide, 8.12, 4.70, 1.54, "update_w2_tile_*", BLUE_FILL, BLUE_LINE)

    tag(slide, 1.02, 5.06, 1.96, "host_load_plan: X/Y/W/B", GREEN_FILL, GREEN_LINE)
    tag(slide, 3.18, 5.06, 2.32, "wait command 不编码进 IMEM", ORANGE_FILL, ORANGE_LINE)
    tag(slide, 5.72, 5.06, 2.46, "sequencer 用 wait_after + vpu_drain", BLUE_FILL, BLUE_LINE)
    tag(slide, 8.44, 5.06, 1.82, "CU 输出 ub_rd_* / ptr_sel", PURPLE_FILL, PURPLE_LINE)
    tag(slide, 10.46, 5.06, 1.84, "vpu_pathway / sys_switch", GREEN_FILL, GREEN_LINE)

    textbox(slide, 6.90, 5.42, 1.28, 0.18, "to core", size=8.8, color=SUB, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(10.98, 4.78), (10.98, 5.42), (8.48, 5.42)], BLUE_LINE, 1.6)

    pill(slide, 6.96, 5.68, 1.24, 0.24, "sys_switch", BLUE_FILL, BLUE_LINE, 9.0)
    pill(slide, 8.40, 5.68, 1.18, 0.24, "ub_rd_*", ORANGE_FILL, ORANGE_LINE, 9.0)
    pill(slide, 9.78, 5.68, 1.00, 0.24, "ptr_sel", GREEN_FILL, GREEN_LINE, 9.0)
    pill(slide, 10.98, 5.68, 1.34, 0.24, "vpu_pathway", PURPLE_FILL, PURPLE_LINE, 9.0)

    textbox(slide, 0.56, 6.74, 5.20, 0.18, "真实链路：schedule.json -> imem.hex -> AXI load -> sequencer -> CU。", size=10.0, color=SUB)
    textbox(slide, 7.38, 6.74, 5.20, 0.18, "测试入口：train_convergence / imem_load / run_one_epoch", size=10.0, color=SUB, align=PP_ALIGN.RIGHT)

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == "__main__":
    main()
