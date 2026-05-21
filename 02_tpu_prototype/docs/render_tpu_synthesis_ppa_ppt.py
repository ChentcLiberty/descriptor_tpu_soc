from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path('/home/jjt/tpu-soc')
PACK = ROOT / 'career' / 'hisilicon_25min_interview_20260410'
OUT_PPT = PACK / '10_TinyTPU_综合时序PPA_可编辑.pptx'

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = 'F8F6F1'
TEXT = '243241'
SUB = '718094'
PANEL = 'D7DFE7'
WHITE = 'FEFEFD'
BLUE_FILL = 'EAF3FA'
BLUE_LINE = '2C7DA8'
ORANGE_FILL = 'F7EBDD'
ORANGE_LINE = 'C9733D'
GREEN_FILL = 'E6F2E9'
GREEN_LINE = '5E946F'
PURPLE_FILL = 'ECE7F8'
PURPLE_LINE = '8770D0'
PILL_BLUE = 'E7F0F6'
GRAY_FILL = 'EFF1F4'
FONT = 'WenQuanYi Zen Hei'


def rgb(v: str) -> RGBColor:
    v = v.replace('#', '')
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def I(v: float):
    return Inches(v)


def set_fill(shape, color: str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.15):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width_pt)


def write_text(shape, text: str, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
               valign=MSO_ANCHOR.TOP, ml=0.08, mr=0.08, mt=0.05, mb=0.03):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def round_box(slide, x, y, w, h, fill, line, width_pt=1.15):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    set_fill(sp, fill)
    set_line(sp, line, width_pt)
    return sp


def textbox(slide, x, y, w, h, text, **kwargs):
    sp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    write_text(sp, text, **kwargs)
    return sp


def pill(slide, x, y, w, h, text, fill, color, size=10.0, bold=False):
    sp = round_box(slide, x, y, w, h, fill, fill, 0.6)
    write_text(sp, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER,
               valign=MSO_ANCHOR.MIDDLE, ml=0.03, mr=0.03, mt=0.01, mb=0.01)
    return sp


def card(slide, x, y, w, h, title, subtitle, fill, accent, *, body_size=10.2, label_w=1.08,
         align=PP_ALIGN.LEFT):
    round_box(slide, x, y, w, h, fill, PANEL, 1.0)
    pill(slide, x + 0.10, y + 0.10, min(w - 0.20, label_w), 0.22, title, fill, accent, 8.8, True)
    textbox(slide, x + 0.14, y + 0.38, w - 0.28, h - 0.44, subtitle, size=body_size,
            align=align, valign=MSO_ANCHOR.MIDDLE, ml=0.05, mr=0.05)


def line(slide, x1, y1, x2, y2, color, width_pt=2.0):
    sp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    sp.line.color.rgb = rgb(color)
    sp.line.width = Pt(width_pt)
    return sp


def arrow_head(slide, x, y, direction, color, size=0.07):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, I(x - size), I(y - size), I(size * 2), I(size * 2))
    set_fill(sp, color)
    sp.line.fill.background()
    sp.rotation = {'up': 0, 'right': 90, 'down': 180, 'left': 270}[direction]


def poly_arrow(slide, pts, color, width_pt=2.0):
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line(slide, x1, y1, x2, y2, color, width_pt)
    (x1, y1), (x2, y2) = pts[-2], pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = 'right' if x2 > x1 else 'left'
    else:
        direction = 'down' if y2 > y1 else 'up'
    arrow_head(slide, x2, y2, direction, color)


def tag(slide, x, y, w, text, fill, color):
    pill(slide, x, y, w, 0.21, text, fill, color, 8.4)


def metric(slide, x, y, w, h, title, value, fill, accent):
    round_box(slide, x, y, w, h, fill, PANEL, 0.9)
    textbox(slide, x + 0.08, y + 0.08, w - 0.16, 0.14, title, size=8.4, color=SUB, align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.08, y + 0.23, w - 0.16, h - 0.26, value, size=12.0, bold=True, color=accent,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def hbar(slide, x, y, w, h, frac, fill, line_color, label, value):
    track = round_box(slide, x, y, w, h, GRAY_FILL, PANEL, 0.7)
    track.line.transparency = 1.0
    round_box(slide, x, y, w * frac, h, fill, fill, 0.4)
    textbox(slide, x, y - 0.18, 1.08, 0.14, label, size=8.8, bold=True)
    textbox(slide, x + w - 0.72, y - 0.18, 0.72, 0.14, value, size=8.6, color=line_color, align=PP_ALIGN.RIGHT)


def main():
    PACK.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    textbox(slide, 0.56, 0.18, 0.42, 0.18, 'P4', size=12.0, color=BLUE_LINE)
    textbox(slide, 0.56, 0.45, 5.80, 0.38, '综合 / 时序 / PPA：VPU→UB 切一拍', size=21.0, bold=True)
    textbox(slide, 0.56, 0.84, 9.70, 0.22, 'SMIC180 / Design Compiler / pipeline_exp / register cut + frequency result', size=11.6, color=SUB)
    pill(slide, 10.08, 0.34, 2.58, 0.34, '一页讲清楚优化证据', PILL_BLUE, BLUE_LINE, 11.0)

    round_box(slide, 0.40, 1.24, 12.52, 5.86, WHITE, PANEL, 1.1)

    # Top-left: critical path and register cut location.
    round_box(slide, 0.64, 1.46, 7.10, 2.05, WHITE, PANEL, 1.0)
    textbox(slide, 0.88, 1.66, 2.46, 0.24, '1. 切在哪里', size=15.2, bold=True)
    textbox(slide, 0.88, 1.92, 5.74, 0.18, 'timing.rpt 指向 VPU 后处理到 UB gradient_descent/fxp_mul 的跨模块长路径', size=8.8, color=SUB)
    card(slide, 0.90, 2.32, 1.34, 0.78, 'start', 'vpu_data_\npathway', BLUE_FILL, BLUE_LINE, body_size=8.6, label_w=0.58, align=PP_ALIGN.CENTER)
    card(slide, 2.60, 2.32, 1.34, 0.78, 'VPU', 'post-process\nraw data/valid', GREEN_FILL, GREEN_LINE, body_size=8.3, label_w=0.54, align=PP_ALIGN.CENTER)
    card(slide, 4.30, 2.25, 1.24, 0.92, 'cut FF', 'vpu_ub_\npipe_stage\n+1 cycle', ORANGE_FILL, ORANGE_LINE, body_size=7.9, label_w=0.66, align=PP_ALIGN.CENTER)
    card(slide, 5.90, 2.32, 1.48, 0.78, 'UB', 'writeback\ngrad_descent', PURPLE_FILL, PURPLE_LINE, body_size=8.3, label_w=0.46, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(2.24, 2.71), (2.60, 2.71)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(3.94, 2.71), (4.30, 2.71)], GREEN_LINE, 1.8)
    poly_arrow(slide, [(5.54, 2.71), (5.90, 2.71)], ORANGE_LINE, 1.8)
    tag(slide, 0.92, 3.20, 1.70, 'top-level 边界切分', ORANGE_FILL, ORANGE_LINE)
    tag(slide, 2.80, 3.20, 1.34, 'valid/data 同打一拍', BLUE_FILL, BLUE_LINE)
    tag(slide, 4.34, 3.20, 1.48, '不改 VPU/UB 内部', GREEN_FILL, GREEN_LINE)
    tag(slide, 6.00, 3.20, 1.04, '吞吐不变', PURPLE_FILL, PURPLE_LINE)

    # Top-right: results display.
    round_box(slide, 7.96, 1.46, 4.70, 2.05, WHITE, PANEL, 1.0)
    textbox(slide, 8.20, 1.66, 2.10, 0.24, '2. 结果展示', size=15.2, bold=True)
    textbox(slide, 8.20, 1.92, 3.76, 0.18, '这里保留为答辩时先指给面试官看的结果区', size=8.8, color=SUB)
    metric(slide, 8.20, 2.22, 1.36, 0.72, 'baseline', '164.10\nMHz', BLUE_FILL, BLUE_LINE)
    metric(slide, 9.76, 2.22, 1.36, 0.72, 'pipeline', '183.91\nMHz', GREEN_FILL, GREEN_LINE)
    metric(slide, 11.32, 2.22, 1.00, 0.72, 'Fmax', '+12.1%', ORANGE_FILL, ORANGE_LINE)
    tag(slide, 8.22, 3.10, 1.18, '6.09ns pass', BLUE_FILL, BLUE_LINE)
    tag(slide, 9.58, 3.10, 1.16, '5.44ns pass', GREEN_FILL, GREEN_LINE)
    tag(slide, 10.92, 3.10, 1.00, 'WNS 0.00', ORANGE_FILL, ORANGE_LINE)
    tag(slide, 12.02, 3.10, 0.44, '0 viol', PURPLE_FILL, PURPLE_LINE)

    # Bottom-left: before code.
    round_box(slide, 0.64, 3.74, 5.54, 1.86, BLUE_FILL, PANEL, 1.0)
    pill(slide, 0.86, 3.93, 2.18, 0.23, '修改前：tpu.sv 直连', BLUE_FILL, BLUE_LINE, 8.6, True)
    textbox(slide, 0.88, 4.25, 5.02, 0.94,
            'src_axi/tpu.sv:72-75\n'
            'ub_wr_data_in[0]  = vpu_data_out_1;\n'
            'ub_wr_valid_in[0] = vpu_valid_out_1;\n'
            'ub_wr_data_in[1]  = vpu_data_out_2;\n'
            'ub_wr_valid_in[1] = vpu_valid_out_2;',
            size=7.7, color=TEXT, valign=MSO_ANCHOR.MIDDLE, ml=0.05, mr=0.05)
    tag(slide, 0.88, 5.26, 1.54, 'VPU 输出直接进 UB', BLUE_FILL, BLUE_LINE)
    tag(slide, 2.58, 5.26, 1.70, '跨模块组合路径未切断', ORANGE_FILL, ORANGE_LINE)

    # Bottom-right: after code.
    round_box(slide, 6.42, 3.74, 6.24, 1.86, GREEN_FILL, PANEL, 1.0)
    pill(slide, 6.64, 3.93, 2.58, 0.23, '修改后：raw -> pipe -> UB', GREEN_FILL, GREEN_LINE, 8.6, True)
    textbox(slide, 6.66, 4.23, 5.70, 0.98,
            'src_axi/pipeline_exp/tpu_pipeline.sv\n'
            'vpu_data_out_1 -> vpu_raw_data_out[0]\n'
            'vpu_ub_pipe_stage(.data_in(raw), .valid_in(raw_valid),\n'
            '                  .data_out(pipe), .valid_out(pipe_valid));\n'
            'ub_wr_data_in[0] = vpu_pipe_data_out[0];',
            size=7.1, color=TEXT, valign=MSO_ANCHOR.MIDDLE, ml=0.05, mr=0.05)
    tag(slide, 6.66, 5.26, 1.92, '新增 vpu_ub_pipe_stage.sv', GREEN_FILL, GREEN_LINE)
    tag(slide, 8.78, 5.26, 1.30, '打一拍寄存器', ORANGE_FILL, ORANGE_LINE)
    tag(slide, 10.26, 5.26, 1.58, 'latency +1 cycle', PURPLE_FILL, PURPLE_LINE)

    pill(slide, 0.90, 5.96, 1.54, 0.24, '看路径：左上', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 2.62, 5.96, 1.74, 0.24, '看结果：右上', GREEN_FILL, GREEN_LINE, 8.8)
    pill(slide, 4.54, 5.96, 1.94, 0.24, '看代码：下方对比', ORANGE_FILL, ORANGE_LINE, 8.8)
    pill(slide, 6.66, 5.96, 2.04, 0.24, '报告：reports/qor.rpt', PURPLE_FILL, PURPLE_LINE, 8.8)
    pill(slide, 8.88, 5.96, 1.64, 0.24, '周期 6.09 -> 5.44ns', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 10.70, 5.96, 1.50, 0.24, 'pre-layout DC', GREEN_FILL, GREEN_LINE, 8.8)

    textbox(slide, 0.56, 6.80, 6.18, 0.18,
            '讲法：先用 timing report 定位 VPU->UB 写回长路径，再在 top-level 边界插一拍 pipe，valid/data 同步延后。',
            size=8.5, color=SUB)
    textbox(slide, 7.08, 6.80, 5.48, 0.18,
            '结果：Fmax 164.10MHz -> 183.91MHz；来源：SYN_CODE_CHANGES + resume_pipeline_ub64_at_5p44_20260412/qor.rpt',
            size=8.5, color=SUB, align=PP_ALIGN.RIGHT)

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == '__main__':
    main()
