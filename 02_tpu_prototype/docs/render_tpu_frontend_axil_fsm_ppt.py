from __future__ import annotations

from math import atan2, degrees, pi, cos, sin
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path('/home/jjt/tpu-soc')
DOCS = ROOT / 'docs'
PACK = ROOT / 'career' / 'hisilicon_25min_interview_20260410'

SEQ_PNG_DOCS = DOCS / 'tpu_frontend_axil_seq_fsm_tidy.png'
WRITE_PNG_DOCS = DOCS / 'tpu_frontend_axil_write_fsm_tidy.png'
SEQ_PNG_PACK = PACK / '05a_tpu_frontend_axil_sequencer_fsm_规整版.png'
WRITE_PNG_PACK = PACK / '05b_tpu_frontend_axil_write_fsm_规整版.png'
OUT_PPT = PACK / '05_tpu_frontend_axil_状态机图_规整版.pptx'
OUT_PPT_EDIT = PACK / '05_tpu_frontend_axil_状态机图_规整版_可编辑副本.pptx'

W = 1920
H = 1080
PX_PER_IN = 144.0
SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = (247, 249, 252)
WHITE = (255, 255, 255)
INK = (32, 39, 48)
MUTED = (90, 103, 118)
BORDER = (203, 212, 223)
NAVY = (31, 58, 96)
BLUE = (39, 113, 185)
BLUE_FILL = (232, 242, 252)
GREEN = (62, 130, 93)
GREEN_FILL = (236, 246, 239)
ORANGE = (174, 106, 55)
ORANGE_FILL = (250, 241, 232)
GRAY_FILL = (244, 247, 250)
RED = (157, 69, 69)
PALE_BLUE = (232, 239, 247)

TITLE_FONT = '/usr/share/fonts/google-noto/NotoSansSC-Medium.otf'
BODY_FONT = '/usr/share/fonts/google-noto/NotoSansSC-Regular.otf'
MONO_FONT = '/usr/share/fonts/dejavu/DejaVuSansMono.ttf'

T54 = ImageFont.truetype(TITLE_FONT, 54)
T24 = ImageFont.truetype(TITLE_FONT, 24)
B20 = ImageFont.truetype(BODY_FONT, 20)
B18 = ImageFont.truetype(BODY_FONT, 18)
B16 = ImageFont.truetype(BODY_FONT, 16)
M18 = ImageFont.truetype(MONO_FONT, 18)
M16 = ImageFont.truetype(MONO_FONT, 16)
M14 = ImageFont.truetype(MONO_FONT, 14)

PPT_FONT = 'Noto Sans SC'
PPT_MONO = 'DejaVu Sans Mono'


def px(v: float):
    return Inches(v / PX_PER_IN)


def rgb(color):
    return RGBColor(color[0], color[1], color[2])


def mix(a, b, ratio: float):
    return tuple(int(a[i] * (1 - ratio) + b[i] * ratio) for i in range(3))


def rounded(draw: ImageDraw.ImageDraw, xy, radius=24, fill=WHITE, outline=BORDER, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def text(draw: ImageDraw.ImageDraw, pos, msg, font, fill=INK, spacing=5, align='left'):
    draw.multiline_text(pos, msg, font=font, fill=fill, spacing=spacing, align=align)


def center(draw: ImageDraw.ImageDraw, cx, cy, msg, font, fill=INK, spacing=5):
    bbox = draw.multiline_textbbox((0, 0), msg, font=font, spacing=spacing, align='center')
    draw.multiline_text((cx - (bbox[2] - bbox[0]) / 2, cy - (bbox[3] - bbox[1]) / 2), msg, font=font, fill=fill, spacing=spacing, align='center')


def label_size(msg: str):
    lines = msg.split('\n')
    longest = max(len(line) for line in lines)
    width = max(132, 13 * longest + 34)
    height = 34 + max(0, len(lines) - 1) * 18
    return width, height


def code_size(msg: str):
    return max(92, 13 * len(msg) + 28), 28


def draw_path(draw: ImageDraw.ImageDraw, points, color, width=7, head=18):
    for p0, p1 in zip(points, points[1:]):
        draw.line((p0[0], p0[1], p1[0], p1[1]), fill=color, width=width)
    x0, y0 = points[-2]
    x1, y1 = points[-1]
    ang = atan2(y1 - y0, x1 - x0)
    p2 = (x1 + head * cos(ang + 5 * pi / 6), y1 + head * sin(ang + 5 * pi / 6))
    p3 = (x1 + head * cos(ang - 5 * pi / 6), y1 + head * sin(ang - 5 * pi / 6))
    draw.polygon([(x1, y1), p2, p3], fill=color)


def draw_header(draw: ImageDraw.ImageDraw, title: str, subtitle: str, badge: str):
    draw.rectangle((0, 0, W, 112), fill=WHITE)
    draw.rectangle((0, 0, 18, H), fill=NAVY)
    text(draw, (74, 28), title, T54, fill=INK)
    text(draw, (76, 90), subtitle, M18, fill=BLUE)
    rounded(draw, (1388, 20, 1848, 78), radius=18, fill=PALE_BLUE, outline=(224, 230, 237), width=1)
    text(draw, (1422, 38), badge, B18, fill=NAVY)


def draw_footer(draw: ImageDraw.ImageDraw, left_text: str, right_text: str):
    draw.line((74, 1018, 1848, 1018), fill=BORDER, width=2)
    text(draw, (78, 1030), left_text, B16, fill=MUTED)
    text(draw, (1438, 1030), right_text, B16, fill=MUTED)


def draw_state_box(draw: ImageDraw.ImageDraw, box):
    x0, y0, x1, y1 = box['rect']
    rounded(draw, (x0, y0, x1, y1), radius=28, fill=box['fill'], outline=mix(BORDER, box['band'], 0.18), width=2)
    draw.rounded_rectangle((x0, y0, x1, y0 + 52), radius=28, fill=box['band'], outline=box['band'], width=1)
    draw.rectangle((x0, y0 + 26, x1, y0 + 52), fill=box['band'])
    text(draw, (x0 + 24, y0 + 12), box['title'], T24, fill=WHITE)
    code_w, code_h = code_size(box['code'])
    rounded(draw, (x1 - code_w - 18, y0 + 12, x1 - 18, y0 + 12 + code_h), radius=11, fill=mix(box['band'], WHITE, 0.18), outline=mix(box['band'], WHITE, 0.45), width=1)
    center(draw, x1 - 18 - code_w / 2, y0 + 12 + code_h / 2, box['code'], M14, fill=WHITE)
    y = y0 + 78
    for idx, line in enumerate(box['lines']):
        font = B20 if idx == 0 else B18
        fill = INK if idx == 0 else MUTED
        text(draw, (x0 + 28, y), f'• {line}', font, fill=fill)
        y += 40 if idx == 0 else 34


def draw_label(draw: ImageDraw.ImageDraw, item):
    w, h = label_size(item['text'])
    x, y = item['pos']
    tone = item['tone']
    fill = mix(tone, WHITE, 0.9)
    outline = mix(tone, WHITE, 0.68)
    rounded(draw, (x - w / 2, y - h / 2, x + w / 2, y + h / 2), radius=13, fill=fill, outline=outline, width=2)
    center(draw, x, y, item['text'], M16, fill=tone, spacing=4)


def draw_note(draw: ImageDraw.ImageDraw, note):
    rounded(draw, note['rect'], radius=16, fill=GRAY_FILL, outline=(224, 230, 237), width=1)
    x0, y0, _, _ = note['rect']
    font = M16 if note.get('mono', False) else B16
    text(draw, (x0 + 28, y0 + 20), note['text'], font, fill=NAVY)


def render_png(diag, out_path: Path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, diag['title'], diag['subtitle'], diag['badge'])
    rounded(draw, diag['panel'], radius=26, fill=WHITE, outline=BORDER, width=2)
    for box in diag['boxes']:
        draw_state_box(draw, box)
    for arrow in diag['arrows']:
        draw_path(draw, arrow['points'], arrow['color'], width=arrow.get('width', 7), head=18)
        draw_label(draw, {'text': arrow['label'], 'pos': arrow['label_pos'], 'tone': arrow['color']})
    if diag.get('note'):
        draw_note(draw, diag['note'])
    draw_footer(draw, diag['footer_left'], diag['footer_right'])
    out_path.parent.mkdir(parents=True, exist_ok=True)
    img.save(out_path)


def add_textbox(slide, rect, msg, font_size, color, bold=False, align=PP_ALIGN.LEFT, font_name=PPT_FONT, vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE):
    x, y, w, h = rect
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = vertical_anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = msg
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    return tb


def add_body_box(slide, rect, lines):
    x, y, w, h = rect
    tb = slide.shapes.add_textbox(px(x), px(y), px(w), px(h))
    tf = tb.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(0)
        p.space_after = Pt(2)
        p.line_spacing = 1.0
        run = p.add_run()
        run.text = f'• {line}'
        run.font.name = PPT_FONT
        run.font.size = Pt(12 if idx == 0 else 10)
        run.font.bold = False
        run.font.color.rgb = rgb(INK if idx == 0 else MUTED)
    return tb
def add_round_rect(slide, rect, fill, line=BORDER, radius_shape=MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, line_width=1.2):
    x0, y0, x1, y1 = rect
    shp = slide.shapes.add_shape(radius_shape, px(x0), px(y0), px(x1 - x0), px(y1 - y0))
    shp.fill.solid()
    shp.fill.fore_color.rgb = rgb(fill)
    shp.line.color.rgb = rgb(line)
    shp.line.width = Pt(line_width)
    return shp


def add_state_box_slide(slide, box):
    x0, y0, x1, y1 = box['rect']
    add_round_rect(slide, box['rect'], box['fill'], line=mix(BORDER, box['band'], 0.18), line_width=1.2)
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, px(x0), px(y0), px(x1 - x0), px(52))
    header.fill.solid()
    header.fill.fore_color.rgb = rgb(box['band'])
    header.line.fill.background()
    add_textbox(slide, (x0 + 22, y0 + 6, x1 - x0 - 150, 36), box['title'], 17, WHITE, bold=True, vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)
    code_w, code_h = code_size(box['code'])
    add_round_rect(slide, (x1 - code_w - 18, y0 + 12, x1 - 18, y0 + 12 + code_h), mix(box['band'], WHITE, 0.18), line=mix(box['band'], WHITE, 0.45), line_width=0.8)
    add_textbox(slide, (x1 - code_w - 18, y0 + 10, code_w, code_h + 4), box['code'], 10, WHITE, align=PP_ALIGN.CENTER, font_name=PPT_MONO, vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)
    add_body_box(slide, (x0 + 28, y0 + 76, x1 - x0 - 56, y1 - y0 - 88), box['lines'])


def add_line_segment(slide, p0, p1, color, width_pt=2.5):
    conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(p0[0]), px(p0[1]), px(p1[0]), px(p1[1]))
    conn.line.color.rgb = rgb(color)
    conn.line.width = Pt(width_pt)
    return conn


def add_arrowhead(slide, prev, tip, color, size=18):
    ang = degrees(atan2(tip[1] - prev[1], tip[0] - prev[0]))
    tri = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, px(tip[0] - size / 2), px(tip[1] - size / 2), px(size), px(size))
    tri.fill.solid()
    tri.fill.fore_color.rgb = rgb(color)
    tri.line.fill.background()
    tri.rotation = ang + 90
    return tri


def add_label_slide(slide, text_value, pos, tone):
    w, h = label_size(text_value)
    w += 36
    h += 8
    x, y = pos
    fill = mix(tone, WHITE, 0.9)
    line = mix(tone, WHITE, 0.68)
    add_round_rect(slide, (x - w / 2, y - h / 2, x + w / 2, y + h / 2), fill, line=line, line_width=1.0)
    add_textbox(slide, (x - w / 2 + 8, y - h / 2 + 2, w - 16, h - 4), text_value, 10, tone, align=PP_ALIGN.CENTER, font_name=PPT_MONO, vertical_anchor=MSO_VERTICAL_ANCHOR.MIDDLE)


def add_note_slide(slide, note):
    add_round_rect(slide, note['rect'], GRAY_FILL, line=(224, 230, 237), line_width=1.0)
    x0, y0, x1, y1 = note['rect']
    add_textbox(slide, (x0 + 22, y0 + 10, x1 - x0 - 44, y1 - y0 - 20), note['text'], 12, NAVY, font_name=PPT_MONO if note.get('mono', False) else PPT_FONT)


def render_editable_slide(prs, diag):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_round_rect(slide, (0, 0, 1920, 1080), BG, line=BG, line_width=0.0)
    left_bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, px(18), px(1080))
    left_bar.fill.solid()
    left_bar.fill.fore_color.rgb = rgb(NAVY)
    left_bar.line.fill.background()
    header_bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, px(1920), px(112))
    header_bg.fill.solid()
    header_bg.fill.fore_color.rgb = rgb(WHITE)
    header_bg.line.fill.background()
    add_textbox(slide, (74, 18, 980, 54), diag['title'], 26, INK, bold=True, vertical_anchor=MSO_VERTICAL_ANCHOR.TOP)
    add_textbox(slide, (76, 78, 620, 28), diag['subtitle'], 12, BLUE, font_name=PPT_MONO, vertical_anchor=MSO_VERTICAL_ANCHOR.TOP)
    add_round_rect(slide, (1388, 20, 1848, 78), PALE_BLUE, line=(224, 230, 237), line_width=1.0)
    add_textbox(slide, (1422, 32, 376, 24), diag['badge'], 13, NAVY)
    add_round_rect(slide, diag['panel'], WHITE, line=BORDER, line_width=1.0)
    for box in diag['boxes']:
        add_state_box_slide(slide, box)
    for arrow in diag['arrows']:
        pts = arrow['points']
        for p0, p1 in zip(pts, pts[1:]):
            add_line_segment(slide, p0, p1, arrow['color'], width_pt=2.5 if arrow.get('width', 7) <= 6 else 3.0)
        add_arrowhead(slide, pts[-2], pts[-1], arrow['color'], size=18)
        add_label_slide(slide, arrow['label'], arrow['label_pos'], arrow['color'])
    if diag.get('note'):
        add_note_slide(slide, diag['note'])
    footer_line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(74), px(1018), px(1848), px(1018))
    footer_line.line.color.rgb = rgb(BORDER)
    footer_line.line.width = Pt(1.2)
    add_textbox(slide, (78, 1026, 320, 20), diag['footer_left'], 11, MUTED)
    add_textbox(slide, (1438, 1026, 360, 20), diag['footer_right'], 11, MUTED)


def render_flat_deck(diagrams):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for png_path in diagrams:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(png_path), 0, 0, width=SLIDE_W, height=SLIDE_H)
    prs.save(OUT_PPT)


def render_editable_deck(diagrams):
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    for diag in diagrams:
        render_editable_slide(prs, diag)
    prs.save(OUT_PPT_EDIT)


def seq_diagram():
    return {
        'title': 'tpu_frontend_axil Sequencer FSM',
        'subtitle': 'src_axi/tpu_frontend_axil.sv:122-207',
        'badge': 'Orthogonal state diagram',
        'panel': (70, 150, 1850, 980),
        'boxes': [
            {
                'title': 'SEQ_IDLE',
                'code': "2'b00",
                'rect': (160, 280, 620, 470),
                'band': NAVY,
                'fill': GRAY_FILL,
                'lines': [
                    'Wait for start_pulse or step_pulse',
                    'seq_running = 0',
                    'Busy cleared after completion',
                ],
            },
            {
                'title': 'SEQ_DISPATCH',
                'code': "2'b01",
                'rect': (1000, 280, 1460, 470),
                'band': BLUE,
                'fill': BLUE_FILL,
                'lines': [
                    'Issue one-cycle seq_instr_pulse',
                    'Drive current instruction to control path',
                    'Branch on seq_needs_wait',
                ],
            },
            {
                'title': 'SEQ_WAIT',
                'code': "2'b10",
                'rect': (1000, 670, 1460, 860),
                'band': ORANGE,
                'fill': ORANGE_FILL,
                'lines': [
                    'Hold until vpu_drain is observed',
                    'vpu_drain = prev(valid) & ~valid',
                    'Single-step returns to IDLE here',
                ],
            },
            {
                'title': 'SEQ_ADVANCE',
                'code': "2'b11",
                'rect': (160, 670, 620, 860),
                'band': GREEN,
                'fill': GREEN_FILL,
                'lines': [
                    'Update pc and fetch next IMEM word',
                    'Continue when pc + 1 < imem_len_reg',
                    'Else clear busy_reg and finish',
                ],
            },
        ],
        'arrows': [
            {'points': [(108, 375), (160, 375)], 'color': NAVY, 'label': 'reset', 'label_pos': (78, 332)},
            {'points': [(620, 375), (1000, 375)], 'color': BLUE, 'label': 'start_pulse', 'label_pos': (810, 314)},
            {'points': [(1230, 470), (1230, 670)], 'color': BLUE, 'label': 'seq_needs_wait', 'label_pos': (1374, 572)},
            {'points': [(1000, 765), (620, 765)], 'color': GREEN, 'label': 'vpu_drain && seq_running', 'label_pos': (810, 904)},
            {'points': [(308, 670), (308, 560), (308, 470)], 'color': RED, 'label': 'program_end', 'label_pos': (202, 568)},
            {'points': [(304, 470), (304, 605), (1210, 605), (1210, 670)], 'color': ORANGE, 'label': 'step_pulse', 'label_pos': (602, 572)},
            {'points': [(1100, 470), (1100, 520), (390, 520), (390, 670)], 'color': GREEN, 'label': '!seq_needs_wait', 'label_pos': (778, 486)},
            {'points': [(1460, 760), (1588, 760), (1588, 205), (390, 205), (390, 280)], 'color': ORANGE, 'label': 'vpu_drain &&\n!seq_running', 'label_pos': (990, 162)},
            {'points': [(390, 860), (390, 948), (940, 948), (940, 430), (1000, 430)], 'color': BLUE, 'label': 'pc + 1 <\nimem_len_reg', 'label_pos': (704, 962)},
        ],
        'note': {'rect': (1320, 892, 1810, 954), 'text': 'seq_needs_wait = seq_instr[23]', 'mono': True},
        'footer_left': 'Slide 1 of 2',
        'footer_right': 'Presentation deck + editable copy',
    }


def write_diagram():
    return {
        'title': 'tpu_frontend_axil AXI-Lite Write FSM',
        'subtitle': 'src_axi/tpu_frontend_axil.sv:214-278',
        'badge': 'Orthogonal state diagram',
        'panel': (70, 150, 1850, 980),
        'boxes': [
            {
                'title': 'W_IDLE',
                'code': "2'b00",
                'rect': (160, 280, 620, 470),
                'band': NAVY,
                'fill': GRAY_FILL,
                'lines': [
                    'Accept AW/W in any order',
                    'Latch both together or partial channel',
                    'Clear s_axil_bvalid',
                ],
            },
            {
                'title': 'W_WAIT_W',
                'code': "2'b01",
                'rect': (1000, 280, 1460, 470),
                'band': BLUE,
                'fill': BLUE_FILL,
                'lines': [
                    'AW captured, wait for WVALID',
                    'Capture wd_lat on arrival',
                    'Assert wr_fire when write completes',
                ],
            },
            {
                'title': 'W_WAIT_AW',
                'code': "2'b10",
                'rect': (160, 670, 620, 860),
                'band': GREEN,
                'fill': GREEN_FILL,
                'lines': [
                    'W captured, wait for AWVALID',
                    'Capture aw_lat on arrival',
                    'Assert wr_fire when write completes',
                ],
            },
            {
                'title': 'W_RESP',
                'code': "2'b11",
                'rect': (1000, 670, 1460, 860),
                'band': ORANGE,
                'fill': ORANGE_FILL,
                'lines': [
                    'Drive write response channel',
                    's_axil_bvalid = 1 and bresp = OKAY',
                    'Return after BVALID & BREADY',
                ],
            },
        ],
        'arrows': [
            {'points': [(108, 375), (160, 375)], 'color': NAVY, 'label': 'reset', 'label_pos': (78, 332)},
            {'points': [(620, 375), (1000, 375)], 'color': BLUE, 'label': 'awvalid_only', 'label_pos': (810, 314)},
            {'points': [(385, 470), (385, 670)], 'color': GREEN, 'label': 'wvalid_only', 'label_pos': (244, 568)},
            {'points': [(505, 470), (505, 605), (1210, 605), (1210, 670)], 'color': ORANGE, 'label': 'awvalid && wvalid /\nwr_fire', 'label_pos': (814, 556)},
            {'points': [(1230, 470), (1230, 670)], 'color': BLUE, 'label': 'wvalid / wr_fire', 'label_pos': (1380, 572)},
            {'points': [(620, 765), (1000, 765)], 'color': GREEN, 'label': 'awvalid / wr_fire', 'label_pos': (810, 904)},
            {'points': [(1460, 760), (1588, 760), (1588, 205), (390, 205), (390, 280)], 'color': RED, 'label': 'bvalid &&\nbready', 'label_pos': (990, 162)},
        ],
        'note': {'rect': (1040, 890, 1810, 954), 'text': 'AWREADY = (W_IDLE | W_WAIT_AW)   |   WREADY = (W_IDLE | W_WAIT_W)', 'mono': True},
        'footer_left': 'Slide 2 of 2',
        'footer_right': 'Presentation deck + editable copy',
    }


def main():
    PACK.mkdir(parents=True, exist_ok=True)
    diagrams = [seq_diagram(), write_diagram()]
    png_pairs = [(diagrams[0], SEQ_PNG_DOCS, SEQ_PNG_PACK), (diagrams[1], WRITE_PNG_DOCS, WRITE_PNG_PACK)]
    for diag, docs_path, pack_path in png_pairs:
        render_png(diag, docs_path)
        render_png(diag, pack_path)
    render_flat_deck([SEQ_PNG_PACK, WRITE_PNG_PACK])
    render_editable_deck(diagrams)


if __name__ == '__main__':
    main()
