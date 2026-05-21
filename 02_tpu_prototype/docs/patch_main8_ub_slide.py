from __future__ import annotations

import shutil
import sys
from pathlib import Path

from pptx import Presentation
from pptx.enum.text import PP_ALIGN

DOCS = Path(__file__).resolve().parent
if str(DOCS) not in sys.path:
    sys.path.insert(0, str(DOCS))

import render_tpu_ppt_hisilicon_deep as deck


SLIDE_INDEX = 4
DEFAULT_TARGET = Path('/tmp/VMwareDnD/e5sTYG/02_TinyTPU_8p_备份主讲.pptx')
FOOTER = '这页最适合回答“为什么叫 Unified Buffer”的问题。'

ADDR_ROWS = [
    ('X', '@0 8w', deck.BLUE_BG, deck.BLUE),
    ('Y', '@8 4w', deck.PURPLE_BG, deck.PURPLE),
    ('W1', '@12 4w', deck.GREEN_BG, deck.GREEN),
    ('B1', '@16 2w', deck.GREEN_BG, deck.GREEN),
    ('W2', '@18 2w', deck.ORANGE_BG, deck.ORANGE),
    ('B2', '@20 1w', deck.ORANGE_BG, deck.ORANGE),
    ('H1', '@21 8w', deck.ORANGE_BG, deck.ORANGE),
    ('dZ2', '@29 4w', deck.GREEN_BG, deck.GREEN),
    ('dZ1', '@33 8w', deck.BLUE_BG, deck.BLUE),
]


def clear_slide(slide) -> None:
    for shape in list(slide.shapes):
        node = shape._element
        node.getparent().remove(node)


def add_section_card(slide, x, y, w, h, tone, title, subtitle) -> None:
    deck._add_panel(slide, x, y, w, h, deck.WHITE)
    deck._add_panel(slide, x + 0.12, y + 0.16, 0.05, h - 0.32, tone, tone)
    deck._add_textbox(slide, x + 0.28, y + 0.18, w - 0.46, 0.22, title, 17, deck.INK, True)
    deck._add_textbox(slide, x + 0.28, y + 0.46, w - 0.48, 0.2, subtitle, 8.8, deck.MUTED)


def add_flow_box(slide, x, y, w, h, title, body, fill, tone) -> None:
    deck._add_panel(slide, x, y, w, h, fill)
    deck._add_textbox(slide, x + 0.08, y + 0.12, w - 0.16, 0.16, title, 10.8, tone, True, PP_ALIGN.CENTER)
    deck._add_textbox(slide, x + 0.08, y + 0.32, w - 0.16, h - 0.38, body, 8.0, deck.MUTED, False, PP_ALIGN.CENTER)


def add_note_panel(slide, x, y, w, h, fill, text_value, tone=deck.INK, size=8.4) -> None:
    deck._add_panel(slide, x, y, w, h, fill)
    deck._add_textbox(slide, x + 0.12, y + 0.07, w - 0.24, h - 0.12, text_value, size, tone)


def render_slide(slide, total_pages: int) -> None:
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = deck._rgb(deck.BG)

    deck.PPT_TOTAL_PAGES = total_pages
    deck._render_header_clean(
        slide,
        SLIDE_INDEX + 1,
        'Unified Buffer 设计',
        '保留地址映射，并把 UB_RD 解析、GD 更新和自动参数回写放到一页讲清。',
    )
    deck._add_clean_frame(slide)
    deck._add_chip_ppt(slide, 0.74, 1.66, 1.04, '地址映射', deck.BLUE_BG, deck.BLUE)
    deck._add_chip_ppt(slide, 1.92, 1.66, 1.54, 'UB + GD 闭环', deck.BLUE_BG, deck.BLUE)
    deck._add_chip_ppt(slide, 3.64, 1.66, 1.22, '读指令解析', deck.WHITE, deck.MUTED)
    deck._add_chip_ppt(slide, 5.02, 1.66, 1.48, '自动参数更新', deck.WHITE, deck.MUTED)

    add_section_card(slide, 0.72, 1.88, 2.68, 4.88, deck.BLUE, '保留地址映射', '静态张量沿用 compiler ub_map。')
    y = 2.82
    for label, note, fill, tone in ADDR_ROWS:
        deck._add_panel(slide, 1.02, y, 0.9, 0.18, fill, fill)
        deck._add_textbox(slide, 1.2, y + 0.01, 0.54, 0.1, label, 8.2, tone, True, PP_ALIGN.CENTER)
        deck._add_textbox(slide, 2.04, y + 0.002, 0.82, 0.1, note, 8.8, tone)
        y += 0.28
    add_note_panel(slide, 0.98, 5.68, 2.16, 0.22, deck.HILITE, 'runtime 写回从 wr_ptr_base 后起步。', deck.MUTED, 7.8)
    deck._add_chip_ppt(slide, 1.0, 6.12, 0.92, 'alloc=41', deck.WHITE, deck.BLUE)
    deck._add_chip_ppt(slide, 2.04, 6.12, 0.92, 'free=87', deck.WHITE, deck.MUTED)

    add_section_card(slide, 3.56, 1.88, 4.72, 4.88, deck.ORANGE, 'UB + GD 数据闭环', '装载、读流、梯度回写和参数更新都在 UB 内闭环。')
    deck._add_chip_ppt(slide, 3.92, 2.74, 1.76, 'CTRL.start -> restore', deck.BLUE_BG, deck.BLUE)
    deck._add_chip_ppt(slide, 5.88, 2.74, 1.7, 'wr_ptr = wr_ptr_base', deck.WHITE, deck.BLUE)

    add_flow_box(slide, 3.92, 3.18, 1.16, 0.72, 'host load', 'AXI / UB_PUSH', deck.BLUE_BG, deck.BLUE)
    add_flow_box(slide, 5.3, 3.12, 1.24, 0.9, 'UB space', 'static + runtime', deck.WHITE, deck.ORANGE)
    add_flow_box(slide, 6.8, 3.18, 1.12, 0.72, 'UB_RD', 'ptr_sel 0..4', deck.ORANGE_BG, deck.ORANGE)
    add_flow_box(slide, 3.84, 4.8, 1.38, 0.74, 'VPU writeback', 'grad in', deck.GREEN_BG, deck.GREEN)
    add_flow_box(slide, 5.44, 4.8, 1.04, 0.74, 'GD', 'old + grad*lr', deck.PURPLE_BG, deck.PURPLE)
    add_flow_box(slide, 6.7, 4.8, 1.3, 0.74, 'param writeback', 'updated -> UB', deck.BLUE_BG, deck.BLUE)

    deck._add_arrow_ppt(slide, 5.08, 3.53, 5.3, 3.53, deck.BLUE, 2.2)
    deck._add_arrow_ppt(slide, 6.54, 3.53, 6.8, 3.53, deck.ORANGE, 2.2)
    deck._add_arrow_ppt(slide, 5.92, 4.02, 5.92, 4.8, deck.ORANGE, 2.2)
    deck._add_arrow_ppt(slide, 5.22, 5.16, 5.44, 5.16, deck.GREEN, 2.2)
    deck._add_arrow_ppt(slide, 6.48, 5.16, 6.7, 5.16, deck.PURPLE, 2.2)
    deck._add_arrow_ppt(slide, 7.28, 4.8, 6.28, 4.02, deck.BLUE, 2.2)

    add_note_panel(slide, 3.92, 5.78, 4.0, 0.22, deck.HILITE, 'ptr_sel=5/6 arm update；bias 等首拍，weight 跟 valid。', deck.INK, 7.7)
    deck._add_chip_ppt(slide, 3.94, 6.16, 1.52, 'ptr_sel 5/6', deck.WHITE, deck.GREEN)
    deck._add_chip_ppt(slide, 5.64, 6.16, 1.1, '0x58 -> lr', deck.WHITE, deck.PURPLE)
    deck._add_chip_ppt(slide, 6.92, 6.16, 1.02, 'in-UB loop', deck.WHITE, deck.BLUE)

    add_section_card(slide, 8.42, 1.88, 4.14, 2.32, deck.BLUE, 'UB_RD 读指令解析', 'opcode=010 拆成地址、尺寸、目标和 pathway。')
    add_note_panel(slide, 8.7, 2.76, 3.56, 0.22, deck.WHITE, 'opcode = 010', deck.INK, 8.1)
    add_note_panel(slide, 8.7, 3.04, 3.56, 0.24, deck.WHITE, '[8:3] addr   [12:9] row   [14:13] col', deck.INK, 7.9)
    add_note_panel(slide, 8.7, 3.34, 3.56, 0.26, deck.WHITE, '[15] transpose   [18:16] ptr_sel   [22:19] pathway', deck.INK, 7.7)
    add_note_panel(slide, 8.7, 3.68, 3.56, 0.22, deck.BLUE_BG, 'ptr_sel=0..4 -> input / weight / bias / Y / H', deck.MUTED, 7.7)

    add_section_card(slide, 8.42, 4.36, 4.14, 2.4, deck.GREEN, '自动启动更新参数', 'start / restore / GD address 连在一起。')
    add_note_panel(slide, 8.7, 5.0, 3.56, 0.28, deck.BLUE_BG, 'CTRL.start -> start_pulse -> ub_wr_ptr_restore_out', deck.INK, 7.8)
    add_note_panel(slide, 8.7, 5.36, 3.56, 0.28, deck.GREEN_BG, 'UB: if (ub_wr_ptr_restore_in) wr_ptr <= wr_ptr_base', deck.INK, 7.8)
    add_note_panel(slide, 8.7, 5.72, 3.56, 0.32, deck.PURPLE_BG, 'ptr_sel=5/6: grad_descent_ptr = ub_rd_addr_in\nvalue_updated_out -> UB', deck.INK, 7.7)
    add_note_panel(slide, 8.7, 6.14, 3.56, 0.22, deck.WHITE, '0 in | 1 wt | 2 bias | 3 Y | 4 H | 5 gB | 6 gW', deck.MUTED, 7.7)
    deck._add_static_footer(slide, FOOTER)


def patch_deck(path: Path) -> Path:
    prs = Presentation(path)
    slide = prs.slides[SLIDE_INDEX]
    clear_slide(slide)
    render_slide(slide, len(prs.slides))
    backup = path.with_name(f'{path.stem}.bak_20260412_before_ub_refresh{path.suffix}')
    if not backup.exists():
        shutil.copy2(path, backup)
    prs.save(path)
    return backup


def main() -> None:
    target = Path(sys.argv[1]) if len(sys.argv) > 1 else DEFAULT_TARGET
    backup = patch_deck(target)
    print(f'updated: {target}')
    print(f'backup:  {backup}')


if __name__ == '__main__':
    main()
