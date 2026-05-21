from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path('/home/jjt/tpu-soc')
PACK = ROOT / 'career' / 'hisilicon_25min_interview_20260410'
OUT_PPT = PACK / '12_TinyTPU_GLS与Signoff路径_可编辑.pptx'

SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)

BG = 'F8F6F1'
TEXT = '243241'
SUB = '718094'
PANEL = 'D7DFE7'
WHITE = 'FEFEFD'
BLUE_FILL = 'EAF3FA'
BLUE_LINE = '2C7DA8'
ORANGE_FILL = 'F7EBDD'
ORANGE_LINE = 'C9733D'
GREEN_FILL = 'E6F2E9'
GREEN_LINE = '5E946F'
PURPLE_FILL = 'ECE7F8'
PURPLE_LINE = '8770D0'
PILL_BLUE = 'E7F0F6'
FONT = 'WenQuanYi Zen Hei'


def rgb(v: str) -> RGBColor:
    v = v.replace('#', '')
    return RGBColor(int(v[:2], 16), int(v[2:4], 16), int(v[4:6], 16))


def I(v: float):
    return Inches(v)


def set_fill(shape, color: str):
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = rgb(color)


def set_line(shape, color: str, width_pt: float = 1.15):
    line = shape.line
    line.color.rgb = rgb(color)
    line.width = Pt(width_pt)


def write_text(shape, text: str, *, size=14, bold=False, color=TEXT, align=PP_ALIGN.LEFT,
               valign=MSO_ANCHOR.TOP, ml=0.08, mr=0.08, mt=0.05, mb=0.03):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = I(ml)
    tf.margin_right = I(mr)
    tf.margin_top = I(mt)
    tf.margin_bottom = I(mb)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = FONT
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def round_box(slide, x, y, w, h, fill, line, width_pt=1.15):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, I(x), I(y), I(w), I(h))
    set_fill(sp, fill)
    set_line(sp, line, width_pt)
    return sp


def textbox(slide, x, y, w, h, text, **kwargs):
    sp = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    write_text(sp, text, **kwargs)
    return sp


def pill(slide, x, y, w, h, text, fill, color, size=10.0, bold=False):
    sp = round_box(slide, x, y, w, h, fill, fill, 0.6)
    write_text(sp, text, size=size, bold=bold, color=color, align=PP_ALIGN.CENTER,
               valign=MSO_ANCHOR.MIDDLE, ml=0.03, mr=0.03, mt=0.01, mb=0.01)
    return sp


def card(slide, x, y, w, h, title, subtitle, fill, accent, *, body_size=10.0, label_w=1.08,
         align=PP_ALIGN.LEFT):
    round_box(slide, x, y, w, h, fill, PANEL, 1.0)
    pill(slide, x + 0.10, y + 0.10, min(w - 0.20, label_w), 0.22, title, fill, accent, 8.8, True)
    textbox(slide, x + 0.14, y + 0.38, w - 0.28, h - 0.44, subtitle, size=body_size,
            align=align, valign=MSO_ANCHOR.MIDDLE, ml=0.05, mr=0.05)


def line(slide, x1, y1, x2, y2, color, width_pt=2.0):
    sp = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2))
    sp.line.color.rgb = rgb(color)
    sp.line.width = Pt(width_pt)
    return sp


def arrow_head(slide, x, y, direction, color, size=0.07):
    sp = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ISOSCELES_TRIANGLE, I(x - size), I(y - size), I(size * 2), I(size * 2))
    set_fill(sp, color)
    sp.line.fill.background()
    sp.rotation = {'up': 0, 'right': 90, 'down': 180, 'left': 270}[direction]


def poly_arrow(slide, pts, color, width_pt=2.0):
    for (x1, y1), (x2, y2) in zip(pts, pts[1:]):
        line(slide, x1, y1, x2, y2, color, width_pt)
    (x1, y1), (x2, y2) = pts[-2], pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        direction = 'right' if x2 > x1 else 'left'
    else:
        direction = 'down' if y2 > y1 else 'up'
    arrow_head(slide, x2, y2, direction, color)


def metric(slide, x, y, w, h, title, value, fill, accent):
    round_box(slide, x, y, w, h, fill, PANEL, 0.9)
    textbox(slide, x + 0.08, y + 0.08, w - 0.16, 0.14, title, size=8.4, color=SUB, align=PP_ALIGN.CENTER)
    textbox(slide, x + 0.08, y + 0.22, w - 0.16, h - 0.24, value, size=11.2, bold=True, color=accent,
            align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


def main():
    PACK.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(BG)

    textbox(slide, 0.56, 0.18, 0.42, 0.18, 'P6', size=12.0, color=BLUE_LINE)
    textbox(slide, 0.56, 0.45, 5.40, 0.38, 'GLS / SDF / Sign-off 路径', size=22.0, bold=True)
    textbox(slide, 0.56, 0.84, 8.90, 0.22, 'gate netlist handoff / SDF semantics / current pre-layout fidelity / next sign-off steps', size=12.0, color=SUB)
    pill(slide, 10.10, 0.34, 2.56, 0.34, '可编辑单页 GLS 图', PILL_BLUE, BLUE_LINE, 11.2)

    round_box(slide, 0.40, 1.28, 12.52, 5.74, WHITE, PANEL, 1.1)

    round_box(slide, 0.64, 1.60, 4.18, 2.18, WHITE, PANEL, 1.0)
    textbox(slide, 0.90, 1.80, 2.30, 0.24, '交接到门级后仿', size=16.0, bold=True)
    textbox(slide, 0.90, 2.06, 2.80, 0.18, '综合后产物已经具备 GLS handoff 入口', size=9.2, color=SUB)
    card(slide, 0.92, 2.42, 1.12, 0.98, 'netlist', 'tpu_syn.v\n2.85 MB', BLUE_FILL, BLUE_LINE, body_size=10.0, label_w=0.66, align=PP_ALIGN.CENTER)
    card(slide, 2.20, 2.42, 1.12, 0.98, 'SDF', 'tpu_syn.sdf\n20.84 MB', GREEN_FILL, GREEN_LINE, body_size=10.0, label_w=0.46, align=PP_ALIGN.CENTER)
    card(slide, 3.48, 2.42, 1.12, 0.98, 'SDC', 'tpu_syn.sdc\n20.99 KB', ORANGE_FILL, ORANGE_LINE, body_size=10.0, label_w=0.46, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(2.04, 2.92), (2.20, 2.92)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(3.32, 2.92), (3.48, 2.92)], GREEN_LINE, 1.8)
    pill(slide, 0.96, 3.50, 1.08, 0.22, 'write -verilog', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 2.18, 3.50, 1.08, 0.22, 'write_sdf', GREEN_FILL, GREEN_LINE, 8.8)
    pill(slide, 3.40, 3.50, 1.12, 0.22, 'write_sdc', ORANGE_FILL, ORANGE_LINE, 8.8)

    round_box(slide, 5.08, 1.60, 3.50, 2.18, WHITE, PANEL, 1.0)
    textbox(slide, 5.34, 1.80, 2.20, 0.24, 'GLS 语义', size=16.0, bold=True)
    textbox(slide, 5.34, 2.06, 2.40, 0.18, 'SDF 给门级网表补单元与连线延迟', size=9.2, color=SUB)
    card(slide, 5.34, 2.42, 1.06, 0.94, 'sim', 'VCS / ModelSim', BLUE_FILL, BLUE_LINE, body_size=10.0, label_w=0.44, align=PP_ALIGN.CENTER)
    card(slide, 6.54, 2.42, 0.94, 0.94, 'hook', '$sdf_annotate', PURPLE_FILL, PURPLE_LINE, body_size=9.6, label_w=0.50, align=PP_ALIGN.CENTER)
    card(slide, 7.60, 2.42, 0.82, 0.94, 'goal', 'RTL vs\nnetlist', ORANGE_FILL, ORANGE_LINE, body_size=10.0, label_w=0.48, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(6.40, 2.90), (6.54, 2.90)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(7.48, 2.90), (7.60, 2.90)], PURPLE_LINE, 1.8)
    pill(slide, 5.38, 3.50, 1.00, 0.22, 'function keep', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 6.54, 3.50, 1.04, 0.22, 'timing aware', GREEN_FILL, GREEN_LINE, 8.8)
    pill(slide, 7.74, 3.50, 0.74, 0.22, 'same tests', ORANGE_FILL, ORANGE_LINE, 8.8)

    round_box(slide, 8.84, 1.60, 3.80, 2.18, WHITE, PANEL, 1.0)
    textbox(slide, 9.10, 1.80, 2.40, 0.24, '当前 SDF 反映了什么', size=16.0, bold=True)
    textbox(slide, 9.10, 2.06, 2.60, 0.18, '直接来自当前 DC pre-layout 输出头部', size=9.2, color=SUB)
    metric(slide, 9.14, 2.42, 0.98, 0.54, 'corner', 'TT\n1.8V 25C', BLUE_FILL, BLUE_LINE)
    metric(slide, 10.24, 2.42, 0.98, 0.54, 'tool', 'DC Ultra', GREEN_FILL, GREEN_LINE)
    metric(slide, 11.34, 2.42, 0.98, 0.54, 'timescale', '1ns', ORANGE_FILL, ORANGE_LINE)
    metric(slide, 9.14, 3.02, 0.98, 0.54, 'design', 'tpu\nSAW=2', PURPLE_FILL, PURPLE_LINE)
    metric(slide, 10.24, 3.02, 0.98, 0.54, 'mode', 'wire\nload', BLUE_FILL, BLUE_LINE)
    metric(slide, 11.34, 3.02, 0.98, 0.54, 'interconnect', '0.000', ORANGE_FILL, ORANGE_LINE)
    pill(slide, 9.18, 3.50, 1.18, 0.22, 'not post-layout RC', PURPLE_FILL, PURPLE_LINE, 8.8)
    pill(slide, 10.52, 3.50, 0.98, 0.22, 'good for GLS', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 11.66, 3.50, 0.72, 0.22, 'not sign-off', ORANGE_FILL, ORANGE_LINE, 8.8)

    round_box(slide, 0.64, 4.10, 5.40, 2.22, WHITE, PANEL, 1.0)
    textbox(slide, 0.90, 4.30, 2.40, 0.24, '面试时怎么讲', size=16.0, bold=True)
    textbox(slide, 0.90, 4.56, 2.80, 0.18, '讲清当前能做什么，不能把 pre-layout 说成 sign-off', size=9.2, color=SUB)
    card(slide, 0.94, 4.92, 1.56, 0.90, 'current', '已有 gate netlist + SDF + SDC\n可做 GLS handoff', BLUE_FILL, BLUE_LINE, body_size=9.4, label_w=0.62, align=PP_ALIGN.CENTER)
    card(slide, 2.72, 4.92, 1.34, 0.90, 'do not overclaim', 'ideal clock\nno CTS\nno SI', ORANGE_FILL, ORANGE_LINE, body_size=9.6, label_w=0.92, align=PP_ALIGN.CENTER)
    card(slide, 4.28, 4.92, 1.50, 0.90, 'answer line', '综合后仿真可做\nsign-off 还要 P&R', GREEN_FILL, GREEN_LINE, body_size=9.4, label_w=0.82, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(2.50, 5.37), (2.72, 5.37)], BLUE_LINE, 1.7)
    poly_arrow(slide, [(4.06, 5.37), (4.28, 5.37)], ORANGE_LINE, 1.7)
    pill(slide, 1.02, 5.94, 1.10, 0.22, 'honest boundary', PURPLE_FILL, PURPLE_LINE, 8.8)
    pill(slide, 2.34, 5.94, 1.12, 0.22, 'artifact exists', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 3.68, 5.94, 1.12, 0.22, 'route not done', ORANGE_FILL, ORANGE_LINE, 8.8)

    round_box(slide, 6.28, 4.10, 6.36, 2.22, WHITE, PANEL, 1.0)
    textbox(slide, 6.54, 4.30, 2.40, 0.24, 'sign-off 下一步', size=16.0, bold=True)
    textbox(slide, 6.54, 4.56, 3.30, 0.18, '从当前 pre-layout GLS 走到真正实现闭环的顺序', size=9.2, color=SUB)
    card(slide, 6.58, 4.94, 1.00, 0.92, '1', 'SS corner\n1.62V 125C', BLUE_FILL, BLUE_LINE, body_size=9.8, label_w=0.28, align=PP_ALIGN.CENTER)
    card(slide, 7.80, 4.94, 1.00, 0.92, '2', 'P&R\nCTS / route', GREEN_FILL, GREEN_LINE, body_size=9.8, label_w=0.28, align=PP_ALIGN.CENTER)
    card(slide, 9.02, 4.94, 1.00, 0.92, '3', 'post-layout\nSTA', ORANGE_FILL, ORANGE_LINE, body_size=9.8, label_w=0.28, align=PP_ALIGN.CENTER)
    card(slide, 10.24, 4.94, 1.08, 0.92, '4', 'VCD +\nPT PX', PURPLE_FILL, PURPLE_LINE, body_size=9.8, label_w=0.28, align=PP_ALIGN.CENTER)
    card(slide, 11.56, 4.94, 0.66, 0.92, '5', 'ECO', BLUE_FILL, BLUE_LINE, body_size=9.6, label_w=0.28, align=PP_ALIGN.CENTER)
    poly_arrow(slide, [(7.58, 5.40), (7.80, 5.40)], BLUE_LINE, 1.8)
    poly_arrow(slide, [(8.80, 5.40), (9.02, 5.40)], GREEN_LINE, 1.8)
    poly_arrow(slide, [(10.02, 5.40), (10.24, 5.40)], ORANGE_LINE, 1.8)
    poly_arrow(slide, [(11.32, 5.40), (11.56, 5.40)], PURPLE_LINE, 1.8)
    pill(slide, 6.70, 5.98, 1.10, 0.22, 'slow corner', BLUE_FILL, BLUE_LINE, 8.8)
    pill(slide, 8.06, 5.98, 1.00, 0.22, 'real RC', GREEN_FILL, GREEN_LINE, 8.8)
    pill(slide, 9.34, 5.98, 1.06, 0.22, 'sign-off STA', ORANGE_FILL, ORANGE_LINE, 8.8)
    pill(slide, 10.66, 5.98, 1.02, 0.22, 'accurate power', PURPLE_FILL, PURPLE_LINE, 8.8)

    textbox(slide, 0.56, 6.82, 6.70, 0.18, '讲法：当前已经有 netlist / SDF / SDC，可做 GLS 交接；但它是 pre-layout，真正 sign-off 要靠 P&R 后的 RC 与 STA。', size=9.4, color=SUB)
    textbox(slide, 8.02, 6.82, 4.62, 0.18, '来源：tpu_syn.v / tpu_syn.sdf / SYN_INTERVIEW_QA / SYN_INTERNALIZE', size=9.4, color=SUB, align=PP_ALIGN.RIGHT)

    prs.save(OUT_PPT)
    print(OUT_PPT)


if __name__ == '__main__':
    main()
