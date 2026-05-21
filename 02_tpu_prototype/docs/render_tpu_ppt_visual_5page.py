from pathlib import Path
import shutil

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

ROOT = Path('/home/jjt/tpu-soc')
DOCS = ROOT / 'docs'
DROP = Path('/tmp/VMwareDnD/71NKe5')
SHARE = Path('/mnt/hgfs/wdchenaic/tpu_interview_ppt')

COVER_IMG = DOCS / 'tpu_cover_16x9_zh.png'
STACK_IMG = DOCS / 'tpu_project_stack_min_16x9_zh.png'
RTL_IMG = DOCS / 'tpu_rtl_core_arch_16x9_zh.png'
LOGIC_IMG = DOCS / 'tpu_rtl_logic_closeups_16x9_zh.png'
RESULT_IMG = DOCS / 'tpu_results_16x9_zh.png'
PPT_OUT = DOCS / 'tpu_project_interview_5page_zh.pptx'

DROP_PPT = DROP / 'tpu_project_interview_5page_zh.pptx'
DROP_MAIN_IMG = DROP / 'tpu.png'
DROP_RTL_IMG = DROP / 'tpu_rtl.png'
DROP_LOGIC_IMG = DROP / 'tpu_logic.png'

AXI_FRONTEND = ROOT / 'src_axi' / 'tpu_frontend_axil.sv'

W, H = 1920, 1080
BG = (250, 248, 244)
PAPER = (255, 255, 255)
INK = (28, 40, 58)
MUTED = (104, 116, 129)
LINE = (222, 227, 233)
ACCENT = (32, 107, 138)
ACCENT2 = (198, 110, 62)
ACCENT3 = (83, 140, 102)
TAG_BG = (236, 243, 247)
TAG_BG2 = (248, 238, 230)
TAG_BG3 = (234, 243, 237)
HILITE = (240, 246, 249)

TITLE_FONT_PATH = '/usr/share/fonts/google-noto/NotoSansSC-Medium.otf'
BODY_FONT_PATH = '/usr/share/fonts/google-noto/NotoSansSC-Regular.otf'
MONO_FONT_PATH = '/usr/share/fonts/dejavu/DejaVuSansMono.ttf'

TITLE_72 = ImageFont.truetype(TITLE_FONT_PATH, 72)
TITLE_58 = ImageFont.truetype(TITLE_FONT_PATH, 58)
TITLE_48 = ImageFont.truetype(TITLE_FONT_PATH, 48)
TITLE_40 = ImageFont.truetype(TITLE_FONT_PATH, 40)
TITLE_34 = ImageFont.truetype(TITLE_FONT_PATH, 34)
TITLE_28 = ImageFont.truetype(TITLE_FONT_PATH, 28)
BODY_28 = ImageFont.truetype(BODY_FONT_PATH, 28)
BODY_24 = ImageFont.truetype(BODY_FONT_PATH, 24)
BODY_22 = ImageFont.truetype(BODY_FONT_PATH, 22)
BODY_20 = ImageFont.truetype(BODY_FONT_PATH, 20)
BODY_18 = ImageFont.truetype(BODY_FONT_PATH, 18)
MONO_22 = ImageFont.truetype(MONO_FONT_PATH, 22)
MONO_20 = ImageFont.truetype(MONO_FONT_PATH, 20)
MONO_18 = ImageFont.truetype(MONO_FONT_PATH, 18)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rounded(draw, xy, radius=28, fill=PAPER, outline=LINE, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def label(draw, pos, msg, font, fill=INK, spacing=6, align='left'):
    draw.multiline_text(pos, msg, font=font, fill=fill, spacing=spacing, align=align)


def center(draw, cx, cy, msg, font, fill=INK, spacing=6):
    bbox = draw.multiline_textbbox((0, 0), msg, font=font, spacing=spacing, align='center')
    w = bbox[2] - bbox[0]
    h = bbox[3] - bbox[1]
    draw.multiline_text((cx - w / 2, cy - h / 2), msg, font=font, fill=fill, spacing=spacing, align='center')


def tag(draw, xy, msg, fill=TAG_BG, text_fill=ACCENT, font=BODY_20):
    rounded(draw, xy, radius=18, fill=fill, outline=fill, width=1)
    l, t, r, b = xy
    center(draw, (l + r) / 2, (t + b) / 2, msg, font, fill=text_fill)


def arrow(draw, pts, fill=ACCENT, width=5, arrow_size=15):
    draw.line(pts, fill=fill, width=width)
    x1, y1 = pts[-2]
    x2, y2 = pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        d = 1 if x2 > x1 else -1
        tri = [(x2, y2), (x2 - d * arrow_size, y2 - arrow_size // 2), (x2 - d * arrow_size, y2 + arrow_size // 2)]
    else:
        d = 1 if y2 > y1 else -1
        tri = [(x2, y2), (x2 - arrow_size // 2, y2 - d * arrow_size), (x2 + arrow_size // 2, y2 - d * arrow_size)]
    draw.polygon(tri, fill=fill)


def snippet_from_segments(path, segments):
    lines = path.read_text().splitlines()
    out = []
    for idx, (start, end) in enumerate(segments):
        for no in range(start, end + 1):
            out.append((no, lines[no - 1]))
        if idx != len(segments) - 1:
            out.append((None, '...'))
    return out


def draw_code_card(draw, xy, title, subtitle, snippet_lines, highlights, tone=ACCENT):
    rounded(draw, xy, radius=32, fill=PAPER, outline=LINE, width=2)
    l, t, r, b = xy
    draw.rounded_rectangle((l + 22, t + 24, l + 30, b - 24), radius=4, fill=tone)
    label(draw, (l + 56, t + 26), title, TITLE_34, fill=INK)
    label(draw, (l + 56, t + 76), subtitle, BODY_18, fill=MUTED)
    y = t + 126
    line_h = 34
    for no, code in snippet_lines:
        is_hl = no in highlights
        if is_hl:
            rounded(draw, (l + 44, y - 4, r - 28, y + 24), radius=12, fill=HILITE, outline=HILITE, width=1)
        ln = '   ...' if no is None else f'{no:>4} '
        label(draw, (l + 58, y), ln, MONO_18, fill=ACCENT2 if is_hl else MUTED)
        label(draw, (l + 126, y), code.replace('\t', '    '), MONO_18, fill=INK)
        y += line_h
        if y > b - 30:
            break


def draw_slide_header(draw, page_no, title, subtitle):
    label(draw, (96, 78), f'P{page_no}', BODY_18, fill=ACCENT)
    label(draw, (96, 112), title, TITLE_48, fill=INK)
    label(draw, (96, 176), subtitle, BODY_22, fill=MUTED)
    rounded(draw, (1460, 72, 1810, 128), radius=24, fill=TAG_BG, outline=TAG_BG, width=1)
    center(draw, 1635, 100, f'5 页 PPT 的第 {page_no} 页', BODY_20, fill=ACCENT)


def render_cover_image():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)

    draw.rectangle((0, 0, 26, H), fill=ACCENT)
    rounded(draw, (1360, 90, 1800, 910), radius=38, fill=(244, 237, 229), outline=(244, 237, 229), width=1)

    label(draw, (110, 120), 'TinyTPU AXI-Lite SoC', BODY_22, fill=ACCENT)
    label(draw, (110, 190), '训练闭环原型', TITLE_72, fill=INK)
    label(draw, (110, 304), '从裸 RTL 到寄存器可控、指令可控的 SoC', BODY_28, fill=MUTED)

    tag(draw, (114, 420, 330, 470), 'AXI-Lite', fill=TAG_BG, text_fill=ACCENT)
    tag(draw, (354, 420, 606, 470), 'IMEM + Sequencer', fill=TAG_BG2, text_fill=ACCENT2)
    tag(draw, (630, 420, 846, 470), '41 / 41 PASS', fill=TAG_BG3, text_fill=ACCENT3)
    tag(draw, (870, 420, 1038, 470), 'XOR', fill=TAG_BG, text_fill=ACCENT)

    label(draw, (110, 562), '核心改动', BODY_22, fill=ACCENT)
    label(draw, (110, 612), '补齐控制链路\n跑通训练闭环', TITLE_58, fill=INK, spacing=12)
    label(draw, (114, 826), '这套 PPT 只讲三件事：系统闭环、RTL 控制、验证收敛。', BODY_22, fill=MUTED)

    metric_titles = ['控制接口', '执行组织', 'RTL 重点', '验证结果']
    metric_values = [('AXI-Lite', ACCENT), ('IMEM + seq', ACCENT2), ('pulse / hold', ACCENT3), ('loss 收敛', ACCENT)]
    y = 168
    for idx, ((value, tone), title) in enumerate(zip(metric_values, metric_titles)):
        rounded(draw, (1412, y, 1750, y + 118), radius=26, fill=PAPER, outline=PAPER, width=1)
        label(draw, (1450, y + 22), title, BODY_20, fill=MUTED)
        label(draw, (1450, y + 56), value, TITLE_34, fill=tone)
        y += 154

    img.save(COVER_IMG)


def render_stack_image():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_slide_header(draw, 2, '项目闭环', '软件 -> SoC 控制 -> 计算核心 -> 验证结果')

    cols = [
        ((96, 270, 424, 792), '软件与驱动', ['scheduler', 'encoder', 'cocotb'], ACCENT, TAG_BG),
        ((490, 230, 870, 832), 'SoC 控制', ['AXI-Lite', '寄存器', 'IMEM', 'sequencer'], ACCENT2, TAG_BG2),
        ((936, 230, 1316, 832), '计算核心', ['UB', 'Systolic', 'VPU', '写回路径'], ACCENT3, TAG_BG3),
        ((1382, 270, 1788, 792), '验证与结果', ['41/41 PASS', '12 epoch', 'loss 收敛'], ACCENT, TAG_BG),
    ]

    for xy, title, items, tone, fill in cols:
        rounded(draw, xy, radius=36, fill=PAPER, outline=LINE, width=2)
        l, t, r, b = xy
        draw.rounded_rectangle((l + 24, t + 26, l + 34, b - 26), radius=4, fill=tone)
        label(draw, (l + 58, t + 42), title, TITLE_40, fill=INK)
        y = t + 150
        for item in items:
            tag(draw, (l + 54, y, r - 54, y + 54), item, fill=fill, text_fill=tone)
            y += 92

    arrow(draw, [(424, 518), (490, 518)], fill=ACCENT)
    arrow(draw, [(870, 518), (936, 518)], fill=ACCENT2)
    arrow(draw, [(1316, 518), (1382, 518)], fill=ACCENT3)
    tag(draw, (610, 892, 836, 944), '指令 + 控制', fill=TAG_BG2, text_fill=ACCENT2)
    tag(draw, (1040, 892, 1264, 944), '数据 + 写回', fill=TAG_BG3, text_fill=ACCENT3)
    label(draw, (96, 980), '建议讲法：先讲闭环，再下钻到 RTL 图和代码特写。', BODY_20, fill=MUTED)

    img.save(STACK_IMG)
    shutil.copyfile(STACK_IMG, DROP_MAIN_IMG)


def render_rtl_arch_image():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_slide_header(draw, 3, 'RTL 架构', '真实模块层次：tpu_soc -> frontend + tpu -> UB / Systolic / VPU')

    frontend = (96, 250, 380, 430)
    seq = (430, 220, 888, 430)
    ub = (126, 520, 610, 934)
    sysa = (968, 320, 1560, 560)
    vpu = (968, 630, 1560, 934)
    obs = (1630, 392, 1816, 876)

    rounded(draw, frontend, radius=30)
    label(draw, (132, 286), 'AXI-Lite Frontend', TITLE_28, fill=INK)
    label(draw, (132, 334), 'tpu_frontend_axil\n寄存器写入 / 状态读回 / UB_PUSH', BODY_20, fill=MUTED)

    rounded(draw, seq, radius=30)
    label(draw, (468, 258), 'IMEM + Sequencer + control_unit', TITLE_34, fill=INK)
    label(draw, (468, 314), 'step/start -> instr -> seq_instr_pulse\nwait_after / vpu_drain / pathway 保持', BODY_20, fill=MUTED)
    tag(draw, (468, 380, 664, 428), 'busy_reg', fill=TAG_BG, text_fill=ACCENT)
    tag(draw, (686, 380, 810, 428), 'pc', fill=TAG_BG2, text_fill=ACCENT2)

    rounded(draw, ub, radius=34)
    label(draw, (170, 566), 'Unified Buffer', TITLE_40, fill=INK)
    label(draw, (170, 624), 'ub_inst\nhost load / read ctrl / writeback', BODY_22, fill=MUTED)
    tag(draw, (170, 752, 378, 804), 'input / weight', fill=TAG_BG, text_fill=ACCENT)
    tag(draw, (170, 824, 378, 876), 'bias / Y / H', fill=TAG_BG3, text_fill=ACCENT3)
    tag(draw, (170, 896, 378, 948), 'grad writeback', fill=TAG_BG2, text_fill=ACCENT2)

    rounded(draw, sysa, radius=34)
    center(draw, 1264, 418, 'Systolic Array\nsystolic_inst (2x2)', TITLE_40, fill=INK)
    tag(draw, (1110, 498, 1420, 548), 'sys_weight_in / sys_data_in', fill=TAG_BG, text_fill=ACCENT)

    rounded(draw, vpu, radius=34)
    label(draw, (1184, 666), 'VPU / vpu_inst', TITLE_34, fill=INK)
    blocks = [
        ((1052, 748, 1220, 836), 'Bias', TAG_BG, ACCENT),
        ((1246, 748, 1414, 836), 'Loss', TAG_BG2, ACCENT2),
        ((1052, 846, 1220, 934), 'LReLU', TAG_BG3, ACCENT3),
        ((1246, 846, 1414, 934), 'dLReLU', (239, 236, 247), ACCENT2),
    ]
    for xy, msg, fill, tone in blocks:
        rounded(draw, xy, radius=22, fill=fill, outline=LINE, width=2)
        center(draw, (xy[0] + xy[2]) / 2, (xy[1] + xy[3]) / 2, msg, TITLE_28, fill=tone)

    rounded(draw, obs, radius=30)
    center(draw, 1723, 446, '输出观测', TITLE_28, fill=INK)
    tag(draw, (1664, 532, 1782, 580), 'sys_out', fill=TAG_BG, text_fill=ACCENT)
    tag(draw, (1664, 608, 1782, 656), 'vpu_out', fill=TAG_BG2, text_fill=ACCENT2)
    tag(draw, (1664, 684, 1782, 732), 'valid', fill=TAG_BG3, text_fill=ACCENT3)

    arrow(draw, [(380, 340), (430, 340)], fill=ACCENT)
    arrow(draw, [(238, 430), (238, 520)], fill=ACCENT2)
    arrow(draw, [(668, 430), (668, 490), (382, 490), (382, 520)], fill=ACCENT2)
    arrow(draw, [(888, 340), (1010, 340), (1010, 630)], fill=ACCENT)
    arrow(draw, [(610, 658), (968, 658), (968, 446)], fill=ACCENT)
    arrow(draw, [(610, 756), (968, 756)], fill=ACCENT2)
    arrow(draw, [(1264, 560), (1264, 630)], fill=ACCENT3)
    arrow(draw, [(968, 904), (842, 904), (842, 934), (610, 934)], fill=ACCENT)
    arrow(draw, [(1560, 448), (1630, 448)], fill=ACCENT3)
    arrow(draw, [(1560, 782), (1630, 782)], fill=ACCENT2)

    tag(draw, (640, 624, 866, 674), 'ub_rd_*', fill=TAG_BG, text_fill=ACCENT)
    tag(draw, (650, 724, 910, 774), 'bias / Y / H', fill=TAG_BG3, text_fill=ACCENT3)
    tag(draw, (1080, 578, 1448, 628), 'sys_data_out -> vpu_data_in', fill=TAG_BG2, text_fill=ACCENT2)
    tag(draw, (688, 960, 1064, 1010), 'vpu_data_out -> UB writeback', fill=(239, 236, 247), text_fill=ACCENT2)
    label(draw, (96, 986), '这页重点讲模块关系和数据回路，不硬讲不存在的 clock-gating cell。', BODY_20, fill=MUTED)

    img.save(RTL_IMG)
    shutil.copyfile(RTL_IMG, DROP_RTL_IMG)


def render_logic_closeups():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_slide_header(draw, 4, '关键 RTL 特写', '只保留 3 处最有价值的控制逻辑。')

    wait_lines = snippet_from_segments(AXI_FRONTEND, [(139, 139), (153, 154), (177, 179)])
    path_lines = snippet_from_segments(AXI_FRONTEND, [(156, 157), (403, 403)])
    mux_lines = snippet_from_segments(AXI_FRONTEND, [(405, 408)])

    draw_code_card(draw, (74, 220, 916, 598), '1. wait_after + vpu_drain', 'src_axi/tpu_frontend_axil.sv', wait_lines, {139, 153, 177, 178, 179}, tone=ACCENT)
    draw_code_card(draw, (1002, 220, 1846, 598), '2. pathway 保持', 'src_axi/tpu_frontend_axil.sv', path_lines, {156, 157, 403}, tone=ACCENT2)
    draw_code_card(draw, (74, 646, 1846, 944), '3. host write mux', 'src_axi/tpu_frontend_axil.sv', mux_lines, {405, 406, 407, 408}, tone=ACCENT3)

    label(draw, (96, 980), '如果被问到 clock-gating：当前版本先把控制正确性和训练闭环做实，低功耗门控不是这版主线。', BODY_20, fill=MUTED)

    img.save(LOGIC_IMG)
    shutil.copyfile(LOGIC_IMG, DROP_LOGIC_IMG)


def render_result_image():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_slide_header(draw, 5, '结果与边界', '结果先报，边界主动讲。')

    metrics = [
        ((96, 252, 428, 400), '回归', '41 / 41 PASS', ACCENT),
        ((466, 252, 798, 400), '训练', '12 epoch', ACCENT2),
        ((836, 252, 1206, 400), 'loss', '0.2529 -> 0.1777', ACCENT3),
        ((1244, 252, 1586, 400), 'XOR', '(0, 1, 1, 0)', ACCENT),
    ]
    for xy, title, value, tone in metrics:
        rounded(draw, xy, radius=28)
        l, t, r, b = xy
        label(draw, (l + 30, t + 26), title, BODY_20, fill=MUTED)
        label(draw, (l + 30, t + 72), value, TITLE_34, fill=tone)

    rounded(draw, (96, 488, 908, 760), radius=34)
    label(draw, (136, 530), '一句话结果', BODY_22, fill=ACCENT)
    label(draw, (136, 582), '单次回归全部通过\n多 epoch loss 持续下降', TITLE_40, fill=INK, spacing=12)

    rounded(draw, (1012, 488, 1824, 760), radius=34)
    label(draw, (1052, 530), '边界口径', BODY_22, fill=ACCENT2)
    label(draw, (1052, 582), '当前是 tiny-tpu 原型\n不宣称 DMA / IRQ / clock-gating', TITLE_34, fill=INK, spacing=10)

    tag(draw, (140, 860, 340, 912), '系统闭环', fill=TAG_BG, text_fill=ACCENT)
    tag(draw, (368, 860, 584, 912), 'RTL 控制', fill=TAG_BG2, text_fill=ACCENT2)
    tag(draw, (612, 860, 844, 912), '验证收敛', fill=TAG_BG3, text_fill=ACCENT3)
    tag(draw, (872, 860, 1140, 912), '边界主动说明', fill=(239, 236, 247), text_fill=ACCENT2)
    label(draw, (96, 976), '结尾建议：项目范围克制，但闭环做实。', BODY_20, fill=MUTED)

    img.save(RESULT_IMG)


def add_full_image_slide(prs, img_path):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.shapes.add_picture(str(img_path), 0, 0, width=SLIDE_W, height=SLIDE_H)


def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = 'TinyTPU AXI-Lite SoC 面试 PPT'
    prs.core_properties.subject = '5 页少字多图版'
    prs.core_properties.author = 'Codex'

    for img_path in [COVER_IMG, STACK_IMG, RTL_IMG, LOGIC_IMG, RESULT_IMG]:
        add_full_image_slide(prs, img_path)

    PPT_OUT.parent.mkdir(parents=True, exist_ok=True)
    DROP.mkdir(parents=True, exist_ok=True)
    prs.save(str(PPT_OUT))
    shutil.copyfile(PPT_OUT, DROP_PPT)


def main():
    render_cover_image()
    render_stack_image()
    render_rtl_arch_image()
    render_logic_closeups()
    render_result_image()
    build_ppt()


if __name__ == '__main__':
    main()
