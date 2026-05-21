from pathlib import Path
import shutil

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path("/home/jjt/tpu-soc")
DOCS = ROOT / "docs"
DROP_DIR = Path("/tmp/VMwareDnD/71NKe5")
OUTPUT = DOCS / "tpu_project_interview_5page_zh.pptx"
DROP_OUTPUT = DROP_DIR / "tpu_project_interview_5page_zh.pptx"
ARCH_IMG = DOCS / "tpu_project_architecture_16x9_zh.png"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(248, 245, 237)
INK = RGBColor(27, 38, 59)
MUTED = RGBColor(97, 110, 126)
ACCENT = RGBColor(22, 113, 135)
ACCENT_SOFT = RGBColor(224, 240, 242)
ACCENT_WARM = RGBColor(198, 92, 53)
ACCENT_GOLD = RGBColor(212, 162, 55)
CARD = RGBColor(255, 252, 245)
LINE = RGBColor(213, 219, 224)

FONT = "Microsoft YaHei"
FONT_BOLD = "Microsoft YaHei"


def set_slide_bg(slide, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        0,
        0,
        SLIDE_W,
        SLIDE_H,
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    slide.shapes._spTree.remove(shape._element)
    slide.shapes._spTree.insert(2, shape._element)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size,
    color=INK,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = valign
    p = frame.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.name = FONT_BOLD if bold else FONT
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(
    slide,
    left,
    top,
    width,
    height,
    bullets,
    font_size=18,
    color=INK,
    line_spacing=1.15,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    for idx, item in enumerate(bullets):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(7)
        p.line_spacing = line_spacing
        p.bullet = True
        for run in p.runs:
            run.font.name = FONT
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
    return box


def add_card(slide, left, top, width, height, title, body_lines, accent=ACCENT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.2)

    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, left, top, Inches(0.10), height)
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()

    add_textbox(slide, left + Inches(0.22), top + Inches(0.16), width - Inches(0.34), Inches(0.38), title, 18, bold=True)
    add_bullets(
        slide,
        left + Inches(0.20),
        top + Inches(0.60),
        width - Inches(0.30),
        height - Inches(0.70),
        body_lines,
        font_size=14,
        color=MUTED,
        line_spacing=1.05,
    )
    return shape


def add_metric(slide, left, top, width, height, title, value, accent=ACCENT):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT_SOFT
    shape.line.fill.background()
    add_textbox(slide, left + Inches(0.18), top + Inches(0.14), width - Inches(0.36), Inches(0.22), title, 12, color=MUTED)
    add_textbox(slide, left + Inches(0.18), top + Inches(0.34), width - Inches(0.36), Inches(0.36), value, 21, color=accent, bold=True)
    return shape


def add_section_title(slide, page_no, title, subtitle):
    add_textbox(slide, Inches(0.62), Inches(0.40), Inches(0.68), Inches(0.30), f"P{page_no}", 14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.62), Inches(0.72), Inches(6.30), Inches(0.58), title, 28, color=INK, bold=True)
    add_textbox(slide, Inches(0.62), Inches(1.22), Inches(8.60), Inches(0.32), subtitle, 14, color=MUTED)

    line = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.62), Inches(1.58), Inches(12.10), Inches(0.02))
    line.fill.solid()
    line.fill.fore_color.rgb = LINE
    line.line.fill.background()


def add_footer(slide, text):
    add_textbox(slide, Inches(0.62), Inches(7.00), Inches(12.0), Inches(0.22), text, 10, color=MUTED)


def fit_picture(slide, img_path, left, top, width, height):
    with Image.open(img_path) as img:
        img_w, img_h = img.size
    box_ratio = width / height
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        pic_w = width
        pic_h = width / img_ratio
        pic_left = left
        pic_top = top + (height - pic_h) / 2
    else:
        pic_h = height
        pic_w = height * img_ratio
        pic_top = top
        pic_left = left + (width - pic_w) / 2
    slide.shapes.add_picture(str(img_path), pic_left, pic_top, width=pic_w, height=pic_h)


def add_flow_box(slide, left, top, width, height, title, detail, accent):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = CARD
    shape.line.color.rgb = accent
    shape.line.width = Pt(1.5)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.16), width - Inches(0.32), Inches(0.34), title, 17, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, left + Inches(0.16), top + Inches(0.58), width - Inches(0.32), Inches(0.52), detail, 12, color=MUTED, align=PP_ALIGN.CENTER)
    return shape


def add_arrow(slide, x1, y1, x2, y2, color=ACCENT):
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line.line.color.rgb = color
    line.line.width = Pt(2.0)
    line.line.end_arrowhead = True


def build_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)

    accent_band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, Inches(0.0), Inches(0.26), SLIDE_H)
    accent_band.fill.solid()
    accent_band.fill.fore_color.rgb = ACCENT
    accent_band.line.fill.background()

    warm_block = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(8.95), Inches(0.72), Inches(3.65), Inches(5.65))
    warm_block.fill.solid()
    warm_block.fill.fore_color.rgb = RGBColor(241, 232, 219)
    warm_block.line.fill.background()

    add_textbox(slide, Inches(0.76), Inches(0.86), Inches(6.6), Inches(0.40), "TinyTPU AXI-Lite SoC", 14, color=ACCENT, bold=True)
    add_textbox(slide, Inches(0.76), Inches(1.34), Inches(7.4), Inches(1.10), "从裸 RTL 到训练闭环原型", 28, color=INK, bold=True)
    add_textbox(
        slide,
        Inches(0.76),
        Inches(2.48),
        Inches(7.3),
        Inches(0.72),
        "系统集成 + 指令化执行 + cocotb 闭环验证",
        17,
        color=MUTED,
    )

    add_bullets(
        slide,
        Inches(0.76),
        Inches(3.34),
        Inches(7.1),
        Inches(2.2),
        [
            "把原始 tiny-tpu 从 testbench 直驱补成 SoC 原型",
            "补齐 AXI-Lite 前端、IMEM、sequencer 和寄存器控制链路",
            "验证范围覆盖 2x2 / Q8.8 / 2-layer MLP / XOR",
        ],
        font_size=18,
        color=INK,
    )

    add_metric(slide, Inches(9.24), Inches(1.10), Inches(3.05), Inches(0.88), "控制接口", "AXI-Lite")
    add_metric(slide, Inches(9.24), Inches(2.15), Inches(3.05), Inches(0.88), "执行组织", "IMEM + Sequencer", accent=ACCENT_WARM)
    add_metric(slide, Inches(9.24), Inches(3.20), Inches(3.05), Inches(0.88), "回归结果", "41 / 41 PASS", accent=ACCENT_GOLD)
    add_metric(slide, Inches(9.24), Inches(4.25), Inches(3.05), Inches(0.88), "训练闭环", "XOR 收敛", accent=ACCENT)

    add_textbox(slide, Inches(0.76), Inches(6.55), Inches(12.0), Inches(0.28), "重点不是单个算子，而是把控制链路和验证闭环补完整。", 12, color=MUTED)
    add_footer(slide, "P1  项目定义")


def build_architecture(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, 2, "整体架构", "软件驱动 -> SoC 控制 -> 计算核心 -> 验证结果")

    tag = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(9.38), Inches(0.58), Inches(3.08), Inches(0.48))
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT_SOFT
    tag.line.fill.background()
    add_textbox(slide, Inches(9.58), Inches(0.70), Inches(2.70), Inches(0.20), "整条闭环，不只是 SoC 连线", 11, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)

    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.70), Inches(1.92), Inches(11.92), Inches(4.72))
    frame.fill.solid()
    frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
    frame.line.color.rgb = LINE
    frame.line.width = Pt(1.2)

    fit_picture(slide, ARCH_IMG, Inches(0.82), Inches(2.02), Inches(11.68), Inches(4.50))
    add_footer(slide, "P2  用这张图讲项目全链路：从软件输入到验证收敛。")


def build_control_chain(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, 3, "控制与执行链路", "AXI-Lite + IMEM + Sequencer 如何驱动 TPU")

    box_w = Inches(2.55)
    box_h = Inches(1.18)
    top = Inches(2.05)
    x_positions = [Inches(0.78), Inches(3.45), Inches(6.12), Inches(8.79)]
    titles = ["AXI-Lite 前端", "寄存器 / IMEM", "Sequencer", "TPU 控制 / 执行"]
    details = [
        "写寄存器\n读状态\n装载 IMEM",
        "控制位配置\n指令缓存\n可被 host 驱动",
        "逐条取指\n生成 pulse\n显式等待",
        "UB / SA / VPU\n参数更新\n结果回读",
    ]
    accents = [ACCENT, ACCENT_WARM, ACCENT_GOLD, ACCENT]

    for left, title, detail, accent in zip(x_positions, titles, details, accents):
        add_flow_box(slide, left, top, box_w, box_h, title, detail, accent)

    center_y = top + box_h / 2
    for idx in range(3):
        add_arrow(slide, x_positions[idx] + box_w, center_y, x_positions[idx + 1], center_y)

    add_card(
        slide,
        Inches(0.86),
        Inches(4.02),
        Inches(3.72),
        Inches(1.85),
        "关键控制点",
        [
            "AXI 写寄存器 / IMEM，不再依赖 testbench 直驱",
            "`wait_after` 明确等待语义，避免隐式时序猜测",
            "`seq_instr_pulse` 把当前指令打一拍送入执行路径",
        ],
        accent=ACCENT,
    )
    add_card(
        slide,
        Inches(4.80),
        Inches(4.02),
        Inches(3.72),
        Inches(1.85),
        "容易被追问的点",
        [
            "`vpu_data_pathway` 为什么要保持",
            "为什么要把等待语义显式化",
            "host 写入如何真正进入 UB / datapath",
        ],
        accent=ACCENT_WARM,
    )
    add_card(
        slide,
        Inches(8.74),
        Inches(4.02),
        Inches(3.72),
        Inches(1.85),
        "一句话解释",
        [
            "打通控制链路之后，这个 TPU 才真正变成寄存器可控、指令可控的 SoC 原型。",
        ],
        accent=ACCENT_GOLD,
    )
    add_footer(slide, "P3  建议面试时顺着左到右讲，不要陷进单个 always block。")


def build_implementation(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, 4, "实现细节与关键 Bug", "模块补齐 + 3 个关键问题收敛")

    add_card(
        slide,
        Inches(0.80),
        Inches(1.95),
        Inches(3.35),
        Inches(2.18),
        "这次真正补齐的内容",
        [
            "AXI-Lite 前端与寄存器映射",
            "IMEM 装载与 sequencer 调度",
            "cocotb e2e / 多 epoch 验证闭环",
        ],
        accent=ACCENT,
    )

    add_card(
        slide,
        Inches(0.80),
        Inches(4.42),
        Inches(3.35),
        Inches(1.72),
        "讲法建议",
        [
            "只挑 3 个问题讲现象、根因、修复、验证。",
            "重点讲你如何定位，而不是堆很多现象。",
        ],
        accent=ACCENT_GOLD,
    )

    issue_x = [Inches(4.42), Inches(7.14), Inches(9.86)]
    issue_titles = ["Host write 无效", "pathway 不保持", "wait 逻辑卡死"]
    issue_bodies = [
        ["现象：AXI 已写入，但 UB 不采样", "根因：host write 没有真正打入数据通路", "修复：补齐写使能与采样路径"],
        ["现象：dispatch 后几拍错误 bypass", "根因：选择信号没有跨流水线延迟保持", "修复：保持 `vpu_data_pathway` 直到结果落地"],
        ["现象：状态机停住，后续指令不推进", "根因：错误消费 drain 事件", "修复：重写等待条件与消费时机"],
    ]
    issue_accents = [ACCENT_WARM, ACCENT, ACCENT_GOLD]

    for left, title, body, accent in zip(issue_x, issue_titles, issue_bodies, issue_accents):
        add_card(slide, left, Inches(2.12), Inches(2.36), Inches(3.78), title, body, accent=accent)

    add_footer(slide, "P4  这一页体现工程判断力：能定位、能收敛、能自证修复有效。")


def build_results(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, BG)
    add_section_title(slide, 5, "验证结果与项目边界", "结果说清楚，边界也讲清楚")

    add_metric(slide, Inches(0.82), Inches(2.02), Inches(2.55), Inches(0.95), "单次 e2e", "41 / 41 PASS", accent=ACCENT)
    add_metric(slide, Inches(0.82), Inches(3.17), Inches(2.55), Inches(0.95), "多 epoch", "0.2529 -> 0.1777", accent=ACCENT_WARM)
    add_metric(slide, Inches(0.82), Inches(4.32), Inches(2.55), Inches(0.95), "XOR 输出", "(0, 1, 1, 0)", accent=ACCENT_GOLD)
    add_metric(slide, Inches(0.82), Inches(5.47), Inches(2.55), Inches(0.95), "pipeline 变体", "164.10 -> 183.91 MHz", accent=ACCENT)

    add_card(
        slide,
        Inches(3.80),
        Inches(2.02),
        Inches(4.05),
        Inches(4.40),
        "项目价值",
        [
            "把 tiny-tpu 从裸 RTL 演示补成了可控、可测、可讲清的 SoC 原型",
            "控制链路、执行链路、验证链路形成闭环",
            "面试时可从系统集成、验证方法、问题定位三条线展开",
        ],
        accent=ACCENT,
    )

    add_card(
        slide,
        Inches(8.06),
        Inches(2.02),
        Inches(4.42),
        Inches(4.40),
        "项目边界",
        [
            "当前是 tiny-tpu 原型，不夸大成完整通用 SoC",
            "验证扎实范围：当前网络规模、Q8.8 配置、XOR 训练闭环",
            "尚未覆盖 DMA / IRQ / 更大编译器能力",
        ],
        accent=ACCENT_WARM,
    )

    add_footer(slide, "P5  结尾先报结果，再主动讲边界，会比单纯强调性能更稳。")


def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = "TinyTPU AXI-Lite SoC 面试 PPT"
    prs.core_properties.subject = "5 页项目面试演示"
    prs.core_properties.author = "Codex"

    build_cover(prs)
    build_architecture(prs)
    build_control_chain(prs)
    build_implementation(prs)
    build_results(prs)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    DROP_DIR.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    shutil.copyfile(OUTPUT, DROP_OUTPUT)


if __name__ == "__main__":
    build_presentation()
