from __future__ import annotations

from pathlib import Path
import shutil

from pptx import Presentation
from pptx.util import Inches

ROOT = Path('/home/jjt/tpu-soc')
DOCS = ROOT / 'docs'
PACK = ROOT / 'career' / 'hisilicon_25min_interview_20260410'

SRC_IMAGES = [
    DOCS / 'tpu_hisilicon_clean_rtl_arch_16x9_zh.png',
    DOCS / 'tpu_rtl_core_arch_16x9_zh.png',
]

DST_IMAGES = [
    PACK / '06a_项目级_RTL架构图.png',
    PACK / '06b_核心数据通路_RTL架构图.png',
]

OUT_PPT = PACK / '06_TinyTPU_RTL架构图_2p.pptx'
SLIDE_W = Inches(13.333333)
SLIDE_H = Inches(7.5)


def main() -> None:
    PACK.mkdir(parents=True, exist_ok=True)

    for src, dst in zip(SRC_IMAGES, DST_IMAGES):
        if not src.exists():
            raise FileNotFoundError(src)
        shutil.copy2(src, dst)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    for img in DST_IMAGES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(img), 0, 0, width=SLIDE_W, height=SLIDE_H)

    prs.save(OUT_PPT)


if __name__ == '__main__':
    main()
