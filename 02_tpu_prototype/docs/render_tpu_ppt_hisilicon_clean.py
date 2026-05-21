from pathlib import Path
import shutil

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.util import Inches

ROOT = Path('/home/jjt/tpu-soc')
DOCS = ROOT / 'docs'
DROP = Path('/tmp/VMwareDnD/71NKe5')
SHARE = Path('/mnt/hgfs/wdchenaic/tpu_interview_ppt')

OUT_PPT = DOCS / 'tpu_project_interview_hisilicon_clean_zh.pptx'
DROP_PPT = DROP / 'tpu_project_interview_hisilicon_clean_zh.pptx'
SHARE_PPT = SHARE / 'tpu_project_interview_hisilicon_clean_zh.pptx'

SLIDE1 = DOCS / 'tpu_hisilicon_clean_cover_16x9_zh.png'
SLIDE2 = DOCS / 'tpu_hisilicon_clean_rtl_arch_16x9_zh.png'
SLIDE3 = DOCS / 'tpu_hisilicon_clean_chain_16x9_zh.png'
SLIDE4 = DOCS / 'tpu_hisilicon_clean_difficulty_16x9_zh.png'
SLIDE5 = DOCS / 'tpu_hisilicon_clean_results_16x9_zh.png'

W, H = 1920, 1080
BG = (248, 246, 241)
WHITE = (255, 255, 255)
INK = (26, 38, 56)
MUTED = (102, 113, 126)
LINE = (220, 226, 232)
BLUE = (36, 110, 142)
ORANGE = (196, 109, 63)
GREEN = (88, 143, 107)
PURPLE = (121, 106, 177)
BLUE_BG = (236, 243, 247)
ORANGE_BG = (248, 238, 229)
GREEN_BG = (235, 243, 237)
PURPLE_BG = (240, 237, 248)
WARM_PANEL = (242, 235, 226)
HILITE = (241, 246, 250)

TITLE_FONT = '/usr/share/fonts/google-noto/NotoSansSC-Medium.otf'
BODY_FONT = '/usr/share/fonts/google-noto/NotoSansSC-Regular.otf'
MONO_FONT = '/usr/share/fonts/dejavu/DejaVuSansMono.ttf'

T72 = ImageFont.truetype(TITLE_FONT, 72)
T60 = ImageFont.truetype(TITLE_FONT, 60)
T52 = ImageFont.truetype(TITLE_FONT, 52)
T42 = ImageFont.truetype(TITLE_FONT, 42)
T34 = ImageFont.truetype(TITLE_FONT, 34)
T28 = ImageFont.truetype(TITLE_FONT, 28)
B30 = ImageFont.truetype(BODY_FONT, 30)
B26 = ImageFont.truetype(BODY_FONT, 26)
B24 = ImageFont.truetype(BODY_FONT, 24)
B22 = ImageFont.truetype(BODY_FONT, 22)
B20 = ImageFont.truetype(BODY_FONT, 20)
B18 = ImageFont.truetype(BODY_FONT, 18)
M22 = ImageFont.truetype(MONO_FONT, 22)
M20 = ImageFont.truetype(MONO_FONT, 20)
M18 = ImageFont.truetype(MONO_FONT, 18)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def rounded(draw, xy, radius=28, fill=WHITE, outline=LINE, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def text(draw, pos, msg, font, fill=INK, spacing=6, align='left'):
    draw.multiline_text(pos, msg, font=font, fill=fill, spacing=spacing, align=align)


def center(draw, cx, cy, msg, font, fill=INK, spacing=6):
    bbox = draw.multiline_textbbox((0, 0), msg, font=font, spacing=spacing, align='center')
    w = bbox[2] - bbox[0]
    h = bbox[3] - bbox[1]
    draw.multiline_text((cx - w / 2, cy - h / 2), msg, font=font, fill=fill, spacing=spacing, align='center')


def chip(draw, xy, msg, fill, text_fill, font=B20):
    rounded(draw, xy, radius=18, fill=fill, outline=fill, width=1)
    l, t, r, b = xy
    center(draw, (l + r) / 2, (t + b) / 2, msg, font, fill=text_fill)


def arrow(draw, pts, fill=BLUE, width=5, arrow_size=16):
    draw.line(pts, fill=fill, width=width)
    x1, y1 = pts[-2]
    x2, y2 = pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        d = 1 if x2 > x1 else -1
        tri = [(x2, y2), (x2 - d * arrow_size, y2 - arrow_size // 2), (x2 - d * arrow_size, y2 + arrow_size // 2)]
    else:
        d = 1 if y2 > y1 else -1
        tri = [(x2, y2), (x2 - arrow_size // 2, y2 - d * arrow_size), (x2 + arrow_size // 2, y2 - d * arrow_size)]
    draw.polygon(tri, fill=fill)


def metric_box(draw, xy, title, value, tone):
    rounded(draw, xy, radius=24)
    l, t, r, b = xy
    text(draw, (l + 26, t + 24), title, B20, fill=MUTED)
    text(draw, (l + 26, t + 68), value, T34, fill=tone)


def issue_row(draw, x, y, label_text, body_text, tone):
    text(draw, (x, y), label_text, B18, fill=tone)
    text(draw, (x + 74, y), body_text, B18, fill=INK)


def code_strip(draw, xy, title, ref, code_lines, tone):
    rounded(draw, xy, radius=22, fill=WHITE, outline=LINE, width=2)
    l, t, r, b = xy
    draw.rounded_rectangle((l + 18, t + 20, l + 26, b - 20), radius=4, fill=tone)
    text(draw, (l + 46, t + 16), title, B20, fill=tone)
    text(draw, (r - 180, t + 18), ref, B18, fill=MUTED)
    yy = t + 64
    for code in code_lines:
        rounded(draw, (l + 36, yy - 2, r - 28, yy + 28), radius=12, fill=HILITE, outline=HILITE, width=1)
        text(draw, (l + 52, yy), code, M18, fill=INK)
        yy += 42


def draw_header(draw, idx, title, subtitle):
    text(draw, (88, 72), f'P{idx}', B18, fill=BLUE)
    text(draw, (88, 108), title, T52, fill=INK)
    text(draw, (88, 172), subtitle, B22, fill=MUTED)
    rounded(draw, (1460, 66, 1812, 122), radius=24, fill=BLUE_BG, outline=BLUE_BG, width=1)
    center(draw, 1636, 94, f'5 页 PPT 的第 {idx} 页', B20, fill=BLUE)


def render_cover():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 24, H), fill=BLUE)
    rounded(draw, (1370, 82, 1808, 924), radius=42, fill=WARM_PANEL, outline=WARM_PANEL, width=1)

    text(draw, (96, 110), 'TinyTPU AXI-Lite SoC', B22, fill=BLUE)
    text(draw, (96, 184), '从裸 RTL 到\n训练闭环原型', T72, fill=INK, spacing=12)
    text(draw, (96, 366), '系统集成、控制执行、验证闭环三条线并行展开。', B26, fill=MUTED)

    chip(draw, (98, 468, 306, 518), '系统集成', BLUE_BG, BLUE)
    chip(draw, (328, 468, 560, 518), '控制执行', ORANGE_BG, ORANGE)
    chip(draw, (582, 468, 820, 518), '验证闭环', GREEN_BG, GREEN)
    chip(draw, (842, 468, 1058, 518), '设计重难点', PURPLE_BG, PURPLE)

    text(draw, (96, 602), '面试主线', B22, fill=BLUE)
    text(draw, (96, 654), 'Why 这个项目值得做\nWhat 系统整体怎么搭起来\nHow 关键控制逻辑怎么收敛\nEvidence 我如何证明它真的跑通', T42, fill=INK, spacing=14)

    text(draw, (96, 968), '答辩口径：范围克制，但系统完整；重点不在“做了个接口壳”，而在“把控制链路和闭环验证补完整”。', B20, fill=MUTED)

    boxes = [
        ('项目定位', 'SoC 原型', BLUE),
        ('执行组织', 'IMEM + Sequencer', ORANGE),
        ('RTL 重点', 'wait / hold / mux', PURPLE),
        ('结果证据', '41 / 41 PASS', GREEN),
    ]
    y = 156
    for title, value, tone in boxes:
        rounded(draw, (1412, y, 1762, y + 126), radius=26, fill=WHITE, outline=WHITE, width=1)
        text(draw, (1450, y + 22), title, B20, fill=MUTED)
        text(draw, (1450, y + 62), value, T34, fill=tone)
        y += 154

    img.save(SLIDE1)


def render_overview():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 2, '项目级 RTL 架构', '从 tpu_soc 顶层看完整 RTL 分层：frontend 控制域 + tpu core 数据域。')

    rounded(draw, (72, 232, 1848, 942), radius=36)
    chip(draw, (98, 250, 328, 298), 'top wrapper', BLUE_BG, BLUE, font=B18)
    text(draw, (344, 258), 'src_axi/tpu_soc.sv', B18, fill=MUTED)

    rounded(draw, (108, 320, 804, 894), radius=30)
    text(draw, (146, 352), 'Frontend 控制域', T34, fill=INK)
    text(draw, (146, 404), 'AXI-Lite / 寄存器映射 / IMEM / sequencer / control_unit', B20, fill=MUTED)
    chip(draw, (146, 452, 334, 498), 'src_axi/tpu_frontend_axil.sv', BLUE_BG, BLUE, font=B18)
    chip(draw, (356, 452, 540, 498), 'control_unit.sv', ORANGE_BG, ORANGE, font=B18)

    frontend_boxes = [
        ((146, 548, 332, 634), 'AXI-Lite IF', BLUE_BG, BLUE),
        ((364, 548, 550, 634), 'Reg Map', BLUE_BG, BLUE),
        ((582, 548, 768, 634), 'IMEM', ORANGE_BG, ORANGE),
        ((146, 676, 410, 762), 'Sequencer FSM', ORANGE_BG, ORANGE),
        ((442, 676, 768, 762), 'Host Write Mux', GREEN_BG, GREEN),
    ]
    for xy, label_s, fill, tone in frontend_boxes:
        rounded(draw, xy, radius=22, fill=fill, outline=LINE, width=1)
        center(draw, (xy[0] + xy[2]) / 2, (xy[1] + xy[3]) / 2, label_s, T28, fill=tone)

    rounded(draw, (848, 320, 1518, 894), radius=30)
    text(draw, (886, 352), 'TPU Core 数据域', T34, fill=INK)
    text(draw, (886, 404), 'src_axi/tpu.sv 负责把 UB / Systolic / VPU 组织成执行闭环', B20, fill=MUTED)

    rounded(draw, (886, 500, 1180, 848), radius=24)
    text(draw, (920, 534), 'Unified Buffer', T28, fill=INK)
    text(draw, (920, 582), 'host load\nread ctrl\nwriteback', B20, fill=MUTED, spacing=10)
    chip(draw, (920, 724, 1082, 770), 'input / weight', BLUE_BG, BLUE, font=B18)
    chip(draw, (920, 782, 1062, 828), 'bias / Y / H', GREEN_BG, GREEN, font=B18)

    rounded(draw, (1220, 500, 1480, 652), radius=24)
    center(draw, 1350, 566, 'Systolic Array\n(2x2)', T28, fill=INK)
    chip(draw, (1250, 606, 1450, 648), 'sys_data / weight', BLUE_BG, BLUE, font=B18)

    rounded(draw, (1220, 696, 1480, 848), radius=24)
    text(draw, (1272, 726), 'VPU', T28, fill=INK)
    chip(draw, (1244, 772, 1336, 814), 'Bias', BLUE_BG, BLUE, font=B18)
    chip(draw, (1360, 772, 1452, 814), 'Loss', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1244, 822, 1336, 864), 'LReLU', GREEN_BG, GREEN, font=B18)
    chip(draw, (1360, 822, 1452, 864), 'dLReLU', PURPLE_BG, PURPLE, font=B18)

    rounded(draw, (1560, 390, 1808, 844), radius=28)
    text(draw, (1600, 426), '对外观测', T28, fill=INK)
    text(draw, (1600, 474), 'host / scoreboard\n可直接读到状态与输出', B20, fill=MUTED, spacing=10)
    chip(draw, (1600, 604, 1766, 648), 'vpu_data_out', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1600, 664, 1766, 708), 'sys_data_out', BLUE_BG, BLUE, font=B18)
    chip(draw, (1600, 724, 1738, 768), 'valid', GREEN_BG, GREEN, font=B18)

    arrow(draw, [(804, 596), (886, 596)], fill=BLUE)
    arrow(draw, [(804, 704), (886, 704)], fill=ORANGE)
    arrow(draw, [(1180, 606), (1220, 606)], fill=BLUE)
    arrow(draw, [(1180, 746), (1220, 746)], fill=GREEN)
    arrow(draw, [(1350, 652), (1350, 696)], fill=ORANGE)
    arrow(draw, [(1220, 828), (1180, 828)], fill=GREEN)
    arrow(draw, [(1518, 566), (1560, 566)], fill=ORANGE)
    arrow(draw, [(1518, 742), (1560, 742)], fill=BLUE)

    chip(draw, (840, 936, 1110, 986), 'ctrl: ub_rd_* / cfg / pathway', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1134, 936, 1328, 986), 'data: UB -> SA/VPU', BLUE_BG, BLUE, font=B18)
    chip(draw, (1352, 936, 1580, 986), 'writeback: VPU -> UB', GREEN_BG, GREEN, font=B18)
    text(draw, (72, 988), '这一页负责回答“模块到底怎么拼在一起”的问题；下一页再讲控制主链路如何实际跑起来。', B20, fill=MUTED)

    img.save(SLIDE2)

def render_chain():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 3, '控制与执行主链路', '这一页负责把“系统能跑”讲成一条严密的执行路径。')

    steps = [
        ('01', 'Host 装载', 'AXI-Lite 写寄存器\n写 UB / 写 IMEM', BLUE, BLUE_BG),
        ('02', '启动方式', 'step / start\n进入 sequencer', ORANGE, ORANGE_BG),
        ('03', '指令分发', 'seq_instr_pulse\n送 control_unit', GREEN, GREEN_BG),
        ('04', 'UB 读路径', 'ub_rd_*\n参数与数据取数', PURPLE, PURPLE_BG),
        ('05', '核心执行', 'Systolic -> VPU\n完成算子处理', BLUE, BLUE_BG),
        ('06', '写回收敛', 'writeback / drain\nbusy 清零', ORANGE, ORANGE_BG),
    ]

    x = 88
    centers = []
    for idx, (num, title, body, tone, fill) in enumerate(steps):
        w = 266
        h = 208
        rounded(draw, (x, 310, x + w, 310 + h), radius=28)
        chip(draw, (x + 24, 330, x + 92, 374), num, fill, tone, font=B18)
        text(draw, (x + 24, 400), title, T28, fill=INK)
        text(draw, (x + 24, 446), body, B20, fill=MUTED, spacing=8)
        centers.append((x + w / 2, 414))
        x += 300

    for i in range(len(centers) - 1):
        arrow(draw, [(centers[i][0] + 120, 414), (centers[i + 1][0] - 120, 414)], fill=BLUE if i % 2 == 0 else ORANGE)

    chip(draw, (164, 616, 424, 666), 'wait_after', BLUE_BG, BLUE)
    chip(draw, (458, 616, 742, 666), 'vpu_pathway_hold', ORANGE_BG, ORANGE)
    chip(draw, (778, 616, 1046, 666), 'host_write_mux', GREEN_BG, GREEN)
    chip(draw, (1078, 616, 1318, 666), 'vpu_drain', PURPLE_BG, PURPLE)

    rounded(draw, (90, 738, 1832, 944), radius=30)
    text(draw, (126, 772), '可追问点', B22, fill=BLUE)
    text(draw, (126, 824), '1. 为什么等待语义不用固定周期，而要显式 `wait_after`？\n2. 为什么 `vpu_data_pathway` 必须保持，而不能只看当前 decode 输出？\n3. 为什么 host write 和 CU write 需要做优先级仲裁？', B24, fill=INK, spacing=16)

    img.save(SLIDE3)


def render_difficulty():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 4, '设计重难点', '按“现象 -> 根因 -> 修复 -> 验证”讲，最容易把深度讲出来。')

    cards = [
        {
            'xy': (72, 230, 608, 944),
            'tone': BLUE,
            'fill': BLUE_BG,
            'title': '显式等待语义',
            'rows': [
                ('现象', 'sequencer 可能提前推进'),
                ('根因', '隐式等待无法覆盖 VPU 延迟'),
                ('修复', '用 wait_after + vpu_drain 收口'),
                ('验证', 'SEQ_WAIT 到 drain 后再 advance'),
            ],
            'ref': 'frontend_axil:139,153,177-179',
            'code': [
                'seq_needs_wait = seq_instr[23];',
                'vpu_drain = vpu_valid_prev && !tpu_vpu_valid_in;',
                'if (seq_needs_wait) seq_state <= SEQ_WAIT;',
            ],
        },
        {
            'xy': (692, 230, 1228, 944),
            'tone': ORANGE,
            'fill': ORANGE_BG,
            'title': '路径保持',
            'rows': [
                ('现象', 'dispatch 后几拍 route 会跑偏'),
                ('根因', 'pipeline latency 与 decode 不同步'),
                ('修复', '把 pathway latch 成持久寄存器'),
                ('验证', '输出阶段仍保持正确模块选择'),
            ],
            'ref': 'frontend_axil:156-157,403',
            'code': [
                'if (seq_instr_pulse && opcode == UB_RD)',
                '    vpu_pathway_reg <= seq_instr[22:19];',
                'vpu_data_pathway_out = vpu_pathway_reg;',
            ],
        },
        {
            'xy': (1312, 230, 1848, 944),
            'tone': GREEN,
            'fill': GREEN_BG,
            'title': 'Host / CU 仲裁',
            'rows': [
                ('现象', 'host write 与 CU write 语义重叠'),
                ('根因', '同一写口需要区分两类驱动源'),
                ('修复', 'UB_PUSH 优先的 host write mux'),
                ('验证', '参数装载和训练写回都能稳定工作'),
            ],
            'ref': 'frontend_axil:405-408',
            'code': [
                'valid0 = push0 ? one : cu_valid0;',
                'valid1 = push1 ? one : cu_valid1;',
                'data0 = push0 ? host0 : cu_data0;',
            ],
        },
    ]

    for c in cards:
        l, t, r, b = c['xy']
        rounded(draw, c['xy'], radius=30)
        draw.rounded_rectangle((l + 20, t + 24, l + 28, b - 24), radius=4, fill=c['tone'])
        text(draw, (l + 52, t + 28), c['title'], T34, fill=INK)
        chip(draw, (l + 52, t + 86, l + 212, t + 132), '可追问', c['fill'], c['tone'], font=B18)
        yy = t + 176
        for label_text, body_text in c['rows']:
            issue_row(draw, l + 52, yy, label_text, body_text, c['tone'])
            yy += 54
        code_strip(draw, (l + 34, b - 226, r - 28, b - 34), 'RTL 特写', c['ref'], c['code'], c['tone'])

    img.save(SLIDE4)


def render_results():
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 5, '验证结果、边界与后续', '结果给证据，边界要主动讲，后续要能展开。')

    metric_box(draw, (72, 240, 394, 392), '回归', '41 / 41 PASS', BLUE)
    metric_box(draw, (424, 240, 746, 392), '训练', '12 epoch', ORANGE)
    metric_box(draw, (776, 240, 1160, 392), 'loss', '0.2529 -> 0.1777', GREEN)
    metric_box(draw, (1190, 240, 1512, 392), 'XOR', '(0, 1, 1, 0)', BLUE)

    rounded(draw, (72, 468, 936, 828), radius=34)
    text(draw, (112, 512), '证据链', B22, fill=BLUE)
    text(draw, (112, 566), '1. 单次 e2e 回归全部通过\n2. 多 epoch 训练 loss 持续下降\n3. XOR 分类结果正确', T42, fill=INK, spacing=16)

    rounded(draw, (984, 468, 1848, 828), radius=34)
    text(draw, (1024, 512), '边界与延展', B22, fill=ORANGE)
    text(draw, (1024, 566), '1. 当前是 tiny-tpu 原型，不夸大成完整通用 SoC\n2. 当前不覆盖 DMA / IRQ / clock-gating\n3. 可自然延展到更大阵列、DMA、低功耗、可观测性', B26, fill=INK, spacing=16)

    chip(draw, (118, 900, 334, 950), '系统闭环', BLUE_BG, BLUE)
    chip(draw, (360, 900, 592, 950), 'RTL 控制', ORANGE_BG, ORANGE)
    chip(draw, (618, 900, 854, 950), '验证收敛', GREEN_BG, GREEN)
    chip(draw, (880, 900, 1164, 950), '边界主动说明', PURPLE_BG, PURPLE)
    text(draw, (72, 986), '结论：这不是“只改了个接口”的项目，而是把控制链路、执行链路和验证链路闭合起来的系统化原型。', B20, fill=MUTED)

    img.save(SLIDE5)


def build_ppt():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = 'TinyTPU AXI-Lite SoC Interview PPT'
    prs.core_properties.subject = 'Hisilicon clean version'
    prs.core_properties.author = 'Codex'

    for p in [SLIDE1, SLIDE2, SLIDE3, SLIDE4, SLIDE5]:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        slide.shapes.add_picture(str(p), 0, 0, width=SLIDE_W, height=SLIDE_H)

    prs.save(str(OUT_PPT))


def main():
    render_cover()
    render_overview()
    render_chain()
    render_difficulty()
    render_results()
    build_ppt()


if __name__ == '__main__':
    main()
