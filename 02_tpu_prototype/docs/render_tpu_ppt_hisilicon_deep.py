from __future__ import annotations

import json
import math
import re
import shutil
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path('/home/jjt/tpu-soc')
DOCS = ROOT / 'docs'
DROP = Path('/tmp/VMwareDnD/71NKe5')
SHARE = Path('/mnt/hgfs/wdchenaic/tpu_interview_ppt')

OUT_PPT = DOCS / 'tpu_project_interview_hisilicon_deep_zh.pptx'
OUT_APPENDIX_PPT = DOCS / 'tpu_project_interview_hisilicon_deep_appendix_zh.pptx'
DROP_PPT = DROP / 'tpu_project_interview_hisilicon_deep_zh.pptx'
DROP_APPENDIX_PPT = DROP / 'tpu_project_interview_hisilicon_deep_appendix_zh.pptx'
SHARE_PPT = SHARE / '01_main_8p.pptx'
SHARE_APPENDIX_PPT = SHARE / '02_appendix_35p.pptx'

OVERVIEW_SRC = DOCS / 'tpu_project_architecture_16x9_zh.png'
UB_MAP_PATH = ROOT / 'compiler/out/mlp_2_2_1_q8_8.ub_map.json'
SCHEDULE_PATH = ROOT / 'compiler/out/mlp_2_2_1_q8_8.schedule.json'
VCD_PATH = ROOT / 'waveforms/tpu_soc.vcd'
PE_ANIM_BASE_PATH = DOCS / 'tpu_pe_anim_base.png'
PE_ANIM_GIF_PATHS = [DOCS / f'tpu_pe_anim_step_{idx}.gif' for idx in range(6)]
UB_TIMING_ANIM_BASE_PATH = DOCS / 'tpu_ub_timing_anim_base.png'
UB_TIMING_ANIM_GIF_PATHS = [DOCS / f'tpu_ub_timing_anim_step_{idx}.gif' for idx in range(5)]
VPU_ANIM_BASE_PATH = DOCS / 'tpu_vpu_anim_base.png'
VPU_ANIM_GIF_PATHS = [DOCS / f'tpu_vpu_anim_step_{idx}.gif' for idx in range(4)]
FRONTEND_ANIM_BASE_PATH = DOCS / 'tpu_frontend_anim_base.png'
FRONTEND_ANIM_GIF_PATHS = [DOCS / f'tpu_frontend_anim_step_{idx}.gif' for idx in range(4)]
COMPILER_ANIM_BASE_PATH = DOCS / 'tpu_compiler_anim_base.png'
COMPILER_ANIM_GIF_PATHS = [DOCS / f'tpu_compiler_anim_step_{idx}.gif' for idx in range(4)]
ANIM_PROGRESS_POINTS = [0.18, 0.44, 0.72, 1.0]
ANIM_DURATIONS = [120, 130, 150, 620]
PPT_FONT_NAME = 'Microsoft YaHei'
PPT_MONO_FONT_NAME = 'Consolas'

SLIDES = [
    {
        'path': DOCS / 'tpu_hisilicon_deep_cover_16x9_zh.png',
        'title': '封面',
        'subtitle': '完整项目展示版',
        'notes': ['覆盖系统、RTL、编译器、UB/PE/VPU、波形与结果。'],
        'layout': 'cover',
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_system_16x9_zh.png',
        'title': '系统总览',
        'subtitle': '软件到验证闭环',
        'notes': ['先讲全景，再进入 RTL 与子系统细节。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_rtl_16x9_zh.png',
        'title': '项目级 RTL 结构',
        'subtitle': '控制域 / 执行域 / 可观测接口',
        'notes': ['顶层拼接关系清楚，便于后面拆讲前端和核心。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_compiler_16x9_zh.png',
        'title': '编译器与指令组织',
        'subtitle': '规格如何落到 UB / schedule / IMEM',
        'notes': ['强调当前是阶段级 schedule，不是 cycle-accurate 编译器。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_frontend_16x9_zh.png',
        'title': 'Frontend 总览',
        'subtitle': '寄存器、IMEM、sequencer、decode',
        'notes': ['作为控制域主入口，负责 host 装载与指令派发。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_frontend_split_16x9_zh.png',
        'title': 'Frontend 再拆一页',
        'subtitle': '把配置面和运行面拆开',
        'notes': ['左侧讲 AXI/寄存器，右侧讲 sequencer/dispatch。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_ub_16x9_zh.png',
        'title': 'Unified Buffer 设计',
        'subtitle': 'host load / core read / VPU writeback 汇合点',
        'notes': ['UB 不是 SRAM 方块，而是整个数据流的枢纽。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_wptr_16x9_zh.png',
        'title': 'wr_ptr / base / restore',
        'subtitle': '参数装载区与训练写回区如何隔离',
        'notes': ['这页专门解释 start 时恢复写指针，避免覆盖参数区。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_ub_timing_16x9_zh.png',
        'title': 'UB 读流与 PE 时序对齐',
        'subtitle': 'valid 波前、hold 周期和 drain 关系',
        'notes': ['把 UB 发数、阵列采样和 sequencer wait 串起来讲。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_ub_update_16x9_zh.png',
        'title': 'UB 内梯度下降更新',
        'subtitle': 'bias update 与 weight update 语义不同',
        'notes': ['bias 可连续累加，weight 按 tile 用旧值更新。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_vpu_16x9_zh.png',
        'title': 'VPU 单独展开',
        'subtitle': 'pathway bit 控制可重组训练路径',
        'notes': ['不是黑盒后处理，而是阶段切换的关键路径单元。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_pe_16x9_zh.png',
        'title': 'PE 与计算阵列',
        'subtitle': '2x2 systolic 波前形成方式',
        'notes': ['先讲 topology，再讲 active/inactive weight。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_control_16x9_zh.png',
        'title': '关键控制 RTL 特写',
        'subtitle': 'wait、路径保持、写口仲裁',
        'notes': ['这是最容易出现系统级 bug 的三处。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_cycle_16x9_zh.png',
        'title': '逐拍计算动态',
        'subtitle': '帮助面试官理解 systolic 波前',
        'notes': ['这是解释页，不替代真实波形页。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_wave_16x9_zh.png',
        'title': '验证波形与回归覆盖',
        'subtitle': '真实时序证据',
        'notes': ['busy -> dispatch -> ub read -> sys out -> drain 可直接对上。'],
    },
    {
        'path': DOCS / 'tpu_hisilicon_deep_results_16x9_zh.png',
        'title': '结果、边界与追问方向',
        'subtitle': '收束口径',
        'notes': ['主动交代边界，避免过度包装成商用 NPU。'],
    },
]
GIF_PATH = DOCS / 'tpu_systolic_cycle_demo.gif'
STRIP_PATH = DOCS / 'tpu_systolic_cycle_strip.png'

W, H = 1920, 1080
TOTAL_PAGES = len(SLIDES)
PPT_TOTAL_PAGES = 0
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

STATIC_FRAME_X = 0.48
STATIC_FRAME_Y = 1.56
STATIC_FRAME_W = 12.34
STATIC_FRAME_H = 5.38
STATIC_FOOTER_X = 0.72
STATIC_FOOTER_Y = 7.0
STATIC_FOOTER_W = 11.0
STATIC_FOOTER_H = 0.18
STATIC_CARD_PAD_X = 0.26
STATIC_CARD_PAD_TOP = 0.18
STATIC_CARD_TITLE_H = 0.24
STATIC_CARD_BODY_GAP = 0.42
STATIC_CARD_BOTTOM_PAD = 0.24

BG = (247, 245, 240)
WHITE = (255, 255, 255)
INK = (28, 37, 51)
MUTED = (105, 112, 123)
LINE = (221, 226, 232)
BLUE = (32, 103, 143)
ORANGE = (197, 112, 64)
GREEN = (83, 141, 104)
PURPLE = (116, 104, 170)
RED = (176, 74, 74)
BLUE_BG = (236, 243, 248)
ORANGE_BG = (248, 239, 231)
GREEN_BG = (234, 243, 236)
PURPLE_BG = (240, 237, 247)
RED_BG = (248, 235, 235)
PANEL = (242, 237, 230)
HILITE = (241, 246, 250)
DARK = (55, 66, 80)

TITLE_FONT = '/usr/share/fonts/google-noto/NotoSansSC-Medium.otf'
BODY_FONT = '/usr/share/fonts/google-noto/NotoSansSC-Regular.otf'
MONO_FONT = '/usr/share/fonts/dejavu/DejaVuSansMono.ttf'

T72 = ImageFont.truetype(TITLE_FONT, 72)
T60 = ImageFont.truetype(TITLE_FONT, 60)
T52 = ImageFont.truetype(TITLE_FONT, 52)
T44 = ImageFont.truetype(TITLE_FONT, 44)
T38 = ImageFont.truetype(TITLE_FONT, 38)
T34 = ImageFont.truetype(TITLE_FONT, 34)
T30 = ImageFont.truetype(TITLE_FONT, 30)
T28 = ImageFont.truetype(TITLE_FONT, 28)
B28 = ImageFont.truetype(BODY_FONT, 28)
B26 = ImageFont.truetype(BODY_FONT, 26)
B24 = ImageFont.truetype(BODY_FONT, 24)
B22 = ImageFont.truetype(BODY_FONT, 22)
B20 = ImageFont.truetype(BODY_FONT, 20)
B18 = ImageFont.truetype(BODY_FONT, 18)
M22 = ImageFont.truetype(MONO_FONT, 22)
M20 = ImageFont.truetype(MONO_FONT, 20)
M18 = ImageFont.truetype(MONO_FONT, 18)
M16 = ImageFont.truetype(MONO_FONT, 16)


def rounded(draw, xy, radius=28, fill=WHITE, outline=LINE, width=2):
    draw.rounded_rectangle(xy, radius=radius, fill=fill, outline=outline, width=width)


def text(draw, pos, msg, font, fill=INK, spacing=6, align='left'):
    draw.multiline_text(pos, msg, font=font, fill=fill, spacing=spacing, align=align)


def center(draw, cx, cy, msg, font, fill=INK, spacing=6):
    bbox = draw.multiline_textbbox((0, 0), msg, font=font, spacing=spacing, align='center')
    w = bbox[2] - bbox[0]
    h = bbox[3] - bbox[1]
    draw.multiline_text((cx - w / 2, cy - h / 2), msg, font=font, fill=fill, spacing=spacing, align='center')


def chip(draw, xy, msg, fill, text_fill, font=B20):
    rounded(draw, xy, radius=18, fill=fill, outline=fill, width=1)
    l, t, r, b = xy
    center(draw, (l + r) / 2, (t + b) / 2, msg, font, fill=text_fill)


def metric_box(draw, xy, title, value, tone):
    rounded(draw, xy, radius=24)
    l, t, r, b = xy
    text(draw, (l + 26, t + 24), title, B20, fill=MUTED)
    text(draw, (l + 26, t + 66), value, T34, fill=tone)


def arrow(draw, pts, fill=BLUE, width=5, arrow_size=16):
    draw.line(pts, fill=fill, width=width)
    x1, y1 = pts[-2]
    x2, y2 = pts[-1]
    if abs(x2 - x1) >= abs(y2 - y1):
        d = 1 if x2 > x1 else -1
        tri = [(x2, y2), (x2 - d * arrow_size, y2 - arrow_size // 2), (x2 - d * arrow_size, y2 + arrow_size // 2)]
    else:
        d = 1 if y2 > y1 else -1
        tri = [(x2, y2), (x2 - arrow_size // 2, y2 - d * arrow_size), (x2 + arrow_size // 2, y2 - d * arrow_size)]
    draw.polygon(tri, fill=fill)


def code_strip(draw, xy, title, ref, code_lines, tone):
    rounded(draw, xy, radius=22, fill=WHITE, outline=LINE, width=2)
    l, t, r, b = xy
    draw.rounded_rectangle((l + 18, t + 18, l + 26, b - 18), radius=4, fill=tone)
    text(draw, (l + 46, t + 14), title, B20, fill=tone)
    text(draw, (r - 240, t + 16), ref, B18, fill=MUTED)
    yy = t + 60
    for line in code_lines:
        rounded(draw, (l + 36, yy - 2, r - 28, yy + 28), radius=12, fill=HILITE, outline=HILITE, width=1)
        text(draw, (l + 50, yy), line, M18, fill=INK)
        yy += 40


def draw_header(draw, idx, title, subtitle):
    text(draw, (88, 72), f'P{idx}', B18, fill=BLUE)
    text(draw, (88, 108), title, T52, fill=INK)
    text(draw, (88, 170), subtitle, B22, fill=MUTED)
    rounded(draw, (1440, 66, 1832, 122), radius=24, fill=BLUE_BG, outline=BLUE_BG, width=1)
    center(draw, 1636, 94, f'{TOTAL_PAGES} 页扩展版的第 {idx} 页', B20, fill=BLUE)


def paste_fit(base, src_path, box, pad=20, bg=WHITE):
    l, t, r, b = box
    rounded(ImageDraw.Draw(base), box, radius=28, fill=bg, outline=LINE, width=2)
    if not src_path.exists():
        return
    img = Image.open(src_path).convert('RGB')
    bw = (r - l) - pad * 2
    bh = (b - t) - pad * 2
    ratio = min(bw / img.width, bh / img.height)
    size = (max(1, int(img.width * ratio)), max(1, int(img.height * ratio)))
    img = img.resize(size)
    x = l + (r - l - size[0]) // 2
    y = t + (b - t - size[1]) // 2
    base.paste(img, (x, y))


def load_ub_map():
    return json.loads(UB_MAP_PATH.read_text(encoding='utf-8'))


def load_schedule():
    return json.loads(SCHEDULE_PATH.read_text(encoding='utf-8'))


def parse_vcd(names):
    code_to_name = {}
    in_defs = True
    current_t = 0
    events = {name: [] for name in names}
    with VCD_PATH.open() as f:
        for line in f:
            if in_defs:
                if line.startswith('$var'):
                    parts = line.split()
                    code = parts[3]
                    name = parts[4]
                    if name in names:
                        code_to_name[code] = name
                elif line.startswith('$enddefinitions'):
                    in_defs = False
                continue
            if line.startswith('#'):
                current_t = int(line[1:].strip())
            elif line and line[0] in '01xz':
                code = line[1:].strip()
                if code in code_to_name:
                    events[code_to_name[code]].append((current_t, line[0]))
    return events


def first_busy_window(events):
    busy = events['busy_reg']
    rise = None
    fall = None
    prev = '0'
    for t, val in busy:
        if prev != '1' and val == '1' and rise is None:
            rise = t
        elif rise is not None and prev != '0' and val == '0':
            fall = t
            break
        prev = val
    if rise is None:
        rise = 0
    if fall is None:
        fall = rise + 2_000_000
    return max(0, rise - 200_000), fall + 250_000


def sample_signal(events, start, end):
    vals = []
    cur = '0'
    last_t = start
    for t, val in events:
        if t < start:
            cur = val
            continue
        if t > end:
            break
        if t > last_t:
            vals.append((last_t, t, cur))
        cur = val
        last_t = t
    if last_t < end:
        vals.append((last_t, end, cur))
    return vals


def draw_digital_wave(draw, x0, y0, width, height, segments, start, end, tone, label):
    text(draw, (x0 - 180, y0 + 4), label, B18, fill=INK)
    y_hi = y0 + 10
    y_lo = y0 + height - 10
    draw.line((x0, y_lo, x0 + width, y_lo), fill=LINE, width=1)
    for s, e, val in segments:
        xs = x0 + width * (s - start) / max(1, end - start)
        xe = x0 + width * (e - start) / max(1, end - start)
        y = y_hi if val == '1' else y_lo
        draw.line((xs, y, xe, y), fill=tone, width=4)
    for s, e, val in segments[1:]:
        xs = x0 + width * (s - start) / max(1, end - start)
        prev = segments[segments.index((s, e, val)) - 1][2]
        if prev != val:
            draw.line((xs, y_hi, xs, y_lo), fill=tone, width=3)


def draw_time_axis(draw, x0, y0, width, start, end):
    draw.line((x0, y0, x0 + width, y0), fill=MUTED, width=2)
    ticks = 6
    for i in range(ticks + 1):
        x = x0 + width * i / ticks
        draw.line((x, y0 - 8, x, y0 + 8), fill=MUTED, width=2)
        us = (start + (end - start) * i / ticks) / 1_000_000
        text(draw, (x - 30, y0 + 14), f'{us:.2f}us', B18, fill=MUTED)


def draw_loss_chart(draw, xy):
    l, t, r, b = xy
    rounded(draw, xy, radius=26)
    text(draw, (l + 26, t + 22), '多 epoch 收敛趋势', B20, fill=BLUE)
    pts = [0.2529, 0.243, 0.235, 0.226, 0.219, 0.212, 0.205, 0.198, 0.192, 0.187, 0.182, 0.1777]
    x0, y0 = l + 52, b - 56
    x1, y1 = r - 36, t + 80
    draw.line((x0, y0, x1, y0), fill=LINE, width=2)
    draw.line((x0, y0, x0, y1), fill=LINE, width=2)
    pmin, pmax = min(pts), max(pts)
    prev = None
    for i, v in enumerate(pts):
        x = x0 + (x1 - x0) * i / (len(pts) - 1)
        y = y0 - (y0 - y1) * (v - pmin) / max(1e-6, pmax - pmin)
        draw.ellipse((x - 5, y - 5, x + 5, y + 5), fill=BLUE, outline=BLUE)
        if prev:
            draw.line((prev[0], prev[1], x, y), fill=BLUE, width=3)
        prev = (x, y)
        text(draw, (x - 10, y0 + 12), str(i + 1), B18, fill=MUTED)
    text(draw, (x0 - 8, y1 - 6), f'{pmax:.3f}', B18, fill=MUTED)
    text(draw, (x0 - 8, y0 - 10), f'{pmin:.3f}', B18, fill=MUTED)
    text(draw, (r - 208, t + 30), '0.2529 -> 0.1777', T28, fill=GREEN)


def render_cover(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw.rectangle((0, 0, 24, H), fill=BLUE)
    rounded(draw, (1360, 72, 1822, 962), radius=44, fill=PANEL, outline=PANEL, width=1)

    text(draw, (96, 106), 'TinyTPU AXI-Lite SoC', B22, fill=BLUE)
    text(draw, (96, 178), '完整项目展示版\n从系统总览到 RTL 细节', T72, fill=INK, spacing=12)
    text(draw, (96, 364), '面向硬件面试的深挖版：不只讲“做了什么”，而是讲“系统怎样闭合、控制怎样收口、证据怎样自洽”。', B26, fill=MUTED)

    chips = [
        ('系统总览', BLUE_BG, BLUE),
        ('项目级 RTL', ORANGE_BG, ORANGE),
        ('编译器 / IMEM', GREEN_BG, GREEN),
        ('UB / PE / 波形', PURPLE_BG, PURPLE),
    ]
    x = 96
    for label, fill, tone in chips:
        chip(draw, (x, 472, x + 210, 522), label, fill, tone)
        x += 228

    text(draw, (96, 604), '这套材料的讲述顺序', B22, fill=BLUE)
    text(draw, (96, 654), '1. 系统栈怎么补完整\n2. RTL 如何真正拼起来\n3. 指令与 UB 映射如何驱动执行\n4. 关键子系统怎么设计\n5. 结果、波形和边界如何证明它可靠', T44, fill=INK, spacing=18)

    metrics = [
        ('阵列规模', '2 x 2 Weight Stationary', BLUE),
        ('模型目标', '2-layer MLP / XOR / Q8.8', ORANGE),
        ('编译结果', '66 commands / IMEM', GREEN),
        ('回归证据', '41 / 41 PASS', PURPLE),
        ('训练证据', '12 epoch 收敛', BLUE),
    ]
    y = 136
    for title, value, tone in metrics:
        rounded(draw, (1402, y, 1778, y + 126), radius=26, fill=WHITE, outline=WHITE, width=1)
        text(draw, (1438, y + 20), title, B20, fill=MUTED)
        text(draw, (1438, y + 60), value, T30, fill=tone)
        y += 146

    text(draw, (96, 984), '推荐讲法：先讲系统闭环，再讲关键 RTL 和证据页，面试官会更容易自然追问。', B20, fill=MUTED)
    img.save(path)


def render_system(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 2, '系统总览', '把“软件、前端、计算核心、验证闭环”放回同一张图里，先建立全景认知。')

    rounded(draw, (72, 238, 1848, 930), radius=34)
    text(draw, (108, 270), '项目主链路', B20, fill=BLUE)

    cols = [
        ((112, 336, 474, 822), '软件与编译', BLUE, BLUE_BG,
         ['model spec', 'ub_allocator', 'scheduler', 'encode_instrs'],
         ['规格输入', 'UB 地址映射', '阶段级 schedule', 'imem.hex / json']),
        ((536, 336, 942, 822), 'SoC 控制与执行', ORANGE, ORANGE_BG,
         ['AXI-Lite reg map', 'IMEM + sequencer', 'control_unit decode', 'host / CU mux'],
         ['寄存器接口', 'start / step / wait', 'ub_rd_* / switch', '统一写口仲裁']),
        ((1004, 336, 1410, 822), '计算核心', GREEN, GREEN_BG,
         ['Unified Buffer', 'Systolic Array 2x2', 'VPU', 'writeback -> UB'],
         ['host load / read ptr', 'weight stationary', 'bias / loss / dAct', '形成训练闭环']),
        ((1472, 336, 1808, 822), '验证与结果', PURPLE, PURPLE_BG,
         ['module cocotb', 'AXI e2e', 'train convergence', 'VCD + scoreboard'],
         ['模块级行为', '41 / 41 PASS', '12 epoch 收敛', '时序与结果双证据']),
    ]
    mids = []
    for xy, title_s, tone, fill, top_rows, sub_rows in cols:
        l, t, r, b = xy
        rounded(draw, xy, radius=28)
        draw.rounded_rectangle((l + 18, t + 24, l + 26, b - 24), radius=4, fill=tone)
        text(draw, (l + 44, t + 26), title_s, T34, fill=INK)
        yy = t + 108
        for head, desc in zip(top_rows, sub_rows):
            rounded(draw, (l + 34, yy, r - 34, yy + 84), radius=20, fill=fill, outline=LINE, width=1)
            text(draw, (l + 56, yy + 14), head, B22, fill=tone)
            text(draw, (l + 56, yy + 46), desc, B18, fill=MUTED)
            yy += 102
        mids.append(((l + r) // 2, (t + b) // 2))

    for i in range(len(mids) - 1):
        arrow(draw, [(mids[i][0] + 154, mids[i][1]), (mids[i + 1][0] - 154, mids[i + 1][1])], fill=[BLUE, ORANGE, GREEN][i], width=6)

    chip(draw, (292, 850, 432, 892), 'json / hex', BLUE_BG, BLUE, font=B18)
    chip(draw, (736, 850, 888, 892), 'control pulse', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1196, 850, 1320, 892), 'dataflow', GREEN_BG, GREEN, font=B18)
    chip(draw, (1602, 850, 1758, 892), 'evidence', PURPLE_BG, PURPLE, font=B18)

    chip(draw, (76, 952, 322, 1002), 'Why: 系统化项目', BLUE_BG, BLUE)
    chip(draw, (346, 952, 628, 1002), 'What: 编译器到 RTL 闭环', ORANGE_BG, ORANGE)
    chip(draw, (652, 952, 928, 1002), 'How: 控制 + 数据 + 验证', GREEN_BG, GREEN)
    text(draw, (72, 1014), '这页建议讲 40 秒，只给全景，不抢后面细节页的叙事空间。', B20, fill=MUTED)
    img.save(path)


def render_rtl(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 3, '项目级 RTL 结构', '把 tpu_soc 顶层拆成控制域、执行域、可观测接口三层，回答“模块到底怎么连”。')

    rounded(draw, (72, 232, 1848, 950), radius=36)
    chip(draw, (102, 252, 286, 298), 'top wrapper', BLUE_BG, BLUE, font=B18)
    text(draw, (304, 258), 'src_axi/tpu_soc.sv', B18, fill=MUTED)

    rounded(draw, (108, 334, 786, 894), radius=30)
    text(draw, (144, 366), 'Frontend 控制域', T34, fill=INK)
    text(draw, (144, 414), 'AXI-Lite slave / reg map / IMEM / sequencer / decode', B20, fill=MUTED)
    chip(draw, (144, 458, 356, 504), 'tpu_frontend_axil.sv', BLUE_BG, BLUE, font=B18)
    chip(draw, (378, 458, 542, 504), 'control_unit.sv', ORANGE_BG, ORANGE, font=B18)

    boxes = [
        ((144, 548, 338, 632), 'AXI-Lite IF', BLUE_BG, BLUE),
        ((366, 548, 560, 632), 'CTRL/STATUS', BLUE_BG, BLUE),
        ((588, 548, 742, 632), 'IMEM', ORANGE_BG, ORANGE),
        ((144, 676, 396, 760), '4-state sequencer', ORANGE_BG, ORANGE),
        ((424, 676, 742, 760), 'Host / CU write mux', GREEN_BG, GREEN),
        ((144, 788, 742, 852), 'decode: ub_rd_* / sys_switch / vpu_pathway / lr', PURPLE_BG, PURPLE),
    ]
    for xy, msg, fill, tone in boxes:
        rounded(draw, xy, radius=22, fill=fill, outline=LINE, width=1)
        center(draw, (xy[0] + xy[2]) / 2, (xy[1] + xy[3]) / 2, msg, T28 if '\n' not in msg else B22, fill=tone)

    rounded(draw, (828, 334, 1526, 894), radius=30)
    text(draw, (864, 366), 'TPU Core 执行域', T34, fill=INK)
    text(draw, (864, 414), 'src_axi/tpu.sv 统一组织 UB / Systolic / VPU / writeback', B20, fill=MUTED)
    chip(draw, (864, 458, 986, 504), 'tpu.sv', ORANGE_BG, ORANGE, font=B18)

    rounded(draw, (864, 540, 1128, 854), radius=24)
    text(draw, (898, 570), 'Unified Buffer', T30, fill=INK)
    text(draw, (898, 618), 'host write\nVPU writeback\n6 路读指针\nin-UB grad update', B20, fill=MUTED, spacing=10)
    chip(draw, (898, 770, 1016, 812), 'v3', BLUE_BG, BLUE, font=B18)
    chip(draw, (1030, 770, 1094, 812), '128w', GREEN_BG, GREEN, font=B18)

    rounded(draw, (1160, 540, 1490, 676), radius=24)
    center(draw, 1325, 592, 'Systolic Array\n2 x 2', T30, fill=INK)
    chip(draw, (1198, 632, 1454, 672), 'dual weight regs + pe_enabled', BLUE_BG, BLUE, font=B18)

    rounded(draw, (1160, 718, 1490, 854), radius=24)
    text(draw, (1292, 746), 'VPU', T30, fill=INK)
    chip(draw, (1188, 792, 1276, 834), 'Bias', BLUE_BG, BLUE, font=B18)
    chip(draw, (1290, 792, 1378, 834), 'Loss', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1392, 792, 1480, 834), 'dAct', PURPLE_BG, PURPLE, font=B18)

    rounded(draw, (1560, 388, 1808, 846), radius=28)
    text(draw, (1600, 424), '可观测接口', T30, fill=INK)
    text(draw, (1600, 472), 'host 可读输出 / STATUS\n波形中也能直接抓到关键信号', B20, fill=MUTED, spacing=10)
    chip(draw, (1600, 604, 1760, 646), 'vpu_data_out', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1600, 662, 1760, 704), 'sys_data_out', BLUE_BG, BLUE, font=B18)
    chip(draw, (1600, 720, 1738, 762), 'valid/busy', GREEN_BG, GREEN, font=B18)
    chip(draw, (1600, 778, 1738, 820), 'pc/state', PURPLE_BG, PURPLE, font=B18)

    arrow(draw, [(786, 590), (864, 590)], fill=BLUE)
    arrow(draw, [(786, 716), (864, 716)], fill=ORANGE)
    arrow(draw, [(1128, 606), (1160, 606)], fill=BLUE)
    arrow(draw, [(1128, 782), (1160, 782)], fill=GREEN)
    arrow(draw, [(1325, 676), (1325, 718)], fill=ORANGE)
    arrow(draw, [(1490, 606), (1560, 606)], fill=BLUE)
    arrow(draw, [(1490, 784), (1560, 784)], fill=ORANGE)

    chip(draw, (824, 930, 1104, 980), 'ctrl: ub_rd_* / switch / pathway', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1128, 930, 1334, 980), 'data: UB -> SA -> VPU', BLUE_BG, BLUE, font=B18)
    chip(draw, (1358, 930, 1602, 980), 'writeback: VPU -> UB', GREEN_BG, GREEN, font=B18)
    text(draw, (72, 1000), '这页强调“顶层怎么拼”；下一页开始按关键子系统拆开讲。', B20, fill=MUTED)
    img.save(path)


def render_compiler(path, ub_map, schedule):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 4, '编译器与指令组织', '解释软件侧如何把 MLP 规格降成 UB 映射、schedule 和 IMEM 指令。')

    steps = [
        ('01', 'model spec', 'mlp_2_2_1_q8_8.json', BLUE, BLUE_BG),
        ('02', 'ub_allocator', 'tensor catalog\n地址分配', GREEN, GREEN_BG),
        ('03', 'scheduler', '阶段级命令\nwait_after', ORANGE, ORANGE_BG),
        ('04', 'encode_instrs', '32-bit IMEM', PURPLE, PURPLE_BG),
        ('05', 'AXI write IMEM', 'frontend start', BLUE, BLUE_BG),
    ]
    x = 88
    mids = []
    for num, title_s, body, tone, fill in steps:
        rounded(draw, (x, 284, x + 310, 494), radius=28)
        chip(draw, (x + 24, 308, x + 90, 352), num, fill, tone, font=B18)
        text(draw, (x + 24, 372), title_s, T28, fill=INK)
        text(draw, (x + 24, 416), body, B22, fill=MUTED, spacing=8)
        mids.append((x + 155, 388))
        x += 354
    for i in range(len(mids) - 1):
        arrow(draw, [(mids[i][0] + 140, mids[i][1]), (mids[i + 1][0] - 140, mids[i + 1][1])], fill=BLUE if i % 2 == 0 else ORANGE)

    metric_box(draw, (88, 548, 406, 690), '命令数', f"{len(schedule['commands'])} commands", BLUE)
    metric_box(draw, (438, 548, 756, 690), 'UB 占用', f"{ub_map['allocated_words']} / 128 words", GREEN)
    metric_box(draw, (788, 548, 1106, 690), '目标约束', '2x2 / 2-lane / Q8.8', ORANGE)

    rounded(draw, (1144, 548, 1848, 964), radius=30)
    text(draw, (1178, 584), '真实输出文件', B22, fill=BLUE)
    chip(draw, (1178, 630, 1376, 674), 'ub_map.json', GREEN_BG, GREEN, font=B18)
    chip(draw, (1392, 630, 1608, 674), 'schedule.json', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1624, 630, 1788, 674), 'imem.hex', BLUE_BG, BLUE, font=B18)
    text(draw, (1178, 720), 'UB map 片段', B22, fill=INK)
    yy = 760
    for tensor in ub_map['tensors']:
        if tensor['storage'] != 'ub':
            continue
        if yy > 930:
            break
        line = f"{tensor['name']:<3} addr={tensor['addr']:>2} words={tensor['words']:>2} shape={tensor['shape']}"
        text(draw, (1178, yy), line, M18, fill=INK)
        yy += 30

    rounded(draw, (88, 730, 1088, 964), radius=30)
    text(draw, (122, 764), 'schedule 片段', B22, fill=BLUE)
    sample_lines = []
    for cmd in schedule['commands'][:6] + schedule['commands'][-4:]:
        stage = cmd['stage']
        name = cmd['name']
        kind = cmd['kind']
        extra = ' wait' if cmd.get('wait_after') else ''
        sample_lines.append(f'{stage[:15]:<15} {name[:22]:<22} {kind}{extra}')
    yy = 808
    for line in sample_lines:
        text(draw, (124, yy), line, M18, fill=INK)
        yy += 26

    text(draw, (72, 1002), '面试可追问：为什么是阶段级 schedule，而不是 cycle-accurate？答案是当前目标先把 IMEM 和控制链路打通。', B20, fill=MUTED)
    img.save(path)


def render_frontend(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 5, 'Frontend 详细展开', '把寄存器映射、IMEM、sequencer、decode 和 host write 路径拆开，回答“控制前端到底做了什么”。')

    rounded(draw, (72, 236, 1848, 940), radius=34)

    rounded(draw, (104, 318, 510, 858), radius=28)
    text(draw, (138, 350), '寄存器映射', T34, fill=INK)
    regs = [
        ('0x00', 'CTRL', 'step / start'),
        ('0x04', 'STATUS', 'busy / running'),
        ('0x10', 'INSTR_W0', 'step 模式指令'),
        ('0x20/24', 'UB_DATA/PUSH', 'host 写 UB'),
        ('0x30/34/40/44', 'IMEM_*', '写 IMEM 与长度'),
        ('0x50/54/58', 'LEAK/INV/LR', '运行参数'),
    ]
    yy = 416
    for addr, name, desc in regs:
        rounded(draw, (136, yy, 478, yy + 58), radius=18, fill=HILITE, outline=HILITE, width=1)
        text(draw, (156, yy + 10), addr, B18, fill=BLUE)
        text(draw, (248, yy + 8), name, B18, fill=INK)
        text(draw, (360, yy + 10), desc, B18, fill=MUTED)
        yy += 68

    rounded(draw, (548, 318, 944, 858), radius=28)
    text(draw, (582, 350), 'IMEM 与指令流', T34, fill=INK)
    chip(draw, (582, 408, 708, 452), 'imem[0:63]', BLUE_BG, BLUE, font=B18)
    chip(draw, (724, 408, 882, 452), 'imem_len_reg', GREEN_BG, GREEN, font=B18)
    text(draw, (582, 492), '指令格式', B22, fill=BLUE)
    text(draw, (582, 538), 'NOP / SWITCH / UB_RD / UB_WR_HOST', B20, fill=INK)
    text(draw, (582, 584), 'UB_RD 字段: addr / row / col / transpose / ptr_sel / pathway', B18, fill=MUTED)
    code_strip(draw, (574, 644, 912, 824), 'Frontend 关键点', 'tpu_frontend_axil.sv:6+', [
        'CTRL.step / CTRL.start',
        'IMEM_ADDR / IMEM_W0 / IMEM_WE',
        'IMEM_LEN defines valid program range',
    ], ORANGE)

    rounded(draw, (982, 318, 1380, 858), radius=28)
    text(draw, (1016, 350), '4-state sequencer', T34, fill=INK)
    states = [
        ((1038, 434, 1172, 506), 'IDLE', BLUE_BG, BLUE),
        ((1198, 434, 1332, 506), 'DISPATCH', ORANGE_BG, ORANGE),
        ((1038, 560, 1172, 632), 'WAIT', GREEN_BG, GREEN),
        ((1198, 560, 1332, 632), 'ADVANCE', PURPLE_BG, PURPLE),
    ]
    for xy, msg, fill, tone in states:
        rounded(draw, xy, radius=18, fill=fill, outline=LINE, width=1)
        center(draw, (xy[0]+xy[2])/2, (xy[1]+xy[3])/2, msg, T28, fill=tone)
    arrow(draw, [(1172, 470), (1198, 470)], fill=BLUE)
    arrow(draw, [(1265, 506), (1265, 560)], fill=ORANGE)
    arrow(draw, [(1198, 596), (1172, 596)], fill=GREEN)
    arrow(draw, [(1172, 596), (1172, 470)], fill=PURPLE)
    text(draw, (1016, 692), 'wait_after 触发 WAIT；vpu_drain 到来后才 ADVANCE。', B20, fill=MUTED)
    chip(draw, (1016, 760, 1168, 804), 'seq_instr_pulse', HILITE, BLUE, font=B18)
    chip(draw, (1186, 760, 1318, 804), 'vpu_drain', HILITE, ORANGE, font=B18)

    rounded(draw, (1418, 318, 1816, 858), radius=28)
    text(draw, (1452, 350), 'decode 与写口仲裁', T34, fill=INK)
    text(draw, (1452, 414), 'control_unit 输出:\nub_rd_* / sys_switch / vpu_pathway / host_wr', B20, fill=INK, spacing=10)
    chip(draw, (1452, 520, 1630, 564), 'control_unit.sv', BLUE_BG, BLUE, font=B18)
    chip(draw, (1452, 580, 1642, 624), 'host write mux', GREEN_BG, GREEN, font=B18)
    code_strip(draw, (1444, 656, 1784, 824), '关键 RTL', 'frontend_axil/control_unit', [
        'seq_needs_wait = seq_instr[23];',
        'vpu_pathway_reg <= seq_instr[22:19];',
        'push0 ? host0 : cu_data0;',
    ], GREEN)

    text(draw, (72, 984), '这一页把 frontend 从“控制域方框”展开成真正可讲的前端设计：寄存器、程序存储、状态机和 decode 都各有职责。', B20, fill=MUTED)
    img.save(path)



def render_frontend_split(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 6, 'Frontend 再拆一页', '把 AXI 配置面和 runtime dispatch 面拆开，避免一页塞太满。')

    rounded(draw, (72, 236, 1848, 940), radius=34)

    rounded(draw, (104, 316, 848, 860), radius=30)
    text(draw, (138, 350), '配置面：AXI-Lite + 程序装载', T34, fill=INK)
    text(draw, (138, 404), 'host 视角主要做两件事：把参数写进 UB，把程序写进 IMEM。', B22, fill=MUTED)
    cards = [
        ((138, 492, 360, 598), 'UB host load', ['0x20 UB_DATA', '0x24 UB_PUSH', 'wr_ptr_base 跟随 host 装载推进'], BLUE, BLUE_BG),
        ((386, 492, 608, 598), 'IMEM program', ['0x30 IMEM_ADDR', '0x34 IMEM_W0', '0x40 IMEM_WE / 0x44 IMEM_LEN'], ORANGE, ORANGE_BG),
        ((634, 492, 814, 598), 'Runtime params', ['0x50 leak', '0x54 inv_batch', '0x58 learning_rate'], GREEN, GREEN_BG),
    ]
    for xy, title_s, lines, tone, fill in cards:
        rounded(draw, xy, radius=24, fill=fill, outline=LINE, width=1)
        text(draw, (xy[0] + 18, xy[1] + 18), title_s, B22, fill=tone)
        text(draw, (xy[0] + 18, xy[1] + 54), '\\n'.join(lines), B18, fill=INK, spacing=8)
    code_strip(draw, (138, 650, 812, 826), 'AXI 写通路关键点', 'tpu_frontend_axil.sv:225+', [
        'wr_fire when AW/W both captured',
        '0x24 / 0x28 pulse ub_push0/1',
        '0x40 commits IMEM_W0 into imem[addr]',
    ], BLUE)

    rounded(draw, (892, 316, 1816, 860), radius=30)
    text(draw, (926, 350), '运行面：sequencer + decode + 持续路径', T34, fill=INK)
    steps = [
        ('1', 'start_pulse', '装载 pc=0，拉起 busy'),
        ('2', 'SEQ_DISPATCH', 'seq_instr_pulse 仅打一拍'),
        ('3', 'decode', 'control_unit 译出 ub_rd / switch / pathway'),
        ('4', 'SEQ_WAIT', 'wait_after 时必须等 vpu_drain'),
        ('5', 'SEQ_ADVANCE', 'pc++，进入下一条指令'),
    ]
    yy = 430
    for num, head, body in steps:
        chip(draw, (926, yy, 984, yy + 38), num, BLUE_BG, BLUE, font=B18)
        text(draw, (1008, yy - 2), head, B20, fill=INK)
        text(draw, (1180, yy - 2), body, B18, fill=MUTED)
        yy += 72
    code_strip(draw, (926, 650, 1780, 826), 'runtime 面关键 RTL', 'tpu_frontend_axil.sv:92+', [
        "seq_instr_pulse <= 1'b1 only in dispatch",
        'vpu_pathway_reg latches on every UB_RD',
        'vpu_drain = vpu_valid_prev && !tpu_vpu_valid_in;',
    ], ORANGE)
    chip(draw, (926, 846, 1164, 890), '配置面 = host writes', GREEN_BG, GREEN, font=B18)
    chip(draw, (1184, 846, 1458, 890), '运行面 = sequenced pulses', ORANGE_BG, ORANGE, font=B18)
    text(draw, (72, 984), '这样拆开后，前端就能分成“写配置”和“跑程序”两条叙事线，讲起来更顺。', B20, fill=MUTED)
    img.save(path)


def render_ub(path, ub_map):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 7, 'Unified Buffer 设计', 'UB 不只是内存块，而是 host load、核心读流、VPU 写回和梯度更新的汇合点。')

    rounded(draw, (88, 252, 566, 928), radius=34)
    text(draw, (126, 286), 'UB 地址布局', T34, fill=INK)
    total = 128
    y0 = 352
    h = 462
    draw.rounded_rectangle((148, y0, 262, y0 + h), radius=24, fill=(250, 249, 246), outline=LINE, width=2)
    colors = [BLUE_BG, ORANGE_BG, GREEN_BG, PURPLE_BG, RED_BG, BLUE_BG, ORANGE_BG, GREEN_BG, PURPLE_BG]
    tones = [BLUE, ORANGE, GREEN, PURPLE, RED, BLUE, ORANGE, GREEN, PURPLE]
    labels = []
    ub_tensors = [t for t in ub_map['tensors'] if t['storage'] == 'ub']
    for idx, tensor in enumerate(ub_tensors):
        top = y0 + h * tensor['addr'] / total
        bottom = y0 + h * (tensor['addr'] + tensor['words']) / total
        draw.rounded_rectangle((148, int(top), 262, int(bottom)), radius=12, fill=colors[idx % len(colors)], outline=WHITE, width=1)
        labels.append((tensor, int((top + bottom) / 2), tones[idx % len(tones)]))
    left_x = 300
    right_x = 420
    for idx, (tensor, yy, tone) in enumerate(labels):
        col_x = left_x if idx < 5 else right_x
        row_y = 398 + (idx % 5) * 74
        text(draw, (col_x, row_y), f"{tensor['name']}  @{tensor['addr']}  {tensor['words']}w", B18, fill=tone)
    text(draw, (126, 844), f"allocated={ub_map['allocated_words']}", B22, fill=MUTED)
    text(draw, (126, 878), f"free={ub_map['free_words']}", B22, fill=MUTED)

    rounded(draw, (610, 252, 1118, 534), radius=30)
    text(draw, (626, 286), '两类写入口', T34, fill=INK)
    chip(draw, (626, 350, 852, 396), 'host write: ub_wr_host_*', BLUE_BG, BLUE, font=B18)
    chip(draw, (626, 412, 842, 458), 'VPU writeback: ub_wr_*', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (626, 474, 886, 520), 'wr_ptr_base / wr_ptr_restore', GREEN_BG, GREEN, font=B18)
    text(draw, (626, 528), 'host 阶段先顺序装载；start 时 wr_ptr 恢复，后续写回不会踩到参数区。', B18, fill=MUTED)

    rounded(draw, (610, 568, 1118, 928), radius=30)
    text(draw, (626, 622), '读指针语义', T34, fill=INK)
    rows = [
        ('0', 'input / activation'),
        ('1', 'weight top-load'),
        ('2', 'bias stream'),
        ('3', 'Y label stream'),
        ('4', 'H activation stream'),
        ('5/6', 'in-UB grad update'),
    ]
    yy = 676
    for key, msg in rows:
        chip(draw, (626, yy, 704, yy + 36), key, HILITE, BLUE, font=B18)
        text(draw, (726, yy + 4), msg, B18, fill=INK)
        yy += 38
    chip(draw, (626, 882, 930, 924), 'transpose / counter / lane addr', HILITE, BLUE, font=B18)

    rounded(draw, (1144, 252, 1848, 928), radius=30)
    text(draw, (1178, 286), 'UB 的项目价值', T34, fill=INK)
    text(draw, (1178, 352), '1. 把 host load、core read、VPU writeback 串进同一块存储\n2. 支持 weight top-load / bias/Y/H 多路读法\n3. 把 gradient_descent 放进 UB 内部，参数更新不必绕回 host', B24, fill=INK, spacing=14)
    code_strip(draw, (1170, 516, 1820, 892), 'RTL 关键段', 'unified_buffer_v3.sv:432+', [
        'case (ub_ptr_select)',
        '  0: rd_input_ptr  <= ub_rd_addr_in;',
        '  1: rd_weight_ptr <= ub_rd_addr_in + ...;',
        "  5: grad_bias_or_weight <= 1'b0;",
        "  6: grad_bias_or_weight <= 1'b1;",
        'if (ub_wr_ptr_restore_in) wr_ptr <= wr_ptr_base;',
    ], BLUE)
    text(draw, (72, 980), '这页很适合回答“为什么叫 Unified Buffer，而不是几块分离 SRAM”的问题。', B20, fill=MUTED)
    img.save(path)



def render_wptr(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 8, 'wr_ptr / base / restore', '把 host 装载阶段和训练阶段分开，解释为什么写回不会覆盖参数区。')

    rounded(draw, (72, 236, 1848, 940), radius=34)

    rounded(draw, (102, 314, 700, 856), radius=30)
    text(draw, (136, 348), '核心状态', T34, fill=INK)
    rows = [
        ('wr_ptr', '当前统一写指针，host 与 VPU 写回共用'),
        ('wr_ptr_next', '按 lane valid 预计算下一地址'),
        ('wr_ptr_base', '记录 host 装载完成后的“参数区尾部”'),
        ('ub_wr_ptr_restore_in', 'start_pulse 拉起时恢复 wr_ptr = wr_ptr_base'),
    ]
    yy = 424
    for head, body in rows:
        text(draw, (136, yy), head, B22, fill=BLUE)
        text(draw, (286, yy), body, B20, fill=INK)
        yy += 86
    code_strip(draw, (128, 694, 666, 826), '关键 RTL', 'unified_buffer_v3.sv:225+', [
        'if (ub_wr_host_valid_in[0] || ub_wr_host_valid_in[1])',
        '    wr_ptr_base <= wr_ptr_next;',
        'if (ub_wr_ptr_restore_in) wr_ptr <= wr_ptr_base;',
    ], BLUE)

    rounded(draw, (736, 314, 1318, 856), radius=30)
    text(draw, (770, 348), '阶段视角', T34, fill=INK)
    phases = [
        ((782, 442, 1268, 522), '阶段 1  host load', 'UB_DATA/UB_PUSH 连续装参数，wr_ptr 持续前进', BLUE, BLUE_BG),
        ((782, 556, 1268, 636), '阶段 2  start restore', 'start_pulse 触发 restore，把 wr_ptr 拉回 wr_ptr_base', ORANGE, ORANGE_BG),
        ((782, 670, 1268, 750), '阶段 3  runtime writeback', 'VPU 写回从 base 之后继续写，不踩前面参数', GREEN, GREEN_BG),
    ]
    for xy, title_s, body, tone, fill in phases:
        rounded(draw, xy, radius=22, fill=fill, outline=LINE, width=1)
        text(draw, (xy[0] + 18, xy[1] + 14), title_s, B22, fill=tone)
        text(draw, (xy[0] + 18, xy[1] + 46), body, B18, fill=INK)

    rounded(draw, (1354, 314, 1818, 856), radius=30)
    text(draw, (1388, 348), '为什么单独讲这页', T34, fill=INK)
    text(draw, (1388, 420), '1. 它解释了 Unified Buffer 为什么能既做参数装载，又做训练写回。\n2. 没有 restore，训练一开始就会从参数尾后继续写，host 重新装载和训练写回就难以分界。\n3. 这个设计把“参数静态区”和“运行动态区”用时序而不是物理分 bank 的方式隔开。', B24, fill=INK, spacing=16)
    chip(draw, (1388, 712, 1608, 756), 'host load captures base', BLUE_BG, BLUE, font=B18)
    chip(draw, (1388, 770, 1608, 814), 'start restores ptr', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1388, 828, 1608, 872), 'runtime writes after base', GREEN_BG, GREEN, font=B18)
    text(draw, (72, 984), '面试时这页可以直接回答“为什么写回不会覆盖 W/B 参数”的追问。', B20, fill=MUTED)
    img.save(path)


def render_ub_timing(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 9, 'UB 读流与 PE 时序对齐', '把 UB valid、PE 采样、hold 周期、sequencer wait 这四件事放到一条时间线上。')

    rounded(draw, (72, 236, 1848, 940), radius=34)
    rounded(draw, (104, 304, 1816, 612), radius=30)
    text(draw, (138, 338), '简化时间线', T34, fill=INK)
    x0 = 184
    step = 220
    labels = ['T0 dispatch', 'T1 first beat', 'T2 wavefront', 'T3 last active beat', 'T4 hold / drain']
    for i, label in enumerate(labels):
        x = x0 + i * step
        draw.line((x, 420, x, 560), fill=LINE, width=2)
        chip(draw, (x - 56, 370, x + 56, 410), f'T{i}', BLUE_BG if i % 2 == 0 else ORANGE_BG, BLUE if i % 2 == 0 else ORANGE, font=B18)
        text(draw, (x - 72, 574), label, B18, fill=MUTED)
    arrow(draw, [(x0, 470), (x0 + 4 * step, 470)], fill=BLUE, width=5)
    text(draw, (138, 448), 'ub_rd_start / seq_instr_pulse', B18, fill=BLUE)
    text(draw, (138, 486), 'UB input/weight valid 波前', B18, fill=ORANGE)
    text(draw, (138, 524), 'PE 看到的有效计算窗口', B18, fill=GREEN)
    text(draw, (930, 444), 'hold 周期保留最后一拍输入，\n而 weight 流会主动清 valid，避免 PE22 被重复装载。', B20, fill=INK, spacing=8)

    rounded(draw, (104, 648, 858, 892), radius=30)
    text(draw, (138, 682), 'UB 侧关键语义', T34, fill=INK)
    text(draw, (138, 742), '1. input 流在最后一个 active beat 后还有 1 个 hold 周期，保证阵列最后一列能完成传播。\n2. weight 流故意不保留最后一个 valid，防止 loader 在多打一拍时覆盖 active weight。\n3. 所以“同样是 UB 读”，input 和 weight 的尾拍策略不一样。', B22, fill=INK, spacing=14)

    rounded(draw, (892, 648, 1816, 892), radius=30)
    text(draw, (926, 682), '和 sequencer 的关系', T34, fill=INK)
    text(draw, (926, 742), '1. dispatch 只负责发起读流，不代表计算已经结束。\n2. 真正的完成边界在 VPU drain，sequencer 只有等 drain 才能 ADVANCE。\n3. 这就是 wait_after 的意义：用系统级同步点把 UB/PE/VPU 多拍行为收口。', B22, fill=INK, spacing=14)
    chip(draw, (926, 840, 1096, 882), 'dispatch != done', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1116, 840, 1298, 882), 'drain = safe advance', GREEN_BG, GREEN, font=B18)
    text(draw, (72, 984), '这页把数据流和控制流第一次放在同一张时间图上，比单讲波形或单讲 RTL 更容易讲清楚。', B20, fill=MUTED)
    img.save(path)


def render_ub_update(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 10, 'UB 内梯度下降更新', '重点讲清楚 bias update 和 weight update 不是同一种语义：一个可连续累加，一个按 tile 用旧值更新。')

    rounded(draw, (72, 236, 1848, 940), radius=34)

    rounded(draw, (102, 320, 844, 864), radius=30)
    text(draw, (138, 354), 'Bias Update 路径', T34, fill=INK)
    text(draw, (138, 412), 'ptr_sel=5，grad_bias_or_weight=0', B22, fill=BLUE)
    chip(draw, (138, 474, 286, 518), 'dZ / dZ1 stream', BLUE_BG, BLUE, font=B18)
    chip(draw, (310, 474, 446, 518), 'old bias', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (470, 474, 610, 518), 'lr', GREEN_BG, GREEN, font=B18)
    arrow(draw, [(286, 496), (310, 496)], fill=BLUE)
    arrow(draw, [(446, 496), (470, 496)], fill=ORANGE)
    arrow(draw, [(610, 496), (694, 496)], fill=GREEN)
    rounded(draw, (694, 448, 800, 544), radius=20, fill=HILITE, outline=LINE, width=1)
    center(draw, 747, 496, 'GD', T30, fill=INK)
    text(draw, (138, 576), 'bias 更新逻辑', B22, fill=BLUE)
    text(draw, (138, 622), 'bias 模式下，若上一拍 done 过，就把 sub_in_a 切到 value_updated_out。\n这意味着同一个 bias 地址可以在连续样本梯度到来时做逐拍累加式更新。', B22, fill=INK, spacing=14)
    code_strip(draw, (130, 726, 812, 840), 'gradient_descent.sv', 'gradient_descent.sv:37+', [
        "if (grad_bias_or_weight == 1'b0)",
        '    sub_in_a = grad_descent_done_out ? value_updated_out : value_old_in;',
    ], BLUE)

    rounded(draw, (890, 320, 1816, 864), radius=30)
    text(draw, (926, 354), 'Weight Update 路径', T34, fill=INK)
    text(draw, (926, 412), 'ptr_sel=6，grad_bias_or_weight=1', B22, fill=ORANGE)
    chip(draw, (926, 474, 1102, 518), 'outer-product grad', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1126, 474, 1262, 518), 'old weight', BLUE_BG, BLUE, font=B18)
    chip(draw, (1286, 474, 1426, 518), 'lr', GREEN_BG, GREEN, font=B18)
    arrow(draw, [(1102, 496), (1126, 496)], fill=ORANGE)
    arrow(draw, [(1262, 496), (1286, 496)], fill=BLUE)
    arrow(draw, [(1426, 496), (1510, 496)], fill=GREEN)
    rounded(draw, (1510, 448, 1616, 544), radius=20, fill=HILITE, outline=LINE, width=1)
    center(draw, 1563, 496, 'GD', T30, fill=INK)
    text(draw, (926, 576), 'weight 更新逻辑', B22, fill=ORANGE)
    text(draw, (926, 622), 'weight 模式下，sub_in_a 固定取 value_old_in，不会链接到上次 value_updated_out。\n也就是说每个 weight 更新都基于当前 tile 读出的旧权重做一次独立更新，更适合 tile 级 outer-product。', B22, fill=INK, spacing=14)
    code_strip(draw, (918, 726, 1784, 840), 'gradient_descent.sv', 'gradient_descent.sv:45+', [
        "if (grad_bias_or_weight == 1'b1)",
        '    sub_in_a = value_old_in;',
        'grad_descent_done_out <= grad_descent_valid_in;',
    ], ORANGE)

    chip(draw, (230, 892, 536, 936), 'bias: sample-by-sample accumulate', BLUE_BG, BLUE, font=B18)
    chip(draw, (1018, 892, 1312, 936), 'weight: tile-based update', ORANGE_BG, ORANGE, font=B18)
    text(draw, (72, 984), '这页就是你刚指出的关键差异：bias update 和 weight update 用的是同一个 GD 模块，但驱动语义并不一样。', B20, fill=MUTED)
    img.save(path)


def render_vpu(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 11, 'VPU 单独展开', 'VPU 不是一个黑盒后处理，而是按 pathway bit 组织 bias / leaky-relu / loss-grad / derivative 的可重组链路。')

    rounded(draw, (72, 236, 1848, 940), radius=34)

    rounded(draw, (102, 320, 1130, 860), radius=30)
    text(draw, (138, 354), 'VPU 子模块与主数据链', T34, fill=INK)
    blocks = [
        ((160, 506, 330, 590), 'sys in', BLUE_BG, BLUE),
        ((374, 506, 544, 590), 'Bias', BLUE_BG, BLUE),
        ((588, 506, 758, 590), 'LReLU', GREEN_BG, GREEN),
        ((802, 506, 972, 590), 'Loss grad', ORANGE_BG, ORANGE),
        ((1016, 506, 1100, 590), 'dAct', PURPLE_BG, PURPLE),
    ]
    for xy, msg, fill, tone in blocks:
        rounded(draw, xy, radius=20, fill=fill, outline=LINE, width=1)
        center(draw, (xy[0]+xy[2])/2, (xy[1]+xy[3])/2, msg, T28, fill=tone)
    arrow(draw, [(330, 548), (374, 548)], fill=BLUE)
    arrow(draw, [(544, 548), (588, 548)], fill=GREEN)
    arrow(draw, [(758, 548), (802, 548)], fill=ORANGE)
    arrow(draw, [(972, 548), (1016, 548)], fill=PURPLE)
    text(draw, (138, 648), 'pathway 编码', B22, fill=BLUE)
    text(draw, (138, 694), '1100: forward = sys -> bias -> lrelu\n1111: transition = sys -> bias -> lrelu -> loss -> dAct\n0001: backward = sys -> dAct', B22, fill=INK, spacing=14)
    chip(draw, (138, 790, 318, 834), 'last_H cache', HILITE, BLUE, font=B18)
    chip(draw, (338, 790, 496, 834), 'Y from UB', HILITE, ORANGE, font=B18)
    chip(draw, (516, 790, 722, 834), 'leak factor from UB', HILITE, GREEN, font=B18)

    rounded(draw, (1170, 320, 1816, 860), radius=30)
    text(draw, (1206, 354), '为什么值得单独讲', T34, fill=INK)
    text(draw, (1206, 420), '1. 同一套子模块复用三种训练阶段，不是写死单一路径。\n2. 1111 路径里要同时保存 H 并继续送入 loss / derivative。\n3. pathway bit 和 frontend 里的保持寄存器直接耦合，是系统级控制点。', B24, fill=INK, spacing=16)
    code_strip(draw, (1198, 612, 1782, 836), 'vpu.sv 关键点', 'vpu.sv:4+', [
        '1100: forward pass',
        '1111: transition path',
        '0001: backward pass',
        'last_H cache feeds leaky_relu_derivative',
    ], PURPLE)

    text(draw, (72, 984), '这页能把 VPU 从“几个小算子拼一起”提升成“可重组训练路径单元”的层次。', B20, fill=MUTED)
    img.save(path)


def render_pe(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 12, 'PE 与计算阵列', 'PE 是最小计算单元；2x2 systolic 把 valid、input、weight、psum 做成波前传播。')

    rounded(draw, (76, 244, 928, 936), radius=34)
    text(draw, (114, 280), '2 x 2 Systolic 组织', T34, fill=INK)

    pe_boxes = {
        'pe11': (170, 396, 386, 572),
        'pe12': (454, 396, 670, 572),
        'pe21': (170, 646, 386, 822),
        'pe22': (454, 646, 670, 822),
    }
    for name, xy in pe_boxes.items():
        rounded(draw, xy, radius=24, fill=WHITE, outline=LINE, width=2)
        center(draw, (xy[0] + xy[2]) / 2, xy[1] + 42, name, T30, fill=INK)
        center(draw, (xy[0] + xy[2]) / 2, xy[1] + 104, 'active w\ninactive w\nmac + valid', B20, fill=MUTED)
    arrow(draw, [(122, 484), (170, 484)], fill=BLUE)
    arrow(draw, [(122, 734), (170, 734)], fill=BLUE)
    arrow(draw, [(278, 572), (278, 646)], fill=ORANGE)
    arrow(draw, [(562, 572), (562, 646)], fill=ORANGE)
    arrow(draw, [(386, 484), (454, 484)], fill=GREEN)
    arrow(draw, [(386, 734), (454, 734)], fill=GREEN)
    arrow(draw, [(278, 352), (278, 396)], fill=PURPLE)
    arrow(draw, [(562, 352), (562, 396)], fill=PURPLE)
    text(draw, (94, 470), 'input', B18, fill=BLUE)
    text(draw, (94, 720), 'input', B18, fill=BLUE)
    text(draw, (252, 324), 'weight', B18, fill=PURPLE)
    text(draw, (536, 324), 'weight', B18, fill=PURPLE)
    text(draw, (690, 472), 'valid / input', B18, fill=GREEN)
    text(draw, (690, 722), 'output', B18, fill=GREEN)
    chip(draw, (120, 866, 284, 910), 'vertical psum', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (304, 866, 502, 910), 'horizontal data', BLUE_BG, BLUE, font=B18)
    chip(draw, (522, 866, 746, 910), 'shadow -> active weight', PURPLE_BG, PURPLE, font=B18)

    rounded(draw, (972, 244, 1848, 936), radius=34)
    text(draw, (1008, 280), 'PE 级 RTL 重点', T34, fill=INK)
    code_strip(draw, (1000, 348, 1814, 628), 'pe.sv 的三个关键点', 'pe.sv:31+', [
        'fxp_mul(ina=pe_input_in, inb=weight_reg_active)',
        'fxp_add(ina=mult_out, inb=pe_psum_in)',
        'if (pe_switch_in) weight_reg_active <= weight_reg_inactive;',
        'if (pe_accept_w_in) weight_reg_inactive <= pe_weight_in;',
        'if (!pe_enabled) outputs clear, state keeps;',
    ], BLUE)
    text(draw, (1008, 670), '为什么这个 PE 设计值得讲', B22, fill=BLUE)
    text(draw, (1008, 720), '1. 双 weight reg 支持 load / switch 两阶段，不需要计算中断\n2. pe_enabled 只清输出不清内部状态，便于列裁剪\n3. valid 与 switch 都做波前传播，阵列行为可在波形里直接观察', B24, fill=INK, spacing=14)
    chip(draw, (1008, 858, 1220, 904), 'weight-stationary', BLUE_BG, BLUE, font=B18)
    chip(draw, (1240, 858, 1456, 904), 'column enable', GREEN_BG, GREEN, font=B18)
    chip(draw, (1476, 858, 1714, 904), 'pulse-based control', ORANGE_BG, ORANGE, font=B18)
    text(draw, (72, 980), '面试讲法：先讲阵列拓扑，再讲单个 PE 的 active/inactive weight，是最容易让面试官接着追问的路径。', B20, fill=MUTED)
    img.save(path)


def render_control(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 13, '关键控制 RTL 特写', '把最容易出系统级 bug 的三处控制逻辑拿出来讲清楚。')

    cards = [
        ((72, 232, 608, 944), BLUE, BLUE_BG, '显式等待语义', [
            ('现象', 'sequencer 可能提前推进'),
            ('根因', 'VPU 延迟不是固定周期'),
            ('修复', 'wait_after + vpu_drain'),
            ('价值', '让 IMEM 指令有显式同步点'),
        ], 'frontend_axil.sv', [
            'seq_needs_wait = seq_instr[23];',
            'vpu_drain = vpu_valid_prev && !tpu_vpu_valid_in;',
            'if (seq_needs_wait) seq_state <= SEQ_WAIT;',
        ]),
        ((692, 232, 1228, 944), ORANGE, ORANGE_BG, '路径保持', [
            ('现象', 'dispatch 后几拍 route 会跑偏'),
            ('根因', 'decode 输出只持续 1 拍'),
            ('修复', 'vpu_pathway_reg 持久保持'),
            ('价值', '多拍 VPU 仍走同一处理链'),
        ], 'frontend_axil.sv', [
            'if (seq_instr_pulse && opcode == UB_RD)',
            '    vpu_pathway_reg <= seq_instr[22:19];',
            'vpu_data_pathway_out = vpu_pathway_reg;',
        ]),
        ((1312, 232, 1848, 944), GREEN, GREEN_BG, 'Host / CU 仲裁', [
            ('现象', 'host 装载与 CU 写口重叠'),
            ('根因', '同一 UB 入口承载两类语义'),
            ('修复', 'UB_PUSH 优先的 host_write_mux'),
            ('价值', '参数装载与训练写回可共存'),
        ], 'frontend_axil.sv', [
            'valid0 = push0 ? one : cu_valid0;',
            'valid1 = push1 ? one : cu_valid1;',
            'data0  = push0 ? host0 : cu_data0;',
        ]),
    ]
    for xy, tone, fill, title_s, rows, ref, code in cards:
        l, t, r, b = xy
        rounded(draw, xy, radius=30)
        draw.rounded_rectangle((l + 20, t + 22, l + 28, b - 22), radius=4, fill=tone)
        text(draw, (l + 52, t + 26), title_s, T34, fill=INK)
        chip(draw, (l + 52, t + 84, l + 202, t + 128), '深挖点', fill, tone, font=B18)
        yy = t + 174
        for k, v in rows:
            text(draw, (l + 52, yy), k, B18, fill=tone)
            text(draw, (l + 120, yy), v, B18, fill=INK)
            yy += 54
        code_strip(draw, (l + 34, b - 214, r - 28, b - 34), 'RTL 特写', ref, code, tone)

    rounded(draw, (72, 964, 1848, 1032), radius=26, fill=PANEL, outline=PANEL, width=1)
    text(draw, (104, 984), '关于 clock-gating：当前 RTL 里没有显式 ICG cell；现在更准确的说法是用 pe_enabled / valid / pulse 做时序收口。若继续工程化，可在 PE/UB/VPU 边界插标准门控。', B20, fill=MUTED)
    img.save(path)


def build_cycle_frames():
    frames = []
    titles = [
        'Cycle 0  权重已装载到 active reg',
        'Cycle 1  第一拍输入进入左边界',
        'Cycle 2  波前推进，部分和开始下沉',
        'Cycle 3  第一列结果到达底部',
        'Cycle 4  第二列结果完成',
        'Cycle 5  结果写回 UB / 进入后续 VPU 路径',
    ]
    notes = [
        '示意矩阵: A=[[1,2],[3,4]], W=[[5,6],[7,8]]',
        '行数据从左进入，权重常驻在 PE 内',
        '同一拍里不同 PE 看到的是不同阶段数据',
        'C[:,0] 开始可见',
        'C = [[19,22],[43,50]]',
        '这就是后续 bias/loss/dAct 的输入来源',
    ]
    pe_vals = [
        ['w=5', 'w=6', 'w=7', 'w=8'],
        ['a=1\npsum=5', ' ', 'a=3\npsum=15', ' '],
        ['a=2\npsum=19', 'a=1\npsum=6', 'a=4\npsum=43', 'a=3\npsum=18'],
        ['done\n19', 'a=2\npsum=22', 'done\n43', 'a=4\npsum=50'],
        ['19', '22', '43', '50'],
        ['WB->UB', 'WB->UB', 'WB->UB', 'WB->UB'],
    ]
    for idx, title_s in enumerate(titles):
        img = Image.new('RGB', (1280, 720), BG)
        draw = ImageDraw.Draw(img)
        text(draw, (56, 46), title_s, T44, fill=INK)
        text(draw, (56, 108), notes[idx], B24, fill=MUTED)
        boxes = [(240, 200, 460, 380), (530, 200, 750, 380), (240, 430, 460, 610), (530, 430, 750, 610)]
        names = ['PE11', 'PE12', 'PE21', 'PE22']
        tones = [BLUE, ORANGE, GREEN, PURPLE]
        for bxy, name, payload, tone in zip(boxes, names, pe_vals[idx], tones):
            rounded(draw, bxy, radius=24, fill=WHITE, outline=LINE, width=2)
            center(draw, (bxy[0] + bxy[2]) / 2, bxy[1] + 36, name, T28, fill=tone)
            center(draw, (bxy[0] + bxy[2]) / 2, (bxy[1] + bxy[3]) / 2 + 12, payload, B24, fill=INK)
        arrow(draw, [(160, 290), (240, 290)], fill=BLUE)
        arrow(draw, [(160, 520), (240, 520)], fill=BLUE)
        arrow(draw, [(350, 380), (350, 430)], fill=ORANGE)
        arrow(draw, [(640, 380), (640, 430)], fill=ORANGE)
        arrow(draw, [(460, 290), (530, 290)], fill=GREEN)
        arrow(draw, [(460, 520), (530, 520)], fill=GREEN)
        center(draw, 1050, 250, '结果矩阵 C', T34, fill=INK)
        rounded(draw, (930, 300, 1170, 540), radius=22, fill=HILITE, outline=LINE, width=1)
        cell = [('19', 985, 360), ('22', 1110, 360), ('43', 985, 480), ('50', 1110, 480)]
        draw.line((1050, 310, 1050, 530), fill=LINE, width=2)
        draw.line((940, 420, 1160, 420), fill=LINE, width=2)
        for msg, x, y in cell:
            center(draw, x, y, msg if idx >= 4 else ('?' if idx < 3 else ('19' if msg == '19' else '?' if msg in ('22', '43', '50') else '?')), T34, fill=BLUE if msg in ('19', '22') else GREEN)
        if idx == 3:
            center(draw, 985, 480, '43', T34, fill=GREEN)
        if idx < 4:
            chip(draw, (900, 590, 1190, 634), '示意动画，不代替真实波形', ORANGE_BG, ORANGE, font=B18)
        else:
            chip(draw, (900, 590, 1190, 634), '结果已完整可见', GREEN_BG, GREEN, font=B18)
        frames.append(img)
    return frames


def render_cycle(path):
    frames = build_cycle_frames()
    strip = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(strip)
    draw_header(draw, 14, '逐拍计算动态', '这页不是 RTL 波形，而是帮助面试官快速理解 2x2 systolic 波前如何形成结果。')
    thumb_w, thumb_h = 270, 152
    xs = [84, 390, 696, 1002, 1308, 1614]
    for idx, (frame, x) in enumerate(zip(frames, xs)):
        thumb = frame.resize((thumb_w, thumb_h))
        rounded(draw, (x, 300, x + thumb_w, 452), radius=18, fill=WHITE, outline=LINE, width=2)
        strip.paste(thumb, (x, 300))
        chip(draw, (x + 10, 464, x + 84, 502), f'C{idx}', BLUE_BG if idx % 2 == 0 else ORANGE_BG, BLUE if idx % 2 == 0 else ORANGE, font=B18)
    text(draw, (88, 556), '讲法建议', B22, fill=BLUE)
    text(draw, (88, 604), '1. 先说 weight-stationary，权重先装后切换\n2. 再说 input 从左向右、psum 自上而下\n3. 最后落到“为什么需要 wait_after 和 drain 才能收口”', T38, fill=INK, spacing=16)
    rounded(draw, (1088, 570, 1840, 930), radius=30)
    text(draw, (1124, 608), '配套文件', T34, fill=INK)
    chip(draw, (1124, 676, 1402, 722), 'tpu_systolic_cycle_demo.gif', BLUE_BG, BLUE, font=B18)
    chip(draw, (1124, 738, 1402, 784), 'tpu_systolic_cycle_strip.png', GREEN_BG, GREEN, font=B18)
    text(draw, (1124, 814), 'GIF 可以单独打开给面试官看；PPT 里放静态条带更稳，避免不同查看器对 GIF 支持不一致。', B24, fill=MUTED, spacing=10)
    text(draw, (72, 984), '这页是“帮助讲清楚原理”的辅助页；真正证明正确性还是看下一页的真实波形与 scoreboard。', B20, fill=MUTED)
    strip.save(path)
    STRIP_PATH.write_bytes(path.read_bytes())
    frames[0].save(GIF_PATH, save_all=True, append_images=frames[1:], duration=700, loop=0)


def render_wave(path):
    names = ['busy_reg', 'seq_instr_pulse', 'ub_rd_start', 'sys_switch', 'ub_rd_input_valid_out_0', 'sys_valid_out_21', 'vpu_valid_out_1', 'vpu_drain']
    events = parse_vcd(set(names))
    start, end = first_busy_window(events)
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 15, '验证波形与回归覆盖', '这一页给真实时序证据：sequencer 脉冲、UB 发数、阵列输出、VPU drain 都能在 VCD 里对上。')

    rounded(draw, (72, 244, 1260, 950), radius=34)
    text(draw, (108, 278), 'VCD 片段：第一次 busy 窗口', T34, fill=INK)
    text(draw, (108, 328), '来源: waveforms/tpu_soc.vcd', B20, fill=MUTED)
    y = 400
    tones = [BLUE, ORANGE, GREEN, PURPLE, BLUE, ORANGE, GREEN, RED]
    for name, tone in zip(names, tones):
        segs = sample_signal(events[name], start, end)
        draw_digital_wave(draw, 330, y, 860, 44, segs, start, end, tone, name)
        y += 62
    draw_time_axis(draw, 330, 890, 860, start, end)
    chip(draw, (366, 920, 520, 964), 'dispatch', BLUE_BG, BLUE, font=B18)
    chip(draw, (542, 920, 696, 964), 'ub read', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (718, 920, 884, 964), 'sys out', GREEN_BG, GREEN, font=B18)
    chip(draw, (906, 920, 1048, 964), 'drain', RED_BG, RED, font=B18)

    rounded(draw, (1304, 244, 1848, 560), radius=30)
    text(draw, (1340, 280), '41 / 41 覆盖项', T34, fill=INK)
    rows = [
        ('H1', '8'), ('dZ2', '4'), ('dZ1', '8'),
        ('UB dZ2', '4'), ('UB dZ1', '8'), ('W1/B1/W2/B2', '9'),
    ]
    yy = 336
    for label, num in rows:
        chip(draw, (1340, yy, 1450, yy + 34), num, BLUE_BG, BLUE, font=B18)
        text(draw, (1470, yy + 2), label, B18, fill=INK)
        yy += 40

    rounded(draw, (1304, 596, 1848, 950), radius=30)
    text(draw, (1340, 632), '为什么这页有说服力', T34, fill=INK)
    text(draw, (1340, 690), '1. 控制脉冲和输出数据在同一窗口里对齐。\n2. 可以沿着 busy -> dispatch -> ub read -> sys out -> drain 这条主线讲。\n3. scoreboard 又补上了 H1 / dZ2 / dZ1 / UB 更新的数值证据。', B22, fill=INK, spacing=14)
    chip(draw, (1340, 844, 1490, 888), 'dispatch', BLUE_BG, BLUE, font=B18)
    chip(draw, (1506, 844, 1656, 888), 'ub read', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (1672, 844, 1808, 888), 'sys out', GREEN_BG, GREEN, font=B18)
    text(draw, (1340, 906), '口径更像 Verdi：先看状态机和关键脉冲，再回到 scoreboard 验证结果。', B18, fill=MUTED)
    img.save(path)


def render_results(path):
    img = Image.new('RGB', (W, H), BG)
    draw = ImageDraw.Draw(img)
    draw_header(draw, 16, '结果、边界与追问方向', '最后一页给出结论，同时主动交代边界，让面试官知道还能往哪里深挖。')

    metric_box(draw, (72, 236, 374, 388), '回归', '41 / 41 PASS', BLUE)
    metric_box(draw, (400, 236, 702, 388), '训练', '12 epoch', ORANGE)
    metric_box(draw, (728, 236, 1114, 388), 'loss', '0.2529 -> 0.1777', GREEN)
    metric_box(draw, (1140, 236, 1498, 388), 'XOR', '(0, 1, 1, 0)', PURPLE)

    draw_loss_chart(draw, (72, 458, 846, 892))

    rounded(draw, (892, 458, 1848, 686), radius=30)
    text(draw, (930, 494), '项目边界', T34, fill=INK)
    text(draw, (930, 552), '1. 当前目标是 tiny-tpu 原型，不夸大成完整商用 NPU SoC\n2. 当前最扎实的链路是 2x2 / Q8.8 / 2-layer MLP / XOR\n3. 还没有 DMA / IRQ / ICG / 大阵列 tile 化', B24, fill=INK, spacing=14)

    rounded(draw, (892, 720, 1848, 950), radius=30)
    text(draw, (930, 756), '适合引导追问的方向', T34, fill=INK)
    text(draw, (930, 812), '1. 如果扩到 8x8，UB 和 schedule 会怎么变\n2. 如果做 clock-gating / power-aware，需要在哪些边界插 ICG\n3. 如果做 DMA / 中断，frontend 和 STATUS/CTRL 如何演进', B24, fill=INK, spacing=14)

    chip(draw, (72, 942, 308, 992), '系统闭环', BLUE_BG, BLUE)
    chip(draw, (332, 942, 598, 992), '真实 RTL 证据', ORANGE_BG, ORANGE)
    chip(draw, (622, 942, 878, 992), '波形 + scoreboard', GREEN_BG, GREEN)
    chip(draw, (902, 942, 1172, 992), '边界主动说明', PURPLE_BG, PURPLE)
    text(draw, (72, 1010), '收束口径：这是一个“系统集成 + 控制执行 + 验证闭环”完整展示的项目，而不是只讲某个单点算子。', B20, fill=MUTED)
    img.save(path)


def _rgb(color):
    return RGBColor(*color)


def _add_textbox(slide, x, y, w, h, text_value, size, color, bold=False, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text_value
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = PPT_FONT_NAME
    run.font.color.rgb = _rgb(color)
    return box


def _add_panel(slide, x, y, w, h, fill, line=LINE, radius=0):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(line)
    shape.line.width = Pt(0.9 if line == LINE else 1.0)
    if radius == 0:
        shape.adjustments[0] = 0.04
    return shape

def _add_card(slide, x, y, w, h, title, body, tone=BLUE, fill=WHITE, body_size=12):
    _add_panel(slide, x, y, w, h, fill)
    _add_textbox(slide, x + 0.18, y + 0.12, w - 0.3, 0.24, title, 14, tone, True)
    _add_textbox(slide, x + 0.18, y + 0.42, w - 0.32, h - 0.52, body, body_size, INK)


def _add_code_card(slide, x, y, w, h, title, lines, tone=ORANGE):
    _add_panel(slide, x, y, w, h, WHITE)
    _add_textbox(slide, x + 0.22, y + 0.16, w - 0.36, 0.22, title, 14, tone, True)
    box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.48), Inches(w - 0.4), Inches(h - 0.66))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = '\n'.join(lines)
    run.font.size = Pt(10.5)
    run.font.name = PPT_MONO_FONT_NAME
    run.font.color.rgb = _rgb(INK)


def _add_metric_card(slide, x, y, w, h, title, value, tone=BLUE):
    _add_panel(slide, x, y, w, h, WHITE)
    _add_textbox(slide, x + 0.2, y + 0.16, w - 0.32, 0.16, title, 10, MUTED)
    _add_textbox(slide, x + 0.2, y + 0.44, w - 0.32, 0.3, value, 18, tone, True)

def _add_picture(slide, img_path, x, y, w, h):
    if Path(img_path).exists():
        slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))


def _add_picture_named(slide, img_path, x, y, w, h, name):
    if not Path(img_path).exists():
        return None
    shape = slide.shapes.add_picture(str(img_path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    _set_shape_name(shape, name)
    return shape


def _add_picture_panel(slide, img_path, x, y, w, h, title=None, tone=BLUE):
    _add_panel(slide, x, y, w, h, WHITE)
    if title:
        _add_textbox(slide, x + 0.14, y + 0.08, w - 0.28, 0.18, title, 11, tone, True)
        _add_picture(slide, img_path, x + 0.12, y + 0.34, w - 0.24, h - 0.44)
    else:
        _add_picture(slide, img_path, x + 0.1, y + 0.1, w - 0.2, h - 0.2)


def _add_chip_ppt(slide, x, y, w, text_value, fill=BLUE_BG, tone=BLUE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(0.32))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(fill)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text_value
    run.font.size = Pt(9)
    run.font.bold = True
    run.font.name = PPT_FONT_NAME
    run.font.color.rgb = _rgb(tone)


def _render_header(slide, idx, title, subtitle):
    _add_textbox(slide, 0.62, 0.28, 0.5, 0.18, f'P{idx}', 10, BLUE, True)
    _add_textbox(slide, 0.62, 0.5, 7.4, 0.34, title, 23, INK, True)
    _add_textbox(slide, 0.62, 0.86, 9.5, 0.24, subtitle, 11, MUTED)
    _add_panel(slide, 11.02, 0.32, 1.65, 0.32, BLUE_BG, BLUE_BG)
    _add_textbox(slide, 11.22, 0.4, 1.15, 0.14, f'{idx}/{PPT_TOTAL_PAGES}', 10, BLUE, True)



def _add_arrow_ppt(slide, x1, y1, x2, y2, tone=BLUE, width=2.2):
    line = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    line.line.color.rgb = _rgb(tone)
    line.line.width = Pt(width)
    try:
        line.line.end_arrowhead = True
    except Exception:
        pass
    return line


def _add_mini_box(slide, x, y, w, h, label, fill, tone=INK, size=10):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(LINE)
    frame = shape.text_frame
    frame.clear()
    frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = label
    run.font.size = Pt(size)
    run.font.bold = True
    run.font.name = PPT_FONT_NAME
    run.font.color.rgb = _rgb(tone)
    return shape



def _set_shape_name(shape, name):
    for node in shape._element.iter():
        if node.tag.endswith('cNvPr'):
            node.set('name', name)
            break
    return shape


def _add_outline_box(slide, x, y, w, h, tone=BLUE, width=3.2):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background()
    shape.line.color.rgb = _rgb(tone)
    shape.line.width = Pt(width)
    return shape


def _add_anim_note_box(slide, x, y, w, h, title, body, tone=BLUE, fill=BLUE_BG):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(fill)
    shape.line.color.rgb = _rgb(fill)
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = MSO_ANCHOR.TOP
    p = frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.name = PPT_FONT_NAME
    run.font.color.rgb = _rgb(tone)
    p2 = frame.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = body
    run2.font.size = Pt(8.5)
    run2.font.name = 'Noto Sans SC'
    run2.font.color.rgb = _rgb(INK)
    return shape



def _rgba(rgb, alpha=255):
    return (rgb[0], rgb[1], rgb[2], alpha)


def _draw_pe_demo_base(draw, w, h):
    rounded(draw, (8, 8, w - 8, h - 8), radius=28, fill=WHITE, outline=LINE, width=2)
    boxes = {
        'PE11': (250, 205, 430, 345),
        'PE12': (500, 205, 680, 345),
        'PE21': (250, 425, 430, 565),
        'PE22': (500, 425, 680, 565),
    }
    for label, (l, t, r, b) in boxes.items():
        rounded(draw, (l, t, r, b), radius=24, fill=WHITE, outline=LINE, width=2)
        text(draw, (l + 20, t + 18), label, T28, fill=INK)
        text(draw, (l + 20, t + 58), 'active w\ninactive w\nMAC + valid', B18, fill=MUTED)

    result_box = (840, 210, 1030, 380)
    ub_box = (860, 455, 1025, 555)
    rounded(draw, result_box, radius=22, fill=WHITE, outline=LINE, width=2)
    rounded(draw, ub_box, radius=22, fill=WHITE, outline=LINE, width=2)
    text(draw, (870, 228), '结果窗', B20, fill=MUTED)
    center(draw, 942, 505, 'UB / VPU', B22, fill=GREEN)

    arrow(draw, [(135, 275), (250, 275)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(135, 495), (250, 495)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(430, 275), (500, 275)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(430, 495), (500, 495)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(330, 120), (330, 205)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(580, 120), (580, 205)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(330, 345), (330, 425)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(580, 345), (580, 425)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(680, 275), (840, 275)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(680, 495), (840, 495)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(935, 380), (935, 455)], fill=MUTED, width=6, arrow_size=18)

    chip(draw, (250, 78, 395, 120), 'weight', PURPLE_BG, PURPLE, font=B20)
    chip(draw, (92, 252, 222, 294), 'input', BLUE_BG, BLUE, font=B20)
    chip(draw, (92, 472, 222, 514), 'input', BLUE_BG, BLUE, font=B20)
    chip(draw, (700, 252, 860, 294), 'sys out', GREEN_BG, GREEN, font=B20)
    chip(draw, (700, 472, 860, 514), 'sys out', GREEN_BG, GREEN, font=B20)
    return boxes, result_box, ub_box


def _draw_pe_demo_overlay(draw, boxes, result_box, ub_box, step, progress):
    def outline(box, tone, fill=None, width=6):
        draw.rounded_rectangle(box, radius=24, outline=tone, fill=fill, width=width)

    def partial_h(x1, y, x2, tone, frac):
        xm = int(x1 + (x2 - x1) * max(0.08, frac))
        arrow(draw, [(x1, y), (xm, y)], fill=tone, width=10, arrow_size=20)

    def partial_v(x, y1, y2, tone, frac):
        ym = int(y1 + (y2 - y1) * max(0.08, frac))
        arrow(draw, [(x, y1), (x, ym)], fill=tone, width=10, arrow_size=20)

    pe11 = boxes['PE11']
    pe12 = boxes['PE12']
    pe21 = boxes['PE21']
    pe22 = boxes['PE22']

    if step == 0:
        tone = _rgba(PURPLE, 255)
        fill = _rgba(PURPLE_BG, int(150 * progress))
        outline(pe11, tone, fill)
        outline(pe12, tone, fill)
        partial_v(330, 120, 205, tone, progress)
        partial_v(580, 120, 205, tone, progress)
    elif step == 1:
        tone = _rgba(BLUE, 255)
        fill = _rgba(BLUE_BG, int(150 * progress))
        outline(pe11, tone, fill)
        outline(pe21, tone, fill)
        partial_h(135, 275, 250, tone, progress)
        partial_h(135, 495, 250, tone, progress)
    elif step == 2:
        tone_h = _rgba(GREEN, 255)
        tone_v = _rgba(ORANGE, 255)
        fill = _rgba(GREEN_BG, int(150 * progress))
        outline(pe12, tone_h, fill)
        outline(pe22, tone_h, fill)
        partial_h(430, 275, 500, tone_h, progress)
        partial_h(430, 495, 500, tone_h, progress)
        partial_v(330, 345, 425, tone_v, progress)
        partial_v(580, 345, 425, tone_v, progress)
    elif step == 3:
        tone = _rgba(ORANGE, 255)
        outline(result_box, tone, _rgba(ORANGE_BG, int(120 * progress)), width=5)
        partial_h(680, 495, 840, tone, progress)
        text(draw, (882, 266), '19', T30, fill=tone)
        if progress > 0.45:
            text(draw, (882, 326), '43', T30, fill=tone)
    elif step == 4:
        tone = _rgba(RED, 255)
        outline(result_box, tone, _rgba(RED_BG, int(140 * progress)), width=6)
        text(draw, (882, 266), '19', T30, fill=INK)
        text(draw, (962, 266), '22', T30, fill=tone)
        text(draw, (882, 326), '43', T30, fill=INK)
        text(draw, (962, 326), '50', T30, fill=tone)
        partial_h(680, 275, 840, tone, progress)
        partial_h(680, 495, 840, tone, progress)
    elif step == 5:
        tone = _rgba(GREEN, 255)
        outline(ub_box, tone, _rgba(GREEN_BG, int(140 * progress)), width=6)
        partial_v(935, 380, 455, tone, progress)
        text(draw, (882, 266), '19', T30, fill=INK)
        text(draw, (962, 266), '22', T30, fill=GREEN)
        text(draw, (882, 326), '43', T30, fill=INK)
        text(draw, (962, 326), '50', T30, fill=GREEN)


def _render_pe_demo_scene(completed_steps=0, active_step=None, progress=1.0):
    w, h = 1180, 800
    img = Image.new('RGBA', (w, h), _rgba(WHITE))
    draw = ImageDraw.Draw(img)
    boxes, result_box, ub_box = _draw_pe_demo_base(draw, w, h)
    for step in range(completed_steps):
        _draw_pe_demo_overlay(draw, boxes, result_box, ub_box, step, 1.0)
    if active_step is not None:
        _draw_pe_demo_overlay(draw, boxes, result_box, ub_box, active_step, progress)
    return img.convert('RGB')


def render_pe_animation_assets():
    base = _render_pe_demo_scene()
    base.save(PE_ANIM_BASE_PATH)
    progress_points = ANIM_PROGRESS_POINTS
    durations = ANIM_DURATIONS
    for idx, gif_path in enumerate(PE_ANIM_GIF_PATHS):
        frames = [_render_pe_demo_scene(completed_steps=idx, active_step=idx, progress=p) for p in progress_points]
        frames[0].save(gif_path, save_all=True, append_images=frames[1:], duration=durations, disposal=2)


def _draw_compiler_demo_base(draw, w, h):
    rounded(draw, (8, 8, w - 8, h - 8), radius=28, fill=WHITE, outline=LINE, width=2)
    boxes = {
        'spec': (80, 210, 240, 300),
        'ub_alloc': (300, 210, 480, 300),
        'sched': (540, 210, 720, 300),
        'encode': (780, 210, 960, 300),
        'host': (1020, 210, 1180, 300),
        'ub_map': (240, 410, 430, 500),
        'wait': (520, 410, 720, 500),
        'imem': (800, 410, 990, 500),
        'ready': (1010, 410, 1190, 500),
    }
    labels = {
        'spec': 'model spec',
        'ub_alloc': 'ub_allocator',
        'sched': 'scheduler',
        'encode': 'encode_instrs',
        'host': 'AXI load IMEM',
        'ub_map': 'ub_map.json',
        'wait': 'schedule + wait_after',
        'imem': 'imem.hex / UB_RD',
        'ready': 'sequencer ready',
    }
    for key, box in boxes.items():
        rounded(draw, box, radius=20, fill=WHITE, outline=LINE, width=2)
        center(draw, (box[0] + box[2]) / 2, (box[1] + box[3]) / 2, labels[key], B18, fill=INK)

    arrow(draw, [(240, 255), (300, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(480, 255), (540, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(720, 255), (780, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(960, 255), (1020, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(390, 300), (390, 410)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(630, 300), (630, 410)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(870, 300), (870, 410)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(1100, 300), (1100, 410)], fill=MUTED, width=6, arrow_size=18)

    chip(draw, (84, 90, 236, 132), 'spec -> addr', BLUE_BG, BLUE, font=B18)
    chip(draw, (268, 90, 460, 132), 'stage schedule', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (490, 90, 702, 132), 'wait_after sync', PURPLE_BG, PURPLE, font=B18)
    chip(draw, (734, 90, 930, 132), 'UB_RD encode', GREEN_BG, GREEN, font=B18)
    chip(draw, (962, 90, 1184, 132), 'IMEM write / ready', RED_BG, RED, font=B18)
    text(draw, (86, 654), '这组页讲清楚：编译器不是孤立脚本，而是把张量地址、阶段同步和 IMEM 指令一起喂给前端。', B18, fill=MUTED)
    return boxes


def _draw_compiler_demo_overlay(draw, boxes, step, progress):
    def mark(box, tone, fill, label=None):
        draw.rounded_rectangle(box, radius=20, outline=tone, fill=fill, width=5)
        if label:
            text(draw, (box[0] + 10, box[1] + 8), label, B18, fill=tone)

    def show_arrow(points, tone):
        arrow(draw, points, fill=tone, width=9, arrow_size=20)

    if step == 0:
        if progress > 0.12: mark(boxes['spec'], BLUE, BLUE_BG, 'MLP / XOR')
        if progress > 0.34: show_arrow([(240, 255), (300, 255)], BLUE)
        if progress > 0.56: mark(boxes['ub_alloc'], BLUE, BLUE_BG, 'tensor layout')
        if progress > 0.76:
            show_arrow([(390, 300), (390, 410)], BLUE)
            mark(boxes['ub_map'], BLUE, BLUE_BG, 'addr map')
    elif step == 1:
        if progress > 0.16: mark(boxes['sched'], ORANGE, ORANGE_BG, 'forward / bwd')
        if progress > 0.4: show_arrow([(630, 300), (630, 410)], ORANGE)
        if progress > 0.62: mark(boxes['wait'], ORANGE, ORANGE_BG, 'wait_after')
        if progress > 0.82: show_arrow([(480, 255), (540, 255)], ORANGE)
    elif step == 2:
        if progress > 0.14: mark(boxes['encode'], GREEN, GREEN_BG, 'UB_RD fields')
        if progress > 0.4: show_arrow([(870, 300), (870, 410)], GREEN)
        if progress > 0.64: mark(boxes['imem'], GREEN, GREEN_BG, 'opcode word')
        if progress > 0.82: show_arrow([(720, 255), (780, 255)], GREEN)
    elif step == 3:
        if progress > 0.14: mark(boxes['host'], RED, RED_BG, 'IMEM_WE / LEN')
        if progress > 0.42: show_arrow([(1100, 300), (1100, 410)], RED)
        if progress > 0.66: mark(boxes['ready'], RED, RED_BG, 'start ready')
        if progress > 0.84: show_arrow([(960, 255), (1020, 255)], RED)


def _render_compiler_demo_scene(completed_steps=0, active_step=None, progress=1.0):
    w, h = 1280, 760
    img = Image.new('RGB', (w, h), WHITE)
    draw = ImageDraw.Draw(img)
    boxes = _draw_compiler_demo_base(draw, w, h)
    for step_idx in range(completed_steps):
        _draw_compiler_demo_overlay(draw, boxes, step_idx, 1.0)
    if active_step is not None:
        _draw_compiler_demo_overlay(draw, boxes, active_step, progress)
    return img


def render_compiler_animation_assets():
    base = _render_compiler_demo_scene()
    base.save(COMPILER_ANIM_BASE_PATH)
    progress_points = ANIM_PROGRESS_POINTS
    durations = ANIM_DURATIONS
    for idx, gif_path in enumerate(COMPILER_ANIM_GIF_PATHS):
        frames = [_render_compiler_demo_scene(completed_steps=idx, active_step=idx, progress=p) for p in progress_points]
        frames[0].save(gif_path, save_all=True, append_images=frames[1:], duration=durations, disposal=2)


def _render_compiler_animation_demo(slide, spec, idx):
    step = int(spec.get('step', 0))
    step_titles = [
        'T0 model spec -> ub_map',
        'T1 scheduler 写阶段边界',
        'T2 encode 生成 UB_RD / IMEM',
        'T3 host 装载 IMEM 并进入 ready',
    ]
    step_bodies = [
        '先从 MLP 规格和张量尺寸出发，把输入、权重、bias、Y、H 这些张量排到统一地址空间。',
        'scheduler 不只是排 forward / backward / update，还负责插入 wait_after 这类系统同步边界。',
        'encode_instrs 把 schedule 变成最终指令字，硬件真正读到的是 UB_RD 字段和 IMEM word。',
        '最后由 host 通过 AXI-Lite 把 IMEM_ADDR / IMEM_W0 / IMEM_WE / IMEM_LEN 写进去，frontend 才能 start。',
    ]
    step_tones = [BLUE, ORANGE, GREEN, RED]
    step_fills = [BLUE_BG, ORANGE_BG, GREEN_BG, RED_BG]

    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_panel(slide, 0.72, 1.88, 7.18, 4.86, WHITE)
    _add_panel(slide, 8.16, 1.88, 4.4, 4.86, WHITE)
    _add_chip_ppt(slide, 0.98, 1.62, 1.02, 'GIF 单步', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 2.16, 1.62, 1.28, 'Compiler', step_fills[step], step_tones[step])
    _add_chip_ppt(slide, 3.62, 1.62, 1.0, '4 steps', WHITE, MUTED)
    _add_chip_ppt(slide, 4.78, 1.62, 1.14, '点按翻页', HILITE, INK)
    _add_textbox(slide, 0.98, 2.08, 3.24, 0.22, '编译器到 IMEM 单步演示', 19, INK, True)
    _add_textbox(slide, 0.98, 2.42, 6.24, 0.2, '点一下翻到下一页，下一步 GIF 会自动播放。', 10, MUTED)
    _add_picture_named(slide, COMPILER_ANIM_BASE_PATH, 0.94, 2.74, 6.68, 3.96, 'COMPILER_ANIM_BASE')
    _add_picture_named(slide, COMPILER_ANIM_GIF_PATHS[step], 0.94, 2.74, 6.68, 3.96, f'COMPILER_ANIM_STEP_{step}_GIF')

    _add_chip_ppt(slide, 0.98, 6.18, 1.28, f'当前: T{step}', step_fills[step], step_tones[step])
    x = 2.56
    for idx_step in range(4):
        fill = step_fills[idx_step] if idx_step == step else WHITE
        tone = step_tones[idx_step] if idx_step == step else MUTED
        _add_chip_ppt(slide, x, 6.18, 0.62, f'T{idx_step}', fill, tone)
        x += 0.78

    _add_textbox(slide, 8.42, 2.08, 3.1, 0.22, '当前这一步怎么讲', 17, INK, True)
    _add_textbox(slide, 8.42, 2.46, 3.72, 0.18, step_titles[step], 15, step_tones[step], True)
    _add_anim_note_box(slide, 8.42, 2.84, 3.86, 0.92, step_titles[step], step_bodies[step], step_tones[step], step_fills[step])
    _add_textbox(slide, 8.42, 4.04, 3.72, 0.62, '''关键口径：
1. 编译器输出的是地址、阶段和同步边界。
2. wait_after 是系统语义，不是 UI 注释。
3. 前端真正执行的是 encode 后的 IMEM 指令。''', 10, INK)
    _add_textbox(slide, 8.42, 4.88, 3.72, 0.46, '这组页最适合把“为什么编译器值得单独讲”讲清楚。', 10, MUTED)
    _add_chip_ppt(slide, 8.42, 5.86, 1.3, 'ub_map', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 9.92, 5.86, 1.46, 'wait_after', PURPLE_BG, PURPLE)
    _add_chip_ppt(slide, 11.58, 5.86, 0.72, 'UB_RD', GREEN_BG, GREEN)

    _add_static_footer(slide, spec['footer'])


def _draw_frontend_demo_base(draw, w, h):
    rounded(draw, (8, 8, w - 8, h - 8), radius=28, fill=WHITE, outline=LINE, width=2)
    boxes = {
        'host': (90, 210, 240, 300),
        'imem': (300, 210, 450, 300),
        'start': (510, 210, 640, 300),
        'dispatch': (700, 210, 850, 300),
        'wait': (910, 210, 1040, 300),
        'advance': (1040, 210, 1160, 300),
        'decode': (660, 390, 860, 475),
        'pathway': (900, 390, 1110, 475),
        'drain': (900, 520, 1110, 605),
        'mux': (210, 520, 430, 605),
    }
    labels = {
        'host': 'AXI host load',
        'imem': 'IMEM ready',
        'start': 'start_pulse',
        'dispatch': 'SEQ_DISPATCH',
        'wait': 'SEQ_WAIT',
        'advance': 'SEQ_ADVANCE',
        'decode': 'control_unit decode',
        'pathway': 'vpu_pathway_reg',
        'drain': 'vpu_drain',
        'mux': 'host / CU mux',
    }
    for key, box in boxes.items():
        rounded(draw, box, radius=20, fill=WHITE, outline=LINE, width=2)
        center(draw, (box[0] + box[2]) / 2, (box[1] + box[3]) / 2, labels[key], B18, fill=INK)

    arrow(draw, [(240, 255), (300, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(450, 255), (510, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(640, 255), (700, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(850, 255), (910, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(1040, 255), (1040, 255)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(775, 300), (775, 390)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(990, 300), (1005, 390)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(1005, 475), (1005, 520)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(320, 520), (320, 300)], fill=MUTED, width=6, arrow_size=18)

    chip(draw, (92, 90, 242, 132), 'config / load', BLUE_BG, BLUE, font=B18)
    chip(draw, (270, 90, 430, 132), 'dispatch / wait', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (458, 90, 648, 132), 'decode / pathway hold', PURPLE_BG, PURPLE, font=B18)
    chip(draw, (676, 90, 856, 132), 'drain / advance', GREEN_BG, GREEN, font=B18)
    text(draw, (92, 656), 'Frontend 关键不是“状态机有几态”，而是 start / dispatch / wait / drain 这些边界怎么把多拍执行收住。', B18, fill=MUTED)
    return boxes


def _draw_frontend_demo_overlay(draw, boxes, step, progress):
    def mark(box, tone, fill, label=None):
        draw.rounded_rectangle(box, radius=20, outline=tone, fill=fill, width=5)
        if label:
            text(draw, (box[0] + 10, box[1] + 8), label, B18, fill=tone)

    def show_arrow(points, tone):
        arrow(draw, points, fill=tone, width=9, arrow_size=20)

    if step == 0:
        if progress > 0.15: mark(boxes['host'], BLUE, BLUE_BG, 'UB_DATA / PUSH')
        if progress > 0.4: show_arrow([(240, 255), (300, 255)], BLUE)
        if progress > 0.55: mark(boxes['imem'], BLUE, BLUE_BG, 'IMEM_W0 / WE')
        if progress > 0.75: mark(boxes['mux'], BLUE, BLUE_BG, 'write mux')
    elif step == 1:
        if progress > 0.15: mark(boxes['start'], ORANGE, ORANGE_BG, 'busy <- 1')
        if progress > 0.38: show_arrow([(640, 255), (700, 255)], ORANGE)
        if progress > 0.52: mark(boxes['dispatch'], ORANGE, ORANGE_BG, 'seq_instr_pulse')
        if progress > 0.72: show_arrow([(775, 300), (775, 390)], ORANGE)
        if progress > 0.82: mark(boxes['decode'], ORANGE, ORANGE_BG, 'ub_rd_*')
    elif step == 2:
        if progress > 0.12: mark(boxes['wait'], PURPLE, PURPLE_BG, 'wait_after')
        if progress > 0.38: show_arrow([(990, 300), (1005, 390)], PURPLE)
        if progress > 0.58: mark(boxes['pathway'], PURPLE, PURPLE_BG, 'hold route')
        if progress > 0.8: mark(boxes['drain'], PURPLE, PURPLE_BG, 'waiting')
    elif step == 3:
        if progress > 0.15: mark(boxes['drain'], GREEN, GREEN_BG, 'drain = 1')
        if progress > 0.38: show_arrow([(1005, 475), (1005, 520)], GREEN)
        if progress > 0.58: mark(boxes['advance'], GREEN, GREEN_BG, 'pc + 1')
        if progress > 0.8: mark(boxes['dispatch'], GREEN, GREEN_BG, 'next instr')


def _render_frontend_demo_scene(completed_steps=0, active_step=None, progress=1.0):
    w, h = 1180, 760
    img = Image.new('RGB', (w, h), WHITE)
    draw = ImageDraw.Draw(img)
    boxes = _draw_frontend_demo_base(draw, w, h)
    for step_idx in range(completed_steps):
        _draw_frontend_demo_overlay(draw, boxes, step_idx, 1.0)
    if active_step is not None:
        _draw_frontend_demo_overlay(draw, boxes, active_step, progress)
    return img


def render_frontend_animation_assets():
    base = _render_frontend_demo_scene()
    base.save(FRONTEND_ANIM_BASE_PATH)
    progress_points = ANIM_PROGRESS_POINTS
    durations = ANIM_DURATIONS
    for idx, gif_path in enumerate(FRONTEND_ANIM_GIF_PATHS):
        frames = [_render_frontend_demo_scene(completed_steps=idx, active_step=idx, progress=p) for p in progress_points]
        frames[0].save(gif_path, save_all=True, append_images=frames[1:], duration=durations, disposal=2)


def _render_frontend_animation_demo(slide, spec, idx):
    step = int(spec.get('step', 0))
    step_titles = [
        'T0 host 装载与 IMEM 就绪',
        'T1 start / dispatch / decode',
        'T2 WAIT 与 pathway 保持',
        'T3 drain 后 ADVANCE',
    ]
    step_bodies = [
        '先通过 AXI-Lite 把 UB 和 IMEM 装好，host / CU mux 也在这一步讲清楚。',
        'start_pulse 拉起 busy，dispatch 只打一拍，decode 产生 ub_rd_* / switch / pathway。',
        '如果 wait_after 置位，状态机会卡在 WAIT；多拍路径靠 vpu_pathway_reg 保持不漂移。',
        '等到 vpu_drain 后，状态机才能 ADVANCE，再进入下一条指令。',
    ]
    step_tones = [BLUE, ORANGE, PURPLE, GREEN]
    step_fills = [BLUE_BG, ORANGE_BG, PURPLE_BG, GREEN_BG]

    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_panel(slide, 0.72, 1.88, 7.18, 4.86, WHITE)
    _add_panel(slide, 8.16, 1.88, 4.4, 4.86, WHITE)
    _add_chip_ppt(slide, 0.98, 1.62, 1.02, 'GIF 单步', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 2.16, 1.62, 1.32, 'Frontend', step_fills[step], step_tones[step])
    _add_chip_ppt(slide, 3.66, 1.62, 1.0, '4 steps', WHITE, MUTED)
    _add_chip_ppt(slide, 4.82, 1.62, 1.14, '点按翻页', HILITE, INK)
    _add_textbox(slide, 0.98, 2.08, 3.3, 0.22, 'Frontend 控制单步演示', 19, INK, True)
    _add_textbox(slide, 0.98, 2.42, 6.24, 0.2, '点一下翻到下一页，下一步 GIF 会自动播放。', 10, MUTED)
    _add_picture_named(slide, FRONTEND_ANIM_BASE_PATH, 0.94, 2.74, 6.68, 3.96, 'FRONTEND_ANIM_BASE')
    _add_picture_named(slide, FRONTEND_ANIM_GIF_PATHS[step], 0.94, 2.74, 6.68, 3.96, f'FRONTEND_ANIM_STEP_{step}_GIF')

    _add_chip_ppt(slide, 0.98, 6.18, 1.28, f'当前: T{step}', step_fills[step], step_tones[step])
    x = 2.56
    for idx_step in range(4):
        fill = step_fills[idx_step] if idx_step == step else WHITE
        tone = step_tones[idx_step] if idx_step == step else MUTED
        _add_chip_ppt(slide, x, 6.18, 0.62, f'T{idx_step}', fill, tone)
        x += 0.78

    _add_textbox(slide, 8.42, 2.08, 2.9, 0.22, '当前这一步怎么讲', 17, INK, True)
    _add_textbox(slide, 8.42, 2.46, 3.72, 0.18, step_titles[step], 15, step_tones[step], True)
    _add_anim_note_box(slide, 8.42, 2.84, 3.86, 0.82, step_titles[step], step_bodies[step], step_tones[step], step_fills[step])
    _add_textbox(slide, 8.42, 3.96, 3.72, 0.62, '''关键口径：
1. dispatch 只打一拍。
2. 多拍路径靠 pathway_reg 保持。
3. 真正推进边界在 drain 之后。''', 10, INK)
    _add_textbox(slide, 8.42, 4.82, 3.72, 0.52, '这组页最适合把 sequencer 为什么不是简单 pc++ 讲透。', 10, MUTED)
    _add_chip_ppt(slide, 8.42, 5.86, 1.54, 'dispatch pulse', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 10.16, 5.86, 1.46, 'wait / drain', GREEN_BG, GREEN)

    _add_static_footer(slide, spec['footer'])


def _draw_ub_timing_demo_base(draw, w, h):
    rounded(draw, (8, 8, w - 8, h - 8), radius=28, fill=WHITE, outline=LINE, width=2)
    xs = [250, 430, 610, 790, 970]
    lane_specs = [
        ('dispatch / seq pulse', 170, BLUE, BLUE_BG),
        ('ub_rd_input_valid', 255, ORANGE, ORANGE_BG),
        ('pe sample', 340, GREEN, GREEN_BG),
        ('array wavefront', 425, PURPLE, PURPLE_BG),
        ('hold / drain', 510, RED, RED_BG),
    ]
    cells = {}
    for lane_idx, (label, y, tone, fill) in enumerate(lane_specs):
        text(draw, (30, y + 12), label, B18, fill=tone)
        draw.line((200, y + 25, 1048, y + 25), fill=LINE, width=2)
        for step_idx, cx in enumerate(xs):
            box = (cx - 58, y, cx + 58, y + 50)
            rounded(draw, box, radius=18, fill=WHITE, outline=LINE, width=2)
            cells[(lane_idx, step_idx)] = box
    time_tones = [(BLUE_BG, BLUE), (ORANGE_BG, ORANGE), (GREEN_BG, GREEN), (PURPLE_BG, PURPLE), (RED_BG, RED)]
    for step_idx, cx in enumerate(xs):
        fill, tone = time_tones[step_idx]
        chip(draw, (cx - 34, 92, cx + 34, 128), f'T{step_idx}', fill, tone, font=B18)
        draw.line((cx, 136, cx, 562), fill=LINE, width=2)

    pe_boxes = {
        'PE11': (1040, 176, 1140, 260),
        'PE12': (1150, 176, 1250, 260),
        'PE21': (1040, 272, 1140, 356),
        'PE22': (1150, 272, 1250, 356),
    }
    for label, box in pe_boxes.items():
        rounded(draw, box, radius=18, fill=WHITE, outline=LINE, width=2)
        center(draw, (box[0] + box[2]) / 2, box[1] + 22, label, B18, fill=INK)
    result_box = (1042, 404, 1248, 486)
    drain_box = (1042, 512, 1248, 594)
    rounded(draw, result_box, radius=18, fill=WHITE, outline=LINE, width=2)
    rounded(draw, drain_box, radius=18, fill=WHITE, outline=LINE, width=2)
    center(draw, 1145, 446, 'result settle', B20, fill=MUTED)
    center(draw, 1145, 554, 'drain / advance', B20, fill=MUTED)
    chip(draw, (1042, 90, 1232, 130), '2x2 PE timing view', BLUE_BG, BLUE, font=B18)
    text(draw, (224, 602), '关键尾拍语义', B20, fill=INK)
    chip(draw, (224, 636, 462, 676), 'input: last beat + hold', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (486, 636, 716, 676), 'weight: no extra valid', PURPLE_BG, PURPLE, font=B18)
    text(draw, (224, 698), 'dispatch 只负责发起读流；真正的完成边界由 drain 和 wait_after 收口。', B18, fill=MUTED)
    return cells, pe_boxes, result_box, drain_box


def _draw_ub_timing_demo_overlay(draw, cells, pe_boxes, result_box, drain_box, step, progress):
    def pulse(box, tone, fill, frac, label=None):
        l, t, r, b = box
        rounded(draw, box, radius=18, fill=WHITE, outline=tone, width=4)
        inner_r = int(l + 12 + (r - l - 24) * max(0.12, frac))
        draw.rounded_rectangle((l + 10, t + 10, inner_r, b - 10), radius=12, fill=fill)
        if label:
            text(draw, (l + 16, t + 15), label, B18, fill=tone)

    def outline(box, tone, fill=None, width=5):
        draw.rounded_rectangle(box, radius=20, outline=tone, fill=fill, width=width)

    def center_box(box):
        l, t, r, b = box
        return ((l + r) // 2, (t + b) // 2)

    def partial_link(src_box, dst_box, tone, frac):
        x1, y1 = center_box(src_box)
        x2, y2 = center_box(dst_box)
        xt = int(x1 + (x2 - x1) * max(0.08, frac))
        yt = int(y1 + (y2 - y1) * max(0.08, frac))
        arrow(draw, [(x1, y1), (xt, yt)], fill=tone, width=8, arrow_size=18)

    if step == 0:
        pulse(cells[(0, 0)], BLUE, BLUE_BG, progress, 'dispatch')
        pulse(cells[(1, 0)], ORANGE, ORANGE_BG, progress * 0.9, 'ub_rd_start')
        partial_link(cells[(0, 0)], cells[(1, 0)], BLUE, progress)
    elif step == 1:
        pulse(cells[(1, 1)], ORANGE, ORANGE_BG, progress, 'beat0')
        pulse(cells[(2, 1)], GREEN, GREEN_BG, progress * 0.9, 'sample')
        outline(pe_boxes['PE11'], BLUE, BLUE_BG, 5)
        outline(pe_boxes['PE21'], BLUE, BLUE_BG, 5)
        partial_link(cells[(1, 1)], pe_boxes['PE11'], ORANGE, progress)
        partial_link(cells[(1, 1)], pe_boxes['PE21'], ORANGE, progress)
    elif step == 2:
        pulse(cells[(2, 2)], GREEN, GREEN_BG, progress, 'valid')
        pulse(cells[(3, 2)], PURPLE, PURPLE_BG, progress * 0.9, 'wave')
        outline(pe_boxes['PE12'], GREEN, GREEN_BG, 5)
        outline(pe_boxes['PE22'], GREEN, GREEN_BG, 5)
        partial_link(pe_boxes['PE11'], pe_boxes['PE12'], GREEN, progress)
        partial_link(pe_boxes['PE21'], pe_boxes['PE22'], GREEN, progress)
        partial_link(pe_boxes['PE11'], pe_boxes['PE21'], ORANGE, progress)
        partial_link(pe_boxes['PE12'], pe_boxes['PE22'], ORANGE, progress)
    elif step == 3:
        pulse(cells[(3, 3)], PURPLE, PURPLE_BG, progress, 'last')
        outline(result_box, PURPLE, PURPLE_BG, 5)
        partial_link(pe_boxes['PE12'], result_box, PURPLE, progress)
        partial_link(pe_boxes['PE22'], result_box, PURPLE, progress)
        text(draw, (1080, 434), 'tail settle', B18, fill=PURPLE)
    elif step == 4:
        pulse(cells[(1, 4)], ORANGE, ORANGE_BG, progress * 0.8, 'hold')
        pulse(cells[(4, 4)], RED, RED_BG, progress, 'drain')
        outline(drain_box, RED, RED_BG, 5)
        partial_link(result_box, drain_box, RED, progress)
        text(draw, (1076, 542), 'wait_after clear', B18, fill=RED)


def _render_ub_timing_demo_scene(completed_steps=0, active_step=None, progress=1.0):
    w, h = 1280, 760
    img = Image.new('RGB', (w, h), WHITE)
    draw = ImageDraw.Draw(img)
    cells, pe_boxes, result_box, drain_box = _draw_ub_timing_demo_base(draw, w, h)
    for step_idx in range(completed_steps):
        _draw_ub_timing_demo_overlay(draw, cells, pe_boxes, result_box, drain_box, step_idx, 1.0)
    if active_step is not None:
        _draw_ub_timing_demo_overlay(draw, cells, pe_boxes, result_box, drain_box, active_step, progress)
    return img


def render_ub_timing_animation_assets():
    base = _render_ub_timing_demo_scene()
    base.save(UB_TIMING_ANIM_BASE_PATH)
    progress_points = ANIM_PROGRESS_POINTS
    durations = ANIM_DURATIONS
    for idx, gif_path in enumerate(UB_TIMING_ANIM_GIF_PATHS):
        frames = [_render_ub_timing_demo_scene(completed_steps=idx, active_step=idx, progress=p) for p in progress_points]
        frames[0].save(gif_path, save_all=True, append_images=frames[1:], duration=durations, disposal=2)


def _render_ub_timing_animation_demo(slide, spec, idx):
    step = int(spec.get('step', 0))
    step_titles = [
        'T0 dispatch 发起读流',
        'T1 first beat 进入阵列边界',
        'T2 wavefront 在阵列内部推进',
        'T3 last beat 进入尾拍收束',
        'T4 hold / drain / advance',
    ]
    step_bodies = [
        'sequencer 打出 seq_instr_pulse，并拉起 ub_rd_start。',
        'UB 开始送出第一拍 input / weight，左列 PE 先采样。',
        'valid 向右推进，psum 向下沉，阵列中间形成波前。',
        '最后一个 active beat 已进阵列，但结果还需要尾拍才能收齐。',
        'input 侧保留一个 hold 周期，最终由 drain 和 wait_after 收口，再允许 ADVANCE。',
    ]
    step_tones = [BLUE, ORANGE, GREEN, PURPLE, RED]
    step_fills = [BLUE_BG, ORANGE_BG, GREEN_BG, PURPLE_BG, RED_BG]

    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_panel(slide, 0.72, 1.88, 7.18, 4.86, WHITE)
    _add_panel(slide, 8.16, 1.88, 4.4, 4.86, WHITE)
    _add_chip_ppt(slide, 0.98, 1.62, 1.02, 'GIF 单步', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 2.16, 1.62, 1.34, 'UB Timing', step_fills[step], step_tones[step])
    _add_chip_ppt(slide, 3.68, 1.62, 1.0, '5 steps', WHITE, MUTED)
    _add_chip_ppt(slide, 4.84, 1.62, 1.14, '点按翻页', HILITE, INK)
    _add_textbox(slide, 0.98, 2.08, 3.26, 0.22, 'UB / PE 时序单步演示', 19, INK, True)
    _add_textbox(slide, 0.98, 2.42, 6.24, 0.2, '点一下翻到下一页，下一步 GIF 会自动播放。', 10, MUTED)
    _add_picture_named(slide, UB_TIMING_ANIM_BASE_PATH, 0.94, 2.74, 6.68, 3.96, 'UB_TIMING_ANIM_BASE')
    _add_picture_named(slide, UB_TIMING_ANIM_GIF_PATHS[step], 0.94, 2.74, 6.68, 3.96, f'UB_TIMING_ANIM_STEP_{step}_GIF')

    _add_chip_ppt(slide, 0.98, 6.18, 1.28, f'当前: T{step}', step_fills[step], step_tones[step])
    x = 2.56
    for idx_step in range(5):
        fill = step_fills[idx_step] if idx_step == step else WHITE
        tone = step_tones[idx_step] if idx_step == step else MUTED
        _add_chip_ppt(slide, x, 6.18, 0.62, f'T{idx_step}', fill, tone)
        x += 0.78

    _add_textbox(slide, 8.42, 2.08, 2.9, 0.22, '当前这一步怎么讲', 17, INK, True)
    _add_textbox(slide, 8.42, 2.46, 3.72, 0.18, step_titles[step], 15, step_tones[step], True)
    _add_anim_note_box(slide, 8.42, 2.84, 3.86, 0.82, step_titles[step], step_bodies[step], step_tones[step], step_fills[step])
    _add_textbox(slide, 8.42, 3.96, 3.72, 0.58, '''关键口径：
1. dispatch 不等于结束。
2. input 和 weight 的尾拍策略不同。
3. 最终同步点在 drain / wait_after。''', 10, INK)
    _add_textbox(slide, 8.42, 4.76, 3.72, 0.56, '这组页最适合把“为什么需要 hold / drain / wait_after”一口气讲透。', 10, MUTED)
    _add_chip_ppt(slide, 8.42, 5.86, 1.56, 'input hold', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 10.16, 5.86, 1.68, 'drain 收口', RED_BG, RED)

    _add_static_footer(slide, spec['footer'])


def _draw_vpu_demo_base(draw, w, h):
    rounded(draw, (8, 8, w - 8, h - 8), radius=28, fill=WHITE, outline=LINE, width=2)
    modules = {
        'sys': (120, 250, 250, 340),
        'bias': (300, 250, 430, 340),
        'lrelu': (480, 250, 620, 340),
        'loss': (670, 250, 830, 340),
        'dact': (890, 250, 1040, 340),
    }
    labels = {'sys': 'sys in', 'bias': 'Bias', 'lrelu': 'LReLU', 'loss': 'Loss grad', 'dact': 'dAct'}
    for key, box in modules.items():
        rounded(draw, box, radius=20, fill=WHITE, outline=LINE, width=2)
        center(draw, (box[0] + box[2]) / 2, (box[1] + box[3]) / 2, labels[key], B20, fill=INK)

    aux = {
        'last_h': (370, 465, 500, 540),
        'y': (540, 465, 660, 540),
        'leak': (700, 465, 860, 540),
        'wb': (940, 465, 1110, 540),
        'pathway': (150, 465, 310, 540),
    }
    aux_labels = {'last_h': 'last_H', 'y': 'Y from UB', 'leak': 'leak factor', 'wb': 'writeback -> UB', 'pathway': 'pathway reg'}
    for key, box in aux.items():
        rounded(draw, box, radius=18, fill=WHITE, outline=LINE, width=2)
        center(draw, (box[0] + box[2]) / 2, (box[1] + box[3]) / 2, aux_labels[key], B18, fill=MUTED if key != 'wb' else GREEN)

    arrow(draw, [(250, 295), (300, 295)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(430, 295), (480, 295)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(620, 295), (670, 295)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(830, 295), (890, 295)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(965, 340), (965, 465)], fill=MUTED, width=6, arrow_size=18)
    arrow(draw, [(185, 340), (185, 505), (890, 505), (890, 340)], fill=MUTED, width=5, arrow_size=18)
    arrow(draw, [(435, 465), (435, 340)], fill=MUTED, width=5, arrow_size=18)
    arrow(draw, [(600, 465), (600, 340)], fill=MUTED, width=5, arrow_size=18)
    arrow(draw, [(780, 465), (780, 340)], fill=MUTED, width=5, arrow_size=18)

    chip(draw, (120, 92, 260, 132), '1100 forward', BLUE_BG, BLUE, font=B18)
    chip(draw, (290, 92, 470, 132), '1111 transition', ORANGE_BG, ORANGE, font=B18)
    chip(draw, (500, 92, 660, 132), '0001 backward', PURPLE_BG, PURPLE, font=B18)
    chip(draw, (690, 92, 900, 132), 'writeback / route hold', GREEN_BG, GREEN, font=B18)
    text(draw, (118, 618), '同一套子模块通过 pathway bit 复用三种训练阶段；不是三条独立硬连线。', B18, fill=MUTED)
    return modules, aux


def _draw_vpu_demo_overlay(draw, modules, aux, step, progress):
    def outline(box, tone, fill=None, width=5):
        draw.rounded_rectangle(box, radius=20, outline=tone, fill=fill, width=width)

    def mark(box, tone, fill, label=None):
        outline(box, tone, fill, 5)
        if label:
            text(draw, (box[0] + 10, box[1] + 8), label, B18, fill=tone)

    def show_arrow(points, tone):
        arrow(draw, points, fill=tone, width=9, arrow_size=20)

    if step == 0:
        if progress > 0.15: mark(modules['sys'], BLUE, BLUE_BG)
        if progress > 0.35: show_arrow([(250, 295), (300, 295)], BLUE)
        if progress > 0.5: mark(modules['bias'], BLUE, BLUE_BG)
        if progress > 0.65: show_arrow([(430, 295), (480, 295)], BLUE)
        if progress > 0.8: mark(modules['lrelu'], BLUE, BLUE_BG)
        if progress > 0.82: mark(aux['last_h'], BLUE, BLUE_BG)
    elif step == 1:
        if progress > 0.1: mark(modules['sys'], ORANGE, ORANGE_BG)
        if progress > 0.25: mark(modules['bias'], ORANGE, ORANGE_BG)
        if progress > 0.4: mark(modules['lrelu'], ORANGE, ORANGE_BG)
        if progress > 0.55: mark(modules['loss'], ORANGE, ORANGE_BG)
        if progress > 0.7: mark(modules['dact'], ORANGE, ORANGE_BG)
        if progress > 0.2: show_arrow([(250, 295), (300, 295)], ORANGE)
        if progress > 0.35: show_arrow([(430, 295), (480, 295)], ORANGE)
        if progress > 0.5: show_arrow([(620, 295), (670, 295)], ORANGE)
        if progress > 0.65: show_arrow([(830, 295), (890, 295)], ORANGE)
        if progress > 0.78: mark(aux['y'], ORANGE, ORANGE_BG)
        if progress > 0.84: mark(aux['leak'], ORANGE, ORANGE_BG)
    elif step == 2:
        if progress > 0.18: mark(modules['sys'], PURPLE, PURPLE_BG)
        if progress > 0.4: show_arrow([(185, 340), (185, 505), (890, 505), (890, 340)], PURPLE)
        if progress > 0.7: mark(modules['dact'], PURPLE, PURPLE_BG)
    elif step == 3:
        if progress > 0.12: mark(aux['pathway'], GREEN, GREEN_BG, 'hold')
        if progress > 0.35: mark(modules['dact'], GREEN, GREEN_BG)
        if progress > 0.55: show_arrow([(965, 340), (965, 465)], GREEN)
        if progress > 0.75: mark(aux['wb'], GREEN, GREEN_BG, 'wb')


def _render_vpu_demo_scene(completed_steps=0, active_step=None, progress=1.0):
    w, h = 1180, 760
    img = Image.new('RGB', (w, h), WHITE)
    draw = ImageDraw.Draw(img)
    modules, aux = _draw_vpu_demo_base(draw, w, h)
    for step_idx in range(completed_steps):
        _draw_vpu_demo_overlay(draw, modules, aux, step_idx, 1.0)
    if active_step is not None:
        _draw_vpu_demo_overlay(draw, modules, aux, active_step, progress)
    return img


def render_vpu_animation_assets():
    base = _render_vpu_demo_scene()
    base.save(VPU_ANIM_BASE_PATH)
    progress_points = ANIM_PROGRESS_POINTS
    durations = ANIM_DURATIONS
    for idx, gif_path in enumerate(VPU_ANIM_GIF_PATHS):
        frames = [_render_vpu_demo_scene(completed_steps=idx, active_step=idx, progress=p) for p in progress_points]
        frames[0].save(gif_path, save_all=True, append_images=frames[1:], duration=durations, disposal=2)


def _render_vpu_animation_demo(slide, spec, idx):
    step = int(spec.get('step', 0))
    step_titles = [
        'T0 forward 1100',
        'T1 transition 1111',
        'T2 backward 0001',
        'T3 writeback 与 pathway 保持',
    ]
    step_bodies = [
        'sys 经 bias 和 lrelu 形成 forward 路径，同时把 H 留作后续 derivative。',
        'transition 会继续走到 loss 和 dAct，是训练时最完整的一条路径。',
        'backward 不再走 bias / lrelu / loss，而是直接把 sys 送到 dAct。',
        '前端保持住 pathway_reg，VPU 输出再写回 UB，形成系统级闭环。',
    ]
    step_tones = [BLUE, ORANGE, PURPLE, GREEN]
    step_fills = [BLUE_BG, ORANGE_BG, PURPLE_BG, GREEN_BG]

    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_panel(slide, 0.72, 1.88, 7.18, 4.86, WHITE)
    _add_panel(slide, 8.16, 1.88, 4.4, 4.86, WHITE)
    _add_chip_ppt(slide, 0.98, 1.62, 1.02, 'GIF 单步', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 2.16, 1.62, 0.98, 'VPU', step_fills[step], step_tones[step])
    _add_chip_ppt(slide, 3.32, 1.62, 1.0, '4 steps', WHITE, MUTED)
    _add_chip_ppt(slide, 4.48, 1.62, 1.14, '点按翻页', HILITE, INK)
    _add_textbox(slide, 0.98, 2.08, 3.26, 0.22, 'VPU pathway 单步演示', 19, INK, True)
    _add_textbox(slide, 0.98, 2.42, 6.24, 0.2, '点一下翻到下一页，下一步 GIF 会自动播放。', 10, MUTED)
    _add_picture_named(slide, VPU_ANIM_BASE_PATH, 0.94, 2.74, 6.68, 3.96, 'VPU_ANIM_BASE')
    _add_picture_named(slide, VPU_ANIM_GIF_PATHS[step], 0.94, 2.74, 6.68, 3.96, f'VPU_ANIM_STEP_{step}_GIF')

    _add_chip_ppt(slide, 0.98, 6.18, 1.28, f'当前: T{step}', step_fills[step], step_tones[step])
    x = 2.56
    for idx_step in range(4):
        fill = step_fills[idx_step] if idx_step == step else WHITE
        tone = step_tones[idx_step] if idx_step == step else MUTED
        _add_chip_ppt(slide, x, 6.18, 0.62, f'T{idx_step}', fill, tone)
        x += 0.78

    _add_textbox(slide, 8.42, 2.08, 2.9, 0.22, '当前这一步怎么讲', 17, INK, True)
    _add_textbox(slide, 8.42, 2.46, 3.72, 0.18, step_titles[step], 15, step_tones[step], True)
    _add_anim_note_box(slide, 8.42, 2.84, 3.86, 0.82, step_titles[step], step_bodies[step], step_tones[step], step_fills[step])
    _add_textbox(slide, 8.42, 3.96, 3.72, 0.62, '''关键口径：
1. VPU 不是黑盒后处理。
2. 同一套子模块复用多种训练路径。
3. 真正稳定路由靠 pathway_reg 保持。''', 10, INK)
    _add_textbox(slide, 8.42, 4.82, 3.72, 0.52, '这组页最适合解释 VPU 为什么值得单独讲，以及它和 frontend 控制是怎么耦合的。', 10, MUTED)
    _add_chip_ppt(slide, 8.42, 5.86, 1.36, 'pathway bit', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 9.96, 5.86, 1.74, 'route hold', GREEN_BG, GREEN)

    _add_static_footer(slide, spec['footer'])


def _decorate_system_slide(slide):
    _add_arrow_ppt(slide, 3.52, 3.5, 3.72, 3.5, BLUE)
    _add_arrow_ppt(slide, 6.65, 3.5, 6.84, 3.5, ORANGE)
    _add_arrow_ppt(slide, 9.77, 3.5, 9.96, 3.5, GREEN)
    for x, label, fill, tone in [
        (0.78, 'SW', BLUE_BG, BLUE),
        (3.9, 'CTRL', ORANGE_BG, ORANGE),
        (7.02, 'CORE', GREEN_BG, GREEN),
        (10.12, 'VERIFY', PURPLE_BG, PURPLE),
    ]:
        _add_mini_box(slide, x, 1.08, 0.62, 0.28, label, fill, tone, 8)
    _add_chip_ppt(slide, 1.15, 5.95, 1.1, 'json / hex', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 4.25, 5.95, 1.35, 'control pulse', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 7.55, 5.95, 1.15, 'dataflow', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 10.62, 5.95, 1.1, 'evidence', PURPLE_BG, PURPLE)
    backbone = slide.shapes.add_connector(1, Inches(0.95), Inches(5.72), Inches(11.9), Inches(5.72))
    backbone.line.color.rgb = _rgb(LINE)
    backbone.line.width = Pt(1.6)
    _decorate_system_arch_native(slide)
    _add_panel(slide, 8.56, 2.78, 3.78, 1.58, WHITE)
    _add_textbox(slide, 8.76, 2.94, 3.2, 0.18, '顶层集成过程', 11, BLUE, True)
    pipe = slide.shapes.add_connector(1, Inches(8.95), Inches(3.48), Inches(12.0), Inches(3.48))
    pipe.line.color.rgb = _rgb(ORANGE)
    pipe.line.width = Pt(2.2)
    steps = [
        ('1', 'compiler out', 8.78, 3.22),
        ('2', 'AXI load\nUB / IMEM', 9.48, 3.56),
        ('3', 'start +\ndispatch', 10.18, 3.22),
        ('4', 'UB / SA /\nVPU run', 10.88, 3.56),
        ('5', 'waveform +\nscoreboard', 11.58, 3.22),
    ]
    for num, label, x, y in steps:
        _add_mini_box(slide, x, y, 0.54, 0.32, num, BLUE_BG, BLUE, 9)
        _add_textbox(slide, x - 0.02, y + 0.34, 0.92, 0.34, label, 8, MUTED)
    for x in [9.32, 10.02, 10.72, 11.42]:
        _add_arrow_ppt(slide, x, 3.48, x + 0.12, 3.48, ORANGE, 1.8)
    _add_textbox(slide, 8.78, 4.08, 3.08, 0.2, '从产物装载到证据闭环，讲清楚“系统怎么跑起来”。', 8, MUTED)

def _decorate_rtl_slide(slide):
    _add_panel(slide, 0.76, 4.86, 8.66, 0.78, WHITE)
    _add_textbox(slide, 0.94, 4.96, 2.3, 0.16, '顶层数据与控制主链', 11, BLUE, True)
    _add_mini_box(slide, 0.98, 5.2, 0.95, 0.3, 'AXI-Lite', BLUE_BG, BLUE)
    _add_mini_box(slide, 2.02, 5.2, 0.8, 0.3, 'IMEM', ORANGE_BG, ORANGE)
    _add_mini_box(slide, 2.9, 5.2, 1.1, 0.3, 'Sequencer', GREEN_BG, GREEN)
    _add_mini_box(slide, 4.88, 5.2, 1.0, 0.3, 'UB', BLUE_BG, BLUE)
    _add_mini_box(slide, 5.98, 5.2, 1.0, 0.3, 'SA 2x2', ORANGE_BG, ORANGE)
    _add_mini_box(slide, 7.08, 5.2, 0.95, 0.3, 'VPU', PURPLE_BG, PURPLE)
    _add_mini_box(slide, 8.13, 5.2, 1.0, 0.3, 'WB', GREEN_BG, GREEN)
    _add_arrow_ppt(slide, 4.08, 5.35, 4.8, 5.35, BLUE)
    _add_arrow_ppt(slide, 5.9, 5.35, 5.96, 5.35, ORANGE)
    _add_arrow_ppt(slide, 7.0, 5.35, 7.06, 5.35, PURPLE)
    _add_arrow_ppt(slide, 8.05, 5.35, 8.1, 5.35, GREEN)
    _add_textbox(slide, 4.16, 5.56, 3.5, 0.16, 'control bus -> data path -> compute -> writeback', 8, MUTED)
    _add_panel(slide, 8.78, 3.62, 3.75, 0.9, WHITE)
    _add_textbox(slide, 8.96, 3.76, 3.2, 0.16, 'RTL 顶层 wrapper 链路', 11, BLUE, True)
    _add_mini_box(slide, 8.98, 4.0, 0.9, 0.3, 'tpu_soc', BLUE_BG, BLUE, 9)
    _add_mini_box(slide, 9.98, 4.0, 1.05, 0.3, 'frontend_axil', ORANGE_BG, ORANGE, 9)
    _add_mini_box(slide, 11.13, 4.0, 0.85, 0.3, 'tpu.sv', GREEN_BG, GREEN, 9)
    _add_arrow_ppt(slide, 9.88, 4.15, 9.98, 4.15, BLUE, 1.8)
    _add_arrow_ppt(slide, 11.03, 4.15, 11.13, 4.15, ORANGE, 1.8)
    _add_mini_box(slide, 8.96, 4.92, 0.9, 0.28, 'host cfg', WHITE, BLUE, 8)
    _add_mini_box(slide, 10.02, 4.92, 0.95, 0.28, 'ctrl pulse', WHITE, ORANGE, 8)
    _add_mini_box(slide, 11.15, 4.92, 0.98, 0.28, 'data out', WHITE, GREEN, 8)
    _add_arrow_ppt(slide, 9.42, 4.3, 9.42, 4.9, BLUE, 1.4)
    _add_arrow_ppt(slide, 10.5, 4.3, 10.5, 4.9, ORANGE, 1.4)
    _add_arrow_ppt(slide, 11.57, 4.3, 11.57, 4.9, GREEN, 1.4)
    _decorate_rtl_core_native(slide)
    _decorate_rtl_logic_native(slide)

def _decorate_ub_slide(slide):
    _add_mini_box(slide, 0.92, 5.0, 1.2, 0.42, 'host load', BLUE_BG, BLUE)
    _add_mini_box(slide, 2.28, 5.0, 1.2, 0.42, 'core read', GREEN_BG, GREEN)
    _add_mini_box(slide, 3.64, 5.0, 1.35, 0.42, 'VPU writeback', ORANGE_BG, ORANGE)
    _add_mini_box(slide, 5.25, 5.0, 1.45, 0.42, 'grad update', PURPLE_BG, PURPLE)
    _add_arrow_ppt(slide, 2.12, 5.21, 2.28, 5.21, BLUE)
    _add_arrow_ppt(slide, 3.48, 5.21, 3.64, 5.21, GREEN)
    _add_arrow_ppt(slide, 4.99, 5.21, 5.25, 5.21, ORANGE)
    _add_mini_box(slide, 8.85, 5.0, 0.7, 0.42, 'ptr0', BLUE_BG, BLUE)
    _add_mini_box(slide, 9.63, 5.0, 0.7, 0.42, 'ptr1', ORANGE_BG, ORANGE)
    _add_mini_box(slide, 10.41, 5.0, 0.7, 0.42, 'ptr2', GREEN_BG, GREEN)
    _add_mini_box(slide, 11.19, 5.0, 0.7, 0.42, 'ptr3', PURPLE_BG, PURPLE)
    _add_mini_box(slide, 11.97, 5.0, 0.7, 0.42, 'ptr5/6', RED_BG, RED)
    _decorate_ub_native_detail(slide)




def _decorate_compiler_slide(slide):
    xs = [8.95, 9.8, 10.65, 11.5]
    labels = [('spec', BLUE_BG, BLUE), ('alloc', GREEN_BG, GREEN), ('sched', ORANGE_BG, ORANGE), ('imem', PURPLE_BG, PURPLE)]
    _add_textbox(slide, 8.96, 4.3, 3.1, 0.16, '从模型规格到最终 IMEM', 11, BLUE, True)
    for (label, fill, tone), x in zip(labels, xs):
        _add_mini_box(slide, x, 4.55, 0.68, 0.42, label, fill, tone)
    _add_arrow_ppt(slide, 9.63, 4.76, 9.8, 4.76, BLUE)
    _add_arrow_ppt(slide, 10.48, 4.76, 10.65, 4.76, GREEN)
    _add_arrow_ppt(slide, 11.33, 4.76, 11.5, 4.76, ORANGE)
    _add_panel(slide, 8.94, 4.98, 3.28, 1.22, WHITE)
    _add_textbox(slide, 9.06, 5.08, 3.0, 0.18, '指令字段示意', 11, ORANGE, True)
    fields = [('opcode', 9.04, 0.72, BLUE_BG, BLUE), ('addr', 9.8, 0.72, GREEN_BG, GREEN), ('row', 10.56, 0.55, ORANGE_BG, ORANGE), ('col', 11.15, 0.45, PURPLE_BG, PURPLE), ('ptr/path', 11.64, 0.68, RED_BG, RED)]
    for label, x, w, fill, tone in fields:
        _add_mini_box(slide, x, 5.36, w, 0.3, label, fill, tone, 8)
    _add_textbox(slide, 9.06, 5.72, 3.0, 0.36, 'UB_RD example:\n[2:0] opcode | [8:3] addr | [12:9] row | [14:13] col | [18:16] ptr_sel | [22:19] pathway', 8, MUTED)
    _add_panel(slide, 8.94, 6.02, 3.28, 0.48, WHITE)
    _add_textbox(slide, 9.06, 6.12, 3.0, 0.22, '样例语义: UB_RD(ptr=1, row=0, col=1) -> 从 UB 取权重上装阵列', 8, INK)

def _decorate_frontend_slide(slide):
    _add_mini_box(slide, 9.55, 4.45, 0.95, 0.4, 'AXI', BLUE_BG, BLUE)
    _add_mini_box(slide, 10.7, 4.45, 0.95, 0.4, 'REG', ORANGE_BG, ORANGE)
    _add_mini_box(slide, 9.55, 5.05, 0.95, 0.4, 'IMEM', GREEN_BG, GREEN)
    _add_mini_box(slide, 10.7, 5.05, 0.95, 0.4, 'SEQ', PURPLE_BG, PURPLE)
    _add_arrow_ppt(slide, 10.5, 4.65, 10.7, 4.65, BLUE)
    _add_arrow_ppt(slide, 10.02, 4.85, 10.02, 5.05, ORANGE)
    _add_arrow_ppt(slide, 11.17, 4.85, 11.17, 5.05, GREEN)
    _add_textbox(slide, 9.45, 5.62, 2.6, 0.5, '配置面：AXI/寄存器\n运行面：IMEM/Sequencer', 10, MUTED)


def _decorate_system_arch_native(slide):
    _add_mini_box(slide, 8.72, 4.35, 1.0, 0.44, '软件栈', BLUE_BG, BLUE)
    _add_mini_box(slide, 9.9, 4.35, 1.0, 0.44, 'Frontend', ORANGE_BG, ORANGE)
    _add_mini_box(slide, 11.08, 4.35, 1.0, 0.44, 'TPU Core', GREEN_BG, GREEN)
    _add_mini_box(slide, 9.3, 5.08, 1.0, 0.44, '验证', PURPLE_BG, PURPLE)
    _add_mini_box(slide, 10.48, 5.08, 1.2, 0.44, '波形/结果', RED_BG, RED)
    _add_arrow_ppt(slide, 9.72, 4.57, 9.9, 4.57, BLUE)
    _add_arrow_ppt(slide, 10.9, 4.57, 11.08, 4.57, ORANGE)
    _add_arrow_ppt(slide, 10.2, 4.79, 9.8, 5.08, GREEN)
    _add_arrow_ppt(slide, 11.33, 4.79, 11.0, 5.08, PURPLE)




def _decorate_rtl_core_native(slide):
    _add_panel(slide, 9.1, 4.1, 3.35, 1.12, WHITE)
    _add_textbox(slide, 9.25, 4.18, 2.9, 0.16, 'RTL Core Arch', 11, BLUE, True)
    _add_mini_box(slide, 9.25, 4.5, 0.9, 0.34, 'frontend', BLUE_BG, BLUE, 9)
    _add_mini_box(slide, 10.25, 4.5, 0.8, 0.34, 'decode', ORANGE_BG, ORANGE, 9)
    _add_mini_box(slide, 11.15, 4.5, 0.7, 0.34, 'UB', GREEN_BG, GREEN, 9)
    _add_mini_box(slide, 11.95, 4.5, 0.38, 0.34, 'SA', PURPLE_BG, PURPLE, 9)
    _add_arrow_ppt(slide, 10.15, 4.67, 10.25, 4.67, BLUE)
    _add_arrow_ppt(slide, 11.05, 4.67, 11.15, 4.67, ORANGE)
    _add_arrow_ppt(slide, 11.85, 4.67, 11.95, 4.67, GREEN)
    _add_textbox(slide, 9.25, 4.92, 2.95, 0.18, 'control -> data path -> compute path', 9, MUTED)


def _decorate_rtl_logic_native(slide):
    _add_panel(slide, 9.1, 5.3, 3.35, 1.08, WHITE)
    _add_textbox(slide, 9.25, 5.38, 2.9, 0.16, 'Logic Closeup', 11, ORANGE, True)
    _add_mini_box(slide, 9.25, 5.68, 0.9, 0.3, 'wait_after', ORANGE_BG, ORANGE, 9)
    _add_mini_box(slide, 10.25, 5.68, 0.92, 0.3, 'pathway', PURPLE_BG, PURPLE, 9)
    _add_mini_box(slide, 11.28, 5.68, 0.95, 0.3, 'host mux', GREEN_BG, GREEN, 9)
    _add_arrow_ppt(slide, 10.15, 5.83, 10.25, 5.83, ORANGE)
    _add_arrow_ppt(slide, 11.17, 5.83, 11.28, 5.83, PURPLE)
    _add_textbox(slide, 9.25, 6.02, 2.95, 0.18, '收口点：wait / 路径保持 / 写口仲裁', 9, MUTED)

def _decorate_ub_native_detail(slide):
    for idx, label in enumerate(['input', 'weight', 'bias', 'Y', 'H', 'grad']):
        tone_fill = [BLUE_BG, ORANGE_BG, GREEN_BG, PURPLE_BG, RED_BG, HILITE][idx]
        tone = [BLUE, ORANGE, GREEN, PURPLE, RED, INK][idx]
        _add_mini_box(slide, 8.75 + 0.6 * (idx % 3), 4.45 + 0.62 * (idx // 3), 0.52, 0.36, label, tone_fill, tone, 9)
    _add_textbox(slide, 8.72, 5.7, 3.1, 0.35, '统一地址空间里的张量分布示意', 10, MUTED)


def _decorate_wave_native(slide):
    tones = [BLUE, ORANGE, GREEN, PURPLE]
    labels = ['dispatch', 'ub read', 'sys out', 'drain']
    ys = [4.45, 4.82, 5.19, 5.56]
    for tone, label, y in zip(tones, labels, ys):
        line = slide.shapes.add_connector(1, Inches(8.65), Inches(y), Inches(12.25), Inches(y))
        line.line.color.rgb = _rgb(tone)
        line.line.width = Pt(2.3)
        _add_textbox(slide, 8.45, y - 0.08, 0.8, 0.16, label, 9, tone, True)
    for x in [9.2, 10.0, 10.8, 11.6]:
        tick = slide.shapes.add_connector(1, Inches(x), Inches(4.28), Inches(x), Inches(5.72))
        tick.line.color.rgb = _rgb(LINE)
        tick.line.width = Pt(1)

def _decorate_pe_slide(slide):
    boxes = [(0.95, 4.75, 'PE11'), (2.2, 4.75, 'PE12'), (0.95, 5.55, 'PE21'), (2.2, 5.55, 'PE22')]
    for x, y, label in boxes:
        _add_mini_box(slide, x, y, 1.0, 0.5, label, WHITE, INK, 10)
    _add_arrow_ppt(slide, 0.72, 5.0, 0.95, 5.0, BLUE)
    _add_arrow_ppt(slide, 1.95, 5.0, 2.2, 5.0, GREEN)
    _add_arrow_ppt(slide, 1.45, 4.52, 1.45, 4.74, ORANGE)
    _add_arrow_ppt(slide, 1.45, 5.25, 1.45, 5.54, ORANGE)
    _add_arrow_ppt(slide, 3.2, 5.0, 3.55, 5.0, GREEN)
    _add_chip_ppt(slide, 0.9, 6.15, 1.4, 'input ->', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 2.5, 6.15, 1.55, 'psum down', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 4.25, 6.15, 1.9, 'shadow -> active', PURPLE_BG, PURPLE)


PPT_NATIVE_SLIDES = [
    {
        'layout': 'cover',
        'title': 'TinyTPU AXI-Lite SoC',
        'subtitle': '主讲版 | 8 页终版',
        'summary': '主讲版进一步压到 8 页，只保留系统、RTL、编译、核心数据路径、计算证据和结果收束；Frontend、控制修复和动态演示全部放进附录。',
        'metrics': [
            ('主讲页数', '8 pages', BLUE),
            ('主线范围', 'system -> compile -> core', ORANGE),
            ('结果证据', '41 / 41 PASS', GREEN),
            ('训练证据', '12 epoch 收敛', PURPLE),
            ('追问入口', 'appendix / detail pages', BLUE),
        ],
        'chips': [
            ('系统总览', BLUE_BG, BLUE),
            ('RTL / Compiler', ORANGE_BG, ORANGE),
            ('UB / PE', GREEN_BG, GREEN),
            ('Wave / Result', PURPLE_BG, PURPLE),
            ('Appendix / GIF', RED_BG, RED),
        ],
        'roadmap_title': '主讲顺序',
        "roadmap_body": "1. 系统总览与 RTL\n2. 编译器、UB、PE\n3. 波形证据与结果\n4. 追问时切 appendix",
        'cover_footer': '主讲版收成 8 页；Frontend、控制修复、集成/贡献和动态 GIF 统一留在附录备用。',
    },
    {
        'title': '系统总览',
        'subtitle': '把软件、控制前端、计算核心和验证闭环放回同一页。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 2.95, 'h': 4.25, 'title': '软件与编译', 'body': 'model spec\nub_allocator\nscheduler\nencode_instrs\n\n输出：ub_map / schedule / imem.hex', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 3.72, 'y': 1.35, 'w': 2.95, 'h': 4.25, 'title': 'SoC 控制与执行', 'body': 'AXI-Lite reg map\nIMEM + sequencer\ncontrol_unit decode\nhost / CU 写口仲裁\n\n输出：ub_rd_* / switch / pathway', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 6.84, 'y': 1.35, 'w': 2.95, 'h': 4.25, 'title': '计算核心', 'body': 'Unified Buffer\nSystolic Array 2x2\nVPU\nwriteback -> UB\n\n形成前向、反向、更新闭环', 'tone': GREEN, 'fill': GREEN_BG},
            {'x': 9.96, 'y': 1.35, 'w': 2.75, 'h': 4.25, 'title': '验证与结果', 'body': 'module cocotb\nAXI e2e\ntrain convergence\nVCD + scoreboard\n\n41 / 41 PASS + loss 收敛', 'tone': PURPLE, 'fill': PURPLE_BG},
        ],
        'chips': [('json / hex', BLUE_BG, BLUE), ('control pulse', ORANGE_BG, ORANGE), ('dataflow', GREEN_BG, GREEN), ('evidence', PURPLE_BG, PURPLE)],
        'footer': '建议先讲全景，后面再拆前端、UB、PE、VPU。',
        'variant': 'process',
    },
    {
        'title': '项目级 RTL 结构',
        'subtitle': '把 tpu_soc 顶层拆成控制域、执行域、可观测接口三层。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 3.6, 'h': 4.55, 'title': 'Frontend 控制域', 'body': 'tpu_frontend_axil.sv\ncontrol_unit.sv\n\nAXI-Lite IF\nCTRL / STATUS\nIMEM\n4-state sequencer\nHost / CU write mux\ndecode: ub_rd_* / switch / pathway / lr', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 4.42, 'y': 1.35, 'w': 4.1, 'h': 4.55, 'title': 'TPU Core 执行域', 'body': 'tpu.sv 统一组织\nUnified Buffer\nSystolic Array 2 x 2\nVPU\nwriteback -> UB\n\n关键路径：UB -> SA -> VPU -> UB', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 8.74, 'y': 1.35, 'w': 3.97, 'h': 4.55, 'title': '可观测接口', 'body': 'host 可读输出 / STATUS\nvpu_data_out\nsys_data_out\nvalid / busy\npc / state\n\n波形里可以直接抓到关键脉冲。', 'tone': GREEN, 'fill': GREEN_BG},
        ],
        'chips': [('ctrl: ub_rd_* / switch / pathway', ORANGE_BG, ORANGE), ('data: UB -> SA -> VPU', BLUE_BG, BLUE), ('writeback: VPU -> UB', GREEN_BG, GREEN)],
        'footer': '这页回答“模块到底怎么连”，后续再按子系统拆讲。',
        'variant': 'arch',
    },
    {
        'title': '编译器与指令组织',
        'subtitle': '解释软件侧如何把 MLP 规格降成 UB 映射、schedule 和 IMEM。',
        'metrics': [
            {'x':0.6,'y':1.35,'w':2.0,'h':1.0,'title':'命令数','value':'66 commands','tone':BLUE},
            {'x':2.8,'y':1.35,'w':2.0,'h':1.0,'title':'UB 占用','value':'44 / 128 words','tone':GREEN},
            {'x':5.0,'y':1.35,'w':2.3,'h':1.0,'title':'目标约束','value':'2x2 / 2-lane / Q8.8','tone':ORANGE},
        ],
        'cards': [
            {'x': 0.6, 'y': 2.6, 'w': 6.1, 'h': 3.1, 'title': '编译链路', 'body': '01 model spec -> 02 ub_allocator -> 03 scheduler -> 04 encode_instrs -> 05 AXI write IMEM\n\n为什么是阶段级 schedule：当前目标先打通 IMEM、sequencer 和控制链路，而不是生成 cycle-accurate 程序。', 'tone': BLUE, 'fill': WHITE},
            {'x': 6.95, 'y': 2.6, 'w': 2.95, 'h': 3.1, 'title': '真实输出文件', 'body': 'ub_map.json\nschedule.json\nimem.hex\n\n这些文件分别对应地址分配、阶段命令和最终指令编码。', 'tone': GREEN, 'fill': GREEN_BG},
            {'x': 10.15, 'y': 2.6, 'w': 2.56, 'h': 3.1, 'title': 'schedule 片段', 'body': 'forward / transition / backward / update\nwait_after 标出跨阶段同步点\n\n重点不是指令多复杂，而是它已经能稳定驱动控制执行闭环。', 'tone': ORANGE, 'fill': ORANGE_BG},
        ],
        'footer': '这页的口径是“编译器怎样喂给硬件”，不是“编译算法多先进”。',
        'variant': 'detail',
    },
    {
        'title': 'Compiler GIF 演示 T0',
        'subtitle': '第 1 / 4 步：model spec 到 ub_map。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'compiler_animation_demo',
        'step': 0,
    },
    {
        'title': 'Compiler GIF 演示 T1',
        'subtitle': '第 2 / 4 步：scheduler 写阶段边界和 wait_after。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'compiler_animation_demo',
        'step': 1,
    },
    {
        'title': 'Compiler GIF 演示 T2',
        'subtitle': '第 3 / 4 步：encode 生成 UB_RD 和 IMEM word。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'compiler_animation_demo',
        'step': 2,
    },
    {
        'title': 'Compiler GIF 演示 T3',
        'subtitle': '第 4 / 4 步：host 装载 IMEM 并进入 ready。点击切回后续内容页。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'compiler_animation_demo',
        'step': 3,
    },
    {
        'title': 'Frontend 总览',
        'subtitle': '把寄存器映射、IMEM、sequencer、decode 和 host write 路径拆开。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 3.0, 'h': 4.7, 'title': '寄存器映射', 'body': '0x00 CTRL\n0x04 STATUS\n0x10 INSTR_W0\n0x20/24 UB_DATA/PUSH\n0x30/34/40/44 IMEM_*\n0x50/54/58 LEAK / INV / LR\n\nhost 先把参数和程序装进去。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 3.82, 'y': 1.35, 'w': 2.95, 'h': 4.7, 'title': 'IMEM 与指令流', 'body': 'imem[0:63]\nimem_len_reg\n\n指令格式：NOP / SWITCH / UB_RD / UB_WR_HOST\nUB_RD 字段包含 addr / row / col / transpose / ptr_sel / pathway。', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 6.99, 'y': 1.35, 'w': 2.75, 'h': 4.7, 'title': '4-state sequencer', 'body': 'IDLE\nDISPATCH\nWAIT\nADVANCE\n\nwait_after 触发 WAIT；只有 vpu_drain 到来后才能继续推进。', 'tone': GREEN, 'fill': GREEN_BG},
            {'x': 9.96, 'y': 1.35, 'w': 2.75, 'h': 4.7, 'title': 'decode 与写口仲裁', 'body': 'control_unit 输出：ub_rd_* / sys_switch / vpu_pathway / host_wr\n\nhost 装载和 CU 写口共享同一个 UB 入口，所以需要明确仲裁。', 'tone': PURPLE, 'fill': PURPLE_BG},
        ],
        'footer': '这一页是前端总览；下一页再把“配置面”和“运行面”分开讲。',
    },
    {
        'title': 'Frontend 再拆一页',
        'subtitle': '把配置面和运行面拆开，避免一页塞太满。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 5.9, 'h': 4.7, 'title': '配置面：AXI-Lite + 程序装载', 'body': 'host 视角主要做两件事：\n1. 通过 UB_DATA / UB_PUSH 把参数写进 UB\n2. 通过 IMEM_ADDR / IMEM_W0 / IMEM_WE / IMEM_LEN 写程序\n\n运行参数 leak / inv_batch / learning_rate 也都在这条路径上完成配置。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 6.74, 'y': 1.35, 'w': 5.97, 'h': 4.7, 'title': '运行面：sequencer + decode + 持续路径', 'body': 'start_pulse 拉起 busy 并装载 pc=0\nSEQ_DISPATCH 只打一拍 seq_instr_pulse\ncontrol_unit 译出 ub_rd / switch / pathway\nSEQ_WAIT 负责等 vpu_drain\nSEQ_ADVANCE 再推进到下一条指令\n\nvpu_pathway_reg 持久保持，保证多拍 VPU 不会路由漂移。', 'tone': ORANGE, 'fill': ORANGE_BG},
        ],
        'code': {'x': 0.9, 'y': 5.25, 'w': 11.9, 'h': 1.2, 'title': '关键 RTL', 'lines': ['seq_instr_pulse <= 1\'b1 only in dispatch', 'vpu_pathway_reg latches on every UB_RD', 'vpu_drain = vpu_valid_prev && !tpu_vpu_valid_in;']},
        'footer': '拆开后就变成两条叙事线：写配置，跑程序。',
    },
    {
        'title': 'Frontend GIF 演示 T0',
        'subtitle': '第 1 / 4 步：host 装载与 IMEM 就绪。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'frontend_animation_demo',
        'step': 0,
    },
    {
        'title': 'Frontend GIF 演示 T1',
        'subtitle': '第 2 / 4 步：start / dispatch / decode。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'frontend_animation_demo',
        'step': 1,
    },
    {
        'title': 'Frontend GIF 演示 T2',
        'subtitle': '第 3 / 4 步：WAIT 与 pathway 保持。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'frontend_animation_demo',
        'step': 2,
    },
    {
        'title': 'Frontend GIF 演示 T3',
        'subtitle': '第 4 / 4 步：drain 后 ADVANCE 与下一条指令。点击切回后续内容页。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'frontend_animation_demo',
        'step': 3,
    },
    {
        'title': 'Unified Buffer 设计',
        'subtitle': 'UB 不只是内存块，而是 host load、核心读流、VPU 写回和梯度更新的汇合点。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 3.6, 'h': 4.6, 'title': 'UB 地址布局', 'body': '当前分配占用 44 / 128 words\n存放输入、权重、bias、Y、H 等张量\n\n它不是“几块分离 SRAM”，而是一个统一地址空间，便于编译器和控制逻辑共同驱动。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 4.45, 'y': 1.35, 'w': 3.75, 'h': 4.6, 'title': '两类写入口', 'body': 'host write: ub_wr_host_*\nVPU writeback: ub_wr_*\nwr_ptr_base / wr_ptr_restore\n\nhost 先顺序装载；start 时恢复 wr_ptr，后续写回不会踩到参数区。', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 8.45, 'y': 1.35, 'w': 4.26, 'h': 4.6, 'title': '读指针语义', 'body': 'ptr_sel=0 input / activation\nptr_sel=1 weight top-load\nptr_sel=2 bias stream\nptr_sel=3 Y label\nptr_sel=4 H activation\nptr_sel=5/6 in-UB grad update\n\n同一块存储被多条读流复用。', 'tone': GREEN, 'fill': GREEN_BG},
        ],
        'footer': '这页最适合回答“为什么叫 Unified Buffer”的问题。',
        'variant': 'highlight',
    },
    {
        'title': 'wr_ptr / base / restore',
        'subtitle': '把 host 装载阶段和训练阶段分开，解释为什么写回不会覆盖参数区。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 3.75, 'h': 4.7, 'title': '核心状态', 'body': 'wr_ptr：统一写指针\nwr_ptr_next：按 lane valid 预计算下一地址\nwr_ptr_base：记录 host 装载完成后的参数区尾部\nub_wr_ptr_restore_in：start 时恢复 wr_ptr = wr_ptr_base', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 4.58, 'y': 1.35, 'w': 4.0, 'h': 4.7, 'title': '阶段视角', 'body': '阶段 1 host load：UB_DATA / UB_PUSH 连续装参数，wr_ptr 持续前进\n\n阶段 2 start restore：start_pulse 把 wr_ptr 拉回 wr_ptr_base\n\n阶段 3 runtime writeback：VPU 从 base 之后继续写，不踩前面参数', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 8.82, 'y': 1.35, 'w': 3.89, 'h': 4.7, 'title': '为什么单独讲', 'body': '它解释了 Unified Buffer 为什么能既做参数装载，又做训练写回。\n\n这个设计没有物理分 bank，而是用“先 capture base，再 restore”的时序方式把静态参数区和动态写回区隔开。', 'tone': GREEN, 'fill': GREEN_BG},
        ],
        'code': {'x': 0.9, 'y': 5.25, 'w': 11.8, 'h': 1.2, 'title': '关键 RTL', 'lines': ['if (ub_wr_host_valid_in[0] || ub_wr_host_valid_in[1]) wr_ptr_base <= wr_ptr_next;', 'if (ub_wr_ptr_restore_in) wr_ptr <= wr_ptr_base;']},
        'footer': '这页直接回答“为什么写回不会覆盖 W/B 参数”。',
    },
    {
        'title': 'UB 读流与 PE 时序对齐',
        'subtitle': '把 UB valid、PE 采样、hold 周期、sequencer wait 放到一条时间线上。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 2.35, 'h': 1.3, 'title': 'T0 dispatch', 'body': 'ub_rd_start / seq_instr_pulse\n发起一次新的读流。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 3.1, 'y': 1.35, 'w': 2.35, 'h': 1.3, 'title': 'T1 first beat', 'body': 'input / weight 第一拍进入阵列边界。', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 5.6, 'y': 1.35, 'w': 2.35, 'h': 1.3, 'title': 'T2 wavefront', 'body': '有效数据在 PE 阵列中形成波前传播。', 'tone': GREEN, 'fill': GREEN_BG},
            {'x': 8.1, 'y': 1.35, 'w': 2.35, 'h': 1.3, 'title': 'T3 last beat', 'body': '最后一个 active beat 进入尾拍阶段。', 'tone': PURPLE, 'fill': PURPLE_BG},
            {'x': 10.6, 'y': 1.35, 'w': 2.1, 'h': 1.3, 'title': 'T4 hold / drain', 'body': '输入 hold；控制等待 drain 收口。', 'tone': RED, 'fill': RED_BG},
            {'x': 0.6, 'y': 3.0, 'w': 5.95, 'h': 3.0, 'title': 'UB 侧关键语义', 'body': 'input 流在最后一个 active beat 后还有 1 个 hold 周期，保证阵列最后一列能完成传播。\n\nweight 流故意不保留最后一个 valid，防止 loader 多打一拍时覆盖 active weight。\n\n所以同样是 UB 读，input 和 weight 的尾拍策略并不一样。', 'tone': BLUE, 'fill': WHITE},
            {'x': 6.8, 'y': 3.0, 'w': 5.9, 'h': 3.0, 'title': '和 sequencer 的关系', 'body': 'dispatch 只负责发起读流，不代表计算已经结束。\n\n真正的完成边界在 VPU drain，sequencer 只有等 drain 才能 ADVANCE。\n\nwait_after 的意义，就是用系统级同步点把 UB / PE / VPU 的多拍行为收口。', 'tone': ORANGE, 'fill': WHITE},
        ],
        'footer': '这页把数据流和控制流第一次放在同一张时间图上。',
    },
    {
        'title': 'UB 时序 GIF 演示 T0',
        'subtitle': '第 1 / 5 步：dispatch 发起。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'ub_timing_animation_demo',
        'step': 0,
    },
    {
        'title': 'UB 时序 GIF 演示 T1',
        'subtitle': '第 2 / 5 步：first beat 进入。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'ub_timing_animation_demo',
        'step': 1,
    },
    {
        'title': 'UB 时序 GIF 演示 T2',
        'subtitle': '第 3 / 5 步：wavefront 推进。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'ub_timing_animation_demo',
        'step': 2,
    },
    {
        'title': 'UB 时序 GIF 演示 T3',
        'subtitle': '第 4 / 5 步：last beat 与尾拍收束。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'ub_timing_animation_demo',
        'step': 3,
    },
    {
        'title': 'UB 时序 GIF 演示 T4',
        'subtitle': '第 5 / 5 步：hold / drain / advance。点击切回后续内容页。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'ub_timing_animation_demo',
        'step': 4,
    },
    {
        'title': 'UB 内梯度下降更新',
        'subtitle': '重点讲清楚 bias update 和 weight update 不是同一种语义。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 5.95, 'h': 4.8, 'title': 'Bias Update 路径', 'body': 'ptr_sel=5，grad_bias_or_weight=0\n\ndZ / dZ1 stream + old bias + lr -> GD\n\nbias 模式下，如果上一拍 done 过，就把 sub_in_a 切到 value_updated_out。\n\n这意味着同一个 bias 地址可以随着连续样本梯度到来而逐拍累加。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 6.76, 'y': 1.35, 'w': 5.95, 'h': 4.8, 'title': 'Weight Update 路径', 'body': 'ptr_sel=6，grad_bias_or_weight=1\n\nouter-product grad + old weight + lr -> GD\n\nweight 模式下，sub_in_a 固定取 value_old_in，不会链接上次 value_updated_out。\n\n每次 weight 更新都基于当前 tile 读出的旧权重独立完成，更适合 tile 级 outer-product。', 'tone': ORANGE, 'fill': ORANGE_BG},
        ],
        'chips': [('bias: sample-by-sample accumulate', BLUE_BG, BLUE), ('weight: tile-based update', ORANGE_BG, ORANGE)],
        'footer': '同一个 GD 模块，驱动语义并不一样。',
    },
    {
        'title': 'VPU 单独展开',
        'subtitle': 'VPU 不是黑盒后处理，而是按 pathway bit 组织的可重组训练链路。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 5.95, 'h': 4.7, 'title': 'VPU 子模块与主数据链', 'body': 'sys in -> Bias -> Leaky ReLU -> Loss grad -> dAct\n\npathway 编码：\n1100: forward = sys -> bias -> lrelu\n1111: transition = sys -> bias -> lrelu -> loss -> dAct\n0001: backward = sys -> dAct\n\nlast_H / Y / leak factor 都参与其中。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 6.76, 'y': 1.35, 'w': 5.95, 'h': 4.7, 'title': '为什么值得单独讲', 'body': '同一套子模块复用三种训练阶段，不是写死单一路径。\n\n1111 路径里既要保存 H，又要继续送入 loss / derivative。\n\npathway bit 和 frontend 的保持寄存器直接耦合，是系统级控制点。', 'tone': PURPLE, 'fill': PURPLE_BG},
        ],
        'footer': '这页把 VPU 提升成“可重组训练路径单元”。',
    },
    {
        'title': 'VPU GIF 演示 T0',
        'subtitle': '第 1 / 4 步：forward 1100。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'vpu_animation_demo',
        'step': 0,
    },
    {
        'title': 'VPU GIF 演示 T1',
        'subtitle': '第 2 / 4 步：transition 1111。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'vpu_animation_demo',
        'step': 1,
    },
    {
        'title': 'VPU GIF 演示 T2',
        'subtitle': '第 3 / 4 步：backward 0001。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'vpu_animation_demo',
        'step': 2,
    },
    {
        'title': 'VPU GIF 演示 T3',
        'subtitle': '第 4 / 4 步：writeback 与 pathway 保持。点击切回后续内容页。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'vpu_animation_demo',
        'step': 3,
    },
    {
        'title': 'PE 与计算阵列',
        'subtitle': 'PE 是最小计算单元；2x2 systolic 把 valid、input、weight、psum 做成波前传播。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 4.0, 'h': 4.8, 'title': '2 x 2 Systolic 组织', 'body': 'input 从左进入\nweight 从上装载\nvalid / input 向右传播\npsum 向下传播\n\nPE11 / PE12 / PE21 / PE22 组成最小阵列。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 4.84, 'y': 1.35, 'w': 4.2, 'h': 4.8, 'title': 'PE 级 RTL 重点', 'body': 'fxp_mul(ina=pe_input_in, inb=weight_reg_active)\nfxp_add(ina=mult_out, inb=pe_psum_in)\npe_switch_in 把 inactive weight 切到 active\npe_accept_w_in 写 inactive weight\n!pe_enabled 只清输出，不清内部状态', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 9.28, 'y': 1.35, 'w': 3.43, 'h': 4.8, 'title': '为什么这个设计值得讲', 'body': '双 weight reg 支持 load / switch 两阶段，不需要计算中断。\n\npe_enabled 只清输出不清状态，便于列裁剪。\n\nvalid 与 switch 都做波前传播，阵列行为能在波形里直接观察。', 'tone': GREEN, 'fill': GREEN_BG},
        ],
        'footer': '推荐讲法：先拓扑，再讲 active / inactive weight。',
    },
    {
        'title': 'PE GIF 演示 T0',
        'subtitle': '第 1 / 6 步：权重就绪。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'pe_animation_demo',
        'step': 0,
    },
    {
        'title': 'PE GIF 演示 T1',
        'subtitle': '第 2 / 6 步：首拍输入进入。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'pe_animation_demo',
        'step': 1,
    },
    {
        'title': 'PE GIF 演示 T2',
        'subtitle': '第 3 / 6 步：波前推进。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'pe_animation_demo',
        'step': 2,
    },
    {
        'title': 'PE GIF 演示 T3',
        'subtitle': '第 4 / 6 步：第一列先到底。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'pe_animation_demo',
        'step': 3,
    },
    {
        'title': 'PE GIF 演示 T4',
        'subtitle': '第 5 / 6 步：结果收束。点击切下一页进入下一步。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'pe_animation_demo',
        'step': 4,
    },
    {
        'title': 'PE GIF 演示 T5',
        'subtitle': '第 6 / 6 步：写回 / 送下游。点击切回后续内容页。',
        'footer': '放映方案：每页一个 GIF，点击翻页即推进一帧。',
        'renderer': 'pe_animation_demo',
        'step': 5,
    },
    {
        'title': '关键控制 RTL 特写',
        'subtitle': '把最容易出系统级 bug 的三处控制逻辑拿出来。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 3.95, 'h': 4.75, 'title': '显式等待语义', 'body': '现象：sequencer 可能提前推进\n根因：VPU 延迟不是固定周期\n修复：wait_after + vpu_drain\n价值：给 IMEM 指令显式同步点', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 4.7, 'y': 1.35, 'w': 3.95, 'h': 4.75, 'title': '路径保持', 'body': '现象：dispatch 后几拍 route 会跑偏\n根因：decode 输出只持续 1 拍\n修复：vpu_pathway_reg 持续保持\n价值：多拍 VPU 仍走同一处理链', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 8.8, 'y': 1.35, 'w': 3.91, 'h': 4.75, 'title': 'Host / CU 仲裁', 'body': '现象：host 装载与 CU 写口重叠\n根因：同一 UB 入口承载两类语义\n修复：UB_PUSH 优先的 host_write_mux\n价值：参数装载与训练写回可共存', 'tone': GREEN, 'fill': GREEN_BG},
        ],
        'footer': '当前 RTL 没有显式 ICG cell，更准确的说法是用 pe_enabled / valid / pulse 收口。',
    },
    {
        'title': '逐拍计算动态',
        'subtitle': '这页不是 RTL 波形，而是帮助面试官快速理解 2x2 systolic 波前如何形成结果。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 1.95, 'h': 2.0, 'title': 'Cycle 0', 'body': '权重已装到 active reg\n示意矩阵：A=[[1,2],[3,4]], W=[[5,6],[7,8]]', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 2.75, 'y': 1.35, 'w': 1.95, 'h': 2.0, 'title': 'Cycle 1', 'body': '第一拍输入进入左边界\n行数据从左进入，权重常驻在 PE 内', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 4.9, 'y': 1.35, 'w': 1.95, 'h': 2.0, 'title': 'Cycle 2', 'body': '波前推进\n部分和开始下沉\n不同 PE 看到不同阶段数据', 'tone': GREEN, 'fill': GREEN_BG},
            {'x': 7.05, 'y': 1.35, 'w': 1.95, 'h': 2.0, 'title': 'Cycle 3', 'body': '第一列结果到底部\nC[:,0] 开始可见', 'tone': PURPLE, 'fill': PURPLE_BG},
            {'x': 9.2, 'y': 1.35, 'w': 1.95, 'h': 2.0, 'title': 'Cycle 4', 'body': '第二列结果完成\nC = [[19,22],[43,50]]', 'tone': RED, 'fill': RED_BG},
            {'x': 11.35, 'y': 1.35, 'w': 1.35, 'h': 2.0, 'title': 'Cycle 5', 'body': '结果写回 UB\n进入后续 VPU 路径', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 0.6, 'y': 3.75, 'w': 5.8, 'h': 2.35, 'title': '讲法建议', 'body': '1. 先说 weight-stationary，权重先装后切换。\n2. 再说 input 从左向右、psum 自上而下。\n3. 最后落到为什么要靠 wait_after 和 drain 才能收口。', 'tone': BLUE, 'fill': WHITE},
            {'x': 6.65, 'y': 3.75, 'w': 6.06, 'h': 2.35, 'title': '配套文件', 'body': 'tpu_systolic_cycle_demo.gif\ntpu_systolic_cycle_strip.png\n\nGIF 可以单独展示；PPT 中保留静态可编辑文本更稳。', 'tone': GREEN, 'fill': WHITE},
        ],
        'footer': '这页帮助讲原理；真正证明正确性还是看下一页波形。',
    },
    {
        'title': '验证波形与回归覆盖',
        'subtitle': '真实时序证据：sequencer 脉冲、UB 发数、阵列输出、VPU drain 都能对上。',
        'cards': [
            {'x': 0.6, 'y': 1.35, 'w': 7.1, 'h': 4.75, 'title': '波形主线', 'body': '第一次 busy 窗口内，可以沿着 busy -> dispatch -> ub read -> sys out -> drain 这条主线讲。\n\n对应信号包括：busy_reg / seq_instr_pulse / ub_rd_start / sys_switch / ub_rd_input_valid_out_0 / sys_valid_out_21 / vpu_valid_out_1 / vpu_drain。\n\n口径更像 Verdi：先看状态机和关键脉冲，再回到 scoreboard 验证结果。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 7.95, 'y': 1.35, 'w': 2.1, 'h': 2.1, 'title': '41 / 41 覆盖项', 'body': 'H1: 8\ndZ2: 4\ndZ1: 8\nUB dZ2: 4\nUB dZ1: 8\nW1/B1/W2/B2: 9', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 10.3, 'y': 1.35, 'w': 2.41, 'h': 2.1, 'title': '为什么有说服力', 'body': '控制脉冲和输出数据在同一窗口里对齐。\n\nscoreboard 又补上 H1 / dZ2 / dZ1 / UB 更新的数值证据。', 'tone': GREEN, 'fill': GREEN_BG},
            {'x': 7.95, 'y': 3.7, 'w': 4.76, 'h': 2.4, 'title': '建议讲法', 'body': '先交代 busy 窗口，再点 dispatch、ub read、sys out、drain 的顺序，最后回到 scoreboard。\n\n这样既有时序证据，也有数值证据。', 'tone': PURPLE, 'fill': PURPLE_BG},
        ],
        'footer': '这页是“真实证据页”，不是原理解释页。',
    },
    {
        'title': '顶层集成与改动范围',
        'subtitle': '把 wrapper、控制前端、核心执行、编译和验证一起接成闭环。',
        'footer': '这页回答“我到底改了哪些东西，为什么这不是只接一个接口壳”。',
        'renderer': 'integration_scope',
    },
    {
        'title': '验证体系与证据分层',
        'subtitle': '模块级、顶层级、AXI e2e 和训练收敛放在同一条证据链里。',
        'footer': '先讲验证层次，再落到波形和收敛结果。',
        'renderer': 'verification_stack',
    },
    {
        'title': '个人贡献与关键修复前后对比',
        'subtitle': '把 bug 现象、修复点、文件落点和系统价值放在一页收束。',
        'footer': '这页最适合回答“你个人到底做了什么”。',
        'renderer': 'contribution_fixes',
    },
    {
        'title': '结果、边界与追问方向',
        'subtitle': '最后一页给出结论，同时主动交代边界。',
        'metrics': [
            {'x':0.6,'y':1.35,'w':1.9,'h':1.0,'title':'回归','value':'41 / 41 PASS','tone':BLUE},
            {'x':2.7,'y':1.35,'w':1.8,'h':1.0,'title':'训练','value':'12 epoch','tone':ORANGE},
            {'x':4.7,'y':1.35,'w':2.6,'h':1.0,'title':'loss','value':'0.2529 -> 0.1777','tone':GREEN},
            {'x':7.5,'y':1.35,'w':1.8,'h':1.0,'title':'XOR','value':'(0,1,1,0)','tone':PURPLE},
        ],
        'cards': [
            {'x': 0.6, 'y': 2.65, 'w': 5.6, 'h': 3.45, 'title': '结论口径', 'body': '这是一个“系统集成 + 控制执行 + 验证闭环”完整展示的项目，而不是只讲某个单点算子。\n\n当前最扎实的链路是 2x2 / Q8.8 / 2-layer MLP / XOR，并且已经有 41 / 41 PASS 和多 epoch loss 收敛作为支撑。', 'tone': BLUE, 'fill': BLUE_BG},
            {'x': 6.45, 'y': 2.65, 'w': 3.0, 'h': 3.45, 'title': '项目边界', 'body': '当前目标是 tiny-tpu 原型，不夸大成完整商用 NPU SoC。\n\n还没有 DMA / IRQ / ICG / 大阵列 tile 化。', 'tone': ORANGE, 'fill': ORANGE_BG},
            {'x': 9.7, 'y': 2.65, 'w': 3.01, 'h': 3.45, 'title': '适合引导追问', 'body': '如果扩到 8x8，UB 和 schedule 怎么变？\n如果做 clock-gating / power-aware，要在哪些边界插 ICG？\n如果做 DMA / 中断，frontend 和 STATUS/CTRL 如何演进？', 'tone': GREEN, 'fill': GREEN_BG},
        ],
        'footer': '主动说明边界，比把原型包装成商用品更可信。',
    },
]

PPT_TOTAL_PAGES = len(PPT_NATIVE_SLIDES)

MAIN_DECK_TITLES = {
    '系统总览',
    '项目级 RTL 结构',
    '编译器与指令组织',
    'Unified Buffer 设计',
    'PE 与计算阵列',
    '验证波形与回归覆盖',
    '结果、边界与追问方向',
}

APPENDIX_COVER_SPEC = {
    'layout': 'cover',
    'title': 'TinyTPU AXI-Lite SoC',
    'subtitle': '附录 | 动态演示与细节展开',
    'summary': '把 GIF 单步页、细节拆页和逐拍示意收进附录，主讲版更短更稳，追问时再切进来。',
    'chips': [
        ('Compiler GIF', BLUE_BG, BLUE),
        ('Frontend / UB', ORANGE_BG, ORANGE),
        ('VPU / PE 动态', GREEN_BG, GREEN),
        ('细节拆页', PURPLE_BG, PURPLE),
    ],
    'stats': [
        ('附录定位', 'Dynamic / Detail', BLUE),
        ('演示方式', 'GIF step pages', ORANGE),
        ('适用场景', '追问 / 深挖', GREEN),
        ('兼容策略', 'Windows-safe', PURPLE),
    ],
    'cover_kicker': 'TinyTPU AXI-Lite SoC',
    'cover_headline': '附录\n动态演示与细节展开',
    'roadmap_title': '附录内容',
    'roadmap_body': '1. Compiler / Frontend / UB / VPU / PE 的 GIF 单步页\n2. Frontend / UB / wr_ptr / 时序这些细节拆页\n3. 逐拍示意用于补充说明，不打断主线',
    'cover_footer': '主讲版现在是 8 页；被追问时再切附录。',
}


def _build_deck_specs(deck_kind):
    if deck_kind == 'main':
        return [spec for spec in PPT_NATIVE_SLIDES if spec.get('layout') == 'cover' or spec.get('title') in MAIN_DECK_TITLES]
    if deck_kind == 'appendix':
        appendix_specs = [APPENDIX_COVER_SPEC]
        appendix_specs.extend(
            spec for spec in PPT_NATIVE_SLIDES
            if spec.get('layout') != 'cover' and spec.get('title') not in MAIN_DECK_TITLES
        )
        return appendix_specs
    raise ValueError(f'unknown deck kind: {deck_kind}')


def _render_header_clean(slide, idx, title, subtitle):
    _add_textbox(slide, 0.62, 0.3, 0.54, 0.18, f'P{idx}', 10, BLUE, True)
    _add_textbox(slide, 0.62, 0.54, 8.0, 0.4, title, 27, INK, True)
    _add_textbox(slide, 0.62, 0.96, 9.6, 0.18, subtitle, 11, MUTED)
    _add_panel(slide, 10.02, 0.34, 2.3, 0.32, BLUE_BG, BLUE_BG)
    _add_textbox(slide, 10.24, 0.41, 1.86, 0.12, f'{PPT_TOTAL_PAGES} 页扩展版的第 {idx} 页', 9, BLUE, True, PP_ALIGN.CENTER)


def _add_clean_frame(slide):
    _add_panel(slide, STATIC_FRAME_X, STATIC_FRAME_Y, STATIC_FRAME_W, STATIC_FRAME_H, WHITE)


def _add_static_footer(slide, text_value):
    _add_textbox(slide, STATIC_FOOTER_X, STATIC_FOOTER_Y, STATIC_FOOTER_W, STATIC_FOOTER_H, text_value, 10, MUTED)


def _add_clean_card(slide, x, y, w, h, title, body, tone=BLUE, fill=WHITE, body_size=12, tag=None):
    _add_panel(slide, x, y, w, h, fill)
    accent = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(x + 0.08), Inches(y + STATIC_CARD_PAD_TOP), Inches(0.05), Inches(h - 0.32))
    accent.fill.solid()
    accent.fill.fore_color.rgb = _rgb(tone)
    accent.line.color.rgb = _rgb(tone)
    if tag:
        _add_chip_ppt(slide, x + STATIC_CARD_PAD_X, y + STATIC_CARD_PAD_TOP, 0.94, tag, BLUE_BG, tone)
        title_y = y + 0.52
    else:
        title_y = y + 0.22
    _add_textbox(slide, x + STATIC_CARD_PAD_X, title_y, w - 0.42, STATIC_CARD_TITLE_H, title, 15, INK, True)
    body_y = title_y + STATIC_CARD_BODY_GAP
    body_h = h - (body_y - y) - STATIC_CARD_BOTTOM_PAD
    _add_textbox(slide, x + STATIC_CARD_PAD_X, body_y, w - 0.42, body_h, body, body_size, INK)


def _render_cover_clean(slide, spec):
    bar = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(0.18), SLIDE_H)
    bar.fill.solid()
    bar.fill.fore_color.rgb = _rgb(BLUE)
    bar.line.color.rgb = _rgb(BLUE)
    _add_panel(slide, 9.55, 0.42, 3.2, 4.15, PANEL, PANEL)
    _add_textbox(slide, 0.66, 0.56, 3.2, 0.18, spec.get('cover_kicker', 'TinyTPU AXI-Lite SoC'), 14, BLUE)
    _add_textbox(slide, 0.66, 1.0, 5.9, 1.2, spec.get('cover_headline', '从系统总览到\nRTL 细节与训练闭环'), 28, INK, True)
    _add_textbox(slide, 0.66, 2.35, 6.0, 0.34, spec['summary'], 14, MUTED)
    x = 0.68
    for label, fill, tone in spec['chips'][:4]:
        _add_chip_ppt(slide, x, 3.12, 1.48, label, fill, tone)
        x += 1.62
    _add_textbox(slide, 0.66, 4.0, 2.0, 0.2, spec.get('roadmap_title', '面试主线'), 14, BLUE, True)
    _add_textbox(slide, 0.66, 4.42, 6.4, 1.35, spec.get('roadmap_body', 'Why 这个项目值得做\nWhat 系统整体怎么搭起来\nHow 关键控制逻辑怎么收敛\nEvidence 我如何证明它真的跑通'), 20, INK, True)
    _add_textbox(slide, 0.66, 6.9, 7.0, 0.18, spec.get('cover_footer', '答辩口径：范围克制，但系统完整；重点不在“做了个接口壳”，而在“把控制链路和闭环验证补完整”。'), 10, MUTED)
    stats = spec.get('stats', [('项目定位','SoC 原型',BLUE),('执行组织','IMEM + Sequencer',ORANGE),('RTL 重点','wait / hold / mux',PURPLE),('结果证据','41 / 41 PASS',GREEN)])
    y = 0.82
    for title, value, tone in stats:
        _add_panel(slide, 9.86, y, 2.55, 0.88, WHITE, WHITE)
        _add_textbox(slide, 10.12, y + 0.12, 1.8, 0.14, title, 10, MUTED)
        _add_textbox(slide, 10.12, y + 0.42, 2.0, 0.22, value, 18, tone, True)
        y += 1.02


def _render_system_overview_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '主线页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 0.92, '4 layers', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 2.94, 1.66, 0.98, '1 chain', WHITE, MUTED)
    _add_panel(slide, 9.38, 1.62, 2.98, 0.34, WHITE)
    _add_textbox(slide, 9.6, 1.72, 2.46, 0.12, '目标: 一页先把层次、主链和证据关系讲清。', 8, MUTED)

    _add_panel(slide, 0.72, 1.88, 7.22, 4.88, WHITE)
    _add_panel(slide, 8.18, 1.88, 4.38, 4.88, WHITE)
    _add_textbox(slide, 0.98, 2.1, 1.6, 0.22, '四层闭环', 19, INK, True)
    _add_textbox(slide, 0.98, 2.44, 4.56, 0.18, '主讲版先把“软件怎么喂、控制怎么发、计算怎么跑、证据怎么收”放回同一页。', 10, MUTED)

    layers = [
        ('软件输入', 'model / quant / schedule 先把阶段语义定住。', BLUE, BLUE_BG, 'spec/json', 'imem.hex'),
        ('控制前端', 'AXI-Lite、IMEM 和 sequencer 把程序真正发起来。', ORANGE, ORANGE_BG, 'reg map', 'dispatch/wait'),
        ('计算核心', 'UB、SA、VPU 负责前向 / 反向 / 更新三段执行。', GREEN, GREEN_BG, 'UB -> SA', 'VPU writeback'),
        ('证据收束', 'wave、scoreboard、convergence 把控制链和数值链钉死。', PURPLE, PURPLE_BG, 'wave/VCD', '41/41 + loss'),
    ]
    y = 2.82
    for title, body, tone, fill, chip1, chip2 in layers:
        _add_panel(slide, 1.0, y, 6.58, 0.76, fill)
        _add_panel(slide, 1.12, y + 0.12, 0.05, 0.52, tone, tone)
        _add_textbox(slide, 1.32, y + 0.12, 1.28, 0.16, title, 15, INK, True)
        _add_textbox(slide, 2.68, y + 0.12, 3.18, 0.18, body, 10, INK)
        _add_chip_ppt(slide, 5.96, y + 0.22, 0.86, chip1, WHITE, tone)
        _add_chip_ppt(slide, 6.94, y + 0.22, 1.06, chip2, WHITE, tone)
        y += 0.92

    _add_textbox(slide, 8.46, 2.1, 1.78, 0.22, '端到端主链', 18, BLUE, True)
    _add_textbox(slide, 8.46, 2.42, 2.92, 0.16, '这条链决定后面 6 页按什么顺序讲。', 10, MUTED)
    steps = [
        ('compile', 'spec -> ub_map -> imem', BLUE, BLUE_BG),
        ('load', 'host 把参数和程序写进 SoC', ORANGE, ORANGE_BG),
        ('dispatch', 'sequencer 打 pulse + control decode', GREEN, GREEN_BG),
        ('run', 'UB / SA / VPU 执行并写回', PURPLE, PURPLE_BG),
        ('evidence', 'wave + scoreboard + convergence', BLUE, BLUE_BG),
    ]
    y = 2.86
    prev_y = None
    for title, body, tone, fill in steps:
        _add_panel(slide, 8.46, y, 3.58, 0.68, WHITE)
        _add_chip_ppt(slide, 8.64, y + 0.18, 0.92, title, fill, tone)
        _add_textbox(slide, 9.74, y + 0.18, 2.06, 0.16, body, 10, INK)
        if prev_y is not None:
            _add_arrow_ppt(slide, 9.02, prev_y + 0.68, 9.02, y, tone, 2.2)
        prev_y = y
        y += 0.82

    _add_panel(slide, 8.46, 6.12, 3.58, 0.42, HILITE, HILITE)
    _add_textbox(slide, 8.66, 6.24, 3.0, 0.12, '讲法: 总览先立结构，后面页再拆 UB / PE，证据页负责收口。', 8, MUTED)
    _add_static_footer(slide, spec['footer'])

def _render_rtl_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '架构页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 0.76, 'RTL', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 2.76, 1.66, 1.18, '3 domains', WHITE, MUTED)
    _add_chip_ppt(slide, 4.12, 1.66, 1.1, 'wrapper', WHITE, MUTED)
    _add_panel(slide, 9.34, 1.62, 3.02, 0.34, WHITE)
    _add_textbox(slide, 9.56, 1.72, 2.5, 0.12, '目标: 用一页把 wrapper、control、core、observe 讲清。', 8, MUTED)

    chain = [
        ('tpu_soc', 0.94, 2.08, 1.62, BLUE_BG, BLUE),
        ('frontend_axil', 3.12, 2.08, 1.86, ORANGE_BG, ORANGE),
        ('tpu.sv', 5.56, 2.08, 1.42, GREEN_BG, GREEN),
        ('observe / host', 9.94, 2.08, 1.9, PURPLE_BG, PURPLE),
    ]
    for label, x, y, w, fill, tone in chain:
        _add_mini_box(slide, x, y, w, 0.58, label, fill, tone, 12)
    _add_arrow_ppt(slide, 2.58, 2.37, 3.12, 2.37, BLUE, 2.4)
    _add_arrow_ppt(slide, 4.98, 2.37, 5.56, 2.37, ORANGE, 2.4)
    _add_arrow_ppt(slide, 6.98, 2.37, 9.94, 2.37, GREEN, 2.4)
    _add_textbox(slide, 7.3, 2.16, 2.2, 0.16, 'wrapper 之后就是 core execute path', 9, MUTED)

    _add_panel(slide, 0.72, 2.88, 3.72, 3.78, WHITE)
    _add_panel(slide, 0.84, 3.04, 0.05, 3.42, BLUE, BLUE)
    _add_textbox(slide, 1.0, 3.06, 1.96, 0.22, 'Frontend 控制域', 18, INK, True)
    _add_textbox(slide, 1.0, 3.4, 2.92, 0.16, 'AXI-Lite / reg map / IMEM / sequencer', 10, MUTED)
    _add_mini_box(slide, 1.04, 3.88, 1.18, 0.48, 'AXI-Lite', BLUE_BG, BLUE, 10)
    _add_mini_box(slide, 2.42, 3.88, 1.02, 0.48, 'Reg Map', WHITE, BLUE, 10)
    _add_mini_box(slide, 1.04, 4.56, 1.02, 0.48, 'IMEM', ORANGE_BG, ORANGE, 10)
    _add_mini_box(slide, 2.28, 4.56, 1.38, 0.48, 'Sequencer', ORANGE_BG, ORANGE, 10)
    _add_mini_box(slide, 1.04, 5.24, 1.34, 0.48, 'wait / drain', WHITE, ORANGE, 10)
    _add_mini_box(slide, 2.58, 5.24, 1.04, 0.48, 'host mux', WHITE, GREEN, 10)
    _add_textbox(slide, 1.04, 5.96, 2.98, 0.22, 'host 装载、程序起停、多拍等待和写口仲裁都在这里收口。', 9, MUTED)

    _add_panel(slide, 4.68, 2.88, 4.04, 3.78, WHITE)
    _add_panel(slide, 4.8, 3.04, 0.05, 3.42, ORANGE, ORANGE)
    _add_textbox(slide, 4.96, 3.06, 1.88, 0.22, 'Core 执行域', 18, INK, True)
    _add_textbox(slide, 4.96, 3.4, 3.18, 0.16, 'tpu.sv 把 UB / SA / VPU 组织成单条执行闭环。', 10, MUTED)
    _add_panel(slide, 5.12, 4.0, 1.2, 0.92, WHITE)
    _add_textbox(slide, 5.42, 4.28, 0.6, 0.18, 'UB', 18, INK, True, PP_ALIGN.CENTER)
    _add_panel(slide, 6.66, 4.0, 1.2, 0.92, WHITE)
    _add_textbox(slide, 6.86, 4.28, 0.8, 0.18, 'SA 2x2', 16, INK, True, PP_ALIGN.CENTER)
    _add_panel(slide, 8.2, 4.0, 1.2, 0.92, WHITE)
    _add_textbox(slide, 8.56, 4.28, 0.48, 0.18, 'VPU', 18, INK, True, PP_ALIGN.CENTER)
    _add_arrow_ppt(slide, 6.32, 4.46, 6.66, 4.46, BLUE, 2.4)
    _add_arrow_ppt(slide, 7.86, 4.46, 8.2, 4.46, GREEN, 2.4)
    _add_arrow_ppt(slide, 8.8, 4.92, 8.8, 5.38, ORANGE, 2.2)
    _add_arrow_ppt(slide, 8.2, 5.38, 6.28, 5.38, GREEN, 2.2)
    _add_chip_ppt(slide, 5.04, 5.86, 1.22, 'input / weight', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 6.46, 5.86, 1.12, 'psum', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 7.78, 5.86, 1.22, 'writeback', GREEN_BG, GREEN)

    _add_panel(slide, 8.96, 2.88, 3.6, 3.78, WHITE)
    _add_panel(slide, 9.08, 3.04, 0.05, 3.42, GREEN, GREEN)
    _add_textbox(slide, 9.24, 3.06, 1.72, 0.22, '对外观测', 18, INK, True)
    _add_textbox(slide, 9.24, 3.4, 2.34, 0.16, 'host / scoreboard 可以直接回读状态和输出。', 10, MUTED)
    _add_chip_ppt(slide, 9.24, 3.92, 0.86, 'busy', WHITE, PURPLE)
    _add_chip_ppt(slide, 10.26, 3.92, 0.86, 'valid', WHITE, GREEN)
    _add_chip_ppt(slide, 11.28, 3.92, 0.92, 'sys_out', WHITE, BLUE)
    _add_chip_ppt(slide, 9.24, 4.42, 0.92, 'vpu_out', WHITE, ORANGE)
    _add_chip_ppt(slide, 10.34, 4.42, 0.62, 'pc', WHITE, BLUE)
    _add_chip_ppt(slide, 11.12, 4.42, 0.9, 'state', WHITE, ORANGE)
    _add_textbox(slide, 9.24, 5.02, 2.6, 0.42, '波形页里真正要看的，就是这些 control/data observe 信号。', 10, INK)
    _add_panel(slide, 9.24, 5.64, 2.98, 0.42, HILITE, HILITE)
    _add_textbox(slide, 9.42, 5.76, 2.54, 0.12, '讲法: wrapper 先定边界，后面页再拆 compile / UB / PE。', 8, MUTED)

    _add_textbox(slide, 0.84, 6.98, 4.72, 0.14, 'wrapper chain: tpu_soc -> frontend_axil -> tpu.sv -> observe/host', 9, MUTED)
def _render_compiler_reference(slide, spec, idx):
    schedule = load_schedule()
    commands = schedule.get('commands', [])
    host_load_plan = schedule.get('host_load_plan', [])
    excerpt = []
    for cmd in commands[:4]:
        excerpt.append(f"{cmd['stage'][:14]:<14} {cmd['name'][:16]:<16} {cmd['kind']}")
    map_lines = []
    for item in host_load_plan[:3]:
        map_lines.append(f"{item['tensor']:<2} addr={item['addr']:<2} words={item['words']:<2} shape={item['shape']}")

    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '软件页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 1.2, 'Compiler', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 3.2, 1.66, 1.14, '5 phases', WHITE, MUTED)
    _add_chip_ppt(slide, 4.52, 1.66, 1.18, 'IMEM 落地', WHITE, MUTED)
    _add_panel(slide, 9.16, 1.62, 3.2, 0.34, WHITE)
    _add_textbox(slide, 9.38, 1.72, 2.68, 0.12, '目标: 讲清软件产物如何真正落成 frontend 可执行指令。', 8, MUTED)

    steps = [
        ('01', 'model spec', 'layer / quant / shape', BLUE, BLUE_BG),
        ('02', 'ub alloc', 'tensor map', GREEN, GREEN_BG),
        ('03', 'scheduler', 'stage + wait', ORANGE, ORANGE_BG),
        ('04', 'encode', '32-bit instr', PURPLE, PURPLE_BG),
        ('05', 'host load', 'write IMEM', BLUE, BLUE_BG),
    ]
    x = 0.82
    mids = []
    for num, title, body, tone, fill in steps:
        _add_panel(slide, x, 1.98, 1.94, 1.16, WHITE)
        _add_chip_ppt(slide, x + 0.14, 2.08, 0.46, num, fill, tone)
        _add_textbox(slide, x + 0.16, 2.42, 1.42, 0.18, title, 14, INK, True)
        _add_textbox(slide, x + 0.16, 2.72, 1.48, 0.16, body, 9, MUTED)
        mids.append(x + 0.97)
        x += 2.26
    for i in range(len(mids) - 1):
        _add_arrow_ppt(slide, mids[i] + 0.78, 2.56, mids[i + 1] - 0.78, 2.56, BLUE if i % 2 == 0 else ORANGE, 2.2)

    _add_metric_card(slide, 0.82, 3.48, 1.82, 0.92, '命令数', '66 commands', BLUE)
    _add_metric_card(slide, 2.88, 3.48, 1.82, 0.92, 'UB 占用', '41 / 128', GREEN)
    _add_metric_card(slide, 4.94, 3.48, 2.08, 0.92, '约束', '2x2 / 2-lane / Q8.8', ORANGE)

    _add_panel(slide, 0.72, 4.62, 6.2, 2.12, WHITE)
    _add_textbox(slide, 0.98, 4.88, 1.6, 0.18, 'schedule 片段', 15, BLUE, True)
    _add_code_card(slide, 0.96, 5.18, 5.68, 1.14, 'phase schedule', excerpt or ['schedule unavailable'], BLUE)
    _add_textbox(slide, 1.0, 6.5, 5.2, 0.14, '重点不是 cycle-accurate，而是先证明 schedule / wait_after / IMEM 已经能驱动前端执行。', 8, MUTED)

    _add_panel(slide, 7.14, 3.48, 5.42, 3.26, WHITE)
    _add_textbox(slide, 7.4, 3.74, 1.8, 0.18, '产物与指令落地', 15, BLUE, True)
    _add_chip_ppt(slide, 7.4, 4.06, 1.12, 'ub_map.json', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 8.72, 4.06, 1.2, 'schedule.json', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 10.12, 4.06, 0.92, 'imem.hex', BLUE_BG, BLUE)
    _add_textbox(slide, 7.4, 4.54, 1.34, 0.16, 'host load plan', 12, INK, True)
    _add_textbox(slide, 7.4, 4.84, 2.34, 0.62, '\n'.join(map_lines) if map_lines else 'host load plan unavailable', 9, INK)
    _add_textbox(slide, 10.0, 4.54, 1.46, 0.16, 'UB_RD 字段', 12, INK, True)
    _add_textbox(slide, 10.0, 4.84, 2.14, 0.68, 'opcode / wait_after\naddr / row / col\ntranspose / ptr_sel / pathway', 10, INK)
    _add_panel(slide, 7.4, 5.92, 4.5, 0.48, HILITE, HILITE)
    _add_textbox(slide, 7.6, 6.06, 4.0, 0.12, '讲法: 先定阶段语义，再编码成 frontend 真正能跑的 IMEM word。', 8, MUTED)
    _add_static_footer(slide, spec['footer'])

def _render_control_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '控制页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 1.0, '3 fixes', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 3.0, 1.66, 1.02, 'system bugs', WHITE, MUTED)
    _add_chip_ppt(slide, 4.2, 1.66, 0.94, 'closure', WHITE, MUTED)
    _add_panel(slide, 9.12, 1.62, 3.24, 0.34, WHITE)
    _add_textbox(slide, 9.34, 1.72, 2.72, 0.12, '目标: 讲清 3 个最关键的系统级收口点。', 8, MUTED)

    _add_textbox(slide, 0.82, 1.94, 3.28, 0.2, '3 个系统级收口点', 18, BLUE, True)
    _add_textbox(slide, 0.82, 2.24, 6.42, 0.18, '这页不讲“代码细节大全”，只讲真正决定系统是否稳定的等待、路径和仲裁。', 10, MUTED)

    focus = [
        (0.68, '等待语义', 'Fix 01', BLUE, ['Before', 'dispatch 后只能猜什么时候结束'], ['After', 'wait_after + vpu_drain + SEQ_WAIT'], ['Value', '多拍 VPU 有显式同步点'], ['seq_needs_wait = seq_instr[23];', 'vpu_drain = vpu_valid_prev && !tpu_vpu_valid_in;', 'if (seq_needs_wait) seq_state <= SEQ_WAIT;'], 'frontend_axil:139,153,177-179'),
        (4.38, '路径保持', 'Fix 02', ORANGE, ['Before', 'decode 只持续 1 拍，route 会漂移'], ['After', 'vpu_pathway_reg 持续保持'], ['Value', '同一条指令跨多拍仍走同一路径'], ['if (seq_instr_pulse && opcode == UB_RD)', '    vpu_pathway_reg <= seq_instr[22:19];', 'vpu_data_pathway_out = vpu_pathway_reg;'], 'frontend_axil:156-157,403'),
        (8.08, '写口仲裁', 'Fix 03', GREEN, ['Before', 'UB_PUSH 与 CU 写口会打架'], ['After', 'host write mux 让 UB_PUSH 优先'], ['Value', '参数装载和训练写回可以共存'], ['valid0 = push0 ? one : cu_valid0;', 'valid1 = push1 ? one : cu_valid1;', 'data0 = push0 ? host0 : cu_data0;'], 'frontend_axil:405-408'),
    ]
    for x, title, tag, tone, row1, row2, row3, lines, ref in focus:
        _add_panel(slide, x, 2.66, 3.52, 3.84, WHITE)
        _add_panel(slide, x + 0.12, 2.82, 0.05, 3.48, tone, tone)
        _add_chip_ppt(slide, x + 0.26, 2.84, 0.78, tag, WHITE, tone)
        _add_textbox(slide, x + 0.26, 3.2, 1.8, 0.22, title, 18, INK, True)
        _add_textbox(slide, x + 0.26, 3.62, 0.56, 0.16, row1[0], 11, tone, True)
        _add_textbox(slide, x + 0.92, 3.62, 2.16, 0.18, row1[1], 10, INK)
        _add_textbox(slide, x + 0.26, 4.08, 0.56, 0.16, row2[0], 11, tone, True)
        _add_textbox(slide, x + 0.92, 4.08, 2.2, 0.18, row2[1], 10, INK)
        _add_textbox(slide, x + 0.26, 4.54, 0.56, 0.16, row3[0], 11, tone, True)
        _add_textbox(slide, x + 0.92, 4.54, 2.18, 0.18, row3[1], 10, INK)
        _add_code_card(slide, x + 0.22, 5.1, 3.08, 0.96, 'RTL 特写', lines, tone)
        _add_textbox(slide, x + 1.86, 5.22, 1.18, 0.14, ref, 8, MUTED)

    _add_chip_ppt(slide, 0.9, 6.32, 1.2, 'wait_after', WHITE, BLUE)
    _add_chip_ppt(slide, 2.28, 6.32, 1.18, 'vpu_drain', WHITE, ORANGE)
    _add_chip_ppt(slide, 3.64, 6.32, 1.28, 'pathway_reg', WHITE, PURPLE)
    _add_chip_ppt(slide, 5.1, 6.32, 1.14, 'host mux', WHITE, GREEN)
    _add_textbox(slide, 0.72, 7.02, 11.0, 0.18, '关于 clock-gating：当前 RTL 里没有显式 ICG cell；更准确的说法是用 pe_enabled / valid / pulse 做收口。若继续工程化，可在 PE / UB / VPU 边界插标准门控。', 10, MUTED)


def _render_results_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '收束页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 0.88, 'Result', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 2.88, 1.66, 1.14, 'evidence', WHITE, MUTED)
    _add_chip_ppt(slide, 4.2, 1.66, 1.14, 'appendix', WHITE, MUTED)
    _add_panel(slide, 9.22, 1.62, 3.14, 0.34, WHITE)
    _add_textbox(slide, 9.44, 1.72, 2.62, 0.12, '目标: 用结论、边界和追问入口把主讲版收住。', 8, MUTED)
    for metric in spec.get('metrics', []):
        _add_metric_card(slide, metric['x'], metric['y'], metric['w'], metric['h'], metric['title'], metric['value'], metric['tone'])

    _add_panel(slide, 0.56, 3.02, 6.0, 2.96, WHITE)
    _add_textbox(slide, 0.86, 3.3, 1.58, 0.18, '一句话结论', 15, BLUE, True)
    _add_textbox(slide, 0.86, 3.76, 5.06, 0.56, '这不是单点算子 demo，而是一条从 compile 到 control、从 execute 到 evidence 的完整闭环。', 18, INK, True)
    _add_chip_ppt(slide, 0.88, 4.62, 1.46, '41 / 41 PASS', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 2.56, 4.62, 1.42, '12 epoch 收敛', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 4.22, 4.62, 1.52, 'XOR = (0,1,1,0)', ORANGE_BG, ORANGE)
    _add_textbox(slide, 0.86, 5.18, 5.12, 0.42, '主讲版到这里为止已经足够回答“项目怎么搭起来、怎么跑起来、怎么证明是真的”。', 10, MUTED)

    _add_panel(slide, 6.84, 3.02, 5.94, 1.26, WHITE)
    _add_textbox(slide, 7.14, 3.28, 1.16, 0.18, '项目边界', 15, ORANGE, True)
    _add_chip_ppt(slide, 7.14, 3.68, 0.98, 'prototype', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 8.3, 3.68, 0.92, '2x2/Q8.8', WHITE, BLUE)
    _add_chip_ppt(slide, 9.42, 3.68, 1.12, 'no DMA/IRQ', WHITE, ORANGE)
    _add_chip_ppt(slide, 10.76, 3.68, 1.02, 'no ICG', WHITE, GREEN)
    _add_textbox(slide, 7.14, 4.0, 5.1, 0.14, '边界说清楚，反而比把原型包装成商用品更可信。', 9, MUTED)

    _add_panel(slide, 6.84, 4.62, 5.94, 1.36, WHITE)
    _add_textbox(slide, 7.14, 4.88, 1.46, 0.18, '追问怎么接', 15, GREEN, True)
    _add_chip_ppt(slide, 7.14, 5.3, 1.04, '8x8 scale', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 8.36, 5.3, 0.96, 'DMA path', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 9.5, 5.3, 1.02, 'IRQ/status', PURPLE_BG, PURPLE)
    _add_chip_ppt(slide, 10.7, 5.3, 1.1, 'open appendix', BLUE_BG, BLUE)
    _add_textbox(slide, 7.14, 5.66, 5.08, 0.14, '被追问时直接切 appendix，继续展开 Frontend、控制修复、GIF 单步页。', 9, MUTED)

    _add_chip_ppt(slide, 0.84, 6.34, 1.22, '系统闭环', WHITE, BLUE)
    _add_chip_ppt(slide, 2.26, 6.34, 1.42, '真实证据', WHITE, ORANGE)
    _add_chip_ppt(slide, 3.92, 6.34, 1.36, '边界主动说明', WHITE, GREEN)
    _add_chip_ppt(slide, 5.52, 6.34, 1.32, 'appendix 可追问', WHITE, PURPLE)
    _add_textbox(slide, 0.5, 7.02, 12.0, 0.18, '收束口径：主讲版只讲最短闭环；需要深挖时，再切附录证明细节。', 10, MUTED)

def _render_frontend_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '控制页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 1.24, 'Frontend', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 3.26, 1.66, 1.42, 'config/runtime', WHITE, MUTED)
    _add_chip_ppt(slide, 4.86, 1.66, 1.18, 'hold/drain', WHITE, MUTED)
    _add_panel(slide, 9.12, 1.62, 3.24, 0.34, WHITE)
    _add_textbox(slide, 9.34, 1.72, 2.72, 0.12, '目标: 证明前端不是寄存器壳，而是控制收口点。', 8, MUTED)

    chain = [
        ('Host write', 0.82, 2.0, 1.72, BLUE_BG, BLUE),
        ('CTRL / STATUS', 2.82, 2.0, 1.84, ORANGE_BG, ORANGE),
        ('IMEM', 4.98, 2.0, 1.4, GREEN_BG, GREEN),
        ('sequencer', 6.64, 2.0, 1.72, PURPLE_BG, PURPLE),
        ('control_unit', 8.64, 2.0, 1.84, BLUE_BG, BLUE),
        ('ub_rd / host mux', 10.74, 2.0, 1.46, GREEN_BG, GREEN),
    ]
    for label, x, y, w, fill, tone in chain:
        _add_mini_box(slide, x, y, w, 0.62, label, fill, tone, 13)
    for x1, x2, tone in [(2.54, 2.82, BLUE), (4.66, 4.98, ORANGE), (6.38, 6.64, GREEN), (8.36, 8.64, PURPLE), (10.48, 10.74, BLUE)]:
        _add_arrow_ppt(slide, x1, 2.31, x2, 2.31, tone, 2.4)
    _add_textbox(slide, 0.82, 2.82, 11.1, 0.16, '把 host 写寄存器、IMEM 程序装载、状态机推进和 decode 输出串成一条主执行链。', 10, MUTED)

    _add_panel(slide, 0.74, 3.36, 3.54, 3.12, WHITE)
    _add_panel(slide, 0.86, 3.52, 0.05, 2.78, BLUE, BLUE)
    _add_textbox(slide, 1.02, 3.54, 1.8, 0.2, '寄存器与状态', 17, INK, True)
    regs = [
        ('0x00', 'CTRL', 'step / start'),
        ('0x04', 'STATUS', 'busy / running'),
        ('0x10', 'INSTR_W0', 'step 模式指令'),
        ('0x20/24', 'UB_DATA/PUSH', 'host 写 UB'),
        ('0x30/34/40/44', 'IMEM_*', '程序写入'),
        ('0x50/54/58', 'LEAK/INV/LR', '运行参数'),
    ]
    y = 4.02
    for addr, name, note in regs:
        _add_panel(slide, 1.02, y, 2.82, 0.34, BLUE_BG, BLUE_BG)
        _add_textbox(slide, 1.14, y + 0.1, 0.7, 0.12, addr, 8, BLUE, True)
        _add_textbox(slide, 1.82, y + 0.1, 0.92, 0.12, name, 9, INK, True)
        _add_textbox(slide, 2.78, y + 0.1, 0.84, 0.12, note, 8, MUTED)
        y += 0.42

    _add_panel(slide, 4.5, 3.36, 4.04, 3.12, WHITE)
    _add_panel(slide, 4.62, 3.52, 0.05, 2.78, ORANGE, ORANGE)
    _add_textbox(slide, 4.78, 3.54, 1.7, 0.2, 'IMEM 与指令', 17, INK, True)
    _add_chip_ppt(slide, 4.8, 3.94, 0.92, 'imem[0:63]', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 5.9, 3.94, 1.12, 'imem_len_reg', GREEN_BG, GREEN)
    _add_textbox(slide, 4.8, 4.4, 3.18, 0.24, 'NOP / SWITCH / UB_RD / UB_WR_HOST', 12, INK)
    _add_textbox(slide, 4.8, 4.86, 3.22, 0.56, 'UB_RD 核心字段：\\naddr / row / col / transpose / ptr_sel / pathway', 10, MUTED)
    _add_code_card(slide, 4.78, 5.46, 3.42, 0.72, 'Frontend 关键点', ['CTRL.step / CTRL.start', 'IMEM_ADDR / IMEM_W0 / IMEM_WE', 'IMEM_LEN defines valid program range'], ORANGE)

    _add_panel(slide, 8.76, 3.36, 3.8, 3.12, WHITE)
    _add_panel(slide, 8.88, 3.52, 0.05, 2.78, GREEN, GREEN)
    _add_textbox(slide, 9.04, 3.54, 2.0, 0.2, '状态机与关键信号', 17, INK, True)
    for label, x, y, fill, tone in [
        ('IDLE', 9.2, 4.06, BLUE_BG, BLUE),
        ('DISPATCH', 10.46, 4.06, ORANGE_BG, ORANGE),
        ('WAIT', 9.2, 4.88, GREEN_BG, GREEN),
        ('ADVANCE', 10.46, 4.88, PURPLE_BG, PURPLE),
    ]:
        _add_mini_box(slide, x, y, 0.94, 0.46, label, fill, tone, 10)
    _add_arrow_ppt(slide, 10.14, 4.28, 10.46, 4.28, BLUE, 2.0)
    _add_arrow_ppt(slide, 10.92, 4.52, 10.92, 4.88, ORANGE, 2.0)
    _add_arrow_ppt(slide, 10.46, 5.1, 10.14, 5.1, GREEN, 2.0)
    _add_arrow_ppt(slide, 9.68, 4.88, 9.68, 4.52, PURPLE, 2.0)
    _add_chip_ppt(slide, 9.16, 5.74, 1.08, 'seq_instr_pulse', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 10.46, 5.74, 0.96, 'vpu_drain', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 11.58, 5.74, 0.82, 'pathway', PURPLE_BG, PURPLE)
    _add_textbox(slide, 9.04, 6.18, 3.1, 0.18, '关键点：dispatch 只打一拍，pathway_reg 把 decode 结果延成多拍有效。', 9, MUTED)
    _add_static_footer(slide, spec['footer'])

def _render_ub_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '数据页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 1.46, 'Unified Buffer', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 3.48, 1.66, 1.12, 'hub', WHITE, MUTED)
    _add_chip_ppt(slide, 4.78, 1.66, 1.26, 'restore/update', WHITE, MUTED)
    _add_panel(slide, 9.16, 1.62, 3.2, 0.34, WHITE)
    _add_textbox(slide, 9.38, 1.72, 2.68, 0.12, '目标: 强调 UB 是汇合点，不是孤立 SRAM。', 8, MUTED)

    _add_panel(slide, 0.72, 1.88, 2.72, 4.88, WHITE)
    _add_panel(slide, 0.84, 2.04, 0.05, 4.54, BLUE, BLUE)
    _add_textbox(slide, 1.0, 2.08, 1.56, 0.22, '地址布局', 18, INK, True)
    _add_textbox(slide, 1.0, 2.42, 1.84, 0.18, '统一地址空间里的张量驻留。', 10, MUTED)
    segs = [
        ('X', 2.78, 0.18, BLUE_BG, BLUE, '@0 8w'),
        ('Y', 3.04, 0.12, PURPLE_BG, PURPLE, '@8 4w'),
        ('W1/B1', 3.24, 0.18, GREEN_BG, GREEN, '@12 6w'),
        ('W2/B2', 3.5, 0.12, ORANGE_BG, ORANGE, '@18 3w'),
        ('H1', 3.7, 0.18, ORANGE_BG, ORANGE, '@21 8w'),
        ('dZ2', 3.96, 0.12, GREEN_BG, GREEN, '@29 4w'),
        ('dZ1', 4.16, 0.18, BLUE_BG, BLUE, '@33 8w'),
    ]
    for label, y, h, fill, tone, note in segs:
        _add_panel(slide, 1.16, y, 0.74, h, fill, fill)
        _add_textbox(slide, 1.32, y + 0.03, 0.4, 0.12, label, 9, tone, True)
        _add_textbox(slide, 2.06, y + 0.02, 0.88, 0.12, note, 9, tone)
    _add_chip_ppt(slide, 1.0, 6.1, 0.94, 'alloc=41', WHITE, BLUE)
    _add_chip_ppt(slide, 2.08, 6.1, 0.86, 'free=87', WHITE, MUTED)

    _add_panel(slide, 3.68, 1.88, 4.72, 4.88, WHITE)
    _add_panel(slide, 3.8, 2.04, 0.05, 4.54, ORANGE, ORANGE)
    _add_textbox(slide, 3.96, 2.08, 2.2, 0.22, 'UB 作为数据汇合点', 18, INK, True)
    _add_textbox(slide, 3.96, 2.42, 3.9, 0.18, '同一块存储同时接 host load、core read、writeback 和 update。', 10, MUTED)
    _add_mini_box(slide, 4.1, 2.96, 1.28, 0.52, 'host load', BLUE_BG, BLUE, 11)
    _add_mini_box(slide, 6.6, 2.96, 1.16, 0.52, 'core read', ORANGE_BG, ORANGE, 11)
    _add_panel(slide, 5.42, 3.62, 1.18, 1.0, WHITE)
    _add_textbox(slide, 5.76, 3.96, 0.5, 0.18, 'UB', 20, INK, True, PP_ALIGN.CENTER)
    _add_mini_box(slide, 4.04, 5.08, 1.58, 0.52, 'VPU writeback', GREEN_BG, GREEN, 11)
    _add_mini_box(slide, 6.32, 5.08, 1.22, 0.52, 'grad update', PURPLE_BG, PURPLE, 11)
    _add_arrow_ppt(slide, 5.38, 3.22, 5.42, 4.02, BLUE, 2.2)
    _add_arrow_ppt(slide, 6.6, 3.22, 6.58, 4.02, ORANGE, 2.2)
    _add_arrow_ppt(slide, 5.06, 5.08, 5.54, 4.62, GREEN, 2.2)
    _add_arrow_ppt(slide, 6.92, 5.08, 6.48, 4.62, PURPLE, 2.2)
    _add_chip_ppt(slide, 3.96, 6.12, 1.62, 'ub_wr_host_*', WHITE, BLUE)
    _add_chip_ppt(slide, 5.78, 6.12, 1.28, 'ptr_sel 0..6', WHITE, ORANGE)
    _add_chip_ppt(slide, 7.24, 6.12, 0.9, 'in-UB', WHITE, GREEN)

    _add_panel(slide, 8.64, 1.88, 3.92, 4.88, WHITE)
    _add_panel(slide, 8.76, 2.04, 0.05, 4.54, GREEN, GREEN)
    _add_textbox(slide, 8.92, 2.08, 2.12, 0.22, '关键语义', 18, INK, True)
    items = [
        ('01', 'wr_ptr_base / restore', 'host 装完先记 base，start 时恢复写指针，避免踩参数区。', BLUE, BLUE_BG),
        ('02', 'ptr_select', 'input / weight / bias / Y / H / grad 由统一指针语义切换。', ORANGE, ORANGE_BG),
        ('03', 'gradient_descent', '更新逻辑放在 UB 内部，不必绕回 host 再改写。', GREEN, GREEN_BG),
    ]
    y = 2.78
    for tag, title, body, tone, fill in items:
        _add_panel(slide, 8.92, y, 3.2, 0.78, fill)
        _add_chip_ppt(slide, 9.08, y + 0.22, 0.44, tag, WHITE, tone)
        _add_textbox(slide, 9.68, y + 0.16, 1.68, 0.14, title, 11, tone, True)
        _add_textbox(slide, 9.68, y + 0.38, 2.18, 0.18, body, 9, INK)
        y += 0.92
    _add_code_card(slide, 8.9, 5.64, 3.26, 0.84, 'RTL 关键段', [
        'case (ub_ptr_select)',
        'if (ub_wr_ptr_restore_in) wr_ptr <= wr_ptr_base;',
        'gradient_descent lives inside UB;'
    ], BLUE)
    _add_static_footer(slide, spec['footer'])

def _render_pe_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '计算页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 0.74, 'PE', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 2.74, 1.66, 1.18, '2x2 array', WHITE, MUTED)
    _add_chip_ppt(slide, 4.1, 1.66, 1.24, 'wavefront', WHITE, MUTED)

    _add_panel(slide, 0.72, 1.9, 6.24, 4.86, WHITE)
    _add_panel(slide, 0.84, 2.06, 0.05, 4.52, BLUE, BLUE)
    _add_textbox(slide, 1.0, 2.08, 2.4, 0.22, '2 x 2 Systolic 组织', 19, INK, True)
    _add_textbox(slide, 1.0, 2.42, 3.86, 0.18, '主讲版只把三件事讲清：输入从左进、权重从上装、psum 向下沉。', 10, MUTED)
    pe_boxes = [
        ('PE11', 1.42, 3.04), ('PE12', 3.62, 3.04),
        ('PE21', 1.42, 4.92), ('PE22', 3.62, 4.92),
    ]
    for label, x, y in pe_boxes:
        _add_clean_card(slide, x, y, 1.62, 1.26, label, 'active w\ninactive w\nmac + valid', BLUE, WHITE, 10)
    _add_arrow_ppt(slide, 1.02, 3.66, 1.42, 3.66, BLUE, 2.4)
    _add_arrow_ppt(slide, 3.2, 3.66, 3.62, 3.66, GREEN, 2.4)
    _add_arrow_ppt(slide, 1.02, 5.54, 1.42, 5.54, BLUE, 2.4)
    _add_arrow_ppt(slide, 3.2, 5.54, 3.62, 5.54, GREEN, 2.4)
    _add_arrow_ppt(slide, 2.24, 2.72, 2.24, 3.04, PURPLE, 2.4)
    _add_arrow_ppt(slide, 4.44, 2.72, 4.44, 3.04, PURPLE, 2.4)
    _add_arrow_ppt(slide, 2.24, 4.32, 2.24, 4.92, ORANGE, 2.4)
    _add_arrow_ppt(slide, 4.44, 4.32, 4.44, 4.92, ORANGE, 2.4)
    _add_textbox(slide, 1.86, 2.5, 0.72, 0.14, 'weight', 9, PURPLE, True)
    _add_textbox(slide, 4.04, 2.5, 0.72, 0.14, 'weight', 9, PURPLE, True)
    _add_textbox(slide, 0.84, 3.56, 0.46, 0.14, 'input', 9, BLUE, True)
    _add_textbox(slide, 0.84, 5.44, 0.46, 0.14, 'input', 9, BLUE, True)
    _add_textbox(slide, 5.42, 3.62, 0.84, 0.14, 'valid', 9, GREEN, True)
    _add_textbox(slide, 5.46, 5.44, 0.56, 0.14, 'out', 9, GREEN, True)
    _add_chip_ppt(slide, 1.02, 6.18, 1.22, 'vertical psum', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 2.52, 6.18, 1.4, 'horizontal data', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 4.14, 6.18, 1.7, 'shadow -> active', PURPLE_BG, PURPLE)

    _add_panel(slide, 7.18, 1.9, 5.1, 4.86, WHITE)
    _add_panel(slide, 7.3, 2.06, 0.05, 4.52, ORANGE, ORANGE)
    _add_textbox(slide, 7.46, 2.08, 2.0, 0.22, '为什么这页值得讲', 19, INK, True)
    _add_code_card(slide, 7.42, 2.68, 4.64, 1.46, 'pe.sv 核心', [
        'mult = input * weight_reg_active',
        'psum = mult + pe_psum_in',
        'switch: inactive -> active',
        'enabled off: clear out, keep state',
    ], BLUE)
    cards = [
        ('双 weight reg', 'load / switch 分离，所以换权重不必打断计算。', BLUE, BLUE_BG),
        ('enabled 语义', 'pe_enabled 只清输出不清状态，列裁剪更自然。', GREEN, GREEN_BG),
        ('波前可观测', 'valid、switch、psum 都能在波形页对应起来。', ORANGE, ORANGE_BG),
    ]
    y = 4.42
    for title, body, tone, fill in cards:
        _add_panel(slide, 7.42, y, 4.64, 0.62, fill)
        _add_textbox(slide, 7.6, y + 0.12, 1.2, 0.14, title, 11, tone, True)
        _add_textbox(slide, 8.96, y + 0.12, 2.84, 0.14, body, 9, INK)
        y += 0.76
    _add_chip_ppt(slide, 7.42, 6.22, 1.28, 'weight-stationary', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 8.92, 6.22, 1.18, 'pulse control', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 10.34, 6.22, 1.1, 'column gate', GREEN_BG, GREEN)
    _add_static_footer(slide, spec['footer'])

def _render_pe_animation_demo(slide, spec, idx):
    step = int(spec.get('step', 0))
    step_titles = [
        'T0 权重就绪',
        'T1 首拍输入进入',
        'T2 波前推进',
        'T3 第一列先到底',
        'T4 结果收束',
        'T5 写回 / 送下游',
    ]
    step_bodies = [
        'top-load / switch 完成，active weight 已稳定。',
        '左边界发入 input，PE11 / PE21 开始第一拍 MAC。',
        'valid 向右、psum 向下，第二列开始接力计算。',
        '第一列结果先到底，阵列继续完成后续传播。',
        '2x2 输出矩阵收齐，结果窗完整可见。',
        'sys output 被收集，后续可写回 UB 或继续送 VPU。',
    ]
    step_tones = [PURPLE, BLUE, GREEN, ORANGE, RED, GREEN]
    step_fills = [PURPLE_BG, BLUE_BG, GREEN_BG, ORANGE_BG, RED_BG, GREEN_BG]

    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)

    _add_panel(slide, 0.72, 1.88, 7.18, 4.86, WHITE)
    _add_panel(slide, 8.16, 1.88, 4.4, 4.86, WHITE)
    _add_chip_ppt(slide, 0.98, 1.62, 1.02, 'GIF 单步', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 2.16, 1.62, 0.92, 'PE', step_fills[step], step_tones[step])
    _add_chip_ppt(slide, 3.24, 1.62, 1.0, '6 steps', WHITE, MUTED)
    _add_chip_ppt(slide, 4.4, 1.62, 1.14, '点按翻页', HILITE, INK)
    _add_textbox(slide, 0.98, 2.08, 3.2, 0.22, 'PE 计算过程单步演示', 19, INK, True)
    _add_textbox(slide, 0.98, 2.42, 6.24, 0.2, '点一下翻到下一页，下一步 GIF 会自动播放。', 10, MUTED)
    _add_picture_named(slide, PE_ANIM_BASE_PATH, 0.94, 2.74, 6.68, 3.96, 'PE_ANIM_BASE')
    _add_picture_named(slide, PE_ANIM_GIF_PATHS[step], 0.94, 2.74, 6.68, 3.96, f'PE_ANIM_STEP_{step}_GIF')

    _add_chip_ppt(slide, 0.98, 6.18, 1.2, f'当前: T{step}', step_fills[step], step_tones[step])
    x = 2.42
    for idx_step in range(6):
        fill = step_fills[idx_step] if idx_step == step else WHITE
        tone = step_tones[idx_step] if idx_step == step else MUTED
        _add_chip_ppt(slide, x, 6.18, 0.62, f'T{idx_step}', fill, tone)
        x += 0.76

    _add_textbox(slide, 8.42, 2.08, 2.6, 0.22, '当前这一步怎么讲', 17, INK, True)
    _add_textbox(slide, 8.42, 2.46, 3.72, 0.18, step_titles[step], 15, step_tones[step], True)
    _add_anim_note_box(slide, 8.42, 2.84, 3.86, 0.72, step_titles[step], step_bodies[step], step_tones[step], step_fills[step])
    _add_textbox(slide, 8.42, 3.78, 3.72, 0.68, '放映建议：\n1. 当前页只讲当前一步。\n2. 讲完后点击一下，直接切下一页。\n3. 下一页的 GIF 会自动播放下一步。', 10, INK)
    _add_textbox(slide, 8.42, 4.82, 3.72, 0.6, '采用分步 GIF 页，兼顾放映效果和 Windows PowerPoint 兼容性。', 10, MUTED)
    _add_chip_ppt(slide, 8.42, 5.86, 1.52, '兼容优先', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 10.4, 5.86, 1.34, 'GIF 自动播放', GREEN_BG, GREEN)

    _add_static_footer(slide, spec['footer'])


def _render_wave_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.94, '证据页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.84, 1.66, 0.82, 'Wave', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 2.82, 1.66, 1.1, 'VCD line', WHITE, MUTED)
    _add_chip_ppt(slide, 4.08, 1.66, 1.22, 'scoreboard', WHITE, MUTED)
    _add_panel(slide, 9.18, 1.62, 3.18, 0.34, WHITE)
    _add_textbox(slide, 9.4, 1.72, 2.66, 0.12, '目标: 用一页真实波形把控制链和数据链钉死。', 8, MUTED)

    _add_panel(slide, 0.58, 1.88, 8.9, 4.88, WHITE)
    _add_textbox(slide, 0.86, 2.12, 3.42, 0.22, 'VCD 片段：第一次 busy 窗口', 18, INK, True)
    _add_textbox(slide, 0.86, 2.46, 2.52, 0.16, '来源: waveforms/tpu_soc.vcd', 10, MUTED)
    signals = [
        ('busy_reg', BLUE, 0),
        ('seq_instr_pulse', ORANGE, 1),
        ('ub_rd_start', GREEN, 2),
        ('sys_switch', PURPLE, 3),
        ('ub_rd_input_valid_out_0', BLUE, 4),
        ('sys_valid_out_21', ORANGE, 5),
        ('vpu_valid_out_1', GREEN, 6),
        ('vpu_drain', RED, 7),
    ]
    x0 = 2.42
    x1 = 8.78
    y0 = 2.96
    row_h = 0.42
    ticks = [2.9, 3.96, 5.02, 6.08, 7.14, 8.2, 9.26]
    for x in ticks[:-1]:
        tick = slide.shapes.add_connector(1, Inches(x), Inches(2.92), Inches(x), Inches(6.28))
        tick.line.color.rgb = _rgb(LINE)
        tick.line.width = Pt(1)
    pulse_sets = {
        'busy_reg': [(0.0, 5.82)],
        'seq_instr_pulse': [(0.58, 0.14), (1.34, 0.14), (2.12, 0.14), (2.88, 0.14), (3.66, 0.14), (4.42, 0.14), (5.18, 0.14)],
        'ub_rd_start': [(0.62, 0.06), (0.98, 0.06), (1.76, 0.06), (2.54, 0.06), (3.32, 0.06), (4.1, 0.06), (4.86, 0.06), (5.56, 0.06)],
        'sys_switch': [(0.9, 0.05), (1.68, 0.05), (2.46, 0.05), (3.24, 0.05), (4.02, 0.05), (4.8, 0.05), (5.52, 0.05)],
        'ub_rd_input_valid_out_0': [(1.0, 0.12), (1.78, 0.12), (2.56, 0.12), (3.34, 0.12), (4.12, 0.12), (4.9, 0.12), (5.6, 0.12)],
        'sys_valid_out_21': [(1.08, 0.1), (1.86, 0.1), (2.64, 0.1), (3.42, 0.1), (4.2, 0.1), (4.98, 0.1), (5.68, 0.1)],
        'vpu_valid_out_1': [(1.14, 0.12), (1.92, 0.12), (2.7, 0.12), (3.48, 0.12), (4.26, 0.12), (5.04, 0.12), (5.74, 0.12)],
        'vpu_drain': [(1.34, 0.05), (2.12, 0.05), (2.9, 0.05), (3.68, 0.05), (4.46, 0.05), (5.24, 0.05), (5.92, 0.05)],
    }
    for label, tone, idx_row in signals:
        y = y0 + idx_row * row_h
        _add_textbox(slide, 0.98, y - 0.02, 1.28, 0.14, label, 8, INK, True)
        base = slide.shapes.add_connector(1, Inches(x0), Inches(y + 0.14), Inches(x1), Inches(y + 0.14))
        base.line.color.rgb = _rgb(tone)
        base.line.width = Pt(1.8)
        for start, width in pulse_sets[label]:
            xl = x0 + start
            xr = xl + width
            top = slide.shapes.add_connector(1, Inches(xl), Inches(y - 0.02), Inches(xr), Inches(y - 0.02))
            top.line.color.rgb = _rgb(tone)
            top.line.width = Pt(1.8)
            left = slide.shapes.add_connector(1, Inches(xl), Inches(y + 0.14), Inches(xl), Inches(y - 0.02))
            right = slide.shapes.add_connector(1, Inches(xr), Inches(y - 0.02), Inches(xr), Inches(y + 0.14))
            left.line.color.rgb = _rgb(tone)
            right.line.color.rgb = _rgb(tone)
            left.line.width = Pt(1.8)
            right.line.width = Pt(1.8)
    for t, x in [('10.64us', 2.34), ('11.01us', 3.4), ('11.37us', 4.46), ('11.73us', 5.52), ('12.10us', 6.58), ('12.46us', 7.64), ('12.83us', 8.7)]:
        _add_textbox(slide, x - 0.12, 6.28, 0.66, 0.14, t, 8, MUTED)
    _add_chip_ppt(slide, 2.72, 6.46, 0.94, 'dispatch', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 4.14, 6.46, 0.98, 'ub read', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 5.58, 6.46, 0.92, 'sys out', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 6.96, 6.46, 0.82, 'drain', RED_BG, RED)

    _add_panel(slide, 9.74, 1.88, 2.84, 2.12, WHITE)
    _add_textbox(slide, 10.0, 2.14, 1.78, 0.2, '覆盖项', 16, ORANGE, True)
    cover_rows = [('H1', '8'), ('dZ2', '4'), ('dZ1', '8'), ('UB dZ2', '4'), ('UB dZ1', '8'), ('W1/B1/W2/B2', '9')]
    y = 2.52
    for label, value in cover_rows:
        _add_mini_box(slide, 10.0, y, 0.7, 0.24, value, BLUE_BG, BLUE, 8)
        _add_textbox(slide, 10.84, y + 0.05, 1.36, 0.12, label, 9, INK)
        y += 0.3

    _add_panel(slide, 9.74, 4.18, 2.84, 2.58, WHITE)
    _add_textbox(slide, 10.0, 4.44, 1.56, 0.2, '讲法顺序', 16, GREEN, True)
    _add_chip_ppt(slide, 10.0, 4.84, 0.72, 'busy', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 10.84, 4.84, 0.9, 'dispatch', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 11.88, 4.84, 0.72, 'drain', GREEN_BG, GREEN)
    _add_textbox(slide, 10.0, 5.28, 2.16, 0.52, '先看 busy 窗口，\n再讲 dispatch -> ub read -> sys out，\n最后回到 scoreboard。', 10, INK)
    _add_textbox(slide, 10.0, 6.02, 2.24, 0.26, '这页证明控制链和数据链在同一时间窗内闭合。', 9, MUTED)
    _add_static_footer(slide, spec['footer'])


def _render_frontend_split_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)

    _add_panel(slide, 0.72, 1.94, 4.08, 4.76, WHITE)
    _add_panel(slide, 0.84, 2.1, 0.05, 4.42, BLUE, BLUE)
    _add_textbox(slide, 1.0, 2.12, 1.9, 0.22, '配置面：host 写入', 18, INK, True)
    _add_mini_box(slide, 1.1, 2.86, 1.56, 0.5, 'UB_DATA / UB_PUSH', BLUE_BG, BLUE, 11)
    _add_mini_box(slide, 1.1, 3.62, 1.72, 0.5, 'IMEM_ADDR / WE', ORANGE_BG, ORANGE, 11)
    _add_mini_box(slide, 1.1, 4.38, 1.7, 0.5, 'IMEM_LEN / START', GREEN_BG, GREEN, 11)
    _add_mini_box(slide, 1.1, 5.14, 1.46, 0.5, 'LEAK / INV / LR', PURPLE_BG, PURPLE, 11)
    for y1, y2, tone in [(3.36, 3.62, BLUE), (4.12, 4.38, ORANGE), (4.88, 5.14, GREEN)]:
        _add_arrow_ppt(slide, 1.86, y1, 1.86, y2, tone, 2.2)
    _add_textbox(slide, 3.06, 2.92, 1.38, 1.9, 'host 先做两件事：\\n1. 连续把参数和样本写进 UB\\n2. 写入 IMEM 程序和执行长度\\n\\n训练参数 leak / inv_batch / learning_rate 也都在这条路径上配置完成。', 11, INK)
    _add_textbox(slide, 3.06, 5.32, 1.34, 0.5, '这页左边只讲“怎么把程序和数据喂进去”，不展开运行时细节。', 9, MUTED)

    _add_panel(slide, 5.02, 1.94, 7.56, 4.76, WHITE)
    _add_panel(slide, 5.14, 2.1, 0.05, 4.42, ORANGE, ORANGE)
    _add_textbox(slide, 5.3, 2.12, 2.5, 0.22, '运行面：sequencer + decode + 持续路径', 18, INK, True)
    flow = [
        ('IMEM', 5.48, 2.92, 1.02, ORANGE_BG, ORANGE),
        ('DISPATCH', 7.06, 2.92, 1.24, BLUE_BG, BLUE),
        ('WAIT', 8.9, 2.92, 0.92, GREEN_BG, GREEN),
        ('ADVANCE', 10.2, 2.92, 1.12, PURPLE_BG, PURPLE),
    ]
    for label, x, y, w, fill, tone in flow:
        _add_mini_box(slide, x, y, w, 0.54, label, fill, tone, 12)
    for x1, x2, tone in [(6.5, 7.06, ORANGE), (8.3, 8.9, BLUE), (9.82, 10.2, GREEN)]:
        _add_arrow_ppt(slide, x1, 3.2, x2, 3.2, tone, 2.2)
    _add_mini_box(slide, 7.34, 4.06, 2.2, 0.58, 'control_unit decode', PURPLE_BG, PURPLE, 11)
    _add_arrow_ppt(slide, 8.16, 3.46, 8.16, 4.06, ORANGE, 2.2)
    _add_chip_ppt(slide, 5.52, 5.02, 1.12, 'seq_instr_pulse', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 6.88, 5.02, 0.96, 'ub_rd_*', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 8.08, 5.02, 1.02, 'pathway hold', PURPLE_BG, PURPLE)
    _add_chip_ppt(slide, 9.38, 5.02, 0.94, 'vpu_drain', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 10.52, 5.02, 1.16, 'host write mux', BLUE_BG, BLUE)
    _add_textbox(slide, 5.5, 5.52, 6.18, 0.64, 'start_pulse 拉起 busy 并把 pc 装到 0。DISPATCH 只打一拍 seq_instr_pulse；如果 wait_after 置位，状态机会在 WAIT 等到 vpu_drain，再进入 ADVANCE。', 10, INK)
    _add_code_card(slide, 5.46, 6.14, 6.66, 0.52, '关键 RTL', ['seq_instr_pulse only in dispatch', 'vpu_pathway_reg latches decode result', 'vpu_drain closes the multi-cycle action'], ORANGE)
    _add_static_footer(slide, spec['footer'])

def _render_wptr_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)

    _add_panel(slide, 0.74, 1.96, 11.84, 1.6, WHITE)
    _add_textbox(slide, 1.0, 2.18, 2.2, 0.2, '把装载阶段和训练阶段分开', 17, INK, True)
    stages = [
        ('阶段 1 host load', 1.06, 2.66, 2.56, 'UB_DATA/UB_PUSH 连续装参数，wr_ptr 持续前进', BLUE_BG, BLUE),
        ('阶段 2 start restore', 4.36, 2.66, 2.68, 'start_pulse 触发 restore，把 wr_ptr 拉回 wr_ptr_base', ORANGE_BG, ORANGE),
        ('阶段 3 runtime writeback', 7.86, 2.66, 3.0, 'VPU 写回从 base 之后继续写，不踩前面参数', GREEN_BG, GREEN),
    ]
    mids = []
    for title, x, y, w, body, fill, tone in stages:
        _add_panel(slide, x, y, w, 0.64, fill, fill)
        _add_textbox(slide, x + 0.16, y + 0.12, w - 0.28, 0.14, title, 12, tone, True)
        _add_textbox(slide, x + 0.16, y + 0.34, w - 0.28, 0.12, body, 9, INK)
        mids.append((x, w, tone))
    _add_arrow_ppt(slide, 3.7, 2.98, 4.36, 2.98, BLUE, 2.2)
    _add_arrow_ppt(slide, 7.12, 2.98, 7.86, 2.98, ORANGE, 2.2)

    _add_panel(slide, 0.74, 3.9, 4.38, 2.78, WHITE)
    _add_panel(slide, 0.86, 4.06, 0.05, 2.44, BLUE, BLUE)
    _add_textbox(slide, 1.02, 4.08, 1.6, 0.2, '核心状态', 17, INK, True)
    states = [
        ('wr_ptr', '当前统一写指针，host 与 VPU 写回共用'),
        ('wr_ptr_next', '按 lane valid 预计算下一地址'),
        ('wr_ptr_base', '记录 host 装载完成后的“参数区尾部”'),
        ('ub_wr_ptr_restore_in', 'pulse 拉起时恢复 wr_ptr = wr_ptr_base'),
    ]
    y = 4.54
    for label, body in states:
        _add_textbox(slide, 1.02, y, 1.12, 0.16, label, 11, BLUE, True)
        _add_textbox(slide, 2.04, y, 2.5, 0.16, body, 10, INK)
        y += 0.46
    _add_code_card(slide, 1.0, 5.76, 3.88, 0.7, '关键 RTL', ['if (ub_wr_host_valid_in[0] || ub_wr_host_valid_in[1])', '    wr_ptr_base <= wr_ptr_next;', 'if (ub_wr_ptr_restore_in) wr_ptr <= wr_ptr_base;'], BLUE)

    _add_panel(slide, 5.38, 3.9, 3.44, 2.78, WHITE)
    _add_panel(slide, 5.5, 4.06, 0.05, 2.44, ORANGE, ORANGE)
    _add_textbox(slide, 5.66, 4.08, 1.8, 0.2, '为什么重要', 17, INK, True)
    _add_textbox(slide, 5.66, 4.54, 2.7, 1.16, '1. 它解释了 Unified Buffer 为什么能既做参数装载又做训练写回。\n2. 没有 restore，训练一开始就会从参数尾后继续写。\n3. 这个设计把“参数静态区”和“运行动态区”用时序切开。', 11, INK)
    _add_chip_ppt(slide, 5.68, 6.12, 1.58, 'host load captures base', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 5.68, 6.44, 1.34, 'start restores ptr', ORANGE_BG, ORANGE)

    _add_panel(slide, 9.06, 3.9, 3.5, 2.78, WHITE)
    _add_panel(slide, 9.18, 4.06, 0.05, 2.44, GREEN, GREEN)
    _add_textbox(slide, 9.34, 4.08, 2.0, 0.2, '面试回答口径', 17, INK, True)
    _add_textbox(slide, 9.34, 4.54, 2.78, 1.02, '为什么写回不会覆盖 W/B 参数？\n因为 host load 结束时先记住 wr_ptr_base；start 时把 wr_ptr restore 回 base；后续 runtime writeback 都从 base 之后继续写。', 11, INK)
    _add_chip_ppt(slide, 9.34, 6.12, 1.76, 'runtime writes after base', GREEN_BG, GREEN)
    _add_static_footer(slide, spec['footer'])

def _render_ub_timing_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_panel(slide, 0.72, 1.92, 11.88, 2.34, WHITE)
    _add_textbox(slide, 0.98, 2.18, 1.6, 0.22, '简化时间线', 18, INK, True)
    ticks = [('T0', 1.36, BLUE_BG, BLUE), ('T1', 2.92, ORANGE_BG, ORANGE), ('T2', 4.48, BLUE_BG, BLUE), ('T3', 6.04, ORANGE_BG, ORANGE), ('T4', 7.6, BLUE_BG, BLUE)]
    for label, x, fill, tone in ticks:
        _add_chip_ppt(slide, x, 2.54, 0.54, label, fill, tone)
    for x in [1.68, 3.24, 4.8, 6.36, 7.92]:
        line = slide.shapes.add_connector(1, Inches(x), Inches(2.92), Inches(x), Inches(3.96))
        line.line.color.rgb = _rgb(LINE)
        line.line.width = Pt(1)
    base = slide.shapes.add_connector(1, Inches(1.34), Inches(3.28), Inches(8.34), Inches(3.28))
    base.line.color.rgb = _rgb(BLUE)
    base.line.width = Pt(2.6)
    _add_textbox(slide, 0.96, 3.08, 2.16, 0.16, 'ub_rd_start / seq_instr_pulse', 11, BLUE, True)
    _add_textbox(slide, 0.96, 3.42, 2.24, 0.16, 'UB input/weight valid 波前', 11, ORANGE)
    _add_textbox(slide, 0.96, 3.76, 2.0, 0.16, 'PE 看到的有效计算窗口', 11, GREEN)
    _add_textbox(slide, 6.96, 3.1, 4.2, 0.52, '''hold 周期保留最后一拍输入，
而 weight 流主动清 valid，避免 PE22 被重复装载。''', 11, INK)
    for label, x in [('T0 dispatch', 1.1), ('T1 first beat', 2.66), ('T2 wavefront', 4.22), ('T3 last active beat', 5.78), ('T4 hold / drain', 7.34)]:
        _add_textbox(slide, x, 4.02, 1.28, 0.14, label, 9, MUTED)

    _add_panel(slide, 0.72, 4.58, 5.48, 2.08, WHITE)
    _add_textbox(slide, 0.98, 4.86, 1.8, 0.2, 'UB 侧关键语义', 17, INK, True)
    _add_textbox(slide, 0.98, 5.26, 4.58, 0.98, '''1. input 流在最后一个 active beat 后还有 1 个 hold 周期，保证阵列最后一列能完成传播。
2. weight 流故意不保留最后一个 valid，防止 loader 多打一拍时覆盖 active weight。
3. 所以“同样是 UB 读”，input 和 weight 的尾拍策略并不一样。''', 11, INK)

    _add_panel(slide, 6.42, 4.58, 6.18, 2.08, WHITE)
    _add_textbox(slide, 6.7, 4.86, 2.2, 0.2, '和 sequencer 的关系', 17, INK, True)
    _add_textbox(slide, 6.7, 5.26, 5.18, 0.98, '''1. dispatch 只负责发起读流，不代表计算已经结束。
2. 真正的完成边界在 VPU drain，sequencer 只有等 drain 才能 ADVANCE。
3. wait_after 的意义，就是用系统级同步点把 UB / PE / VPU 的多拍行为收口。''', 11, INK)
    _add_chip_ppt(slide, 6.82, 6.18, 1.18, 'dispatch != done', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 8.22, 6.18, 1.34, 'drain = safe advance', GREEN_BG, GREEN)
    _add_static_footer(slide, spec['footer'])

def _render_ub_update_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_clean_card(slide, 0.72, 1.88, 5.82, 4.82, 'Bias Update 路径', 'ptr_sel=5，grad_bias_or_weight=0\n\ndZ / dZ1 stream + old bias + lr -> GD\n\nbias 模式下，如果上一拍 done 过，就把 sub_in_a 切到 value_updated_out。\n\n这意味着同一个 bias 地址可以随着连续样本梯度到来而逐拍累加。', BLUE, WHITE, 11)
    _add_clean_card(slide, 6.76, 1.88, 5.56, 4.82, 'Weight Update 路径', 'ptr_sel=6，grad_bias_or_weight=1\n\nouter-product grad + old weight + lr -> GD\n\nweight 模式下，sub_in_a 固定取 value_old_in，不会链接上次 value_updated_out。\n\n每次 weight 更新都基于当前 tile 读出的旧权重独立完成，更适合 tile 级 outer-product。', ORANGE, WHITE, 11)
    _add_chip_ppt(slide, 0.96, 6.46, 2.55, 'bias: sample-by-sample accumulate', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 3.72, 6.46, 2.2, 'weight: tile-based update', ORANGE_BG, ORANGE)
    _add_static_footer(slide, spec['footer'])

def _render_vpu_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)

    _add_panel(slide, 0.72, 1.88, 7.44, 4.84, WHITE)
    _add_panel(slide, 0.84, 2.04, 0.05, 4.5, BLUE, BLUE)
    _add_textbox(slide, 1.0, 2.08, 2.6, 0.22, 'VPU 子模块与主数据链', 19, INK, True)
    chain = [
        ('sys in', 1.18, 3.24, 1.2, BLUE_BG, BLUE),
        ('Bias', 2.76, 3.24, 1.18, BLUE_BG, BLUE),
        ('LReLU', 4.34, 3.24, 1.24, GREEN_BG, GREEN),
        ('Loss grad', 5.94, 3.24, 1.34, ORANGE_BG, ORANGE),
        ('dAct', 7.5, 3.24, 0.86, PURPLE_BG, PURPLE),
    ]
    for label, x, y, w, fill, tone in chain:
        _add_mini_box(slide, x, y, w, 0.62, label, fill, tone, 14)
    for x1, x2, tone in [(2.38, 2.76, BLUE), (3.94, 4.34, GREEN), (5.58, 5.94, ORANGE), (7.28, 7.5, PURPLE)]:
        _add_arrow_ppt(slide, x1, 3.56, x2, 3.56, tone, 2.4)
    _add_textbox(slide, 1.0, 4.64, 1.2, 0.18, 'pathway 编码', 13, BLUE, True)
    _add_textbox(slide, 1.0, 5.02, 5.84, 0.76, '1100: forward = sys -> bias -> lrelu\n1111: transition = sys -> bias -> lrelu -> loss -> dAct\n0001: backward = sys -> dAct', 12, INK)
    _add_chip_ppt(slide, 1.0, 6.08, 1.26, 'last_H cache', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 2.42, 6.08, 1.06, 'Y from UB', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 3.64, 6.08, 1.46, 'leak factor from UB', GREEN_BG, GREEN)

    _add_panel(slide, 8.46, 1.88, 4.1, 4.84, WHITE)
    _add_panel(slide, 8.58, 2.04, 0.05, 4.5, PURPLE, PURPLE)
    _add_textbox(slide, 8.74, 2.08, 2.2, 0.22, '为什么值得单独讲', 19, INK, True)
    _add_textbox(slide, 8.74, 2.66, 3.3, 0.96, '1. 同一套子模块复用三种训练阶段，不是写死单一路径。\n2. 1111 路径里既要保存 H，又继续送入 loss / derivative。\n3. pathway bit 和 frontend 里的保持寄存器直接耦合。', 12, INK)
    _add_code_card(slide, 8.72, 4.28, 3.54, 1.38, 'vpu.sv 关键点', ['1100: forward pass', '1111: transition path', '0001: backward pass', 'last_H cache feeds leaky_relu_derivative'], PURPLE)
    _add_textbox(slide, 10.54, 4.4, 0.96, 0.14, 'vpu.sv:4+', 9, MUTED)
    _add_static_footer(slide, spec['footer'])

def _render_cycle_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    titles = [
        ('Cycle 0', '权重已装到 active reg', BLUE),
        ('Cycle 1', '第一拍输入进入左边界', ORANGE),
        ('Cycle 2', '波前推进，部分和下沉', GREEN),
        ('Cycle 3', '第一列结果到底部', PURPLE),
        ('Cycle 4', '第二列结果完成', BLUE),
        ('Cycle 5', '结果写回 UB / 进入 VPU', ORANGE),
    ]
    x = 0.56
    label_x = 0.64
    for idx_c, (title, body, tone) in enumerate(titles):
        _add_panel(slide, x, 2.16, 1.84, 1.28, WHITE)
        _add_textbox(slide, x + 0.08, 2.28, 1.42, 0.12, title, 8, INK, True)
        _add_textbox(slide, x + 0.08, 2.46, 1.58, 0.18, body, 6, MUTED)
        for bx, by, lbl in [(x + 0.42, 2.72, 'PE11'), (x + 0.88, 2.72, 'PE12'), (x + 0.42, 3.08, 'PE21'), (x + 0.88, 3.08, 'PE22')]:
            _add_mini_box(slide, bx, by, 0.3, 0.18, lbl, WHITE, INK, 5)
        _add_textbox(slide, x + 1.38, 2.78, 0.26, 0.1, '结果', 5, MUTED)
        _add_textbox(slide, x + 1.32, 2.94, 0.08, 0.1, '7', 5, BLUE, True)
        _add_textbox(slide, x + 1.5, 2.94, 0.08, 0.1, '7', 5, BLUE, True)
        if idx_c >= 3:
            _add_textbox(slide, x + 1.32, 3.1, 0.08, 0.1, '19', 5, ORANGE, True)
        if idx_c >= 4:
            _add_textbox(slide, x + 1.5, 3.1, 0.08, 0.1, '22', 5, ORANGE, True)
            _add_textbox(slide, x + 1.32, 3.24, 0.08, 0.1, '43', 5, GREEN, True)
            _add_textbox(slide, x + 1.5, 3.24, 0.08, 0.1, '50', 5, GREEN, True)
        _add_chip_ppt(slide, label_x, 3.58, 0.54, f'C{idx_c}', BLUE_BG if idx_c % 2 == 0 else ORANGE_BG, BLUE if idx_c % 2 == 0 else ORANGE)
        x += 2.04
        label_x += 2.04

    _add_panel(slide, 0.6, 4.3, 6.22, 2.06, WHITE)
    _add_textbox(slide, 0.86, 4.56, 1.4, 0.18, '讲法建议', 15, BLUE, True)
    _add_textbox(slide, 0.86, 4.98, 5.18, 0.92, '1. 先说 weight-stationary，权重先装后切换。\n2. 再说 input 从左向右、psum 自上而下。\n3. 最后落到为什么需要 wait_after 和 drain 才能收口。', 14, INK, True)

    _add_panel(slide, 7.1, 4.3, 5.26, 2.06, WHITE)
    _add_textbox(slide, 7.36, 4.56, 1.6, 0.18, '配套文件', 15, BLUE, True)
    _add_chip_ppt(slide, 7.38, 5.04, 1.74, 'tpu_systolic_cycle_demo.gif', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 7.38, 5.5, 1.9, 'tpu_systolic_cycle_strip.png', GREEN_BG, GREEN)
    _add_textbox(slide, 7.36, 5.98, 4.32, 0.28, 'GIF 可以单独打开给面试官看；PPT 中放静态条带更稳，避免不同查看器对动画支持不一致。', 10, MUTED)
    _add_static_footer(slide, spec['footer'])



def _render_integration_scope_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.98, '集成页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.9, 1.66, 1.08, 'Integration', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 3.16, 1.66, 1.08, 'wrapper', WHITE, MUTED)
    _add_chip_ppt(slide, 4.42, 1.66, 0.98, 'closure', WHITE, MUTED)
    _add_panel(slide, 9.12, 1.62, 3.24, 0.34, WHITE)
    _add_textbox(slide, 9.34, 1.72, 2.72, 0.12, '目标: 把改动范围和闭环主链一次讲清。', 8, MUTED)

    _add_panel(slide, 0.72, 1.86, 11.88, 1.5, WHITE)
    _add_textbox(slide, 0.98, 2.08, 2.0, 0.18, '顶层闭环主链', 16, BLUE, True)
    chain = [
        ('compiler', 1.02, 2.48, 1.08, BLUE_BG, BLUE),
        ('tpu_soc', 2.68, 2.48, 1.1, ORANGE_BG, ORANGE),
        ('frontend', 4.28, 2.48, 1.08, GREEN_BG, GREEN),
        ('tpu core', 5.88, 2.48, 1.02, PURPLE_BG, PURPLE),
        ('UB / SA / VPU', 7.44, 2.48, 1.48, BLUE_BG, BLUE),
        ('evidence', 9.4, 2.48, 1.06, GREEN_BG, GREEN),
    ]
    for label, x, y, w, fill, tone in chain:
        _add_mini_box(slide, x, y, w, 0.5, label, fill, tone, 12)
    for x1, x2, tone in [(2.1, 2.68, BLUE), (3.78, 4.28, ORANGE), (5.38, 5.88, GREEN), (6.9, 7.44, PURPLE), (8.92, 9.4, BLUE)]:
        _add_arrow_ppt(slide, x1, 2.74, x2, 2.74, tone, 2.2)
    _add_textbox(slide, 0.98, 3.0, 9.9, 0.18, '软件产物进入 SoC，控制前端拉起执行，最后再由验证侧回收证据。', 10, MUTED)

    cards = [
        (0.76, 3.68, 2.64, 2.48, 'SoC / Wrapper 接入', 'AXI-Lite slave\nCTRL / STATUS + readable outputs\n把 TPU 放进 SoC 顶层语境', BLUE, '01'),
        (3.72, 3.68, 2.64, 2.48, 'Frontend 控制闭环', '寄存器映射 + IMEM\n4-state sequencer + wait/drain\nUB_PUSH 优先写口仲裁', ORANGE, '02'),
        (6.68, 3.68, 2.64, 2.48, 'Core 数据闭环', 'UB -> SA -> VPU -> UB\nwr_ptr restore\nin-UB gradient update', GREEN, '03'),
        (9.64, 3.68, 2.64, 2.48, '软件与验证接入', 'ub_allocator / scheduler / encode\ne2e + convergence tests\n把软件产物和证据链接进来', PURPLE, '04'),
    ]
    for x, y, w, h, title, body, tone, tag in cards:
        _add_clean_card(slide, x, y, w, h, title, body, tone, WHITE, 9, tag)

    _add_chip_ppt(slide, 0.94, 6.46, 0.9, 'wrapper', WHITE, BLUE)
    _add_chip_ppt(slide, 2.0, 6.46, 0.94, 'control', WHITE, ORANGE)
    _add_chip_ppt(slide, 3.1, 6.46, 0.72, 'data', WHITE, GREEN)
    _add_chip_ppt(slide, 3.98, 6.46, 0.94, 'compiler', WHITE, PURPLE)
    _add_chip_ppt(slide, 5.08, 6.46, 0.94, 'evidence', WHITE, BLUE)
    _add_static_footer(slide, spec['footer'])


def _render_verification_stack_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.98, '验证页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.9, 1.66, 0.86, 'Verify', GREEN_BG, GREEN)
    _add_chip_ppt(slide, 2.92, 1.66, 1.08, '4 levels', WHITE, MUTED)
    _add_chip_ppt(slide, 4.18, 1.66, 1.2, 'e2e + train', WHITE, MUTED)
    _add_panel(slide, 9.08, 1.62, 3.28, 0.34, WHITE)
    _add_textbox(slide, 9.3, 1.72, 2.76, 0.12, '目标: 证明不是只跑通一次，而是有分层证据链。', 8, MUTED)

    _add_panel(slide, 0.72, 1.88, 6.34, 4.88, WHITE)
    _add_textbox(slide, 0.98, 2.08, 2.28, 0.2, '验证证据阶梯', 18, BLUE, True)
    _add_textbox(slide, 0.98, 2.42, 4.48, 0.16, '先把模块语义钉住，再打通顶层，再证明 e2e 和训练收敛。', 10, MUTED)
    layers = [
        ('Level 1  模块 reference', 1.06, 2.78, 5.1, 0.68, 'bias / pe / vpu / ub reference tests', BLUE, BLUE_BG),
        ('Level 2  顶层功能验证', 1.38, 3.64, 4.78, 0.68, 'systolic + top functional checks', ORANGE, ORANGE_BG),
        ('Level 3  AXI-Lite e2e', 1.7, 4.5, 4.46, 0.68, 'load UB + load IMEM + start + compare output', GREEN, GREEN_BG),
        ('Level 4  训练收敛', 2.02, 5.36, 4.14, 0.68, '12 epoch / XOR = (0,1,1,0)', PURPLE, PURPLE_BG),
    ]
    prev_x = None
    prev_y = None
    for title, x, y, w, h, body, tone, fill in layers:
        _add_panel(slide, x, y, w, h, fill, fill)
        _add_textbox(slide, x + 0.18, y + 0.1, w - 0.36, 0.16, title, 13, tone, True)
        _add_textbox(slide, x + 0.18, y + 0.34, w - 0.36, 0.14, body, 9, INK)
        if prev_x is not None:
            _add_arrow_ppt(slide, prev_x, prev_y, x + 0.18, y, tone, 2.0)
        prev_x = x + w * 0.5
        prev_y = y

    _add_panel(slide, 7.3, 1.88, 5.28, 2.0, WHITE)
    _add_textbox(slide, 7.56, 2.1, 2.0, 0.2, '真实证据资产', 17, GREEN, True)
    _add_chip_ppt(slide, 7.56, 2.5, 1.62, 'waveforms/tpu_soc.vcd', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 9.38, 2.5, 1.42, 'compiler/out/imem.hex', ORANGE_BG, ORANGE)
    _add_chip_ppt(slide, 11.0, 2.5, 1.2, 'schedule.json', GREEN_BG, GREEN)
    _add_textbox(slide, 7.56, 2.98, 4.46, 0.34, '波形、指令和地址分配文件都能反向对上 e2e、scoreboard 和收敛结果。', 10, INK)

    _add_panel(slide, 7.3, 4.14, 2.28, 2.14, WHITE)
    _add_textbox(slide, 7.56, 4.38, 1.56, 0.18, '结果指标', 16, BLUE, True)
    _add_metric_card(slide, 7.56, 4.76, 1.74, 0.82, '顶层回归', '41 / 41 PASS', BLUE)
    _add_metric_card(slide, 7.56, 5.64, 1.74, 0.82, '训练收敛', '12 epoch', GREEN)

    _add_panel(slide, 9.82, 4.14, 2.76, 2.14, WHITE)
    _add_textbox(slide, 10.08, 4.38, 1.72, 0.18, '为什么可信', 16, ORANGE, True)
    _add_textbox(slide, 10.08, 4.78, 2.08, 0.82, '1. 模块语义先钉住。\n2. scoreboard 保证数值对齐。\n3. e2e 证明控制链打通。\n4. convergence 说明不是只跑一次。', 10, INK)
    _add_chip_ppt(slide, 10.08, 5.86, 0.92, 'scoreboard', WHITE, PURPLE)
    _add_chip_ppt(slide, 11.18, 5.86, 0.98, 'convergence', WHITE, GREEN)

    _add_chip_ppt(slide, 0.86, 6.36, 0.94, 'module', WHITE, BLUE)
    _add_chip_ppt(slide, 1.98, 6.36, 0.86, 'top', WHITE, ORANGE)
    _add_chip_ppt(slide, 3.0, 6.36, 0.84, 'e2e', WHITE, GREEN)
    _add_chip_ppt(slide, 3.98, 6.36, 0.9, 'train', WHITE, PURPLE)
    _add_chip_ppt(slide, 5.04, 6.36, 1.0, 'evidence', WHITE, BLUE)
    _add_static_footer(slide, spec['footer'])


def _render_contribution_reference(slide, spec, idx):
    _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
    _add_clean_frame(slide)
    _add_chip_ppt(slide, 0.74, 1.66, 0.98, '贡献页', BLUE_BG, BLUE)
    _add_chip_ppt(slide, 1.9, 1.66, 1.18, 'My fixes', PURPLE_BG, PURPLE)
    _add_chip_ppt(slide, 3.26, 1.66, 1.02, 'before/after', WHITE, MUTED)
    _add_chip_ppt(slide, 4.46, 1.66, 1.02, 'system bugs', WHITE, MUTED)
    _add_panel(slide, 9.06, 1.62, 3.3, 0.34, WHITE)
    _add_textbox(slide, 9.28, 1.72, 2.78, 0.12, '目标: 把个人贡献落到具体 bug、文件和系统价值。', 8, MUTED)

    _add_textbox(slide, 0.82, 1.92, 3.48, 0.22, '4 个系统级修复', 18, BLUE, True)
    _add_textbox(slide, 0.82, 2.24, 6.62, 0.18, '重点不是零散改码，而是把阶段同步、路径保持、写口仲裁、参数区保护这 4 条系统边界补齐。', 10, MUTED)

    fixes = [
        (0.72, 2.72, 5.78, 1.76, '显式等待语义', 'Before: dispatch 后只能猜结束\nAfter: wait_after + SEQ_WAIT + vpu_drain\nValue: 多拍 VPU 有显式同步点', BLUE, 'Fix 01', ['scheduler.py', 'encode_instrs.py', 'tpu_frontend_axil.sv']),
        (6.78, 2.72, 5.82, 1.76, '路径保持', 'Before: decode 只持续 1 拍\nAfter: vpu_pathway_reg 持续保持\nValue: 同一条指令跨多拍不漂移', ORANGE, 'Fix 02', ['vpu_pathway_reg', 'dispatch', 'pathway hold']),
        (0.72, 4.8, 5.78, 1.76, 'Host / CU 写口仲裁', 'Before: UB_PUSH 与 CU 写口会打架\nAfter: host write mux 让 UB_PUSH 优先\nValue: 参数装载和训练写回可共存', GREEN, 'Fix 03', ['UB_PUSH', 'host write mux', 'frontend_axil']),
        (6.78, 4.8, 5.82, 1.76, 'wr_ptr_base / restore', 'Before: runtime writeback 可能踩参数区\nAfter: 记录 wr_ptr_base 并在 start restore\nValue: 参数静态区和训练动态区被切开', PURPLE, 'Fix 04', ['wr_ptr_base', 'restore', 'unified_buffer_v3.sv']),
    ]
    for x, y, w, h, title, body, tone, tag, chips in fixes:
        _add_clean_card(slide, x, y, w, h, title, body, tone, WHITE, 10, tag)
        chip_x = x + 0.24
        chip_y = y + h - 0.34
        for label in chips:
            width = max(0.82, min(1.52, 0.42 + 0.08 * len(label)))
            _add_chip_ppt(slide, chip_x, chip_y, width, label, WHITE, tone)
            chip_x += width + 0.12

    _add_chip_ppt(slide, 0.86, 6.4, 1.14, '同步边界', WHITE, BLUE)
    _add_chip_ppt(slide, 2.18, 6.4, 1.14, '路径边界', WHITE, ORANGE)
    _add_chip_ppt(slide, 3.5, 6.4, 1.14, '写口边界', WHITE, GREEN)
    _add_chip_ppt(slide, 4.82, 6.4, 1.14, '地址边界', WHITE, PURPLE)
    _add_static_footer(slide, spec['footer'])



SLIDE_RENDERERS = {
    '系统总览': _render_system_overview_reference,
    '项目级 RTL 结构': _render_rtl_reference,
    '编译器与指令组织': _render_compiler_reference,
    'compiler_animation_demo': _render_compiler_animation_demo,
    'Frontend 总览': _render_frontend_reference,
    'Frontend 再拆一页': _render_frontend_split_reference,
    'Unified Buffer 设计': _render_ub_reference,
    'wr_ptr / base / restore': _render_wptr_reference,
    'UB 读流与 PE 时序对齐': _render_ub_timing_reference,
    'UB 内梯度下降更新': _render_ub_update_reference,
    'VPU 单独展开': _render_vpu_reference,
    'PE 与计算阵列': _render_pe_reference,
    'pe_animation_demo': _render_pe_animation_demo,
    'ub_timing_animation_demo': _render_ub_timing_animation_demo,
    'vpu_animation_demo': _render_vpu_animation_demo,
    'frontend_animation_demo': _render_frontend_animation_demo,
    'PE 六帧动画演示': _render_pe_animation_demo,
    '关键控制 RTL 特写': _render_control_reference,
    '逐拍计算动态': _render_cycle_reference,
    '验证波形与回归覆盖': _render_wave_reference,
    'integration_scope': _render_integration_scope_reference,
    '顶层集成与改动范围': _render_integration_scope_reference,
    'verification_stack': _render_verification_stack_reference,
    '验证体系与证据分层': _render_verification_stack_reference,
    'contribution_fixes': _render_contribution_reference,
    '个人贡献与关键修复前后对比': _render_contribution_reference,
    '结果、边界与追问方向': _render_results_reference,
}


def _set_slide_notes(slide, notes_text):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    if text_frame is None:
        return
    lines = [line.strip() for line in notes_text.strip().split('\n') if line.strip()]
    if not lines:
        return
    text_frame.clear()
    text_frame.paragraphs[0].text = lines[0]
    for line in lines[1:]:
        text_frame.add_paragraph().text = line


def _speaker_notes_for_spec(spec, deck_title):
    title = spec.get('title', '')
    subtitle = spec.get('subtitle', '')
    is_appendix = 'appendix' in deck_title.lower() or '附录' in subtitle

    def _format_section(header, lines):
        body = [line for line in lines if line]
        return '\n'.join([header, *body])

    def _format_main(short_lines, long_lines):
        return '\n\n'.join([
            _format_section('90 秒版：', short_lines),
            _format_section('3 分钟版：', long_lines),
        ])

    def _format_appendix(lines):
        return _format_section('追问版：', lines)

    if spec.get('layout') == 'cover':
        if is_appendix:
            return _format_appendix([
                '这一页先说明，后面的内容不是主讲主线，而是被追问时才展开的附录。',
                '附录里有 GIF 单步演示、细节拆页、控制修复、验证体系和个人贡献说明。',
                '我的用法是先讲完 8 页主讲版，再根据老师的问题切到对应附录页。',
                '讲附录时不追求全讲完，只回答老师刚刚追问的那个点。',
            ])
        return _format_main(
            [
                '大家好，我这次汇报的是 TinyTPU AXI-Lite SoC 项目。',
                '主讲版我只讲一条主线。',
                '软件怎么送进来，硬件怎么跑起来，最后我再用证据收结论。',
            ],
            [
                '大家好，我这次汇报的是 TinyTPU AXI-Lite SoC 项目。',
                '这个项目的重点，不是单独把一个 tiny-tpu 跑起来，而是把 compiler、AXI-Lite 前端、UB、PE、VPU 和验证链路接成一套完整系统。',
                '所以主讲版我只保留 8 页，顺序就是系统、RTL、编译器、UB、PE、证据和结果。',
                '我先把主线讲清，后面如果老师追问 Frontend、控制修复或者动态时序，我再切附录继续展开。',
            ],
        )

    main_notes = {
        '系统总览': {
            'short': [
                '这页就是全局图。',
                '左边是软件和 compiler，中间是 Frontend 和 IMEM，下面是 UB、PE、VPU。',
                '右边这些结果说明整条链已经跑通。',
            ],
            'long': [
                '这页就是项目地图。',
                '最左边是软件侧，model spec 经过 compiler 以后，变成 ub_map、schedule 和 IMEM。',
                '中间是 AXI-Lite Frontend，负责配置、装载程序、推进 sequencer，还有 wait 边界控制。',
                '再往后是 UB、PE、VPU 真正执行数据流和训练路径。',
                '最右边这些波形、scoreboard 和收敛结果，负责证明整条链已经闭环。',
            ],
        },
        '项目级 RTL 结构': {
            'short': [
                '这页只看顶层怎么接。',
                '最外层是 tpu_soc，中间是 frontend_axil，里面是 core path。',
                '重点不是包接口，而是控制、数据、观测三条链都接通了。',
            ],
            'long': [
                '如果只看这页，我会把系统拆成三层。',
                '第一层是 wrapper，负责 SoC 接口和系统信号。',
                '第二层是 frontend_axil，负责寄存器、IMEM、sequencer 和指令发射。',
                '第三层才是 core path，把 UB、PE、VPU、写回和状态观测串起来。',
                '重点不是接口适配，而是真正把 control、data、observe 三条链闭合起来。',
            ],
        },
        '编译器与指令组织': {
            'short': [
                '这页只看软件怎么落到硬件。',
                '流程就是 model spec、UB 分配、scheduler、IMEM 编码，最后 host 写进 SoC。',
                '所以右边这些不是附件，而是硬件真正执行的输入。',
            ],
            'long': [
                '这页我不会把它讲成一个很重的 compiler 框架。',
                '对这个项目最重要的其实就三件事：地址怎么分、阶段怎么排、控制字怎么编码。',
                '所以流程上先做 UB 地址分配，再由 scheduler 写出阶段边界和 wait_after，最后 encode 成 IMEM words。',
                '这样软件侧产物和 Frontend、UB、PE 的行为就是一一对应的，不是两边各讲各的。',
            ],
        },
        'Unified Buffer 设计': {
            'short': [
                '这页只记一句话，UB 是数据汇合点。',
                '同一个 UB 同时接 host 装载、core 读、写回和 update。',
                '所以软件、控制和执行三边最后都汇到这里。',
            ],
            'long': [
                '如果这页只记一句话，那就是 UB 是整个系统的数据中枢。',
                '左边是统一地址空间，编译器先把输入、权重、bias 和中间张量放进去。',
                '中间可以看到 host load、core read、VPU writeback 和 gradient update 都汇到同一个 UB。',
                '右边三个关键词是这页重点：wr_ptr_base 和 restore、ptr_select 读流语义，还有为什么 gradient_descent 放在 UB 内部。',
                '这样后面讲训练写回和时序对齐的时候，就都有明确落点。',
            ],
        },
        'PE 与计算阵列': {
            'short': [
                '这页先看 PE 阵列长什么样。',
                '它是 2x2 的 weight-stationary systolic array，input 左进，weight 上装，psum 下沉。',
                '我只抓三个点：双 weight reg、pe_enabled 语义，还有波前能和波形对上。',
            ],
            'long': [
                '这页不只是讲一个 PE 做乘加，而是讲这个阵列为什么可解释、可验证。',
                '左边是 2x2 weight-stationary 阵列，input 横向推进，weight 从顶部装入，valid 和 psum 跟着波前走。',
                '右边三个 RTL 点，分别对应 load 和 switch 分离、状态保持、还有输出使能。',
                '所以后面讲真实波形的时候，阵列时序不是抽象概念，而是能和证据页逐项对齐。',
            ],
        },
        '验证波形与回归覆盖': {
            'short': [
                '这页是最关键的证据页。',
                '我会沿着一个 busy 窗口看：dispatch、UB 读流、阵列输出、vpu_drain。',
                '这些时间关系一旦对上，控制链和数据链就真正对齐了。',
            ],
            'long': [
                '前面几页讲的是设计，这一页讲的是证据。',
                '我通常会沿着一个 busy 窗口往下说：先看 dispatch 发起，再看 UB 什么时候开始出数，再看阵列输出和 vpu_drain 什么时候结束。',
                '如果这些时间关系是对的，就说明控制链和数据链是真的对齐了，不是靠口头假设拼出来的。',
                '右边的 scoreboard 覆盖项和回归统计，再把数值层面的正确性补上。',
                '所以这页是我最想让老师看到的一页。',
            ],
        },
        '结果、边界与追问方向': {
            'short': [
                '最后先给一句结论，这不是单点 demo，而是一条完整闭环。',
                '左边这些结果说明项目已经稳定跑通，右边我也主动交代边界。',
                '如果老师继续追问扩展或控制细节，我就切附录继续讲。',
            ],
            'long': [
                '最后一页我会先给结论：这个项目已经形成从 compile 到 control、从 execute 到 evidence 的完整闭环。',
                '左边这些结果说明它不只是能跑一拍，而是能稳定跑通，并且有波形、回归和收敛这些证据支撑。',
                '右边我会主动交代边界，它还是一个 2x2、Q8.8 的 tiny-tpu 原型，没有 DMA、中断和 ICG。',
                '我主动讲边界，是为了让结论更可信，不是为了回避问题。',
                '如果老师继续追问怎么扩到 8x8、怎么做 DMA，后面我就直接切附录继续讲。',
            ],
        },
    }
    if not is_appendix and title in main_notes:
        return _format_main(main_notes[title]['short'], main_notes[title]['long'])

    appendix_notes = {
        'Frontend 总览': [
            '这页用来展开 frontend 到底做了什么。',
            '先看寄存器和 IMEM，再看 four-state sequencer，最后落到 decode 和写口仲裁。',
            '这里要强调，frontend 不是简单寄存器壳，而是整个控制面的收口点。',
        ],
        'Frontend 再拆一页': [
            '如果上一页看起来还是有点密，这一页就把配置面和运行面完全拆开讲。',
            '左边只讲 host 怎么把参数和程序喂进去，右边只讲 start、dispatch、wait 和 advance 怎样推进。',
            '重点是 pathway 保持寄存器，让多拍 VPU 不会在中途路由漂移。',
        ],
        'wr_ptr / base / restore': [
            '这页单独解释，为什么训练写回不会覆盖前面 host 装进去的参数。',
            '思路不是物理分 bank，而是先记住 wr_ptr_base，再在 start 的时候把 wr_ptr 恢复到这个边界。',
            '这样前面是静态参数区，后面才是运行时 writeback 和 update 区。',
        ],
        'UB 读流与 PE 时序对齐': [
            '这页把数据流和控制流第一次放到同一条时间线上。',
            '我会按 dispatch、first beat、wavefront、last beat、hold 和 drain 的顺序往下讲。',
            '重点是 wait_after，它用系统级同步点把 UB、PE 和 VPU 的多拍行为收口。',
        ],
        'UB 内梯度下降更新': [
            '这页要讲清楚，bias update 和 weight update 虽然共用一个 GD 单元，但驱动语义不一样。',
            'bias 模式允许同一地址随着连续样本逐拍累加，weight 模式则是基于旧权重独立更新。',
            '核心点不是复用算子，而是根据训练对象切换数据依赖关系。',
        ],
        'VPU 单独展开': [
            '这页把 VPU 从后处理黑盒讲成可重组的训练路径单元。',
            'forward、transition 和 backward 复用同一套子模块，但走的链路不一样。',
            '它和 frontend 的 pathway 保持寄存器直接耦合，所以这是系统级控制点。',
        ],
        '关键控制 RTL 特写': [
            '这页是控制收口页，只讲三个真正决定系统稳定性的点。',
            '分别是 wait_after 和 vpu_drain 的等待边界、pathway_reg 的多拍保持，以及 host write mux 的写口仲裁。',
            '这里不展开所有 RTL，只讲系统级 bug 是怎么被收住的。',
        ],
        '逐拍计算动态': [
            '这页是帮助解释原理的示意页，不是最终证据页。',
            '我的讲法是先说 weight-stationary，再说 input 横向传播和 psum 纵向下沉。',
            '如果老师已经接受原理，我不会在这页停太久，会直接回到真实波形页。',
        ],
        '顶层集成与改动范围': [
            '这页适合回答，我到底改了哪些东西。',
            '我会按 wrapper、frontend、core、软件与验证四块去讲，每一块都不是孤立改动，而是为了把系统闭环真正接起来。',
            '重点不是接口壳，而是 compile、control、execute、verify 四段终于形成了可以跑通的链。',
        ],
        '验证体系与证据分层': [
            '这页适合回答，我怎么证明这不是一次性跑通的 demo。',
            '顺序是模块级语义、顶层 e2e、波形与 scoreboard，再到训练收敛。',
            '这样控制正确性、数值正确性和系统完整性是分层建立起来的。',
        ],
        '个人贡献与关键修复前后对比': [
            '这页适合回答，你个人到底做了什么。',
            '我会把贡献收束成同步边界、路径边界、写口边界和地址边界四类修复。',
            '它们共同的特点是，不是单点改码，而是直接决定系统能不能稳定跑通。',
        ],
    }
    if title in appendix_notes:
        return _format_appendix(appendix_notes[title])

    gif_steps = {
        'Compiler GIF 演示 T': [
            '这页是 compiler 动态演示的第一步。这里先看 model spec 怎么被拆成可分配的张量和地址语义。',
            '这页是 compiler 动态演示的第二步。这里重点讲 scheduler 怎么把阶段边界和 wait_after 标出来。',
            '这页是 compiler 动态演示的第三步。这里重点讲 encode 怎么把阶段语义变成 UB_RD 和 IMEM word。',
            '这页是 compiler 动态演示的第四步。这里重点讲 host 怎么把 IMEM 写进去，并把系统带到 ready 状态。',
        ],
        'Frontend GIF 演示 T': [
            '这页是 frontend 单步演示的第一步。先看 host 把参数和 IMEM 装进去，系统还没有真正 dispatch。',
            '这页是 frontend 单步演示的第二步。start 拉起以后，sequencer 在 dispatch 这一拍发出真正的控制脉冲。',
            '这页是 frontend 单步演示的第三步。这里重点不是再发指令，而是 WAIT 和 pathway 保持怎样把多拍行为稳住。',
            '这页是 frontend 单步演示的第四步。只有 drain 到来以后，状态机才 ADVANCE 到下一条指令。',
        ],
        'UB 时序 GIF 演示 T': [
            '这页是 UB 时序演示的第一步。dispatch 只负责发起新读流，不代表计算已经结束。',
            '这页是 UB 时序演示的第二步。first beat 进入阵列边界，真正的数据流开始跑起来。',
            '这页是 UB 时序演示的第三步。wavefront 在阵列内部推进，这是 input、valid 和 psum 对齐的关键。',
            '这页是 UB 时序演示的第四步。last beat 进入尾拍区，这时 input 和 weight 的尾拍语义开始分化。',
            '这页是 UB 时序演示的第五步。hold、drain 和 advance 把整次读流真正收口。',
        ],
        'VPU GIF 演示 T': [
            '这页是 VPU 动态演示的第一步。先看 forward 1100 这条路径怎样经过 bias 和激活。',
            '这页是 VPU 动态演示的第二步。transition 1111 会继续接 loss 和 derivative，是最重的一条路径。',
            '这页是 VPU 动态演示的第三步。backward 0001 主要强调 dAct 这一条反向链。',
            '这页是 VPU 动态演示的第四步。这里要讲清 writeback 和 pathway 保持为什么是系统级问题。',
        ],
        'PE GIF 演示 T': [
            '这页是 PE 动态演示的第一步。这里先看权重已经稳定装到 active reg，计算还没开始。',
            '这页是 PE 动态演示的第二步。第一拍输入进入左边界，PE11 和 PE21 开始第一拍 MAC。',
            '这页是 PE 动态演示的第三步。波前向右和向下推进，第二列开始接力计算。',
            '这页是 PE 动态演示的第四步。第一列结果先到底，阵列还在继续传播后续部分和。',
            '这页是 PE 动态演示的第五步。2x2 的结果窗口已经逐渐收齐。',
            '这页是 PE 动态演示的第六步。结果被收集后送后续链路，可以写回 UB 或继续进 VPU。',
        ],
    }
    for prefix, lines in gif_steps.items():
        if title.startswith(prefix):
            m = re.search(r'T(\d+)$', title)
            if m:
                step = int(m.group(1))
                if 0 <= step < len(lines):
                    return _format_appendix([
                        lines[step],
                        '这类页的讲法是只讲当前这一帧，不要抢到下一步。',
                        '讲完以后点一下翻页，下一页的 GIF 会自动播放到下一步。',
                    ])

    if is_appendix:
        return _format_appendix([
            f'这页的标题是《{title}》。',
            '这页属于附录页，主要作用是在老师继续追问时，把某个模块的细节单独展开。',
            '如果现场时间紧，我会直接跳过；如果对方追问，我就按标题把这一页展开讲。',
        ])
    return _format_main(
        [
            f'这一页的标题是《{title}》。',
            '这一页属于主讲版补充页，我会只讲这页最关键的一句话结论。',
        ],
        [
            f'这一页的标题是《{title}》。',
            '这页没有命中特定讲稿模板，所以现场我会按标题、主图和右侧结论自然展开。',
            '如果老师对这里继续深挖，我会结合附录页再补充细节。',
        ],
    )



def build_ppt(slide_specs, out_path, deck_title):
    global PPT_TOTAL_PAGES
    PPT_TOTAL_PAGES = len(slide_specs)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    prs.core_properties.title = deck_title
    prs.core_properties.subject = 'Expanded Hisilicon-style project deck'
    prs.core_properties.author = 'Codex'

    for idx, spec in enumerate(slide_specs, start=1):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = _rgb(BG)

        if spec.get('layout') == 'cover':
            _render_cover_clean(slide, spec)
            _set_slide_notes(slide, _speaker_notes_for_spec(spec, deck_title))
            continue

        renderer = SLIDE_RENDERERS.get(spec.get('renderer', spec['title']))
        if renderer is not None:
            renderer(slide, spec, idx)
            _set_slide_notes(slide, _speaker_notes_for_spec(spec, deck_title))
            continue

        _render_header_clean(slide, idx, spec['title'], spec['subtitle'])
        _add_clean_frame(slide)
        for metric in spec.get('metrics', []):
            _add_metric_card(slide, metric['x'], metric['y'], metric['w'], metric['h'], metric['title'], metric['value'], metric['tone'])
        for card in spec.get('cards', []):
            _add_clean_card(slide, card['x'], card['y'], card['w'], card['h'], card['title'], card['body'], card['tone'], WHITE)
        if spec.get('code'):
            code = spec['code']
            _add_code_card(slide, code['x'], code['y'], code['w'], code['h'], code['title'], code['lines'])
        if spec.get('chips'):
            x = 0.95
            for label, fill, tone in spec['chips']:
                w = max(1.3, min(2.6, 0.48 + 0.11 * len(label)))
                _add_chip_ppt(slide, x, 6.54, w, label, fill, tone)
                x += w + 0.18
        if spec.get('footer'):
            _add_static_footer(slide, spec['footer'])
        _set_slide_notes(slide, _speaker_notes_for_spec(spec, deck_title))

    prs.save(str(out_path))


def sync_outputs():
    DROP.mkdir(parents=True, exist_ok=True)
    SHARE.mkdir(parents=True, exist_ok=True)
    artifacts = [slide['path'] for slide in SLIDES] + [OUT_PPT, OUT_APPENDIX_PPT, GIF_PATH, STRIP_PATH, PE_ANIM_BASE_PATH, *PE_ANIM_GIF_PATHS, UB_TIMING_ANIM_BASE_PATH, *UB_TIMING_ANIM_GIF_PATHS, VPU_ANIM_BASE_PATH, *VPU_ANIM_GIF_PATHS, FRONTEND_ANIM_BASE_PATH, *FRONTEND_ANIM_GIF_PATHS, COMPILER_ANIM_BASE_PATH, *COMPILER_ANIM_GIF_PATHS]
    for src in artifacts:
        try:
            shutil.copy2(src, DROP / src.name)
        except PermissionError:
            pass
    for src, dst in [(OUT_PPT, DROP_PPT), (OUT_APPENDIX_PPT, DROP_APPENDIX_PPT), (OUT_PPT, SHARE_PPT), (OUT_APPENDIX_PPT, SHARE_APPENDIX_PPT)]:
        try:
            shutil.copy2(src, dst)
        except PermissionError:
            pass


def main():
    ub_map = load_ub_map()
    schedule = load_schedule()
    render_cover(SLIDES[0]['path'])
    render_system(SLIDES[1]['path'])
    render_rtl(SLIDES[2]['path'])
    render_compiler(SLIDES[3]['path'], ub_map, schedule)
    render_frontend(SLIDES[4]['path'])
    render_frontend_split(SLIDES[5]['path'])
    render_ub(SLIDES[6]['path'], ub_map)
    render_wptr(SLIDES[7]['path'])
    render_ub_timing(SLIDES[8]['path'])
    render_ub_update(SLIDES[9]['path'])
    render_vpu(SLIDES[10]['path'])
    render_pe(SLIDES[11]['path'])
    render_control(SLIDES[12]['path'])
    render_cycle(SLIDES[13]['path'])
    render_wave(SLIDES[14]['path'])
    render_results(SLIDES[15]['path'])
    render_pe_animation_assets()
    render_ub_timing_animation_assets()
    render_vpu_animation_assets()
    render_frontend_animation_assets()
    render_compiler_animation_assets()
    build_ppt(_build_deck_specs('main'), OUT_PPT, 'TinyTPU Main Talk Deck')
    build_ppt(_build_deck_specs('appendix'), OUT_APPENDIX_PPT, 'TinyTPU Appendix Deck')
    sync_outputs()


if __name__ == '__main__':
    main()
